from __future__ import annotations

import csv
import io
import json
import re
import subprocess
import tempfile
import zipfile
from collections import defaultdict
from collections.abc import Iterable
from pathlib import Path
from typing import Any

from docx import Document
from lxml import etree
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from .models import FileRecord, TextUnit
from .normalize import nfkc, norm, overlap_score
from .roles import infer_area, infer_project, infer_role, infer_version

SUPPORTED = {".pdf", ".docx", ".pptx", ".xlsx", ".xlsm", ".csv", ".tsv", ".json", ".py", ".ipynb", ".md", ".txt", ".yaml", ".yml", ".toml", ".png", ".jpg", ".jpeg"}
TEXT_EXTS = {".csv", ".tsv", ".json", ".py", ".md", ".txt", ".yaml", ".yml", ".toml"}


class DocumentStore:
    def __init__(self, root: Path):
        self.root = root.resolve()
        self.records = self._build_manifest()
        self.by_relative = {norm(r.relative_path): r for r in self.records}
        self.projects = sorted({r.project for r in self.records if r.project})
        self._text_cache: dict[str, list[TextUnit]] = {}
        self.aliases = self._discover_aliases()

    def _build_manifest(self) -> list[FileRecord]:
        out: list[FileRecord] = []
        for path in sorted(self.root.rglob("*")):
            if not path.is_file() or path.name == ".extracted_ok" or path.name.startswith("~$") or path.name.endswith("#") or path.suffix.lower() not in SUPPORTED:
                continue
            rel = path.relative_to(self.root).as_posix()
            out.append(FileRecord(path, rel, path.suffix.lower(), infer_project(rel), infer_area(rel), infer_role(rel), path.name, path.stem, infer_version(rel), path.stat().st_size))
        return out

    def _discover_aliases(self) -> dict[str, str]:
        aliases: dict[str, str] = {}
        for project in self.projects:
            aliases[norm(project)] = project
            compact = re.sub(r"株式会社|医療法人社団|合同会社|有限会社", "", project).strip()
            aliases[norm(compact)] = project
            initials = "".join(part[0].upper() for part in re.findall(r"[A-Za-z]+", project) if part)
            if initials:
                aliases[norm(initials)] = project
        # TERM-BOOK / glossary often contains project abbreviations and internal aliases.
        # Parse tabular rows directly so comma-separated alias candidates are also registered.
        for rec in self.records:
            if rec.role != "internal" and "用語" not in rec.filename:
                continue
            try:
                for unit in self.extract_text_units(rec):
                    rows = unit.metadata.get("rows") or []
                    if rows:
                        for row in rows:
                            if len(row) < 2:
                                continue
                            official = str(row[0]).strip()
                            project = self.resolve_project(official, strict=False)
                            if not project:
                                continue
                            aliases[norm(official)] = project
                            for cell in row[1:3]:
                                for alias in re.split(r"[,、/\s]+", str(cell)):
                                    alias = alias.strip()
                                    if len(alias) >= 2:
                                        aliases[norm(alias)] = project
                    for line in unit.text.splitlines():
                        cols = [x.strip() for x in re.split(r"\t|\s{2,}|\|", line) if x.strip()]
                        if len(cols) < 2:
                            continue
                        project = self.resolve_project(cols[0], strict=False)
                        if project:
                            aliases[norm(cols[1])] = project
            except Exception:
                continue
        # Discover aliases from project-owned documents. This keeps new projects data-driven.
        alias_pattern = re.compile(r"(?:主略称|プロジェクト略称|案件略称|略称|project\s*(?:code|alias))\s*[:：=]\s*([A-Z][A-Z0-9_-]{1,15})", re.I)
        for rec in self.records:
            if not rec.project or rec.role not in {"proposal", "schedule", "final_report", "internal"}:
                continue
            try:
                for unit in self.extract_text_units(rec)[:20]:
                    for match in alias_pattern.findall(unit.text):
                        aliases[norm(match)] = rec.project
            except Exception:
                continue
        return aliases

    def resolve_project(self, hint: str, *, strict: bool = True) -> str:
        nh = norm(hint)
        if not nh:
            return ""
        if nh in getattr(self, "aliases", {}):
            return self.aliases[nh]
        exact = [p for p in self.projects if nh == norm(p) or nh in norm(p) or norm(p) in nh]
        if len(exact) == 1:
            return exact[0]
        scored = sorted(((overlap_score(hint, p), p) for p in self.projects), reverse=True)
        if scored and scored[0][0] >= (0.45 if strict else 0.25):
            if len(scored) == 1 or scored[0][0] > scored[1][0] + 0.08:
                return scored[0][1]
        return ""

    def find(self, *, project_hint: str = "", filename_hint: str = "", extensions: Iterable[str] | None = None, roles: Iterable[str] = (), selected_sources: Iterable[str] = (), limit: int | None = None) -> list[FileRecord]:
        extset = {e.lower() if e.startswith(".") else f".{e.lower()}" for e in extensions} if extensions else None
        project = self.resolve_project(project_hint, strict=False) if project_hint else ""
        selected_norms = [norm(x) for x in selected_sources if str(x).strip()]
        role_set = set(roles)
        scored: list[tuple[float, FileRecord]] = []
        for rec in self.records:
            if extset and rec.extension not in extset:
                continue
            if role_set and rec.role not in role_set:
                continue
            score = 0.0
            if project:
                if rec.project == project:
                    score += 50
                else:
                    continue
            elif project_hint:
                score += 20 * overlap_score(project_hint, rec.project or rec.relative_path)
            if filename_hint:
                nf = norm(filename_hint)
                if nf == norm(rec.filename):
                    score += 45
                elif nf in norm(rec.filename) or norm(rec.filename) in nf:
                    score += 30
                else:
                    score += 8 * overlap_score(filename_hint, rec.filename)
            if selected_norms:
                path_norm = norm(rec.relative_path)
                best = max((1.0 if s == path_norm else 0.9 if s.endswith(path_norm) or path_norm.endswith(s) else overlap_score(s, path_norm) for s in selected_norms), default=0)
                score += 40 * best
            if rec.version == "current":
                score += 2
            score -= len(rec.relative_path) / 10000
            if score > 0 or (not project_hint and not filename_hint and not selected_norms and not role_set):
                scored.append((score, rec))
        scored.sort(key=lambda x: (-x[0], x[1].relative_path))
        records = [rec for _, rec in scored]
        return records[:limit] if limit else records

    def extract_text_units(self, rec: FileRecord) -> list[TextUnit]:
        if rec.relative_path in self._text_cache:
            return self._text_cache[rec.relative_path]
        try:
            units = self._extract_text_units_uncached(rec)
        except Exception as exc:
            units = [TextUnit(rec.relative_path, "parse_error", "", {"error": repr(exc)})]
        self._text_cache[rec.relative_path] = units
        return units

    def _extract_text_units_uncached(self, rec: FileRecord) -> list[TextUnit]:
        ext = rec.extension
        if ext == ".pdf":
            reader = PdfReader(str(rec.path))
            return [TextUnit(rec.relative_path, f"page:{i+1}", page.extract_text() or "", {"page": i + 1}) for i, page in enumerate(reader.pages)]
        if ext == ".docx":
            return self._docx_units(rec)
        if ext == ".pptx":
            return self._pptx_units(rec)
        if ext in {".xlsx", ".xlsm"}:
            return self._xlsx_units(rec)
        if ext in TEXT_EXTS:
            return [TextUnit(rec.relative_path, "file", self._read_text(rec.path), {})]
        if ext == ".ipynb":
            data = json.loads(self._read_text(rec.path))
            units: list[TextUnit] = []
            for i, cell in enumerate(data.get("cells", []), 1):
                source = "".join(cell.get("source", []))
                outputs: list[str] = []
                for output in cell.get("outputs", []):
                    outputs.extend(output.get("text", []))
                    for value in output.get("data", {}).values():
                        if isinstance(value, list):
                            outputs.extend(map(str, value))
                        elif isinstance(value, str):
                            outputs.append(value)
                units.append(TextUnit(rec.relative_path, f"cell:{i}", source + "\n" + "\n".join(outputs), {"cell_type": cell.get("cell_type")}))
            return units
        return []

    def _docx_units(self, rec: FileRecord) -> list[TextUnit]:
        doc = Document(str(rec.path))
        units: list[TextUnit] = []
        for i, p in enumerate(doc.paragraphs, 1):
            if not p.text.strip():
                continue
            runs = [{"text": r.text, "bold": bool(r.bold), "italic": bool(r.italic), "underline": bool(r.underline), "style": r.style.name if r.style else ""} for r in p.runs if r.text]
            units.append(TextUnit(rec.relative_path, f"paragraph:{i}", p.text, {"runs": runs, "style": p.style.name if p.style else ""}))
        for ti, table in enumerate(doc.tables, 1):
            rows = [[cell.text for cell in row.cells] for row in table.rows]
            units.append(TextUnit(rec.relative_path, f"table:{ti}", "\n".join("\t".join(row) for row in rows), {"rows": rows}))
        comments = self._office_comments(rec.path, "word/comments.xml")
        units.extend(TextUnit(rec.relative_path, f"comment:{i}", text, {"comment": True}) for i, text in enumerate(comments, 1))
        return units

    def _pptx_units(self, rec: FileRecord) -> list[TextUnit]:
        prs = Presentation(str(rec.path))
        units: list[TextUnit] = []
        for si, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            runs: list[dict[str, Any]] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for p in shape.text_frame.paragraphs:
                        if p.text:
                            texts.append(p.text)
                        for r in p.runs:
                            if r.text:
                                runs.append({"text": r.text, "bold": bool(r.font.bold), "italic": bool(r.font.italic), "underline": bool(r.font.underline), "size": r.font.size.pt if r.font.size else None})
                if getattr(shape, "has_table", False):
                    rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                    texts.extend("\t".join(row) for row in rows)
            units.append(TextUnit(rec.relative_path, f"slide:{si}", "\n".join(texts), {"runs": runs, "slide": si}))
        comments = self._office_comments(rec.path, "ppt/comments/comment")
        units.extend(TextUnit(rec.relative_path, f"comment:{i}", text, {"comment": True}) for i, text in enumerate(comments, 1))
        return units

    def _xlsx_units(self, rec: FileRecord) -> list[TextUnit]:
        wb = load_workbook(rec.path, data_only=False, read_only=False)
        units: list[TextUnit] = []
        for ws in wb.worksheets:
            rows: list[list[Any]] = []
            nonempty = 0
            for row in ws.iter_rows():
                values = [cell.value for cell in row]
                if any(v not in (None, "") for v in values):
                    rows.append(values)
                    nonempty += 1
                if nonempty >= 300:
                    break
            text = "\n".join("\t".join("" if v is None else str(v) for v in row) for row in rows)
            units.append(TextUnit(rec.relative_path, f"sheet:{ws.title}", text, {"sheet": ws.title, "max_row": ws.max_row, "max_column": ws.max_column}))
        return units

    def render_to_pdf_pages(self, rec: FileRecord, cache_dir: Path) -> list[TextUnit]:
        cache_dir.mkdir(parents=True, exist_ok=True)
        pdf_path = cache_dir / f"{rec.path.stem}_{abs(hash(rec.relative_path)) & 0xffffffff:x}.pdf"
        if not pdf_path.exists():
            with tempfile.TemporaryDirectory() as td:
                subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", td, str(rec.path)], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=90)
                generated = Path(td) / f"{rec.path.stem}.pdf"
                if not generated.exists():
                    raise RuntimeError(f"LibreOffice did not create PDF for {rec.relative_path}")
                pdf_path.write_bytes(generated.read_bytes())
        reader = PdfReader(str(pdf_path))
        # LibreOffice PDFs can contain text that pypdf cannot recover even though
        # the page renders normally.  Use PyMuPDF only for those empty pages so
        # page-based lookups remain portable across LibreOffice releases.
        pypdf_text = [page.extract_text() or "" for page in reader.pages]
        if any(not text.strip() for text in pypdf_text):
            import fitz
            doc = fitz.open(pdf_path)
            pypdf_text = [text or doc[i].get_text() for i, text in enumerate(pypdf_text)]
        return [TextUnit(rec.relative_path, f"page:{i+1}", text, {"rendered_from": rec.extension}) for i, text in enumerate(pypdf_text)]

    @staticmethod
    def _read_text(path: Path) -> str:
        for encoding in ("utf-8-sig", "utf-8", "cp932", "shift_jis"):
            try:
                return path.read_text(encoding=encoding)
            except UnicodeDecodeError:
                continue
        return path.read_text(encoding="utf-8", errors="replace")

    @staticmethod
    def _office_comments(path: Path, prefix: str) -> list[str]:
        out: list[str] = []
        try:
            with zipfile.ZipFile(path) as zf:
                names = [n for n in zf.namelist() if n.startswith(prefix) and n.endswith(".xml")]
                for name in names:
                    root = etree.fromstring(zf.read(name))
                    for node in root.xpath("//*[local-name()='comment']"):
                        text = "".join(node.xpath(".//*[local-name()='t']/text()"))
                        if text.strip():
                            out.append(text.strip())
        except (zipfile.BadZipFile, etree.XMLSyntaxError):
            pass
        return out

    def read_csv(self, rec: FileRecord):
        import pandas as pd
        for encoding in ("utf-8-sig", "utf-8", "cp932", "shift_jis"):
            try:
                return pd.read_csv(rec.path, encoding=encoding)
            except UnicodeDecodeError:
                continue
        return pd.read_csv(rec.path, encoding_errors="replace")

    def load_workbook(self, rec: FileRecord, *, data_only: bool = False):
        return load_workbook(rec.path, data_only=data_only, read_only=False, keep_vba=rec.extension == ".xlsm")
