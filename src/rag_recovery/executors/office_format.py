from __future__ import annotations

import re
import zipfile
from dataclasses import dataclass
from typing import Any

from docx import Document
from lxml import etree
from pptx import Presentation

from ..models import Evidence, ExecutionResult, QueryPlan, Question
from ..normalize import nfkc, norm
from ..store import DocumentStore
from .base import Executor
from .table import _is_color, _rgb


@dataclass
class OfficeFormatCommentExecutor(Executor):
    name: str = "office_format"

    def execute(self, question: Question, plan: QueryPlan, store: DocumentStore) -> ExecutionResult:
        q = nfkc(question.text)
        project = plan.project_hints[0] if plan.project_hints else ""
        filename = plan.filename_hints[0] if plan.filename_hints else ""
        records = store.find(project_hint=project, filename_hint=filename, extensions={".docx", ".pptx", ".xlsx", ".xlsm", ".pdf"}, selected_sources=question.selected_sources, limit=16)
        if not records:
            return ExecutionResult.abstain("Office資料を特定できない")
        if "コメント" in q or "注釈" in q:
            return self._comments(q, records, store)
        if any(k in q for k in ("太字", "下線", "イタリック", "斜体")):
            return self._run_format(q, records)
        if any(k in q for k in ("ハイライト", "赤字", "黄色", "青色", "文字色")):
            return self._color_format(q, records, store)
        if "ページ" in q:
            # The document executor handles rendered pagination; leave route fallback intact.
            return ExecutionResult.abstain("ページ位置はDocumentLookupExecutorへ委譲")
        return ExecutionResult.abstain("要求書式を解析できない")

    def _run_format(self, q, records):
        require_bold = "太字" in q
        require_underline = "下線" in q
        require_italic = "イタリック" in q or "斜体" in q
        hits = []
        for rec in records:
            if rec.extension == ".docx":
                doc = Document(rec.path)
                for pi, paragraph in enumerate(doc.paragraphs, 1):
                    for ri, run in enumerate(paragraph.runs, 1):
                        if not run.text.strip(): continue
                        if require_bold and not bool(run.bold): continue
                        if require_underline and not bool(run.underline): continue
                        if require_italic and not bool(run.italic): continue
                        hits.append((rec.relative_path, f"paragraph:{pi}/run:{ri}", run.text.strip()))
                for ti, table in enumerate(doc.tables, 1):
                    for rowi, row in enumerate(table.rows, 1):
                        for ci, cell in enumerate(row.cells, 1):
                            for pi, paragraph in enumerate(cell.paragraphs, 1):
                                for ri, run in enumerate(paragraph.runs, 1):
                                    if not run.text.strip(): continue
                                    if require_bold and not bool(run.bold): continue
                                    if require_underline and not bool(run.underline): continue
                                    if require_italic and not bool(run.italic): continue
                                    hits.append((rec.relative_path, f"table:{ti}/r:{rowi}/c:{ci}/p:{pi}/run:{ri}", run.text.strip()))
            elif rec.extension == ".pptx":
                prs = Presentation(rec.path)
                for si, slide in enumerate(prs.slides, 1):
                    for shapei, shape in enumerate(slide.shapes, 1):
                        if not getattr(shape, "has_text_frame", False): continue
                        for pi, paragraph in enumerate(shape.text_frame.paragraphs, 1):
                            for ri, run in enumerate(paragraph.runs, 1):
                                if not run.text.strip(): continue
                                if require_bold and not bool(run.font.bold): continue
                                if require_underline and not bool(run.font.underline): continue
                                if require_italic and not bool(run.font.italic): continue
                                hits.append((rec.relative_path, f"slide:{si}/shape:{shapei}/p:{pi}/run:{ri}", run.text.strip()))
        unique = []
        for item in hits:
            if item[2] not in [x[2] for x in unique]: unique.append(item)
        if unique:
            return ExecutionResult(True, "、".join(x[2] for x in unique), .97, "office_run_format_intersection", [Evidence(*x) for x in unique[:30]])
        return ExecutionResult.abstain("指定書式を同時に満たす文字列がない")

    def _comments(self, q, records, store):
        hits = []
        for rec in records:
            if rec.extension in {".xlsx", ".xlsm"}:
                wb = store.load_workbook(rec, data_only=False)
                for ws in wb.worksheets:
                    for row in ws.iter_rows():
                        for cell in row:
                            if cell.comment and cell.comment.text.strip():
                                hits.append((rec.relative_path, f"{ws.title}!{cell.coordinate}", cell.comment.text.strip()))
            else:
                for unit in store.extract_text_units(rec):
                    if unit.metadata.get("comment") and unit.text.strip():
                        hits.append((rec.relative_path, unit.locator, unit.text.strip()))
        if hits:
            return ExecutionResult(True, "、".join(dict.fromkeys(x[2] for x in hits)), .96, "office_comment_extract", [Evidence(*x) for x in hits])
        return ExecutionResult.abstain("コメントを検出できない")

    def _color_format(self, q, records, store):
        color = next((name for jp, name in (("黄色", "yellow"), ("青色", "blue"), ("赤", "red"), ("オレンジ", "orange")) if jp in q), "")
        hits = []
        for rec in records:
            if rec.extension in {".xlsx", ".xlsm"}:
                wb = store.load_workbook(rec, data_only=True)
                for ws in wb.worksheets:
                    for row in ws.iter_rows():
                        for cell in row:
                            if color and _is_color(_rgb(cell), color) and cell.value not in (None, ""):
                                hits.append((rec.relative_path, f"{ws.title}!{cell.coordinate}", str(cell.value)))
            elif rec.extension == ".pptx":
                prs = Presentation(rec.path)
                for si, slide in enumerate(prs.slides, 1):
                    for shapei, shape in enumerate(slide.shapes, 1):
                        if not getattr(shape, "has_text_frame", False): continue
                        shape_fill = self._pptx_fill_rgb(shape)
                        for pi, paragraph in enumerate(shape.text_frame.paragraphs, 1):
                            for ri, run in enumerate(paragraph.runs, 1):
                                text = run.text.strip()
                                if not text: continue
                                font_rgb = str(run.font.color.rgb) if run.font.color.type is not None and run.font.color.rgb else ""
                                if color == "red" and font_rgb and _is_color(font_rgb[-6:], "red"):
                                    hits.append((rec.relative_path, f"slide:{si}/shape:{shapei}/p:{pi}/run:{ri}", text))
                                elif color and shape_fill and _is_color(shape_fill, color):
                                    hits.append((rec.relative_path, f"slide:{si}/shape:{shapei}", text))
            elif rec.extension == ".docx":
                doc = Document(rec.path)
                for pi, paragraph in enumerate(doc.paragraphs, 1):
                    for ri, run in enumerate(paragraph.runs, 1):
                        text = run.text.strip()
                        if not text: continue
                        rgb = str(run.font.color.rgb) if run.font.color.rgb else ""
                        highlight = str(run.font.highlight_color or "").lower()
                        if color == "red" and rgb and _is_color(rgb[-6:], "red"):
                            hits.append((rec.relative_path, f"paragraph:{pi}/run:{ri}", text))
                        if color == "yellow" and "yellow" in highlight:
                            hits.append((rec.relative_path, f"paragraph:{pi}/run:{ri}", text))
        unique = []
        for item in hits:
            if item[2] not in [x[2] for x in unique]: unique.append(item)
        if unique:
            return ExecutionResult(True, "、".join(x[2] for x in unique), .91, "office_color_format_extract", [Evidence(*x) for x in unique])
        return ExecutionResult.abstain("指定色の文字・セル・図形を検出できない")

    @staticmethod
    def _pptx_fill_rgb(shape) -> str:
        try:
            fill = shape.fill
            if fill.type is not None and fill.fore_color.rgb:
                return str(fill.fore_color.rgb)[-6:]
        except Exception:
            pass
        return ""
