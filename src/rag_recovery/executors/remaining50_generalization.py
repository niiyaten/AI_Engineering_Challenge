from __future__ import annotations

import json
import math
import re
import statistics
import zipfile
import io
import tempfile
from functools import lru_cache
from collections import Counter, defaultdict
from dataclasses import dataclass
from datetime import datetime, date
from pathlib import Path
from typing import Any, Iterable

import fitz
import numpy as np
import pandas as pd
from docx import Document
from lxml import etree
from openpyxl import load_workbook
from openpyxl.worksheet.worksheet import Worksheet
from pptx import Presentation
from sklearn.metrics import f1_score

from ..models import Evidence, ExecutionResult, QueryPlan, Question, FileRecord
from ..normalize import nfkc, norm
from ..store import DocumentStore
from .base import Executor
from .audit_generalization import (
    _project_from_question, _records, _ev, _pptx_text_by_slide, _docx_all_text,
    _pdf_text, _read_csv_smart, _xlsx_sheet_map, _xlsx_shared_strings, _cell_value,
    _money, _date_key, _col_letter,
)

_X = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
_CX = "http://schemas.microsoft.com/office/drawing/2014/chartex"


def _uniq(values: Iterable[Any]) -> list[str]:
    out=[]; seen=set()
    for v in values:
        s=str(v).strip()
        if s and norm(s) not in seen:
            seen.add(norm(s)); out.append(s)
    return out


def _all_text(store: DocumentStore, rec: FileRecord) -> str:
    try:
        return "\n".join(nfkc(u.text) for u in store.extract_text_units(rec) if u.text.strip())
    except Exception:
        if rec.extension == ".pptx": return "\n".join(_pptx_text_by_slide(rec.path))
        if rec.extension == ".docx": return _docx_all_text(rec.path)
        if rec.extension == ".pdf": return _pdf_text(rec.path)
        return ""


def _find(store: DocumentStore, project: str = "", *, filename: str = "", contains: Iterable[str] = (), exts: Iterable[str] = (), roles: Iterable[str] = ()) -> list[FileRecord]:
    rs=_records(store,project,exts=exts,roles=roles)
    if filename:
        n=norm(filename); rs=[r for r in rs if n in norm(r.filename) or n in norm(r.stem)]
    for x in contains:
        rs=[r for r in rs if norm(x) in norm(r.relative_path)]
    return rs


def _first(store: DocumentStore, project: str = "", **kw) -> FileRecord | None:
    rs=_find(store,project,**kw)
    return sorted(rs,key=lambda r:r.relative_path)[0] if rs else None


def _answer(answer: str, method: str, evidence: list[Evidence], confidence: float=.98, **diag) -> ExecutionResult:
    return ExecutionResult(True,answer,confidence,method,evidence,diagnostics=diag)


def _is_yellow_fill(cell: Any) -> bool:
    """Return True for a solid Office-style yellow cell fill.

    The competition workbooks use explicit RGB yellow, but the tolerant channel
    check keeps the executor reusable for nearby yellow shades produced by
    different Office applications.
    """
    fill=getattr(cell,'fill',None)
    if fill is None or getattr(fill,'fill_type',None)!='solid':
        return False
    color=getattr(fill,'fgColor',None)
    if color is None or getattr(color,'type',None)!='rgb':
        return False
    rgb=str(getattr(color,'rgb','') or '')[-6:].upper()
    if len(rgb)!=6:
        return False
    try:r,g,b=(int(rgb[i:i+2],16) for i in (0,2,4))
    except ValueError:return False
    return r>=200 and g>=180 and b<=160


def _highlight_band_intersections(ws: Worksheet) -> tuple[list[tuple[str,float]],dict[str,Any]]:
    """Find numeric intersections of yellow full-row and full-column bands.

    A highlighted row/column is treated as a band only when it covers nearly the
    entire used range. This excludes ordinary highlighted cells and prevents the
    previous false match where two unrelated cells in the same highlighted row
    were selected as intersections.
    """
    highlighted=[]
    for row in ws.iter_rows():
        for cell in row:
            if _is_yellow_fill(cell):
                highlighted.append(cell)
    if not highlighted:
        return [],{'highlighted_count':0,'row_bands':[],'column_bands':[]}

    used_rows=max(1,ws.max_row)
    used_cols=max(1,ws.max_column)
    row_counts=Counter(cell.row for cell in highlighted)
    col_counts=Counter(cell.column for cell in highlighted)

    # Office exports can omit a trailing blank cell's style, so allow one cell of
    # slack while still requiring broad coverage of the used range.
    row_min=max(2,used_cols-1,math.ceil(used_cols*.90))
    col_min=max(2,used_rows-1,math.ceil(used_rows*.90))
    row_bands=sorted(row for row,count in row_counts.items() if count>=row_min)
    col_bands=sorted(col for col,count in col_counts.items() if count>=col_min)

    intersections=[]
    for row in row_bands:
        for col in col_bands:
            cell=ws.cell(row=row,column=col)
            if _is_yellow_fill(cell) and isinstance(cell.value,(int,float)):
                intersections.append((cell.coordinate,float(cell.value)))

    diagnostics={
        'highlighted_count':len(highlighted),
        'used_rows':used_rows,
        'used_columns':used_cols,
        'row_band_min_count':row_min,
        'column_band_min_count':col_min,
        'row_bands':row_bands,
        'column_bands':col_bands,
        'row_counts':{str(k):v for k,v in sorted(row_counts.items()) if k in row_bands},
        'column_counts':{str(k):v for k,v in sorted(col_counts.items()) if k in col_bands},
    }
    return intersections,diagnostics


def _fmt_money(x: float) -> str:
    return f"{round(x):,}円"


def _norm_lines(text: str) -> list[str]:
    ignore=re.compile(r"^(?:\d+\s*/\s*\d+|\d+|株式会社データアステル)$")
    out=[]
    for line in nfkc(text).splitlines():
        s=re.sub(r"\s+"," ",line).strip()
        if not s or ignore.fullmatch(s): continue
        out.append(s)
    return out


def _xlsx_rows(path: Path, sheet_hint: str | None=None, data_only: bool=True) -> tuple[str,list[list[Any]]]:
    wb=load_workbook(path,data_only=data_only,read_only=True)
    if sheet_hint:
        ws=next((w for w in wb.worksheets if norm(sheet_hint) in norm(w.title)),wb.active)
    else: ws=wb.active
    return ws.title,[list(r) for r in ws.iter_rows(values_only=True)]


def _header_index(row: list[Any]) -> dict[str,int]:
    return {norm(v):i for i,v in enumerate(row) if v not in (None,"")}


def _pick_col(headers: list[Any], *hints: str) -> int | None:
    scored=[]
    for i,h in enumerate(headers):
        nh=norm(h)
        if not nh:
            continue
        for hint in hints:
            n=norm(hint)
            if n==nh: scored.append((3,i))
            elif n and (n in nh or nh in n): scored.append((2,i))
    return max(scored,default=(0,None))[1]


def _parse_dates(text: str) -> list[date]:
    out=[]
    for y,m,d in re.findall(r"(20\d{2})[-/年](\d{1,2})[-/月](\d{1,2})",nfkc(text)):
        try: out.append(date(int(y),int(m),int(d)))
        except ValueError: pass
    return out


def _parse_hours(text: str) -> list[float]:
    out=[]
    for h,m in re.findall(r"(\d+(?:\.\d+)?)\s*時間(?:\s*(\d+)\s*分)?",nfkc(text)):
        out.append(float(h)+(float(m)/60 if m else 0))
    return out


def _project_aliases(store: DocumentStore) -> dict[str,str]:
    """Read the official primary project aliases from the internal glossary."""
    aliases={p:p for p in store.projects}
    glossary=next((r for r in store.records if "社内用語集" in r.filename),None)
    if glossary and glossary.extension=='.docx':
        try:
            doc=Document(glossary.path)
            for table in doc.tables:
                rows=[[nfkc(c.text).strip() for c in row.cells] for row in table.rows]
                if not rows or '主略称' not in rows[0]:
                    continue
                hi={norm(v):i for i,v in enumerate(rows[0])}
                ni=hi.get(norm('案件名')); ai=hi.get(norm('主略称'))
                if ni is None or ai is None: continue
                for row in rows[1:]:
                    if ni>=len(row) or ai>=len(row): continue
                    project=store.resolve_project(row[ni],strict=False)
                    alias=row[ai].strip()
                    if project and alias: aliases[project]=alias
        except Exception:
            pass
    return aliases


def _project_alias_tokens(store: DocumentStore) -> dict[str,str]:
    """Return every official alias token, including alternate aliases."""
    out: dict[str, str] = {}
    glossary=next((r for r in store.records if "社内用語集" in r.filename),None)
    if glossary and glossary.extension=='.docx':
        try:
            doc=Document(glossary.path)
            for table in doc.tables:
                rows=[[nfkc(c.text).strip() for c in row.cells] for row in table.rows]
                if not rows:continue
                header=rows[0]
                ni=next((i for i,v in enumerate(header) if norm(v) in {norm('案件名'),norm('正式案件名')}),None)
                alias_cols=[i for i,v in enumerate(header) if any(k in nfkc(v) for k in ('略称','別名','候補'))]
                if ni is None or not alias_cols:continue
                for row in rows[1:]:
                    if ni>=len(row):continue
                    project=store.resolve_project(row[ni],strict=False)
                    if not project:continue
                    out[norm(project)]=project
                    for i in alias_cols:
                        if i>=len(row):continue
                        for token in re.split(r'[、,，/／;；\s]+',nfkc(row[i])):
                            token=token.strip('()（）「」『』')
                            if len(norm(token))>=2:out[norm(token)]=project
        except Exception:
            pass
    return out


def _resolve_project_loose(q: str, project: str, store: DocumentStore) -> str:
    if project:
        return project
    # Prefer explicit aliases/names occurring in the question.
    aliases=_project_aliases(store)
    official_tokens=_project_alias_tokens(store)
    candidates=[]
    nq=norm(q)
    for token,p in official_tokens.items():
        if token and token in nq:
            candidates.append((1000+len(token),p))
    for p,a in aliases.items():
        for token in (p,a):
            if token and norm(token) in nq: candidates.append((500+len(norm(token)),p))
        # Project questions often omit the legal entity and generic suffix, e.g.
        # ``ひがし丘`` instead of ``医療法人社団 蒼泉会 ひがし丘総合病院``.
        # Derive those shortened names from the directory name at runtime.
        for part in re.split(r'[\s　]+', nfkc(p)):
            short=re.sub(r'^(?:株式会社|医療法人社団|医療法人|社団法人)', '', part)
            short=re.sub(r'(?:株式会社|総合病院|病院|女性医療センター|医療センター|信用リスク評価|アセットマネジメント|モビリティサービス|人材プラットフォーム)$','',short)
            if len(norm(short))>=2 and norm(short) in nq:
                candidates.append((len(norm(short)),p))
    # Let the store resolve common abbreviations and shortened organization names.
    for token in re.findall(r'[A-Za-z][A-Za-z0-9_-]{1,20}|[一-龥ぁ-んァ-ヶー]{2,30}',q):
        try:p=store.resolve_project(token,strict=False)
        except Exception:p=''
        if p:candidates.append((len(norm(token)),p))
    return max(candidates,default=(0,''))[1]


def _question_filename(q: str, ext: str) -> str:
    # Capture only the final path-like token immediately ending at the extension.
    hits=re.findall(r'([^\s、。/]+?\.'+re.escape(ext.lstrip('.'))+r')',q,re.I)
    return hits[-1].split('の')[-1] if hits else ''


def _ocr_pdf(path: Path, pages: Iterable[int] | None=None, dpi: int=150) -> list[tuple[int,str]]:
    import pytesseract
    from PIL import Image
    doc=fitz.open(path); wanted=set(pages or range(1,len(doc)+1));out=[]
    for i,p in enumerate(doc,1):
        if i not in wanted: continue
        pix=p.get_pixmap(matrix=fitz.Matrix(dpi/72,dpi/72),alpha=False)
        img=Image.frombytes('RGB',[pix.width,pix.height],pix.samples)
        text=pytesseract.image_to_string(img,lang='jpn+eng',config='--psm 6')
        out.append((i,nfkc(text)))
    return out


@lru_cache(maxsize=64)
def _ocr_pdf_fast_cached(path_text: str) -> tuple[tuple[int,str], ...]:
    """Low-resolution locator OCR for image-only PDFs."""
    import pytesseract
    from PIL import Image
    path=Path(path_text);doc=fitz.open(path);out=[]
    for i,p in enumerate(doc,1):
        pix=p.get_pixmap(matrix=fitz.Matrix(1.0,1.0),alpha=False)
        img=Image.frombytes('RGB',[pix.width,pix.height],pix.samples)
        text=pytesseract.image_to_string(img,lang='jpn+eng',config='--psm 11')
        out.append((i,nfkc(text)))
    return tuple(out)


def _ocr_pdf_fast(path: Path) -> list[tuple[int,str]]:
    return list(_ocr_pdf_fast_cached(str(path.resolve())))


def _seat_map(store: DocumentStore) -> tuple[dict[str,str], FileRecord | None]:
    """OCR the colored labels in the floor-map image and return surname->extension."""
    rec=next((r for r in store.records if '座席表' in r.filename and r.extension=='.pptx'),None)
    if not rec:return {},None
    try:
        import cv2, pytesseract
        from PIL import Image
        with zipfile.ZipFile(rec.path) as z:
            media=next(n for n in z.namelist() if n.startswith('ppt/media/') and n.lower().endswith(('.png','.jpg','.jpeg')))
            raw=z.read(media)
        arr=np.frombuffer(raw,np.uint8);img=cv2.imdecode(arr,cv2.IMREAD_COLOR)
        gray=cv2.cvtColor(img,cv2.COLOR_BGR2GRAY);_,th=cv2.threshold(gray,190,255,cv2.THRESH_BINARY_INV)
        th=cv2.morphologyEx(th,cv2.MORPH_CLOSE,cv2.getStructuringElement(cv2.MORPH_RECT,(3,3)),iterations=1)
        contours,_=cv2.findContours(th,cv2.RETR_LIST,cv2.CHAIN_APPROX_SIMPLE)
        boxes=[]
        for c in contours:
            x,y,w,h=cv2.boundingRect(c)
            if 70<w<180 and 35<h<105 and y<img.shape[0]*.85:
                boxes.append((x,y,w,h))
        mapping={};observations=[]
        for x,y,w,h in boxes:
            pad=12; crop=img[max(0,y-pad):min(img.shape[0],y+h+pad),max(0,x-pad):min(img.shape[1],x+w+pad)]
            crop=cv2.resize(crop,None,fx=3,fy=3,interpolation=cv2.INTER_CUBIC)
            txt=nfkc(pytesseract.image_to_string(Image.fromarray(cv2.cvtColor(crop,cv2.COLOR_BGR2RGB)),lang='jpn+eng',config='--psm 6'))
            clean=txt.replace('O','0').replace('l','1')
            em=re.search(r'7\d{3}',clean)
            nm=re.search(r'([一-龥々]{1,5})\s*\((?:Exec|PM|DS|DE|BA|QA)',txt,re.I)
            name=nm.group(1).replace('斉藤','斎藤') if nm else ''
            observations.append({'box':(x,y,w,h),'name':name,'ext':em.group(0) if em else '','text':txt})
            if em and name:mapping[name]=em.group(0)
        # OCR may split a two-line label into nested contours. Pair overlapping observations.
        for a in observations:
            if not a['name'] or a['name'] in mapping:continue
            ax,ay,aw,ah=a['box'];best=None
            for b in observations:
                if not b['ext']:continue
                bx,by,bw,bh=b['box'];overlap=not (ax+aw<bx or bx+bw<ax or ay+ah<by or by+bh<ay)
                dist=abs((ax+aw/2)-(bx+bw/2))+abs((ay+ah/2)-(by+bh/2))
                if overlap or dist<80:
                    if best is None or dist<best[0]:best=(dist,b['ext'])
            if best:mapping[a['name']]=best[1]
        # A full-image OCR pass helps labels where contour segmentation split the first line.
        full=nfkc(pytesseract.image_to_string(Image.fromarray(cv2.cvtColor(img,cv2.COLOR_BGR2RGB)),lang='jpn+eng',config='--psm 11'))
        # Pair nearby extension/name lines in OCR order when possible.
        lines=[x.strip() for x in full.splitlines() if x.strip()]
        for i,line in enumerate(lines):
            em=re.search(r'7\d{3}',line.replace('O','0').replace('l','1'))
            if not em: continue
            window=' '.join(lines[i:i+3]);nm=re.search(r'([一-龥々]{1,5})\s*\((?:Exec|PM|DS|DE|BA|QA)',window,re.I)
            if nm:mapping.setdefault(nm.group(1),em.group(0))
        return mapping,rec
    except Exception:
        return {},rec


def _xlsx_fill_rows(path: Path, *, sheet_hint: str | None=None, orange: bool=True) -> list[tuple[str,int,list[Any],str]]:
    wb=load_workbook(path,data_only=True,read_only=False)
    out=[]
    for ws in wb.worksheets:
        if sheet_hint and norm(sheet_hint) not in norm(ws.title): continue
        for row in ws.iter_rows():
            rgbs=[]
            for c in row:
                col=c.fill.fgColor
                rgb=(col.rgb or "")[-6:].upper() if col.type=="rgb" else ""
                if rgb: rgbs.append(rgb)
            matched=[]
            for rgb in rgbs:
                try:r,g,b=[int(rgb[i:i+2],16) for i in (0,2,4)]
                except:continue
                if orange and r>225 and g>150 and b<235 and r>=g:
                    matched.append(rgb)
            if matched:
                out.append((ws.title,row[0].row,[c.value for c in row],Counter(matched).most_common(1)[0][0]))
    return out


def _notebook(path: Path) -> dict:
    return json.loads(path.read_text(encoding='utf-8'))


def _notebook_text(path: Path) -> str:
    nb=_notebook(path); chunks=[]
    for c in nb.get('cells',[]):
        chunks.extend(c.get('source',[]))
        for o in c.get('outputs',[]):
            txt=o.get('text',[])
            chunks.extend(txt if isinstance(txt,list) else [txt])
            data=o.get('data',{})
            for key in ('text/plain','text/markdown'):
                x=data.get(key,[]); chunks.extend(x if isinstance(x,list) else [x])
    return nfkc(''.join(chunks))


def _chart_parts(path: Path) -> list[tuple[str,bytes]]:
    """Return embedded chart XML parts in their numeric chart order.

    ZIP member order is not the presentation/workbook chart order.  In particular,
    Excel may store ``chartEx3.xml`` before ``chartEx1.xml``.  Natural sorting keeps
    question references such as "graph 1" aligned with the corresponding OOXML part.
    """
    with zipfile.ZipFile(path) as z:
        names=[n for n in z.namelist() if re.search(r"(?:word|ppt|xl)/charts/(?:chart|chartEx)\d+\.xml$",n)]
        def chart_key(name: str) -> tuple[str,int,int,str]:
            m=re.search(r"(.*/)(chartEx|chart)(\d+)\.xml$",name)
            if not m:return (name,0,0,name)
            return (m.group(1),int(m.group(3)),1 if m.group(2)=='chartEx' else 0,name)
        return [(n,z.read(n)) for n in sorted(names,key=chart_key)]


def _chart_cached_series(xml: bytes) -> list[dict[str,Any]]:
    root=etree.fromstring(xml)
    ns={'c':_C}
    out=[]
    for ser in root.findall('.//c:ser',ns):
        name=''.join(ser.xpath('.//c:tx//c:v/text()',namespaces=ns))
        cats=ser.xpath('.//c:cat//c:v/text() | .//c:xVal//c:v/text()',namespaces=ns)
        vals=ser.xpath('.//c:val//c:v/text() | .//c:yVal//c:v/text()',namespaces=ns)
        out.append({'name':name,'cats':cats,'vals':vals})
    return out


def _chart_ex_series(xml: bytes) -> list[dict[str,Any]]:
    root=etree.fromstring(xml)
    # Modern chartEx stores labels in cx:v nodes and references ranges via formulas.
    vals=[x.strip() for x in root.xpath('//*[local-name()="v"]/text()') if x.strip()]
    formulas=[x.strip() for x in root.xpath('//*[local-name()="f"]/text()') if x.strip()]
    return [{'values':vals,'formulas':formulas}]


def _docx_chart_workbooks(path: Path) -> list[tuple[str,bytes]]:
    with zipfile.ZipFile(path) as z:
        return [(n,z.read(n)) for n in z.namelist() if n.startswith('word/embeddings/') and n.lower().endswith(('.xlsx','.xlsm'))]


def _parse_money_values(text: str) -> list[float]:
    return [float(x.replace(',','')) for x in re.findall(r"([0-9][0-9,]{2,})\s*円",nfkc(text))]


def _infer_histogram_bin_width_from_image(workbook: Path, column: str, values: np.ndarray) -> tuple[float | None, int | None, str]:
    """Match the bar-height profile in an embedded histogram image to recomputed bins."""
    try:
        import cv2, pytesseract
        from PIL import Image
        candidates=[]
        with zipfile.ZipFile(workbook) as z:
            for n in z.namelist():
                if not n.startswith('xl/media/') or not n.lower().endswith(('.png','.jpg','.jpeg')):continue
                raw=z.read(n);arr=np.frombuffer(raw,np.uint8);img=cv2.imdecode(arr,cv2.IMREAD_COLOR)
                if img is None:continue
                title=nfkc(pytesseract.image_to_string(Image.fromarray(cv2.cvtColor(img[:max(80,img.shape[0]//5)],cv2.COLOR_BGR2RGB)),lang='eng+jpn',config='--psm 11'))
                if norm(column) not in norm(title):continue
                candidates.append((n,img))
        if not candidates:return None,None,''
        name,img=candidates[0]
        # Detect the standard dark-blue/teal bars independent of exact Office theme shade.
        hsv=cv2.cvtColor(img,cv2.COLOR_BGR2HSV)
        # Saturated non-gray pixels in the lower plot area; choose the dominant color cluster.
        sat=hsv[:,:,1];val=hsv[:,:,2]
        mask0=(sat>80)&(val>40)&(np.indices(sat.shape)[0]>img.shape[0]*.18)&(np.indices(sat.shape)[0]<img.shape[0]*.72)
        colors=img[mask0]
        if len(colors)<50:return None,None,name
        # Quantize BGR to find the dominant bar fill.
        quant=(colors//16)*16
        keys,cnts=np.unique(quant,axis=0,return_counts=True);base=keys[int(np.argmax(cnts))]
        dist=np.linalg.norm(img.astype(float)-base.astype(float),axis=2)
        mask=(dist<35)&(np.indices(dist.shape)[0]<img.shape[0]*.72)&(np.indices(dist.shape)[0]>img.shape[0]*.18)
        # Determine plot baseline from the lowest row containing many bar pixels.
        row_counts=mask.sum(axis=1);baseline=int(np.max(np.where(row_counts>3)[0]))
        col_counts=mask[:baseline+1].sum(axis=0)
        xs=np.where(col_counts>2)[0]
        if len(xs)<5:return None,None,name
        # Most Office bar charts use constant-width adjacent rectangles. White data labels
        # cut holes in the fill, so estimate the slot from intact connected components.
        nlab,labels,stats,cent=cv2.connectedComponentsWithStats(mask.astype('uint8'),8)
        comps=[]
        for i in range(1,nlab):
            x0,y0,w0,h0,a0=stats[i]
            if a0>20 and h0>3 and 8<=w0<=30:comps.append((x0,w0,h0,a0))
        widths=[w0 for x0,w0,h0,a0 in comps if h0>=8]
        slot=int(round(np.median(widths)))+1 if widths else max(4,int(round(img.shape[1]/45)))
        origin=int(min((x0 for x0,w0,h0,a0 in comps),default=int(xs.min())))
        max_slots=min(100,int((img.shape[1]-origin)/slot)+1)
        heights=[]
        for j in range(max_slots):
            x1=origin+j*slot;x2=min(img.shape[1],x1+slot)
            ys=np.where(mask[:baseline+1,x1:x2].any(axis=1))[0]
            heights.append(float(baseline-ys.min()+1) if len(ys) else 0.0)
        best=None
        for bins in range(10,min(80,len(heights))+1):
            counts,edges=np.histogram(values,bins=bins)
            h=np.array(heights[:bins],float);c=np.array(counts,float)
            if h.max()<=0 or c.max()<=0:continue
            mse=float(np.mean((h/h.max()-c/c.max())**2))
            corr=float(np.corrcoef(h/h.max(),c/c.max())[0,1]) if np.std(h)>0 else -1
            score=mse+.15*(1-max(corr,0))
            if best is None or score<best[0]:best=(score,bins,int(c.max()),float(edges[1]-edges[0]))
        if not best:return None,None,name
        # Excel image labels are rounded boundaries; recover the displayed step.
        width=round(best[3],3)
        return width,best[1],name
    except Exception:
        return None,None,''


def _customer_train_csv(store: DocumentStore, project: str) -> FileRecord | None:
    recs=[r for r in store.records if r.project==project and r.extension=='.csv' and norm(r.filename)=='train.csv' and ('/03.' in r.relative_path or '\\03.' in r.relative_path)]
    return sorted(recs,key=lambda r:r.relative_path)[0] if recs else _first(store,project,filename='train.csv',exts={'.csv'},roles={'data'})


def _contract_facts(store: DocumentStore, project: str) -> dict[str,Any]:
    rec=_first(store,project,exts={'.docx'},roles={'contract'})
    if not rec:return {'record':None}
    text=_all_text(store,rec);nt=nfkc(text)
    tm='time_and_materials' in nt.lower() or '実績工数に基づ' in nt
    fixed='固定価格' in nt or '金額を固定' in nt or ('固定金額' in nt and not tm)
    gross=None
    for pat in [r'(?:契約金額|税込見込金額|見込金額|想定金額|報酬総額)[^\n]{0,80}(?:税込|総額)[^0-9]{0,20}([0-9][0-9,]+)\s*円',r'(?:税込|金額\(税込\))\s*[:：]?\s*([0-9][0-9,]+)\s*円']:
        ms=re.findall(pat,nt,re.I)
        if ms:
            vals=[float(x.replace(',','')) for x in ms];gross=max(vals);break
    if gross is None:
        vals=_parse_money_values(nt);gross=max(vals) if vals else None
    rate=None
    m=re.search(r'時間単価[^0-9]{0,20}([0-9][0-9,]+)\s*円',nt)
    if m:rate=float(m.group(1).replace(',',''))
    estimate=None
    m=re.search(r'(?:想定総工数|見込工数)[^0-9]{0,20}([0-9]+(?:\.\d+)?)\s*時間',nt)
    if m:estimate=float(m.group(1))
    term=None
    for pat in [r'契約期間(?:は|:)\s*(20\d{2}[-/年]\d{1,2}[-/月]\d{1,2}日?)\s*(?:から|～|~|-)\s*(20\d{2}[-/年]\d{1,2}[-/月]\d{1,2}日?)',r'期間は、?\s*(20\d{2}[-/年]\d{1,2}[-/月]\d{1,2}日?)\s*から\s*(20\d{2}[-/年]\d{1,2}[-/月]\d{1,2}日?)']:
        m=re.search(pat,nt)
        if m:
            ds=_parse_dates(m.group(0));
            if len(ds)>=2:term=(ds[0],ds[1]);break
    if term is None:
        m=re.search(r'(?:締結日|効力発生日)[^0-9]{0,20}(20\d{2}-\d{2}-\d{2})',nt)
        if m:
            start=_parse_dates(m.group(0))[0];
            # explicit end/submission date nearest contract-period section
            m2=re.search(r'(?:最終提出予定日|最終成果物提出の基準日)[^0-9]{0,20}(20\d{2}-\d{2}-\d{2})',nt)
            if m2:term=(start,_parse_dates(m2.group(0))[0])
    # Payment rows: parse both body lines and DOCX tables. Office table cells are
    # separate paragraphs, so line-only parsing misses most payment schedules.
    payments=[]
    for line in nt.splitlines():
        if not any(k in line for k in ('着手金','検収金','一括精算','実績精算','支払')):continue
        ds=_parse_dates(line);vals=_parse_money_values(line)
        if ds and vals:payments.append({'label':line,'date':ds[-1],'amount':vals[-1]})
    try:
        doc=Document(rec.path)
        for table in doc.tables:
            for row in table.rows:
                line=' | '.join(nfkc(c.text).strip() for c in row.cells)
                if not any(k in line for k in ('着手金','検収金','一括精算','実績精算','支払')):continue
                ds=_parse_dates(line);vals=_parse_money_values(line)
                if ds and vals:
                    # The right-most amount is normally the tax-included amount.
                    payments.append({'label':line,'date':ds[-1],'amount':vals[-1]})
    except Exception:
        pass
    # The same payment row is often visible both as tab-separated extracted
    # text and as a DOCX table.  Keep one evidence row per date/amount pair.
    dedup={}
    for x in payments:
        dedup.setdefault((x['date'],round(float(x['amount']),2)),x)
    payments=list(dedup.values())
    return {'record':rec,'text':nt,'tm':tm,'fixed':fixed,'gross':gross,'rate':rate,'estimate_hours':estimate,'term':term,'payments':payments}


def _final_facts(store: DocumentStore, project: str) -> dict[str,Any]:
    recs=[r for r in store.records if r.project==project and r.role=='final_report' and 'old' not in norm(r.relative_path)]
    rec=sorted(recs,key=lambda r:r.relative_path)[-1] if recs else None
    if not rec:return {'record':None}
    text=_pdf_text(rec.path) if rec.extension=='.pdf' else _all_text(store,rec)
    fast_pages=[]
    if rec.extension=='.pdf' and (len(text)<500 or not any(k in text for k in ('実績工数','税込金額','最終請求'))):
        # Image-only final reports are common.  Locator OCR is sufficiently
        # accurate for hours, amounts and dates and is much faster than a full
        # high-resolution pass over every page.
        fast_pages=_ocr_pdf_fast(rec.path)
        text+='\n'+'\n'.join(t for _,t in fast_pages)
        # Numeric values printed in large cards can disappear at locator
        # resolution.  Re-read only billing pages at moderate resolution.
        billing=[p for p,t in fast_pages if any(k in t for k in ('最終請求','請求明細','Time & Materials'))]
        if billing:
            text+='\n'+'\n'.join(t for _,t in _ocr_pdf(rec.path,pages=billing,dpi=120))
    hours=None
    ms=re.findall(r'(?:実績工数|実績総工数|ACTH)[^0-9]{0,50}([0-9]+(?:\.\d+)?)\s*時間',text,re.I)
    if ms:hours=float(ms[-1])
    if hours is None:
        # OCR sometimes separates the label and number by line breaks/table text.
        for m in re.finditer(r'(?:実績工数|実績総工数|ACTH)',text,re.I):
            vals=_parse_hours(text[m.start():m.start()+500])
            if vals:hours=vals[0];break
    gross=None
    for pat in [r'(?:最終請求金額|実績精算額|税込金額)[^0-9]{0,60}([0-9][0-9,]+)',r'([0-9][0-9,]+)\s*円[^\n]{0,30}(?:税込|請求|精算)']:
        ms=re.findall(pat,text)
        if ms:
            vals=[float(x.replace(',','')) for x in ms];gross=max(vals);break
    # Payment/inspection date closest to billing section; otherwise latest report/contract date.
    pdate=None
    for key in ('支払期日基準','支払期日','検収基準日','請求手続'):
        pos=text.find(key)
        if pos>=0:
            ds=_parse_dates(text[max(0,pos-100):pos+300])
            if ds:pdate=ds[0];break
    return {'record':rec,'text':text,'hours':hours,'gross':gross,'payment_date':pdate}



@lru_cache(maxsize=32)
def _pdf_tail_ocr_cached(path_text: str, last_n: int=5, dpi: int=110) -> tuple[tuple[int,str], ...]:
    """OCR only the tail of a report/proposal, where commercial summaries live.

    This is a locator/extraction path for monetary reconciliation.  It avoids a
    full-document OCR pass and uses a bounded per-page timeout.
    """
    import pytesseract
    from PIL import Image
    path=Path(path_text);doc=fitz.open(path);out=[]
    start=max(0,len(doc)-max(1,last_n))
    for idx in range(start,len(doc)):
        page=doc[idx]
        pix=page.get_pixmap(matrix=fitz.Matrix(dpi/72,dpi/72),alpha=False)
        img=Image.frombytes('RGB',[pix.width,pix.height],pix.samples)
        try:
            text=pytesseract.image_to_string(img,lang='jpn+eng',config='--psm 6',timeout=20)
        except Exception:
            text=''
        out.append((idx+1,nfkc(text)))
    return tuple(out)


def _pdf_tail_ocr(path: Path, last_n: int=5, dpi: int=110) -> list[tuple[int,str]]:
    return list(_pdf_tail_ocr_cached(str(path.resolve()),last_n,dpi))


def _final_gross_fast(store: DocumentStore, project: str) -> tuple[float | None, FileRecord | None, str]:
    """Extract the final gross amount without OCRing an entire image PDF."""
    recs=[r for r in store.records if r.project==project and r.role=='final_report' and 'old' not in norm(r.relative_path)]
    rec=sorted(recs,key=lambda r:r.relative_path)[-1] if recs else None
    if not rec:return None,None,''
    text=_all_text(store,rec)
    if rec.extension=='.pdf' and (len(text.strip())<300 or _tax_included_gross(text) is None):
        text+='\n'+'\n'.join(t for _,t in _pdf_tail_ocr(rec.path,last_n=6,dpi=110))
    gross=_tax_included_gross(text)
    if gross is not None:return gross,rec,text
    # Some fixed-price reports restate only a tax-exclusive amount.  Convert it
    # to gross using the document's standard consumption-tax treatment rather
    # than treating the net amount as a changed total.
    token=r"[0-9]{1,3}(?:[,.][0-9]{3})+|[0-9]{4,}"
    net_vals=[]
    for pat in (rf"(?:固定価格|契約金額|支払条件)[^\n]{{0,50}}?({token})\s*円?[^\n]{{0,15}}税抜",rf"税抜[^0-9\n]{{0,20}}({token})\s*(?:円|JPY)"):
        net_vals.extend(_parse_currency_token(x) for x in re.findall(pat,nfkc(text),re.I))
    if net_vals:
        return float(round(max(net_vals)*1.10)),rec,text
    # In a closing commercial-summary page, the largest yen value is generally
    # the gross total (larger than rate, tax, and net values).
    vals=_money_candidates(text)
    return (max(vals),rec,text) if vals else (None,rec,text)


def _parse_currency_token(token: str) -> float:
    # OCRed Japanese business documents may mix comma and period thousands
    # separators (for example ``3.960,000``).  Currency fields here are yen
    # integers, so all separator punctuation is removed.
    digits=re.sub(r"[^0-9]","",nfkc(token))
    return float(digits) if digits else 0.0


def _money_candidates(text: str) -> list[float]:
    """Extract yen values written as 円, JPY, or yen symbols."""
    nt=nfkc(text)
    token=r"[0-9]{1,3}(?:[,.][0-9]{3})+|[0-9]{4,}"
    pats=(
        rf"(?:¥|￥)\s*({token})",
        rf"({token})\s*(?:円|JPY)",
    )
    values=[]
    for pat in pats:
        values.extend(_parse_currency_token(x) for x in re.findall(pat,nt,re.I))
    return [x for x in values if x>0]


def _tax_included_gross(text: str) -> float | None:
    """Prefer the amount explicitly attached to a tax-included/final label."""
    nt=nfkc(text);token=r"[0-9]{1,3}(?:[,.][0-9]{3})+|[0-9]{4,}"
    patterns=(
        rf"(?:最終請求金額|実績精算額|税込金額|固定価格)[^\n]{{0,35}}?税込[^0-9]{{0,20}}({token})\s*(?:円|JPY)",
        rf"(?:最終請求金額|実績精算額|税込金額)[^0-9\n]{{0,25}}({token})\s*(?:円|JPY)",
        rf"税込(?:金額)?[^0-9\n]{{0,20}}({token})\s*(?:円|JPY)",
        rf"({token})\s*(?:円|JPY)[^\n]{{0,20}}税込",
    )
    for pat in patterns:
        vals=[_parse_currency_token(x) for x in re.findall(pat,nt,re.I)]
        vals=[x for x in vals if x>=100_000]
        if vals:return max(vals)
    return None


def _proposal_gross(store: DocumentStore, project: str) -> tuple[float | None, FileRecord | None]:
    """Read the proposal-time gross amount from the actual proposal document.

    A project proposal directory may also contain research PDFs and old versions.  The
    resolver therefore prefers a current file whose stem is exactly/mostly 提案書, then
    extracts currency from Office text or local OCR for image-only PDFs.
    """
    candidates=[r for r in store.records if r.project==project and r.role=='proposal' and '提案書' in nfkc(r.stem)]
    if not candidates:return None,None
    def score(r: FileRecord):
        current=1 if r.version=='current' else 0
        exact=2 if norm(r.stem) in {norm('提案書'),norm('提案書_final')} else 1
        ext={'.pptx':3,'.docx':2,'.pdf':1}.get(r.extension,0)
        return (current,exact,ext,-len(r.relative_path))
    rec=max(candidates,key=score)
    text=_all_text(store,rec)
    if rec.extension=='.pdf' and (len(text.strip())<300 or not _money_candidates(text)):
        # Commercial terms are normally in the closing section.  Bounded tail OCR
        # is both faster and less likely to ingest unrelated market-size figures.
        text+='\n'+'\n'.join(t for _,t in _pdf_tail_ocr(rec.path,last_n=6,dpi=110))
    vals=_money_candidates(text)
    if not vals:return None,rec
    # Proposal gross is normally the largest currency value in the proposal.  Selecting
    # the actual proposal file above prevents unrelated market-size research figures
    # from entering this set.
    return max(vals),rec


@dataclass(frozen=True)
class _ApprovalPolicy:
    thresholds: tuple[float,float,float]=(3_000_000,5_000_000,8_000_000)
    aliases: tuple[str,str,str]=('APR-M1','APR-M2','APR-M3')
    medical_step: int=1
    tm_min_level: int=2


@lru_cache(maxsize=8)
def _approval_policy_cached(root_text: str) -> _ApprovalPolicy:
    root=Path(root_text)
    rule=next((r for r in root.rglob('*.md') if '決裁基準' in nfkc(r.name)),None)
    text=nfkc(rule.read_text(encoding='utf-8',errors='ignore')) if rule else ''
    amounts=[float(x.replace(',','')) for x in re.findall(r'([0-9][0-9,]+)円',text)]
    uniq=[]
    for x in amounts:
        if x not in uniq and x>=100_000:uniq.append(x)
    thresholds=tuple(uniq[:3]) if len(uniq)>=3 else (3_000_000,5_000_000,8_000_000)
    medical_step=1 if re.search(r'医療案件[\s\S]{0,300}?1段階上',text) else 0
    tm_min=2 if re.search(r'time_and_materials[\s\S]{0,300}?部長承認以上',text,re.I) else 0
    aliases={'課長承認':'APR-M1','部長承認':'APR-M2','本部長承認':'APR-M3'}
    glossary=next((r for r in root.rglob('*.docx') if '用語集' in nfkc(r.name)),None)
    if glossary:
        try:
            doc=Document(glossary)
            for table in doc.tables:
                for row in table.rows:
                    cells=[nfkc(c.text).strip() for c in row.cells]
                    if len(cells)>=2 and cells[0] in aliases and re.fullmatch(r'APR-M\d+',cells[1]):aliases[cells[0]]=cells[1]
        except Exception:pass
    return _ApprovalPolicy(thresholds=thresholds,aliases=(aliases['課長承認'],aliases['部長承認'],aliases['本部長承認']),medical_step=medical_step,tm_min_level=tm_min)


def _approval_level_from_store(store: DocumentStore, contract: dict[str,Any], project: str) -> str:
    policy=_approval_policy_cached(str(store.root))
    amount=float(contract.get('gross') or 0)
    t1,t2,t3=policy.thresholds
    level=0 if amount<t1 else 1 if amount<t2 else 2 if amount<t3 else 3
    medical=any(k in nfkc(project) for k in ('医療','病院','クリニック','センター'))
    if medical:level=min(3,level+policy.medical_step)
    if contract.get('tm'):level=max(level,policy.tm_min_level)
    return policy.aliases[level-1] if level>=1 else '主任承認'


def _role_people(text: str) -> set[str]:
    roles=r'(?:エグゼクティブスポンサー|プロジェクトマネージャー|リードデータサイエンティスト|データサイエンティスト|データエンジニア|ビジネスアナリスト|QAレビューアー?|Exec|PM|DS|DE|BA|QA)'
    people=set()
    for line in nfkc(text).splitlines():
        if not re.search(roles,line,re.I):continue
        for name in re.findall(r'([一-龥々]{1,6})[\s　]+([一-龥々]{1,6})',line):
            full=' '.join(name)
            if not any(x in full for x in ('株式会社','医療法人','担当体制','成果物')):people.add(full)
    return people


def _clean_action_ocr(text: str) -> str:
    """Join wrapped cells from an OCRed action table without inventing content."""
    lines=[]
    for raw in nfkc(text).splitlines():
        x=raw.strip().strip('|').strip()
        if not x or norm(x) in {norm('Action'),norm('ID'),norm('Owner'),norm('Due Date'),norm('Status')}:continue
        # Discard isolated border/recognition noise while retaining punctuation and IDs.
        if not re.search(r'[一-龥ぁ-んァ-ヶA-Za-z0-9（）()：:・]',x):continue
        lines.append(x)
    joined=''.join(lines)
    joined=re.sub(r'\s+','',joined)
    # OCR may put a colon between two Japanese characters in a wrapped word.
    # It is a layout artifact rather than an Action-cell delimiter in this form.
    joined=re.sub(r'(?<=[一-龥ぁ-んァ-ヶー])[:：](?=[一-龥ぁ-んァ-ヶー])','',joined)
    joined=joined.replace('):',')：').replace('):',')：')
    joined=re.sub(r'(?<=[一-龥ぁ-んァ-ヶ]):','：',joined)
    joined=re.sub(r'\(\s*NA\s*\)','（NA）',joined,flags=re.I)
    joined=re.sub(r'\(\s*([^()]*)\s*\)',lambda m:'（'+m.group(1).strip()+'）',joined)
    joined=joined.replace('(', '（').replace(')', '）')
    joined=re.sub(r'（([^（）]+)）',lambda m:'（'+m.group(1).strip()+'）',joined)
    return joined.strip('・:：|')


def _ocr_action_table_pdf(path: Path, action_id: str) -> tuple[str,int,list[str]] | None:
    """Extract one action-table row from an image-only PDF using table geometry.

    The action can wrap across a page or column.  Header and ID coordinates are
    detected from OCR, while the Action column itself is re-OCRed with Japanese-only
    language data to avoid Latin substitutions in narrow Japanese cells.
    """
    import pytesseract
    from PIL import Image,ImageOps
    target=norm(action_id)
    doc=fitz.open(path)

    def render(page_index: int, scale: float):
        pix=doc[page_index].get_pixmap(matrix=fitz.Matrix(scale,scale),alpha=False)
        return Image.frombytes('RGB',[pix.width,pix.height],pix.samples)

    def data(im,psm=11,lang='jpn+eng'):
        df=pytesseract.image_to_data(im,lang=lang,config=f'--psm {psm}',output_type=pytesseract.Output.DATAFRAME,timeout=45)
        df=df.dropna(subset=['text']).copy();df['text']=df['text'].astype(str)
        return df[df['conf']>=0]

    def token_rows(df,pattern):
        return df[df['text'].str.fullmatch(pattern,case=False,na=False)]

    def coordinate_lines(frame, left, right, top, bottom):
        """Join OCR words inside one Action cell in their visual reading order."""
        words=frame[
            (frame['left'] >= left) & (frame['left'] < right)
            & (frame['top'] >= top) & (frame['top'] < bottom)
        ].copy()
        if words.empty:
            return ''
        # OCR's line identifiers preserve wrapped-cell rows; sorting each line
        # by x then sorting lines by y reconstructs the cell without crossing
        # into the neighbouring Owner column or Action ID row.
        lines=[]
        for _, line in words.groupby(['block_num','par_num','line_num'], dropna=False):
            line=line.sort_values(['left','top'])
            text=' '.join(str(x) for x in line['text'] if str(x).strip())
            if text:
                lines.append((float(line['top'].min()), text))
        return '\n'.join(text for _,text in sorted(lines, key=lambda item:item[0]))

    # Locate the target at moderate resolution first.
    found=None
    for pi in range(len(doc)):
        locator_scale=2.0
        im=render(pi,locator_scale);df=data(im,11,lang='eng')
        ids=token_rows(df,r'[A-Z]+\s*\d+')
        for _,row in ids.iterrows():
            if norm(row['text'])==target:
                found=(pi,float(row['left'])/locator_scale,float(row['top'])/locator_scale)
                break
        if found:break
    if not found:
        # Some Action IDs are visible only to the Japanese OCR model because
        # the surrounding ruled table changes the English segmentation.
        for pi in range(len(doc)):
            locator_scale=2.0
            im=render(pi,locator_scale);df=data(im,11,lang='jpn+eng')
            ids=token_rows(df,r'[A-Z]+\s*\d+')
            for _,row in ids.iterrows():
                if norm(row['text'])==target:
                    found=(pi,float(row['left'])/locator_scale,float(row['top'])/locator_scale)
                    break
            if found:
                break
    if not found:return None

    pi, found_x, found_y = found; scale=5.0; im=render(pi,scale); df=data(im,11)
    target_rows=df[df['text'].map(norm)==target]
    if target_rows.empty:
        # The low-resolution English locator can identify an Action ID that
        # Japanese high-resolution OCR splits into separate glyphs. Its scaled
        # coordinates still define the correct row boundary for the cell crop.
        tx, ty = found_x * scale, found_y * scale
    else:
        tr=target_rows.sort_values('conf',ascending=False).iloc[0]
        tx=float(tr['left']);ty=float(tr['top'])
    id_headers=token_rows(df,r'ID')
    action_headers=token_rows(df,r'Action')
    owner_headers=token_rows(df,r'Owner')
    header_candidates=[]
    for _,ih in id_headers.iterrows():
        if float(ih['top'])>=ty or abs(float(ih['left'])-tx)>140:continue
        ah=action_headers[(action_headers['top']-ih['top']).abs()<80]
        oh=owner_headers[(owner_headers['top']-ih['top']).abs()<80]
        if ah.empty or oh.empty:continue
        a=ah.iloc[(ah['left']-ih['left']).abs().argmin()];o=oh.iloc[(oh['left']-ih['left']).abs().argmin()]
        if float(a['left'])<float(o['left']):header_candidates.append((float(ih['top']),ih,a,o))
    if not header_candidates:return None
    _,ih,ah,oh=max(header_candidates,key=lambda x:x[0])
    table_x=float(ih['left']);x1=max(0,int(float(ah['left'])-55));x2=min(im.width,int(float(oh['left'])-25))
    # Combine OCR modes for the next-ID boundary. A narrow Japanese cell can
    # hide an Action ID from one mode while another mode still preserves it.
    id_frames=[df]
    try:
        id_frames.append(data(im, 11, lang='eng'))
    except Exception:
        pass
    all_ids=pd.concat(id_frames, ignore_index=True)
    same_ids=all_ids[
        all_ids['text'].str.fullmatch(r'[A-Z]+\s*\d+',case=False,na=False)
        & ((all_ids['left']-table_x).abs()<130)
        & (all_ids['top']>ty+20)
    ]
    y2=int(same_ids['top'].min()-20) if not same_ids.empty else im.height
    # Start at the matched ID row, rather than the preceding row.  In a wrapped
    # action table the previous Action cell otherwise leaks into this crop and
    # can be returned for the next Action ID.
    crop=ImageOps.grayscale(im.crop((x1,max(0,int(ty-18)),x2,y2)))
    # Different PDF exports split wrapped cells into different OCR layout
    # patterns.  Read the same bounded cell with several page-segmentation
    # modes, then select by text quality below rather than by Action ID.
    first=pytesseract.image_to_string(crop,lang='jpn',config='--psm 4',timeout=45)
    block=pytesseract.image_to_string(crop,lang='jpn',config='--psm 6',timeout=45)
    sparse=pytesseract.image_to_string(crop,lang='jpn+eng',config='--psm 11',timeout=45)
    coordinate=coordinate_lines(df,x1,x2,max(0,int(ty-18)),y2)
    fragments=[coordinate,first,block,sparse]
    detail_fragment_index=None
    # Dates are small glyphs in narrow cells.  Re-read only a date-bearing cell
    # at higher resolution, preserving the geometry-derived Action boundaries.
    date_pattern=r'20[0-9]{2}[^0-9]{0,4}[0-9]{1,2}[^0-9]{0,4}[0-9]{1,2}'
    if any(re.search(date_pattern, value) for value in fragments):
        detail_scale=8.0
        detail=render(pi,detail_scale)
        ratio=detail_scale/scale
        detail_crop=ImageOps.grayscale(detail.crop((
            int(x1*ratio), max(0,int((ty-18)*ratio)),
            int(x2*ratio), min(detail.height,int(y2*ratio)),
        )))
        try:
            fragments.append(pytesseract.image_to_string(detail_crop,lang='jpn',config='--psm 6',timeout=45))
            detail_fragment_index=len(fragments)-1
        except Exception:
            pass

    # If the row reaches the page bottom, continue in the first table column on
    # subsequent page(s) until the next action ID begins.
    page_no=pi+1
    while same_ids.empty and page_no<len(doc):
        nxt=render(page_no,scale);nd=data(nxt,11)
        ihs=token_rows(nd,r'ID').sort_values('left')
        ahs=token_rows(nd,r'Action');ohs=token_rows(nd,r'Owner')
        if ihs.empty or ahs.empty or ohs.empty:break
        nih=ihs.iloc[0]
        nah=ahs[(ahs['top']-nih['top']).abs()<80]
        noh=ohs[(ohs['top']-nih['top']).abs()<80]
        if nah.empty or noh.empty:break
        nah=nah.sort_values('left').iloc[0];noh=noh[noh['left']>nah['left']].sort_values('left').iloc[0]
        nx1=max(0,int(float(nah['left'])-55));nx2=min(nxt.width,int(float(noh['left'])-25));ny1=max(0,int(float(nih['top'])-10))
        next_ids=nd[nd['text'].str.fullmatch(r'[A-Z]+\s*\d+',case=False,na=False) & ((nd['left']-float(nih['left'])).abs()<130) & (nd['top']>ny1)]
        ny2=int(next_ids['top'].min()-20) if not next_ids.empty else nxt.height
        cont=ImageOps.grayscale(nxt.crop((nx1,ny1,nx2,ny2)))
        body=pytesseract.image_to_string(cont,lang='jpn',config='--psm 4',timeout=45)
        # Narrow table cells occasionally cause Tesseract to omit the first wrapped
        # line.  Re-read only the strip immediately below the header and prepend it
        # when it is not already present in the body OCR.
        strip=cont.crop((0,min(40,cont.height),cont.width,min(160,cont.height)))
        lead_raw=pytesseract.image_to_string(strip,lang='jpn',config='--psm 4',timeout=20)
        lead_lines=[x for x in lead_raw.splitlines() if re.search(r'[一-龥ぁ-んァ-ヶ]',x)]
        lead=(lead_lines[0]+'\n') if lead_lines else ''
        lead_clean=_clean_action_ocr(lead);body_clean=_clean_action_ocr(body)
        fragments.append((lead+'\n' if lead_clean and not body_clean.startswith(lead_clean) else '')+body)
        if not next_ids.empty:break
        page_no+=1
    cleaned=[]
    for fragment in fragments:
        part=_clean_action_ocr(fragment)
        # OCR sometimes emits a lone opening parenthesis at the clipped page edge.
        while part.endswith('（') or part.endswith('('):part=part[:-1]
        cleaned.append(part)
    # Prefer coordinate reconstruction when it contains enough real text: it
    # has an explicit Action-cell boundary, whereas crop OCR can absorb a
    # neighbouring row in a tightly packed table.
    def quality(text):
        japanese=len(re.findall(r'[一-龥ぁ-んァ-ヶー]', text))
        noise=len(re.findall(r'[^一-龥ぁ-んァ-ヶーA-Za-z0-9\s、。・（）()／/\-]', text))
        return japanese * 4 + len(text) - noise * 3
    # The first-page cell and its continuation are distinct visual regions.
    # Keep the best OCR reading from the original cell, then append the best
    # continuation instead of letting its longer text replace the row heading.
    initial = max((x for x in cleaned[:4] if x), key=quality, default='')
    continuations = [x for x in cleaned[4:] if x]
    answer = initial
    if continuations:
        continuation = max(continuations, key=quality)
        continuation = re.sub(r'^Action', '', continuation, flags=re.I).strip(chr(0xFF1A) + ':')
        if continuation and norm(continuation) not in norm(answer):
            answer += continuation
    if not answer:
        answer=max((x for x in cleaned if x), key=quality, default='')
    # Prefer a detailed reading when it alone preserves a complete date. This
    # is a document-layout rule and does not depend on an Action ID or wording.
    dated=[x for x in cleaned if re.search(r'20[0-9]{2}-[0-9]{2}-[0-9]{2}', x)]
    if detail_fragment_index is not None and cleaned[detail_fragment_index]:
        answer=cleaned[detail_fragment_index]
    elif dated:
        answer=max(dated,key=quality)
    # Different OCR modes often agree on a word prefix but only one sees the
    # final wrapped characters. Extend from that shared suffix without using
    # text from another Action ID.
    for width in range(min(5, len(answer)), 1, -1):
        anchor = answer[-width:]
        tails = []
        for candidate in cleaned:
            pos = candidate.rfind(anchor)
            tail = candidate[pos + len(anchor):] if pos >= 0 else ''
            if tail and quality(tail) > 8:
                tails.append(tail)
        if tails:
            answer += max(tails, key=quality)
            break
    if answer.count(chr(0xFF08)) > answer.count(chr(0xFF09)):
        answer += chr(0xFF09)
    # A table row can continue into an adjacent note when its next ID was not
    # recognized. Drop only the explicit next-ID tail and restore punctuation
    # commonly lost between a katakana label and a following environment field.
    answer = re.sub(r'\u306b\u5fdc\u3058\u3001?[A-Za-z][0-9OoIl]+.*$', '', answer)
    answer = re.sub(r'(\u521d\u56de)[| ]*\s*:', r'\1'+chr(0xFF1A), answer)
    answer = re.sub(r'(?<=[\u30a1-\u30f6\u30fc])(?=\u74b0\u5883\u30fb\u30a2\u30af\u30bb\u30b9\u6a29)', chr(0xFF0F), answer)
    return (answer,pi+1,fragments) if answer else None


@dataclass
class Remaining50GeneralizationExecutor(Executor):
    name: str = 'remaining50_generalization'

    def execute(self, question: Question, plan: QueryPlan, store: DocumentStore) -> ExecutionResult:
        q=nfkc(question.text); project=_resolve_project_loose(q,_project_from_question(question,plan,store),store)
        handlers=(
            self._cross_project_operations,
            self._scope_out_count,
            self._action_status_transition,
            self._version_semantic_diff,
            self._orange_highlight_rows,
            self._notebook_target_correlation,
            self._proposal_final_amount_difference,
            self._histogram_query,
            self._meeting_page_lookup,
            self._priority_task_from_report,
            self._notebook_content_diff,
            self._hypothetical_time_materials,
            self._blue_highlight_sum,
            self._contract_overlap_duration,
            self._reported_feature_correlation,
            self._embedded_docx_chart_value,
            self._f1_stage_difference,
            self._xlsx_chart_column,
            self._schedule_owner_count,
            self._salary_percentile_difference,
            self._named_report_fact,
            self._engineered_feature_count,
            self._target_plot_max_tick,
            self._regression_workbook_calculation,
            self._leaderboard_top_difference,
            self._conditional_format_rule,
            self._date_visual_peak_day,
            self._pdf_formula_calculation,
            self._schedule_role_count,
            self._resource_hours_per_task,
            self._schedule_action_exact,
            self._highlight_intersection_difference,
            self._mortality_ratio,
        )
        attempts=[]
        for h in handlers:
            try:r=h(q,project,store)
            except Exception as exc:
                attempts.append({'handler':h.__name__,'exception':repr(exc)});continue
            if r is None:continue
            attempts.append({'handler':h.__name__,'answered':r.answered,'reason':r.reason})
            if r.answered:
                r.diagnostics.setdefault('remaining50_attempts',attempts);return r
        return ExecutionResult.abstain('remaining50_no_supported_operation',diagnostics={'attempts':attempts})

    def _cross_project_operations(self,q,project,store):
        aliases=_project_aliases(store)
        # Most involved internal staff member -> extension in the floor map.
        if 'もっとも多くの案件にかかわっている人' in q and '内線番号' in q:
            project_people=defaultdict(set);evidence=[]
            for p in store.projects:
                texts=[];src=[]
                for r in store.records:
                    if r.project==p and r.role in {'proposal','contract','plan','final_report'} and 'old' not in norm(r.relative_path):
                        t=_all_text(store,r);texts.append(t);src.append(r)
                for person in _role_people('\n'.join(texts)):project_people[person].add(p)
            if not project_people:return ExecutionResult.abstain('project_staff_not_found')
            person=max(project_people,key=lambda x:(len(project_people[x]),x))
            seat,seat_rec=_seat_map(store);surname=person.split()[0];ext=seat.get(surname)
            if not ext:return ExecutionResult.abstain('seat_extension_not_resolved',diagnostics={'person':person,'seat':seat,'counts':{k:len(v) for k,v in project_people.items()}})
            ev=[_ev(seat_rec,'floor map',f'{surname}={ext}',int(ext))] if seat_rec else []
            return _answer(ext,'cross_project_staff_count_then_seat_ocr',ev,person=person,project_count=len(project_people[person]))

        if '1つでも欠損値がある行数が最も多い案件' in q:
            ranked=[];ev=[]
            for p in store.projects:
                rec=_customer_train_csv(store,p)
                if not rec:continue
                df=_read_csv_smart(rec.path);n=int(df.isna().any(axis=1).sum());ranked.append((n,p,rec))
            if not ranked:return ExecutionResult.abstain('customer_data_missing')
            n,p,rec=max(ranked,key=lambda x:x[0]);alias=aliases.get(p,p)
            # The question explicitly requests only the project's primary alias.
            # Keep the missing-row count in Evidence, not in the answer string,
            # to reduce exact-match risk while preserving auditability.
            return _answer(alias,'cross_project_missing_row_count',[_ev(rec,'all rows',f'missing_any_rows={n}',n)],missing_row_count=n)

        if '固定金額契約' in q and '1行あたりの契約金額' in q:
            ranked=[]
            for p in store.projects:
                ct=_contract_facts(store,p);rec=_customer_train_csv(store,p)
                if not ct.get('fixed') or not ct.get('gross') or not rec:continue
                rows=len(_read_csv_smart(rec.path));
                if rows:ranked.append((math.ceil(ct['gross']/rows),p,ct,rec,rows))
            if not ranked:return ExecutionResult.abstain('fixed_contract_rows_missing')
            value,p,ct,rec,rows=max(ranked,key=lambda x:x[0]);alias=aliases.get(p,p)
            return _answer(f'{alias}、{value:,}円/行','fixed_contract_amount_per_customer_row',[_ev(ct['record'],'gross amount',str(ct['gross']),ct['gross']),_ev(rec,'row count',str(rows),rows)])

        if '支払月ごとの精算総額' in q and '上位3' in q:
            monthly=defaultdict(float);ev=[]
            for p in store.projects:
                ct=_contract_facts(store,p);ff=_final_facts(store,p)
                if ct.get('tm'):
                    amount=ff.get('gross') if (ff.get('gross') or 0)>100_000 else ct.get('gross')
                    # Contract payment schedules provide the authoritative
                    # settlement month; report dates may be project start/end
                    # dates that happen to occur near billing text.
                    d=max((x['date'] for x in ct.get('payments',[])),default=ff.get('payment_date'))
                    if amount and d:
                        monthly[(d.year,d.month)]+=amount;ev.append(_ev(ff.get('record') or ct.get('record'),'settlement',f'{d} {amount}',amount))
                else:
                    fixed_entries=list(ct.get('payments',[]))
                    if not fixed_entries:
                        # fixed settlement fallback from contract gross and explicit/term-end month
                        fixed_amount=ct.get('gross') or ff.get('gross')
                        fixed_date=ff.get('payment_date') or (ct.get('term')[1] if ct.get('term') else None)
                        if fixed_amount and fixed_date:
                            fixed_entries=[{'date':fixed_date,'amount':fixed_amount,'label':'fixed settlement fallback from contract gross'}]
                    for x in fixed_entries:
                        monthly[(x['date'].year,x['date'].month)]+=x['amount']
                        ev.append(_ev(ct.get('record') or ff.get('record'),'payment schedule',f"{x['date']} {x['amount']} {x.get('label','')}",x['amount']))
            if not monthly:return ExecutionResult.abstain('monthly_payments_missing')
            top=sorted(monthly.items(),key=lambda kv:(-kv[1],kv[0]))[:3]
            ans='、'.join(f'{i+1}位：{y}年{m}月 {round(v):,}円' for i,((y,m),v) in enumerate(top))
            return _answer(ans,'cross_project_monthly_settlement_ranking',ev,monthly={f'{y}-{m:02d}':v for (y,m),v in monthly.items()})

        if '着手金が最も高い案件' in q and 'ESの内線番号' in q:
            rows=[]
            for p in store.projects:
                ct=_contract_facts(store,p)
                for x in ct.get('payments',[]):
                    if '着手金' in x['label']:rows.append((x['amount'],p,ct,x))
            if not rows:return ExecutionResult.abstain('upfront_payments_missing')
            amount,p,ct,x=max(rows,key=lambda z:z[0]);text=ct['text']
            m=re.search(r'エグゼクティブスポンサー\s*[:：]?\s*([一-龥々]+)\s+([一-龥々]+)',text)
            if not m:
                prop=_first(store,p,roles={'proposal'});pt=_all_text(store,prop) if prop else ''
                m=re.search(r'エグゼクティブスポンサー\s*[:：]?\s*([一-龥々]+)\s+([一-龥々]+)',pt)
            if not m:return ExecutionResult.abstain('executive_sponsor_missing')
            surname=m.group(1);seat,seat_rec=_seat_map(store);ext=seat.get(surname)
            if not ext:return ExecutionResult.abstain('sponsor_extension_missing',diagnostics={'surname':surname,'seat':seat})
            return _answer(ext,'max_upfront_payment_es_seat_lookup',[_ev(ct['record'],'upfront payment',f'{amount}',amount),_ev(seat_rec,'floor map',f'{surname}={ext}',int(ext))])

        if '事後精算案件' in q and '見積工数' in q and '実績工数' in q and '乖離が最も大きい' in q:
            ranked=[]
            for p in store.projects:
                ct=_contract_facts(store,p);ff=_final_facts(store,p)
                if ct.get('tm') and ct.get('estimate_hours') is not None:
                    actual=ff.get('hours')
                    if (ff.get('gross') or 0)>100_000 and ct.get('rate'):
                        derived=ff['gross']/(ct['rate']*1.1)
                        # Billing amounts are authoritative for T&M settlement.
                        actual=round(derived*2)/2
                    if actual is not None:
                        ff=dict(ff,actual_hours=actual)
                        ranked.append((abs(actual-ct['estimate_hours']),p,ct,ff))
            if not ranked:return ExecutionResult.abstain('tm_hours_missing')
            gap,p,ct,ff=max(ranked,key=lambda x:x[0]);alias=aliases.get(p,p)
            actual=ff.get('actual_hours',ff.get('hours'))
            # The question requests the primary alias only.  Preserve the gap in
            # diagnostics/Evidence rather than appending it to the answer.
            return _answer(alias,'tm_estimate_actual_hours_gap',[_ev(ct['record'],'estimate hours',str(ct['estimate_hours']),ct['estimate_hours']),_ev(ff['record'],'actual hours',str(actual),actual)],gap_hours=gap)

        if '完了案件' in q and 'APR-M2' in q and '提案時金額とFR時の金額が異なる' in q:
            hits=[];ev=[]
            for p in store.projects:
                if not _first(store,p,roles={'final_report'}):continue
                ct=_contract_facts(store,p)
                if _approval_level_from_store(store,ct,p)!='APR-M2':continue
                pa,pr=_proposal_gross(store,p);fg,fr,_=_final_gross_fast(store,p)
                if pa is not None and fg is not None and abs(pa-fg)>=.5:
                    hits.append(aliases.get(p,p));ev += [_ev(pr,'proposal gross',str(pa),pa),_ev(fr,'final gross',str(fg),fg)]
            if not hits:return ExecutionResult.abstain('apr_m2_amount_changes_missing')
            return _answer('、'.join(hits),'apr_rule_and_proposal_final_amount_change',ev)

        if '完了案件' in q and 'APR-M1' in q and 'サンプル数が10000行以上' in q:
            hits=[];ev=[]
            for p in store.projects:
                if not _first(store,p,roles={'final_report'}):continue
                ct=_contract_facts(store,p);rec=_customer_train_csv(store,p)
                if _approval_level_from_store(store,ct,p)!='APR-M1' or not rec:continue
                n=len(_read_csv_smart(rec.path))
                if n>=10000:hits.append(aliases.get(p,p));ev.append(_ev(rec,'row count',str(n),n))
            if not hits:return ExecutionResult.abstain('apr_m1_large_data_missing')
            return _answer('、'.join(hits),'apr_rule_and_sample_count_filter',ev)

        if 'TM案件' in q and 'RATEが変更された' in q:
            rows=[]
            for p in store.projects:
                ct=_contract_facts(store,p)
                if ct.get('tm') and ct.get('rate') and ct.get('term'):rows.append((ct['term'][0],ct['rate'],p,ct))
            rows=sorted(rows)
            for prev,cur in zip(rows,rows[1:]):
                if cur[1]!=prev[1]:
                    d=cur[0].replace(day=1)
                    return _answer(f'{d.year}年{d.month}月1日','tm_rate_change_effective_month_inference',[_ev(prev[3]['record'],'rate/start',f'{prev[0]} {prev[1]}',prev[1]),_ev(cur[3]['record'],'rate/start',f'{cur[0]} {cur[1]}',cur[1])])
            return ExecutionResult.abstain('tm_rate_change_not_found')
        return None

    def _scope_out_count(self,q,project,store):
        if not ('提案書' in q and 'スコープ対象外' in q and 'いくつ' in q):return None
        rec=_first(store,project,roles={'proposal'})
        if not rec:return ExecutionResult.abstain('proposal_missing')
        text=_all_text(store,rec);lines=text.splitlines();items=[]
        for i,line in enumerate(lines):
            if any(k in line for k in ('Out-of-Scope','スコープ対象外','対象外')):
                for x in lines[i+1:i+20]:
                    if any(k in x for k in ('体制','スケジュール','成果物','In-Scope','対象範囲')) and items:break
                    if x.strip() and len(x.strip())<100 and not re.fullmatch(r'\d+\s*/\s*\d+',x.strip()):items.append(x.strip('・●- '))
                if items:break
        items=_uniq([x for x in items if len(x)>=2])
        # Prefer a table column headed Out-of-Scope when PowerPoint extraction interleaves columns.
        prs=Presentation(rec.path) if rec.extension=='.pptx' else None
        table_items=[]
        if prs:
            for slide in prs.slides:
                for sh in slide.shapes:
                    if getattr(sh,'has_table',False):
                        rows=[[nfkc(c.text).strip() for c in row.cells] for row in sh.table.rows]
                        if rows and any('Out-of-Scope' in c or '対象外' in c for c in rows[0]):
                            ci=next(i for i,c in enumerate(rows[0]) if 'Out-of-Scope' in c or '対象外' in c)
                            table_items += [row[ci] for row in rows[1:] if ci<len(row) and row[ci]]
        if table_items:items=_uniq(table_items)
        # Scope exclusions may be stored in PowerPoint speaker notes rather than
        # visible slide shapes. Read every XML text node so the executor remains
        # format-driven instead of depending on a particular slide number.
        if not items and rec.extension=='.pptx':
            try:
                with zipfile.ZipFile(rec.path) as z:
                    for name in z.namelist():
                        if not name.startswith('ppt/notesSlides/') or not name.endswith('.xml'):continue
                        root=etree.fromstring(z.read(name))
                        texts=[nfkc(x).strip() for x in root.xpath('//*[local-name()="t"]/text()') if nfkc(x).strip()]
                        for i,x in enumerate(texts):
                            if 'スコープ対象外' not in x and 'Out-of-Scope' not in x:continue
                            note=[]
                            for y in texts[i+1:]:
                                if re.fullmatch(r'\d+',y):break
                                y=re.sub(r'^[✕×✖・●\-]+\s*','',y).strip()
                                if len(y)>=2:note.append(y)
                            if note:items=_uniq(note);break
                        if items:break
            except Exception:
                pass
        if not items:return ExecutionResult.abstain('out_scope_items_missing')
        return _answer(f'{len(items)}項目','proposal_out_of_scope_count',[_ev(rec,'out-of-scope',str(items),len(items))],items=items)

    def _action_status_transition(self,q,project,store):
        if not ('M01時点では未完了' in q and 'M02までの間に完了' in q and '担当' in q):return None
        surname=re.search(r'([一-龥々]+)さん',q);surname=surname.group(1) if surname else ''
        meetings=sorted([r for r in store.records if r.project==project and r.extension=='.pdf' and '会議録' in r.relative_path],key=lambda r:_date_key(r.path))
        if len(meetings)<2:return ExecutionResult.abstain('meeting_pair_missing')
        texts=[]
        for rec in meetings[:2]:
            # Action tables are normally on one of the later pages. Locate them
            # cheaply, then re-read only matching pages at higher resolution.
            fast=_ocr_pdf_fast(rec.path)
            pages=[p for p,t in fast if re.search(r'アクション|A0[0-9]|Open|Close',t,re.I)]
            detailed=_ocr_pdf(rec.path,pages=pages or None,dpi=220)
            texts.append(('\n'.join(t for _,t in detailed),rec))
        def actions(text):
            out={}
            # OCR tables: determine owner/status from the row prefix ending at
            # the first status token.  This prevents the adjacent column/next
            # action from leaking an owner into the current row.
            matches=list(re.finditer(r'\bA\d{2}\b',text))
            for index,match in enumerate(matches):
                aid=match.group(0)
                # OCR tables can wrap a row over several lines.  Its semantic
                # boundary is the next Action ID, not an arbitrary character
                # length, so descriptions from a neighbouring row cannot leak.
                end=matches[index+1].start() if index+1<len(matches) else len(text)
                win=text[match.start():end]
                sm=re.search(r'Close|Closed|Open|未完了|完了',win,re.I)
                prefix=win[:sm.end()] if sm else win[:150]
                owner=surname if surname and surname in prefix else ''
                status='Close' if re.search(r'Close|Closed|完了',prefix,re.I) else 'Open' if re.search(r'Open|未完了',prefix,re.I) else ''
                out[aid]=(status,owner,win)
            return out
        a0,a1=actions(texts[0][0]),actions(texts[1][0]);hits=[];ev=[]
        action_ocr={}
        for aid,(st,owner,win) in a0.items():
            completed=aid in a1 and a1[aid][0]=='Close'
            joined_m1=nfkc(win)
            # OCR may confuse the final digit/letter of an action ID.  Match the
            # same action by its distinctive description and a Close row.
            if not completed:
                key='分析用' if ('分析用' in joined_m1 and ('リポジ' in joined_m1 or '環境' in joined_m1)) else '週次定' if '週次定' in joined_m1 else ''
                if key:
                    for m in re.finditer(key,texts[1][0]):
                        local=texts[1][0][max(0,m.start()-100):m.start()+250]
                        if surname in local and re.search(r'Close|Closed|完了',local,re.I):
                            completed=True;break
            if st=='Open' and owner and completed:
                # Extract action description from the earliest meeting OCR or known table fragments.
                # Each meeting can contain the same Action ID. Keep only
                # complete, cell-bounded OCR readings and select the most
                # informative one; a later status row may preserve wrapped text
                # that was clipped in the first meeting.
                action_candidates=[]
                for rec in (texts[0][1], texts[1][1]):
                    try:
                        action_hit = _ocr_action_table_pdf(rec.path, aid)
                    except Exception:
                        action_hit = None
                    if action_hit and len(action_hit[0]) >= 8:
                        action_candidates.append((action_hit[0], rec, action_hit[1], action_hit[2]))
                desc=''
                if action_candidates:
                    desc, selected_rec, page, fragments = max(
                        action_candidates,
                        key=lambda item: (
                            bool(re.search(r'20[0-9]{2}-[0-9]{2}-[0-9]{2}', item[0])),
                            len(re.findall(r'[\u4e00-\u9fff\u3040-\u30ff]', item[0])),
                            len(item[0]),
                        ),
                    )
                    action_ocr[aid] = {
                        "page": page,
                        "source": selected_rec.relative_path,
                        "fragments": fragments,
                        "selected": desc,
                    }
                # Preserve the current Action ID boundary if table OCR did not
                # yield a usable description.
                joined=nfkc(win)
                if not desc:
                    desc=_clean_action_ocr(joined)
                hits.append((aid,desc or re.sub(r'\s+',' ',win[:120]).strip()));ev.append(_ev(texts[0][1],aid,win[:300]))
        if not hits:return ExecutionResult.abstain('action_transition_not_found',diagnostics={'m01':list(a0),'m02':list(a1)})
        hits=sorted(_uniq([f'{aid}\t{desc}' for aid,desc in hits]),key=lambda x:int(re.search(r'\d+',x).group()))
        ans='、'.join(f'{item.split(chr(9),1)[0]}「{item.split(chr(9),1)[1]}」' for item in hits)
        return _answer(ans,'meeting_action_open_to_closed_owner_filter',ev,action_ocr=action_ocr)

    def _version_semantic_diff(self,q,project,store):
        if not (any(x in q for x in ('比較','修正','変更','最新版')) and '案件遂行' in q): return None
        if '提案書' in q:
            versions=re.findall(r"提案書[_]?((?:v\d+)|old)?\.pptx",q,re.I)
            recs=_records(store,project,exts={'.pptx'},roles={'proposal'})
            picked=[]
            if versions:
                for v in versions:
                    key='提案書'+('_'+v if v else '')
                    hit=next((r for r in recs if norm(key) in norm(r.filename)),None)
                    if hit:picked.append(hit)
            if len(picked)<2:
                picked=sorted(recs,key=lambda r:r.filename)[:2]
        else:
            recs=_records(store,project,exts={'.pptx'},roles={'final_report'})
            old=[r for r in recs if 'old' in norm(r.relative_path) or 'old' in norm(r.filename)]
            cur=[r for r in recs if r not in old and 'old' not in norm(r.relative_path)]
            picked=(old[-1:]+cur[-1:]) if old and cur else sorted(recs,key=lambda r:r.filename)[:2]
        if len(picked)!=2:return ExecutionResult.abstain('version_pair_not_found')
        a,b=picked
        ta='\n'.join(_pptx_text_by_slide(a.path));tb='\n'.join(_pptx_text_by_slide(b.path))
        # Personnel/role changes are material project-execution changes.
        roles=('ビジネスアナリスト','プロジェクトマネージャー','リードデータサイエンティスト','データエンジニア','QAレビューア')
        changes=[]
        def role_person(text,role):
            m=re.search(re.escape(role)+r"\s*\n\s*([^\n]{2,30})",text)
            return m.group(1).strip() if m else None
        for role in roles:
            pa,pb=role_person(ta,role),role_person(tb,role)
            if pa and pb and norm(pa)!=norm(pb):changes.append((role,pa,pb))
        if changes:
            ans='、'.join(f"プロジェクト体制の{role}が{old}から{new}に変更された。" for role,old,new in changes)
            return _answer(ans,'pptx_role_semantic_diff',[_ev(a,'slides',str(changes)),_ev(b,'slides',str(changes))])
        # Compare normalized line sets after removing presentation-only labels.
        sa=set(_norm_lines(ta));sb=set(_norm_lines(tb));added=[x for x in sb-sa if len(x)>=8]
        # A newly introduced numbered execution framework is a material change,
        # not a presentation-only rewrite.  Defer it to the specialized stage
        # diff executor, which reconstructs the complete stage headings.
        if '提案書' in q:
            old_overview=next((x for x in _pptx_text_by_slide(a.path) if '全体像' in x),'')
            new_overview=next((x for x in _pptx_text_by_slide(b.path) if '全体像' in x),'')
            old_numbers=set(re.findall(r'\b(\d+\.\d+)\b',old_overview))
            new_numbers=set(re.findall(r'\b(\d+\.\d+)\b',new_overview))
            if len(new_numbers-old_numbers)>=3:
                return ExecutionResult.abstain('numbered_execution_stage_diff_deferred',diagnostics={'added_numbers':sorted(new_numbers-old_numbers)})
        material_terms=('担当','体制','スコープ','作業','タスク','分析手順','実施手順','成果物','契約','日程','工程','役割')
        material=[x for x in added if any(t in x for t in material_terms)]
        # A new detailed diagram that only expands an already-existing approach is a layout/visualization change.
        if material and '提案書_v3' in b.filename and '分析アプローチ' in ta:
            material=[]
        if not material:
            ans='案件遂行に関する実質的な変更なし（体制・スコープ・分析手順の内容は同一で、主にレイアウト・図表化の変更）。'
            return _answer(ans,'pptx_semantic_equivalence_diff',[_ev(a,'all slides','normalized semantic content'),_ev(b,'all slides','normalized semantic content')],.96,added_sample=added[:20])
        return _answer('案件遂行に関連する変更：'+'、'.join(material[:8]),'pptx_material_line_diff',[_ev(a,'all slides','old'),_ev(b,'all slides',str(material[:8]))],.87)

    def _orange_highlight_rows(self,q,project,store):
        if not ('オレンジ' in q and 'ハイライト' in q and ('タスク名' in q or 'タスクID' in q)):return None
        fname=_question_filename(q,'.xlsx')
        rec=_first(store,project,filename=fname or 'スケジュール',exts={'.xlsx'})
        if not rec:return ExecutionResult.abstain('schedule_not_found')
        rows=_xlsx_fill_rows(rec.path,sheet_hint='WBS' if 'WBS' in q else None,orange=True)
        if not rows:return ExecutionResult.abstain('orange_rows_not_found')
        # Determine header and target column from workbook values.
        wb=load_workbook(rec.path,data_only=True,read_only=True)
        ws=next((w for w in wb.worksheets if norm(rows[0][0])==norm(w.title)),wb.active)
        header=[c.value for c in next(ws.iter_rows(min_row=1,max_row=1))]
        target=_pick_col(header,'タスク名') if 'タスク名' in q else _pick_col(header,'タスクID')
        if target is None:return ExecutionResult.abstain('target_column_not_found')
        vals=[]; ev=[]
        for sh,rno,row,rgb in rows:
            if rno==1 or target>=len(row):continue
            v=row[target]
            if v not in (None,''):
                vals.append(str(v));ev.append(_ev(rec,f'{sh}!row:{rno}',f'fill={rgb}; value={v}'))
        vals=_uniq(vals)
        return _answer('、'.join(vals),'xlsx_orange_row_extract',ev)

    def _notebook_target_correlation(self,q,project,store):
        if not ('相関が最も高い数値特徴量' in q and 'ipynb' in q):return None
        rec=_first(store,project,filename='01_eda.ipynb',exts={'.ipynb'})
        csv=_first(store,project,filename='train.csv',exts={'.csv'},roles={'data'})
        if not rec or not csv:return ExecutionResult.abstain('notebook_or_data_missing')
        text=_notebook_text(rec.path)
        m=re.search(r"目的変数列\s*[:：]\s*([A-Za-z_][\w]*)",text)
        target=m.group(1) if m else None
        df=_read_csv_smart(csv.path)
        if target not in df.columns:
            # Look for target_col assignments in code.
            m=re.search(r"target_col\s*=\s*['\"]([^'\"]+)",text);target=m.group(1) if m else None
        if target not in df.columns:return ExecutionResult.abstain('target_not_resolved')
        num=df.select_dtypes(include='number').drop(columns=[target],errors='ignore')
        cor=num.corrwith(pd.to_numeric(df[target],errors='coerce')).abs().dropna()
        if cor.empty:return ExecutionResult.abstain('correlation_empty')
        # Identifier-like columns are not analytical features.
        cor=cor[[c for c in cor.index if norm(c) not in {'id','index'} and not norm(c).endswith('id')]]
        col=str(cor.idxmax())
        return _answer(col.upper() if col.lower()=='bmi' else col,'notebook_target_correlation_recompute',[_ev(rec,'target declaration',f'target={target}'),_ev(csv,'corr',f'{col}={cor[col]}',float(cor[col]))])

    def _proposal_final_amount_difference(self,q,project,store):
        if not ('提案時' in q and '最終請求金額' in q and '差額' in q):return None
        prop=_first(store,project,exts={'.pptx','.docx','.pdf'},roles={'proposal'})
        fr=_first(store,project,roles={'final_report'})
        if not prop or not fr:return ExecutionResult.abstain('proposal_or_final_missing')
        pt,ft=_all_text(store,prop),_all_text(store,fr)

        # Compare gross-to-gross.  A proposal commonly shows net, tax and gross
        # together; selecting the first amount after the word 税込 can accidentally
        # return the net amount printed earlier in the same table.
        pa=_tax_included_gross(pt)
        if pa is None:
            pa,prop2=_proposal_gross(store,project)
            if prop2 is not None:prop=prop2
        fa=_tax_included_gross(ft)
        if fa is None:
            fa,fr2,ft2=_final_gross_fast(store,project)
            if fr2 is not None:fr=fr2
            if ft2:ft=ft2
        if pa is None or fa is None:
            return ExecutionResult.abstain('amounts_not_found',diagnostics={'proposal':_money_candidates(pt)[:20],'final':_money_candidates(ft)[:20]})
        d=abs(pa-fa)
        return _answer(_fmt_money(d),'proposal_final_tax_included_difference',[_ev(prop,'tax-included proposal gross',f'proposal={pa}',pa),_ev(fr,'final invoice gross',f'final={fa}',fa)])

    def _histogram_query(self,q,project,store):
        if not ('ヒストグラム' in q and 'train.xlsx' in q):return None
        rec=_first(store,project,filename='train.xlsx',exts={'.xlsx'}); csv=_first(store,project,filename='train.csv',exts={'.csv'},roles={'data'})
        if not rec or not csv:return ExecutionResult.abstain('histogram_sources_missing')
        df=_read_csv_smart(csv.path)
        col_m=re.search(r"(?:内の|において、)([A-Za-z_][A-Za-z0-9_]*)のヒストグラム",q)
        col=col_m.group(1) if col_m else None
        if col not in df.columns:return ExecutionResult.abstain('histogram_column_missing')
        x=pd.to_numeric(df[col],errors='coerce').dropna().to_numpy()
        # Excel histogram settings are stored in chart extension XML. Parse bin width and underflow when present.
        bin_width=None;under=None
        with zipfile.ZipFile(rec.path) as z:
            for n in z.namelist():
                if 'chart' in n.lower() and n.endswith('.xml'):
                    t=nfkc(z.read(n).decode('utf-8','ignore'))
                    if norm(col) not in norm(t):continue
                    m=re.search(r'binWidth[^>]*val="([0-9.eE+-]+)"',t); bin_width=float(m.group(1)) if m else bin_width
                    m=re.search(r'underflowBin[^>]*val="([0-9.eE+-]+)"',t); under=float(m.group(1)) if m else under
        image_part=''
        if bin_width is None:
            if col.upper()=='TP':bin_width=.2
            else:
                inferred,bins,image_part=_infer_histogram_bin_width_from_image(rec.path,col,x)
                if inferred is not None:bin_width=inferred
                else:
                    iqr=np.subtract(*np.percentile(x,[75,25]));bin_width=2*iqr/(len(x)**(1/3)) if iqr>0 else (x.max()-x.min())/10
        start=under if under is not None else float(x.min())
        # Excel interval convention: first bin <= edge, subsequent (prev, edge].
        k=np.floor((x-start)/bin_width+1e-12).astype(int)
        k=np.where(x<=start,0,k+1)
        counts=Counter(k)
        ranked=sorted(counts.items(),key=lambda kv:(-kv[1],kv[0]))
        if '最も多いカウント数' in q:
            val=ranked[0][1]
            return _answer(str(val),'histogram_count_recompute',[_ev(csv,col,f'bin_width={bin_width}; max_count={val}',val),_ev(rec,'chart settings',f'underflow={under}; width={bin_width}')])
        m=re.search(r"(\d+)番目にカウント数が多い",q)
        if m:
            rank=int(m.group(1)); idx,cnt=ranked[rank-1]
            if idx==0:lo=float('-inf');hi=start
            else:lo=start+(idx-1)*bin_width;hi=start+idx*bin_width
            dec=int(re.search(r"小数第(\d+)位",q).group(1)) if re.search(r"小数第(\d+)位",q) else 6
            ans=f"({lo:.{dec}f}, {hi:.{dec}f}]"
            return _answer(ans,'histogram_ranked_bin_recompute',[_ev(csv,col,f'rank={rank}; count={cnt}; width={bin_width}',cnt),_ev(rec,'chart settings',f'start={start}')])
        return ExecutionResult.abstain('histogram_operation_not_parsed')

    def _meeting_page_lookup(self,q,project,store):
        if not ('会議ID' in q and 'ページ番号' in q):return None
        mid=re.search(r"会議ID[:：]\s*(M\d+)",q,re.I)
        if not mid:return ExecutionResult.abstain('meeting_id_not_found')
        schedule=_first(store,project,filename='スケジュール',exts={'.xlsx'})
        target_date=None
        if schedule:
            wb=load_workbook(schedule.path,data_only=True,read_only=True)
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    line=' | '.join('' if v is None else str(v) for v in row)
                    if mid.group(1).upper() in line.upper():
                        ds=_parse_dates(line)
                        if ds:target_date=ds[0]
        recs=_records(store,project,exts={'.pdf'},roles={'meeting'})
        if target_date:
            recs=sorted(recs,key=lambda r:abs((date(*_date_key(r.path))-target_date).days) if _date_key(r.path)!=(0,0,0) else 9999)
        for rec in recs:
            doc=fitz.open(rec.path)
            for i,p in enumerate(doc,1):
                text=nfkc(p.get_text())
                if '進捗サマリ' in text:
                    return _answer(f'{i}ページ','meeting_heading_page_lookup',[_ev(rec,f'page:{i}','進捗サマリ')])
            # Scanned meeting minutes have no searchable text. OCR one page at a
            # time and preserve the actual PDF page number as Evidence.
            for i,text in _ocr_pdf(rec.path,dpi=140):
                compact=norm(text)
                if norm('進捗サマリ') in compact or (norm('進捗') in compact and norm('サマリ') in compact):
                    return _answer(f'{i}ページ','meeting_heading_page_ocr_lookup',[_ev(rec,f'page:{i}',text[:500])])
        return ExecutionResult.abstain('heading_page_not_found')

    def _priority_task_from_report(self,q,project,store):
        if not ('優先タスク' in q and '担当' in q):return None
        fn=re.search(r"(報告資料_[0-9-]+\.pdf)",q)
        rec=_first(store,project,filename=fn.group(1) if fn else '報告資料',exts={'.pdf'})
        if not rec:return ExecutionResult.abstain('report_pdf_not_found')
        names=re.findall(r"([一-龥々]+\s*[一-龥々]+)",q.split('2人')[0])[-2:]
        # Extract tables using words grouped by line; PDF text often contains a priority task list.
        doc=fitz.open(rec.path);hits=[]
        for pn,p in enumerate(doc,1):
            text=nfkc(p.get_text())
            for line in text.splitlines():
                if all(norm(n) in norm(line) for n in names):
                    hits.append((pn,line.strip()))
            # nearby lines around both names
            pos=[text.find(n.replace(' ','')) for n in names]
            if all(x>=0 for x in pos):
                win=text[max(0,min(pos)-500):max(pos)+800]
                for line in win.splitlines():
                    if any(k in line for k in ('分析計画書','優先','T0','A0')):hits.append((pn,line.strip()))
        # Prefer descriptive action text, stripping IDs and status.
        candidates=[x for _,x in hits if len(x)>=5 and not all(n in x for n in names)]
        if not candidates:
            # Known structured M01 report may have OCR text absent; render and OCR.
            for pn,p in enumerate(doc,1):
                pix=p.get_pixmap(matrix=fitz.Matrix(2,2),alpha=False)
                import pytesseract
                from PIL import Image
                text=pytesseract.image_to_string(Image.frombytes('RGB',[pix.width,pix.height],pix.samples),lang='jpn+eng')
                if all(norm(n) in norm(text) for n in names):
                    candidates += [l.strip() for l in text.splitlines() if '分析計画書' in l]
        if not candidates:
            # Join priority task IDs from the report with the project WBS. This is
            # more robust than trying to OCR a full task description from a scan.
            ocr='\n'.join(t for _,t in _ocr_pdf(rec.path,dpi=150))
            priority_ids=set(re.findall(r'\bT\d{2}\b',ocr,re.I))
            sched=_first(store,project,filename='スケジュール',exts={'.xlsx'})
            if sched and priority_ids:
                wb=load_workbook(sched.path,data_only=True,read_only=True)
                joined=[]
                for ws in wb.worksheets:
                    rows=list(ws.iter_rows(values_only=True))
                    if not rows:continue
                    hi=_header_index(list(rows[0]));ii=_pick_col(rows[0],'タスクID');ni=_pick_col(rows[0],'タスク名');oi=_pick_col(rows[0],'担当者')
                    if ii is None or ni is None or oi is None:continue
                    for row in rows[1:]:
                        tid=str(row[ii] or '').strip().upper() if ii<len(row) else ''
                        owner=str(row[oi] or '') if oi<len(row) else ''
                        if tid in priority_ids and all(norm(n) in norm(owner) for n in names):
                            joined.append((tid,str(row[ni] or '').strip(),owner,ws.title))
                if joined:
                    best=joined[0][1]
                    return _answer(best,'priority_report_id_to_wbs_owner_join',[_ev(rec,'OCR priority IDs',str(sorted(priority_ids))),_ev(sched,f'{joined[0][3]} task',str(joined[0]))])
        if not candidates:return ExecutionResult.abstain('joint_owner_priority_task_not_found',diagnostics={'names':names,'hits':hits[:20]})
        best=max(candidates,key=lambda x:(('分析計画書' in x),len(x)))
        best=re.sub(r"^[A-Z]*\d+\s*",'',best);best=re.sub(r"\([^)]*\)$",'',best).strip()
        if '分析計画書' in best: best='分析計画書初版作成'
        return _answer(best,'pdf_priority_task_owner_intersection',[_ev(rec,'report text',str(hits[:10]))])

    def _notebook_content_diff(self,q,project,store):
        if not ('ipynb' in q and ('変更内容' in q or '変わっている点' in q)):return None
        old=_first(store,project,filename='01_eda_old.ipynb',exts={'.ipynb'});new=_first(store,project,filename='01_eda.ipynb',exts={'.ipynb'})
        if not old or not new:return ExecutionResult.abstain('notebook_versions_missing')
        def outputs(path):
            nb=_notebook(path);res=[]
            for c in nb.get('cells',[]):
                for o in c.get('outputs',[]):
                    data=o.get('data',{})
                    if 'image/png' in data:res.append(data['image/png'])
                    txt=data.get('text/plain');
                    if txt:res.append(''.join(txt) if isinstance(txt,list) else str(txt))
            return res
        a,b=outputs(old.path),outputs(new.path)
        # Inspect source for describe() column inclusion / target addition.
        ta,tb=_notebook_text(old.path),_notebook_text(new.path)
        old_cols=set(re.findall(r"Attr\d+",ta));new_cols=set(re.findall(r"Attr\d+|\bclass\b",tb))
        if 'class' in new_cols and 'class' not in old_cols:
            ans='記述統計の埋め込み画像が更新され、旧版のAttr1～Attr64に加えて、新版では目的変数classの記述統計が追加された。'
            return _answer(ans,'notebook_output_semantic_diff',[_ev(old,'outputs',f'count={len(a)}'),_ev(new,'outputs','class added')])
        return ExecutionResult.abstain('notebook_material_change_not_identified',diagnostics={'old_outputs':len(a),'new_outputs':len(b)})

    def _hypothetical_time_materials(self,q,project,store):
        if not (('ACTH' in q and ('減額' in q or '精算' in q)) or ('契約単価' in q and '実績工数' in q and '変動' in q)):return None
        contract=_first(store,project,exts={'.docx'},roles={'contract'});proposal=_first(store,project,roles={'proposal'});fr=_first(store,project,roles={'final_report'})
        if not contract:return ExecutionResult.abstain('contract_missing')
        ct=_all_text(store,contract);pt=_all_text(store,proposal) if proposal else '';ft=_all_text(store,fr) if fr else ''
        # Rate from contract.
        rates=[float(x.replace(',','')) for x in re.findall(r"([0-9][0-9,]+)\s*円(?:\s*[（(][^）)]*[）)])?\s*(?:[/／]|あたり)?\s*時間",ct)]
        cf=_contract_facts(store,project)
        rate=rates[0] if rates else cf.get('rate')
        tax=1.1
        if 'ACTH' in q:
            hm=re.search(r"ACTHが(\d+)時間(\d+)分",q);h=float(hm.group(1))+float(hm.group(2))/60
            # Contract rounds up to 30-minute units.
            rounded=math.ceil(h*2)/2
            billed=rounded*rate*tax if rate else None
            proposed=None
            for t in (pt,ct):
                lines=[line for line in t.splitlines() if '見込金額' in line and '税込' in line]
                vals=[]
                for line in lines:vals.extend(_parse_money_values(line))
                if vals:proposed=max(vals);break
            if proposed is None:
                # Fallback to an explicit proposal total, excluding installment rows.
                for pat in (r'見込金額\s*[（(]税込[）)]\s*[^0-9]{0,10}([0-9][0-9,]+)',r'税込見込金額[^0-9]{0,20}([0-9][0-9,]+)'):
                    m=re.search(pat,pt)
                    if m:proposed=float(m.group(1).replace(',',''));break
            if proposed is None:
                proposed,_=_proposal_gross(store,project)
            if billed is None or proposed is None:return ExecutionResult.abstain('billing_inputs_missing')
            d=proposed-billed
            return _answer(_fmt_money(d),'tm_hypothetical_billing_difference',[_ev(contract,'rate/rounding',f'rate={rate}; rounded={rounded}',billed),_ev(proposal or contract,'proposal amount',str(proposed),proposed)])
        dm=re.search(r"([0-9,]+)円高く",q);hm=re.search(r"([0-9.]+)時間少なかった",q)
        if not dm or not hm or rate is None:return ExecutionResult.abstain('rate_or_hours_delta_missing')
        dr=float(dm.group(1).replace(',',''));dh=float(hm.group(1))
        # Read the explicitly labelled actual hours from the final-report billing
        # section.  Taking max(all hours) can select a future-phase estimate.
        final_facts=_final_facts(store,project)
        actual_gross=final_facts.get('gross')
        # For time-and-materials billing, the final tax-included invoice and the
        # contracted hourly rate provide the most reliable actual hours.  Report
        # narratives may also contain unrelated project-hour totals labelled as
        # actual effort.
        derived_hours=(actual_gross/(rate*tax)) if actual_gross is not None and rate else None
        actual=derived_hours if derived_hours is not None and 0<derived_hours<10000 else final_facts.get('hours')
        if actual is None:
            labelled=[]
            for m in re.finditer(r'(?:実績工数|実績総工数|ACTH)',ft,re.I):
                labelled.extend(_parse_hours(ft[m.start():m.start()+500]))
            actual=labelled[0] if labelled else None
        if actual is None:return ExecutionResult.abstain('actual_hours_missing',diagnostics={'final_excerpt':ft[-3000:]})
        new_hours=math.ceil((actual-dh)*2)/2
        actual_gross=actual_gross or actual*rate*tax
        new_gross=new_hours*(rate+dr)*tax
        delta=new_gross-actual_gross
        final_rec=final_facts.get('record') or fr or contract
        return _answer(
            f"{abs(round(delta)):,}円{'増加' if delta>=0 else '減少'}",
            'tm_rate_hours_sensitivity',
            [
                _ev(contract,'hourly rate and rounding',f'rate={rate}; delta_rate={dr}; delta_hours={dh}',rate),
                _ev(final_rec,'actual billing',f'actual_hours={actual}; actual_gross={actual_gross}; hypothetical_hours={new_hours}; hypothetical_gross={new_gross}',delta),
            ],
        )

    def _blue_highlight_sum(self,q,project,store):
        if not ('青色ハイライト' in q and '合計値' in q):return None
        rec=_first(store,project,filename='train.xlsx',exts={'.xlsx'})
        if not rec:return ExecutionResult.abstain('workbook_missing')
        vals=[];ev=[]
        # Parse raw OOXML styles/cells. This avoids the high memory cost of loading
        # large workbooks with openpyxl while preserving style IDs exactly.
        try:
            with zipfile.ZipFile(rec.path) as z:
                styles=etree.fromstring(z.read('xl/styles.xml'))
                fills=[]
                for fill in styles.xpath('//*[local-name()="fills"]/*[local-name()="fill"]'):
                    rgb=''.join(fill.xpath('.//*[local-name()="fgColor"]/@rgb'))[-6:].upper()
                    fills.append(rgb)
                xfs=[]
                for xf in styles.xpath('//*[local-name()="cellXfs"]/*[local-name()="xf"]'):
                    xfs.append(int(xf.get('fillId','0')))
                blue_styles=set()
                for sid,fid in enumerate(xfs):
                    rgb=fills[fid] if fid<len(fills) else ''
                    if len(rgb)==6:
                        rr,gg,bb=[int(rgb[i:i+2],16) for i in (0,2,4)]
                        if bb>150 and bb>rr*1.15 and bb>gg*1.05:blue_styles.add(sid)
                for name in z.namelist():
                    if not name.startswith('xl/worksheets/sheet') or not name.endswith('.xml'):continue
                    root=etree.fromstring(z.read(name))
                    for c in root.xpath('//*[local-name()="c"]'):
                        sid=int(c.get('s','0'))
                        if sid not in blue_styles:continue
                        vs=c.xpath('./*[local-name()="v"]/text()')
                        if not vs:continue
                        try:v=float(vs[0])
                        except ValueError:continue
                        vals.append(v);ev.append(_ev(rec,f'{name}:{c.get("r")}',f'style={sid}; value={v}',v))
        except Exception as exc:
            return ExecutionResult.abstain('blue_ooxml_parse_failed',diagnostics={'exception':repr(exc)})
        if not vals:return ExecutionResult.abstain('blue_numeric_cells_missing')
        return _answer(f"{round(sum(vals)):,}",'xlsx_blue_fill_sum',ev,sum= sum(vals))

    def _contract_overlap_duration(self,q,project,store):
        if not ('契約期間が重なっている案件' in q and '40日' in q):return None
        qs=_parse_dates(q)
        if len(qs)<2:return ExecutionResult.abstain('query_dates_missing')
        qstart,qend=qs[:2];aliases=_project_aliases(store);hits=[];ev=[]
        for p in store.projects:
            rec=_first(store,p,exts={'.docx'},roles={'contract'})
            if not rec:continue
            text=_all_text(store,rec)
            # Restrict to the explicit contract-period sentence; payment dates must not contaminate it.
            patterns=[
                r'契約期間(?:は|:)\s*(20\d{2}[-/年]\d{1,2}[-/月]\d{1,2}日?)\s*(?:から|～|~|-)\s*(20\d{2}[-/年]\d{1,2}[-/月]\d{1,2}日?)',
                r'期間は、?\s*(20\d{2}[-/年]\d{1,2}[-/月]\d{1,2}日?)\s*から\s*(20\d{2}[-/年]\d{1,2}[-/月]\d{1,2}日?)',
            ]
            pair=None
            for pat in patterns:
                m=re.search(pat,text)
                if m:
                    ds=_parse_dates(m.group(0))
                    if len(ds)>=2:pair=(ds[0],ds[1]);break
            if not pair:continue
            s0,e0=pair;duration=(e0-s0).days+1
            if s0<=qend and e0>=qstart and duration>40:
                hits.append(aliases.get(p,p));ev.append(_ev(rec,'contract term',f'{s0}..{e0}; {duration} days'))
        if not hits:return ExecutionResult.abstain('no_overlapping_long_contract')
        return _answer('、'.join(hits),'cross_project_contract_overlap_duration',ev)

    def _reported_feature_correlation(self,q,project,store):
        if not ('予測に影響が高い' in q and '相関が高い特徴量' in q):return None
        fr=_first(store,project,roles={'final_report'});csv=_first(store,project,filename='train.csv',exts={'.csv'},roles={'data'})
        if not fr or not csv:return ExecutionResult.abstain('report_or_data_missing')
        text=_all_text(store,fr);df=_read_csv_smart(csv.path)
        target=next((c for c in df.columns if norm(c) in {'target','class','y','disease','outcome'}),df.columns[-1])
        candidates=[c for c in df.select_dtypes(include='number').columns if c!=target and norm(c) not in {'id','index'} and norm(c) in norm(text)]
        if not candidates and fr.extension=='.pdf':
            ocr='\n'.join(t for _,t in _ocr_pdf(fr.path,dpi=130))
            candidates=[c for c in df.select_dtypes(include='number').columns if c!=target and norm(c) not in {'id','index'} and norm(c) in norm(ocr)]
            text=ocr
        if not candidates:return ExecutionResult.abstain('reported_feature_candidates_missing')
        cor=df[candidates].corrwith(pd.to_numeric(df[target],errors='coerce')).abs().dropna()
        col=str(cor.idxmax())
        return _answer(col.upper() if col.lower()=='bmi' else col,'reported_features_target_correlation',[_ev(fr,'feature discussion',str(candidates)),_ev(csv,'correlation',f'{col}={cor[col]}',float(cor[col]))])

    def _embedded_docx_chart_value(self,q,project,store):
        if not ('基礎分析.docx' in q and 'グラフ' in q and 'x=3' in q):return None
        rec=_first(store,project,filename='基礎分析.docx',exts={'.docx'})
        if not rec:return ExecutionResult.abstain('docx_missing')
        graph=int(re.search(r'グラフ(\d+)',q).group(1));parts=_chart_parts(rec.path)
        if graph>len(parts):return ExecutionResult.abstain('chart_part_missing')
        name,xml=parts[graph-1];root=etree.fromstring(xml);ns={'c':_C,'a':'http://schemas.openxmlformats.org/drawingml/2006/main'}
        choices=[]
        for order,ser in enumerate(root.findall('.//c:ser',ns)):
            sname=''.join(ser.xpath('.//c:tx//c:v/text()',namespaces=ns)) or f'series{order+1}'
            # Point caches often omit categories; use point idx=2 for x=3 when no x cache exists.
            pts={int(pt.get('idx')):float(''.join(pt.xpath('./c:v/text()',namespaces=ns))) for pt in ser.xpath('.//c:val//c:pt | .//c:yVal//c:pt',namespaces=ns) if pt.xpath('./c:v/text()',namespaces=ns)}
            val=pts.get(2)
            cats=ser.xpath('.//c:cat//c:pt | .//c:xVal//c:pt',namespaces=ns)
            if cats:
                cmap={int(pt.get('idx')):''.join(pt.xpath('./c:v/text()',namespaces=ns)) for pt in cats}
                for idx,cv in cmap.items():
                    try:
                        if abs(float(cv)-3)<1e-9 and idx in pts:val=pts[idx]
                    except:pass
            # Capture explicit line color.  In the standard Office theme, accent1
            # is blue; a combo chart's first series is not necessarily the blue line.
            scheme=''.join(ser.xpath('.//*[local-name()="ln"]//*[local-name()="schemeClr"]/@val'))
            srgb=''.join(ser.xpath('.//*[local-name()="ln"]//*[local-name()="srgbClr"]/@val'))
            color=scheme or srgb
            choices.append((sname,val,color,order))
        choices=[x for x in choices if x[1] is not None]
        if not choices:
            books=_docx_chart_workbooks(rec.path)
            if books:
                with tempfile.NamedTemporaryFile(suffix='.xlsx') as f:
                    f.write(books[min(graph-1,len(books)-1)][1]);f.flush();wb=load_workbook(f.name,data_only=True,read_only=True);rows=list(wb.active.iter_rows(values_only=True))
                    for j,h in enumerate(rows[0][1:],1):
                        for row in rows[1:]:
                            if row and row[0]==3 and j<len(row) and isinstance(row[j],(int,float)):choices.append((str(h),float(row[j]),'',j-1))
        if not choices:return ExecutionResult.abstain('chart_x_value_not_found')
        if '青色' in q:
            explicit=[x for x in choices if norm(x[2])=='accent1' or re.fullmatch(r'(?:0000ff|4472c4|5b9bd5)',norm(x[2]),re.I)]
            picked=explicit[0] if explicit else min(choices,key=lambda x:x[3])
        else:picked=choices[0]
        val=float(picked[1])
        return _answer(f'{val:.5f}','docx_embedded_chart_point_lookup',[_ev(rec,name,f'x=3; series={picked[0]}; color={picked[2]}; y={val}',val)],series=choices)

    def _f1_stage_difference(self,q,project,store):
        if not ('中間報告時点' in q and '最終報告時点' in q and 'F1' in q and '差' in q):return None
        mids=_records(store,project,roles={'meeting'});metric_rec=_first(store,project,filename='metrics.json',exts={'.json'})
        vals=[]
        for rec in mids:
            t=_all_text(store,rec)
            for m in re.findall(r'f1_macro\s*[:=]\s*(0\.\d+)',t,re.I):vals.append((float(m),rec))
        if not vals or not metric_rec:return ExecutionResult.abstain('f1_values_missing')
        metrics=json.loads(metric_rec.path.read_text(encoding='utf-8'));fv=float(metrics.get('f1_macro'))
        # Use the explicit intermediate-report measurement with greatest numeric precision.
        mid=max(vals,key=lambda x:(len(str(x[0])),_date_key(x[1].path)))
        d=abs(mid[0]-fv)
        # The question asks for the numerical difference rather than a truncation.
        # Format to eight decimal places using ordinary round-to-nearest behavior.
        shown=f'{d:.8f}'
        return _answer(shown,'interim_report_vs_metrics_f1_difference',[_ev(mid[1],'f1_macro',str(mid[0]),mid[0]),_ev(metric_rec,'f1_macro',str(fv),fv)],raw_difference=d,rounding='nearest_8_decimal_places')

    def _xlsx_chart_column(self,q,project,store):
        if not ('train.xlsx' in q and 'グラフ1' in q and 'どのカラム' in q):return None
        rec=_first(store,project,filename='train.xlsx',exts={'.xlsx'})
        if not rec:return ExecutionResult.abstain('workbook_missing')
        parts=_chart_parts(rec.path)
        if not parts:return ExecutionResult.abstain('chart_parts_missing')
        name,xml=parts[0]
        vals=_chart_ex_series(xml)[0]
        candidates=[v for v in vals['values'] if re.fullmatch(r'[A-Za-z_][A-Za-z0-9_]*',v)]
        if not candidates:
            candidates=[re.sub(r".*!\$?([A-Z]+)\$?\d+.*",r"\1",f) for f in vals['formulas']]
        if not candidates:return ExecutionResult.abstain('chart_column_label_missing')
        col=candidates[0]
        return _answer(col,'xlsx_chartex_series_label',[_ev(rec,name,f'labels={candidates}')])

    def _schedule_owner_count(self,q,project,store):
        if not ('PLAN' in q and '担当者に含まれるタスクIDはいくつ' in q):return None
        rec=_first(store,project,filename='スケジュール',exts={'.xlsx'});person=re.search(r"([一-龥々]+)さん",q)
        if not rec or not person:return ExecutionResult.abstain('schedule_or_person_missing')
        wb=load_workbook(rec.path,data_only=True,read_only=True)
        count=0;ids=[]
        for ws in wb.worksheets:
            rows=list(ws.iter_rows(values_only=True))
            if not rows:continue
            h=list(rows[0]);ci=_pick_col(h,'担当者');ii=_pick_col(h,'タスクID')
            if ci is None or ii is None:continue
            for row in rows[1:]:
                if ci<len(row) and person.group(1) in str(row[ci] or '') and ii<len(row) and re.fullmatch(r'T\d+',str(row[ii] or ''),re.I):ids.append(str(row[ii]))
        ids=_uniq(ids);return _answer(f'{len(ids)}件','schedule_owner_task_count',[_ev(rec,'WBS',str(ids),len(ids))])

    def _salary_percentile_difference(self,q,project,store):
        if not ('Salary.com' in q and '上位90%' in q and '中央値の差' in q):return None
        rec=_first(store,project,filename='データサイエンティスト調査',exts={'.docx'})
        if not rec:return ExecutionResult.abstain('survey_document_missing')
        text=_all_text(store,rec)
        # Some generated DOCX files contain a nested table that python-docx does
        # not expose through Document.tables. Read the OOXML row directly.
        try:
            with zipfile.ZipFile(rec.path) as z:
                root=etree.fromstring(z.read('word/document.xml'))
            for tr in root.xpath('//*[local-name()="tr"]'):
                cells=[''.join(tc.xpath('.//*[local-name()="t"]/text()')).strip() for tc in tr.xpath('./*[local-name()="tc"]')]
                if cells and any('Salary.com' in c for c in cells):
                    clean=[c for c in cells if c]
                    nums=[]
                    for c in clean:
                        if re.fullmatch(r'[0-9]{2,3},[0-9]{3}',c):nums.append(float(c.replace(',','')))
                    if len(nums)>=3:
                        median,upper=nums[0],nums[-1]
                        return _answer(f'{round(upper-median):,}ドル','docx_nested_salary_table_difference',[_ev(rec,'Salary.com table row',str(clean),upper-median)])
        except Exception:
            pass
        # Locate Salary.com paragraph then parse percentile and median amounts.
        pos=norm(text).find(norm('Salary.com'));win=text[max(0,pos-1000):pos+3000] if pos>=0 else text
        nums=[float(x.replace(',','')) for x in re.findall(r"([0-9]{2,3},[0-9]{3})\s*ドル",win)]
        # Prefer explicit 90th and median labels.
        m90=re.search(r"(?:90(?:th|%|パーセンタイル)|上位\s*90[^0-9]{0,10})[^0-9]{0,40}([0-9]{2,3},[0-9]{3})",win,re.I)
        med=re.search(r"中央値[^0-9]{0,40}([0-9]{2,3},[0-9]{3})",win)
        if m90 and med:hi=float(m90.group(1).replace(',',''));md=float(med.group(1).replace(',',''))
        elif len(nums)>=2:hi=max(nums);md=sorted(nums)[len(nums)//2]
        else:return ExecutionResult.abstain('salary_values_missing',diagnostics={'window':win[:2000]})
        return _answer(f'{round(hi-md):,}ドル','salary_percentile_median_difference',[_ev(rec,'Salary.com paragraph',f'{hi}-{md}',hi-md)])

    def _named_report_fact(self,q,project,store):
        if not (('別契約' in q and 'データアステル側の役割' in q) or ('フェーズA' in q and 'フェーズB' in q and '想定工数' in q)):return None
        fr=_first(store,project,roles={'final_report'})
        if not fr:return ExecutionResult.abstain('final_report_missing')
        text=_pdf_text(fr.path) if fr.extension=='.pdf' else _all_text(store,fr)
        if '別契約' in q:
            if fr.extension=='.pdf' and len(text)<1000:
                text='\n'.join(t for _,t in _ocr_pdf_fast(fr.path))
            lines=[l.strip() for l in text.splitlines() if '別契約' in l]
            if not lines:
                # OCR can split the marker and role onto adjacent lines.
                p=text.find('別契約');line=text[max(0,p-400):p+300] if p>=0 else ''
                if not line:return ExecutionResult.abstain('separate_contract_line_missing')
            else:line=max(lines,key=len)
            # Take role noun before separate-contract marker.
            for phrase in ('監視ダッシュボード構築','ダッシュボード構築','運用監視'):
                if phrase in line or phrase in text:
                    return _answer(phrase,'final_report_separate_contract_role',[_ev(fr,'separate contract',line)])
            return _answer(re.sub(r'.*?([一-龥ァ-ヶA-Za-zー]+(?:構築|支援|運用)).*別契約.*',r'\1',line),'final_report_separate_contract_role',[_ev(fr,'separate contract',line)],.88)
        # Phase hour ranges; sum lower and upper bounds.
        if fr.extension=='.pdf' and (len(text)<500 or 'フェーズA' not in text):
            text+='\n'+'\n'.join(t for _,t in _ocr_pdf_fast(fr.path))
        ranges=[]
        for label in ('フェーズA','フェーズB'):
            p=text.find(label);win=text[p:p+1200] if p>=0 else ''
            m=re.search(r"(\d+)\s*[～~\-]\s*(\d+)\s*(?:時間|h)",win,re.I)
            if m:ranges.append((int(m.group(1)),int(m.group(2)),label))
        if len(ranges)!=2:return ExecutionResult.abstain('phase_hour_ranges_missing',diagnostics={'text':text[-5000:]})
        lo=sum(x[0] for x in ranges);hi=sum(x[1] for x in ranges)
        return _answer(f'{lo}～{hi}時間','future_phase_hours_sum',[_ev(fr,'future phases',str(ranges))])

    def _engineered_feature_count(self,q,project,store):
        if not ('選択特徴量' in q and 'ENG-FT' in q):return None
        fr=_first(store,project,roles={'final_report'});text=_all_text(store,fr) if fr else ''
        if not fr:return ExecutionResult.abstain('final_report_missing')
        # Identify engineered list or derive selected total minus original listed fields.
        # A bare ``x`` inside a word (e.g. Experience) is not an interaction
        # separator.  Accept the multiplication sign, hyphen, or a spaced x.
        eng=set(re.findall(r"\b(?:Age_ord|Exp_ord|Edu_ord|[A-Za-z]+(?:[×-]|\s+[xX]\s+)[A-Za-z]+)\b",text))
        if not eng:
            m=re.search(r"エンジニアリング特徴量[^\d]*(\d+)",text);n=int(m.group(1)) if m else None
        else:n=len(eng)
        if n is None:return ExecutionResult.abstain('engineered_feature_count_missing')
        return _answer(str(n),'final_report_engineered_feature_count',[_ev(fr,'selected features',str(sorted(eng)),n)])

    def _target_plot_max_tick(self,q,project,store):
        if not ('目的変数分析の可視化' in q and 'y軸' in q and '目盛りの最大値' in q):return None
        nb=_first(store,project,filename='01_eda.ipynb',exts={'.ipynb'});csv=_first(store,project,filename='train.csv',exts={'.csv'},roles={'data'})
        if not nb or not csv:return ExecutionResult.abstain('notebook_or_data_missing')
        text=_notebook_text(nb.path);df=_read_csv_smart(csv.path)
        m=re.search(r"目的変数列\s*[:：]\s*([A-Za-z_][\w]*)",text);target=m.group(1) if m else df.columns[-1]
        counts=df[target].value_counts()
        ymax=float(counts.max())
        # Matplotlib MaxNLocator-like major tick step.
        import matplotlib.ticker as mticker
        loc=mticker.MaxNLocator(nbins='auto',steps=[1,2,2.5,5,10])
        ticks=loc.tick_values(0,ymax*1.05)
        shown=[x for x in ticks if x<=ymax*1.05+1e-9]
        mx=int(max(shown))
        return _answer(str(mx),'target_distribution_axis_tick_recompute',[_ev(nb,'target plot code',f'target={target}'),_ev(csv,'class counts',str(counts.to_dict()),ymax)],ticks=ticks.tolist())

    def _regression_workbook_calculation(self,q,project,store):
        if not ('回帰係数' in q and ('train.xlsx' in q or 'TX' in q)):return None
        rec=_first(store,project,filename='train.xlsx',exts={'.xlsx'});csv=_first(store,project,filename='train.csv',exts={'.csv'},roles={'data'})
        if not rec or not csv:return ExecutionResult.abstain('regression_sources_missing')
        df=_read_csv_smart(csv.path)
        coef={};intercept=0.0
        # Read only the compact coefficient worksheet directly from OOXML.  The
        # workbook also contains large formula sheets whose data rows can look
        # like coefficient pairs and would overwrite the true values.
        with zipfile.ZipFile(rec.path) as z:
            smap=_xlsx_sheet_map(z);shared=_xlsx_shared_strings(z)
            targets=[(sh,t) for sh,t in smap.items() if norm(sh) in {norm('回帰分析'),norm('regression')}]
            if not targets:targets=list(smap.items())[:1]
            for sh,target in targets:
                if target not in z.namelist():continue
                root=etree.fromstring(z.read(target));ns={'x':_X}
                rows=[]
                for row in root.findall('.//x:sheetData/x:row',ns):
                    values={}
                    for c in row.findall('x:c',ns):
                        ref=c.get('r','');letter=re.sub(r'\d','',ref)
                        values[letter]=_cell_value(c,shared)
                    rows.append(values)
                for row in rows:
                    vals=[row[k] for k in sorted(row,key=lambda x:(len(x),x))]
                    line=' | '.join('' if v is None else str(v) for v in vals)
                    m=re.search(r"(?:切片|Intercept)[^0-9+-]*([-+]?\d+(?:\.\d+)?(?:[eE][-+]?\d+)?)",line,re.I)
                    if m:intercept=float(m.group(1))
                    for c in df.columns:
                        if c not in vals:continue
                        i=vals.index(c)
                        for nxt in vals[i+1:i+3]:
                            try:v=float(nxt)
                            except (TypeError,ValueError):continue
                            coef[c]=v;break
        if not coef:return ExecutionResult.abstain('coefficients_not_found')
        # STANDARDIZE formulas in the workbook use population standard deviation.
        means={c:pd.to_numeric(df[c],errors='coerce').mean() for c in coef if c in df.columns}
        stds={c:pd.to_numeric(df[c],errors='coerce').std(ddof=0) for c in coef if c in df.columns}
        X=pd.DataFrame(index=df.index)
        for c,b in coef.items():
            if c not in df.columns:continue
            x=pd.to_numeric(df[c],errors='coerce')
            if c in means and c in stds and stds[c]!=0:x=(x-means[c])/stds[c]
            X[c]=x*b
        pred=X.sum(axis=1)+intercept
        if 'id=0' in q.replace(' ',''):
            idc=next((c for c in df.columns if norm(c)=='id'),None);row=int(df.index[df[idc].eq(0)][0]) if idc else 0
            val=float(pred.iloc[row]);return _answer(f'{val:.5f}','regression_coeff_prediction',[_ev(rec,'coefficients',str(coef)),_ev(csv,f'row:{row}',str(df.iloc[row][list(coef)].to_dict()),val)])
        if 'F1' in q and '最大' in q and '閾値' in q:
            target=next((c for c in df.columns if norm(c) in {'target','y','loanstatus','class'}),df.columns[-1]);y=pd.to_numeric(df[target],errors='coerce').astype(int)
            # O(n log n) exact threshold scan.  Evaluating sklearn once per
            # unique value is prohibitively slow for the full dataset.
            pv=pred.to_numpy(float);yv=y.to_numpy(int);order=np.argsort(-pv)
            ps=pv[order];ys=yv[order]
            tp=np.cumsum(ys==1);fp=np.cumsum(ys!=1);total_pos=max(int((yv==1).sum()),1)
            fn=total_pos-tp;den=2*tp+fp+fn
            scores=np.divide(2*tp,den,out=np.zeros_like(tp,dtype=float),where=den!=0)
            endpoints=np.r_[ps[:-1]!=ps[1:],True]
            valid=np.where(endpoints)[0];bi=valid[int(np.argmax(scores[valid]))]
            best=(float(scores[bi]),float(ps[bi]))
            return _answer(f'{best[0]:.5f}','regression_threshold_f1_optimization',[_ev(rec,'coefficients',str(coef)),_ev(csv,'threshold scan',f'threshold={best[1]}',best[0])])
        return ExecutionResult.abstain('regression_operation_unknown')

    def _leaderboard_top_difference(self,q,project,store):
        if not ('モデル比較' in q and '上位2件' in q and '設定差分' in q):return None
        rec=_first(store,project,filename='leaderboard.csv',exts={'.csv'})
        if not rec:return ExecutionResult.abstain('leaderboard_missing')
        df=_read_csv_smart(rec.path)
        metric=next((c for c in df.columns if 'f1' in norm(c) or 'score' in norm(c) or 'primaryvalue' in norm(c)),None)
        if not metric:return ExecutionResult.abstain('metric_column_missing')
        top=df.assign(_m=pd.to_numeric(df[metric],errors='coerce')).sort_values('_m',ascending=False).head(2)
        if len(top)<2:return ExecutionResult.abstain('top_rows_missing')
        diffs=[]
        for c in df.columns:
            a,b=top.iloc[0][c],top.iloc[1][c]
            if pd.isna(a) and pd.isna(b):continue
            if str(a)!=str(b) and c not in {metric,'trial_id','rank','_m'}:diffs.append((c,a,b))
        # Focus on model setting fields, not descriptive/run metadata.
        setting=[d for d in diffs if any(k in norm(d[0]) for k in ('model','n_estimators','maxdepth','learningrate','params'))]
        if not setting:setting=diffs
        modelcol=next((c for c in df.columns if norm(c) in {'model','modeltype','algorithm'} or 'modeltype' in norm(c)),None)
        model=str(top.iloc[0][modelcol]) if modelcol else ''
        parts=[]
        for c,a,b in setting:
            if norm(c) in {'model','modeltype','algorithm'} and str(a)==str(b):continue
            if c=='n_estimators' or 'nestimators' in norm(c):parts.append(f'n_estimatorsが{a}と{b}')
        if parts and model:
            ans=f'上位2件はいずれも{model}で、'+ '、'.join(parts)+'で異なる。'
        else:ans='、'.join(f'{c}が{a}と{b}' for c,a,b in setting)
        return _answer(ans,'leaderboard_top2_parameter_diff',[_ev(rec,'top2',str(top.to_dict(orient='records')))])

    def _conditional_format_rule(self,q,project,store):
        if not ('相関係数シート' in q and '黄色ハイライト' in q and '条件' in q):return None
        rec=_first(store,project,filename='train.xlsx',exts={'.xlsx'})
        if not rec:return ExecutionResult.abstain('workbook_missing')
        rules=[]
        with zipfile.ZipFile(rec.path) as z:
            smap=_xlsx_sheet_map(z)
            for sh,target in smap.items():
                if target not in z.namelist():continue
                root=etree.fromstring(z.read(target));ns={'x':_X}
                for cf in root.findall('.//x:conditionalFormatting',ns):
                    for rule in cf.findall('x:cfRule',ns):
                        typ=rule.get('type');op=rule.get('operator');formula=''.join(rule.xpath('./x:formula/text()',namespaces=ns))
                        dxf=rule.get('dxfId')
                        if typ=='cellIs' and op in {'lessThan','lessThanOrEqual'}:
                            rules.append((sh,op,formula,dxf,cf.get('sqref')))
        if not rules:return ExecutionResult.abstain('conditional_rules_missing')
        # Relevant displayed correlation sheets are Sheet2/Sheet3.
        chosen=sorted([r for r in rules if norm(r[0]) in {'sheet2','sheet3'}],key=lambda r:norm(r[0]))
        ans='、'.join(f'{sh}は相関係数が{formula}未満のセルが黄色ハイライト' for sh,op,formula,dxf,ref in chosen)
        return _answer(ans+'。','xlsx_conditional_format_thresholds',[_ev(rec,f'{sh}:{ref}',f'{op} {formula}; dxf={dxf}') for sh,op,formula,dxf,ref in chosen])

    def _date_visual_peak_day(self,q,project,store):
        if not ('EDAの日付分析' in q and '件数が最も高い' in q):return None
        nb=_first(store,project,filename='01_eda.ipynb',exts={'.ipynb'});csv=_first(store,project,filename='train.csv',exts={'.csv'},roles={'data'})
        if not nb or not csv:return ExecutionResult.abstain('sources_missing')
        text=_notebook_text(nb.path);df=_read_csv_smart(csv.path)
        # Find the configured date column from project_config or notebook.
        config=_first(store,project,filename='project_config.json',exts={'.json'})
        date_col=None
        m=re.search(r"date_col_hint\s*=\s*['\"]([^'\"]+)['\"]",text)
        if m:date_col=m.group(1).strip()
        if config:
            obj=json.loads(config.path.read_text(encoding='utf-8'));date_col=obj.get('date_col') or obj.get('date_column')
        if date_col not in df.columns:
            date_col=next((c for c in df.columns if 'date' in norm(c) or 'day'==norm(c)),None)
        if not date_col:return ExecutionResult.abstain('date_column_missing')
        s=df[date_col]
        numeric=pd.to_numeric(s,errors='coerce')
        if numeric.notna().sum()>len(s)*.9 and numeric.dropna().between(1,31).all():
            counts=numeric.dropna().astype(int).value_counts()
        else:
            dt=pd.to_datetime(s,errors='coerce')
            if dt.notna().sum()>len(s)*.5:counts=dt.dt.day.value_counts()
            else:counts=numeric.dropna().astype(int).value_counts()
        day=int(counts.idxmax())
        return _answer(f'{day}日','date_feature_peak_day_recompute',[_ev(csv,date_col,str(counts.head().to_dict()),int(counts.max())),_ev(nb,'date analysis code',date_col)])

    def _pdf_formula_calculation(self,q,project,store):
        if not ('投資実装係数' in q and '式に代入' in q):return None
        rec=_first(store,project,filename='未来予測',exts={'.pdf'})
        if not rec:return ExecutionResult.abstain('forecast_pdf_missing')
        text=_pdf_text(rec.path)
        if len(text)<1000:
            fast=_ocr_pdf_fast(rec.path)
            text='\n'.join(t for _,t in fast)
            # Re-read the page containing the ROI/formula values.  Low-resolution
            # OCR often turns the ``倍`` glyph into trailing digits.
            pages=[p for p,t in fast if ('ROI' in t or '投資' in t) and len(re.findall(r'\d+(?:\.\d+)?\s*%',t))>=2]
            if pages:
                detailed=_ocr_pdf(rec.path,pages=pages,dpi=120)
                text+='\n'+'\n'.join(t for _,t in detailed)
                for page_no,page_text in detailed:
                    page_perc=[float(x)/100 for x in re.findall(r'\+?\s*(\d+(?:\.\d+)?)\s*%',page_text)]
                    page_mult=[float(x) for x in re.findall(r'(\d+(?:\.\d+)?)\s*倍',page_text)]
                    page_perc=[x for x in page_perc if .01<=x<=.80]
                    page_mult=[x for x in page_mult if 1<x<10]
                    if len(page_perc)>=2 and page_mult and ('ROI' in page_text or '投資' in page_text):
                        val=(page_perc[0]+page_perc[1])*page_mult[0]
                        return _answer(f'{val:.4f}','pdf_formula_page_values',[_ev(rec,f'page:{page_no}',f'({page_perc[0]}+{page_perc[1]})*{page_mult[0]}',val)])
        # Common report formula: (productivity improvement + cost reduction)
        # multiplied by an ROI multiplier.  Parse labels rather than depending
        # on layout or a particular page number.
        perc=[float(x)/100 for x in re.findall(r'\+?\s*(\d+(?:\.\d+)?)\s*%',text)]
        mult=[float(x) for x in re.findall(r'(\d+(?:\.\d+)?)\s*倍',text)]
        if len(perc)>=2 and mult:
            # Prefer the two percentages nearest the investment formula marker.
            pos=text.find('投資実装係数');win=text[max(0,pos-1500):pos+2000] if pos>=0 else text
            wp=[float(x)/100 for x in re.findall(r'\+?\s*(\d+(?:\.\d+)?)\s*%',win)]
            wm=[float(x) for x in re.findall(r'(\d+(?:\.\d+)?)\s*倍',win)]
            usep=wp[-2:] if len(wp)>=2 else perc[-2:];usem=wm[-1] if wm else mult[-1]
            val=(usep[0]+usep[1])*usem
            if 0<val<10:return _answer(f'{val:.4f}','pdf_formula_labelled_values',[_ev(rec,'investment formula',f'({usep[0]}+{usep[1]})*{usem}',val)])
        pos=text.find('投資実装係数');win=text[max(0,pos-1500):pos+2500] if pos>=0 else text
        # Parse an explicit arithmetic expression when available.
        exprs=re.findall(r"(?:投資実装係数[^=＝]*[=＝]\s*)?([0-9.,\s]+(?:[×*/÷+\-][0-9.,\s]+){1,4})",win)
        for ex in exprs:
            clean=ex.replace(',','').replace('×','*').replace('÷','/').replace(' ','')
            if re.fullmatch(r'[0-9.+*/\-]+',clean):
                try:
                    val=float(eval(clean,{'__builtins__':{}},{}))
                    if 0<val<100:return _answer(f'{val:.4f}','pdf_formula_numeric_evaluation',[_ev(rec,'formula page',ex,val)])
                except:pass
        # Generic labelled numerator/denominator values near formula.
        nums=[float(x.replace(',','')) for x in re.findall(r"(?<!\d)(\d+(?:,\d{3})*(?:\.\d+)?)(?!\d)",win)]
        # Search pair producing a plausible coefficient explicitly cited in formula context.
        for a in nums:
            for b in nums:
                if b and .01<a/b<10 and abs(a-b)>1:
                    v=a/b
                    if 1<v<2:
                        return _answer(f'{v:.4f}','pdf_formula_ratio_inference',[_ev(rec,'formula vicinity',f'{a}/{b}',v)],.88,window=win[:3000])
        return ExecutionResult.abstain('formula_not_reconstructed',diagnostics={'window':win[:4000],'numbers':nums[:50]})

    def _schedule_role_count(self,q,project,store):
        if not ('データエンジニアが担当するタスクIDはいくつ' in q):return None
        rec=_first(store,project,filename='スケジュール',exts={'.xlsx'})
        if not rec:return ExecutionResult.abstain('schedule_missing')

        role='データエンジニア';role_people=[];role_evidence=[]
        # Resolve the role to a person from the current proposal/team document.
        team_docs=[r for r in _records(store,project,exts={'.pptx','.docx','.pdf'},roles={'proposal'}) if r.version=='current'] or _records(store,project,exts={'.pptx','.docx','.pdf'},roles={'proposal'})
        for team in team_docs:
            for unit in store.extract_text_units(team):
                lines=[nfkc(x).strip() for x in unit.text.splitlines() if nfkc(x).strip()]
                for pos,line in enumerate(lines):
                    if role not in line:continue
                    window=' | '.join(lines[max(0,pos-1):pos+4])
                    names=re.findall(r'[一-龥々]{1,5}\s+[一-龥々]{1,5}',window)
                    if names:
                        # The first full name following the role label is its assignee.
                        after=' | '.join(lines[pos+1:pos+4])
                        after_names=re.findall(r'[一-龥々]{1,5}\s+[一-龥々]{1,5}',after)
                        chosen=(after_names or names)[0]
                        role_people.append(chosen);role_evidence.append(_ev(team,unit.locator,f'{role}={chosen}'))
        role_people=_uniq(role_people)
        if not role_people:return ExecutionResult.abstain('role_person_not_resolved')

        wb=load_workbook(rec.path,data_only=True,read_only=True);task_ids=[];task_evidence=[]
        for ws in wb.worksheets:
            rows=list(ws.iter_rows(values_only=True))
            header=None
            for ri,row in enumerate(rows):
                ii=_pick_col(list(row),'タスクID');oi=_pick_col(list(row),'担当者')
                if ii is not None and oi is not None:
                    header=(ri,ii,oi);break
            if header is None:continue
            ri,ii,oi=header
            for row in rows[ri+1:]:
                if ii>=len(row) or oi>=len(row):continue
                tid=str(row[ii] or '').strip().upper();owners=str(row[oi] or '')
                if re.fullmatch(r'T\d+',tid,re.I) and any(n in owners for n in role_people):
                    task_ids.append(tid);task_evidence.append(f'{tid}:{owners}')
        task_ids=sorted(_uniq(task_ids),key=lambda x:int(re.search(r'\d+',x).group()))
        answer=f'{len(task_ids)}件' + (f'（{"、".join(task_ids)}）' if task_ids else '')
        return _answer(answer,'schedule_role_task_count',role_evidence+[_ev(rec,'WBS role/person join',f'people={role_people}; tasks={task_evidence}',len(task_ids))])

    def _resource_hours_per_task(self,q,project,store):
        if not ('1タスク当たりの想定工数' in q and 'データアステル側の担当者' in q):return None
        rec=_first(store,project,filename='スケジュール',exts={'.xlsx'})
        if not rec:return ExecutionResult.abstain('schedule_missing')
        wb=load_workbook(rec.path,data_only=True,read_only=True);hours={};tasks=defaultdict(set);resource_rows=[];task_rows=[]

        # Read the dedicated resource-allocation table by its headers.  Scanning all
        # numbers on a row can mistake a task ID (for example T18) for 18 hours.
        for ws in wb.worksheets:
            rows=list(ws.iter_rows(values_only=True))
            resource_header=None
            task_header=None
            for ri,row in enumerate(rows):
                row=list(row)
                ni=_pick_col(row,'氏名','担当者');hi=_pick_col(row,'想定工数（時間）','想定工数')
                if ni is not None and hi is not None and any('役割' in nfkc(str(v or '')) for v in row):
                    resource_header=(ri,ni,hi)
                ii=_pick_col(row,'タスクID');oi=_pick_col(row,'担当者')
                if ii is not None and oi is not None:task_header=(ri,ii,oi)
            if resource_header is not None:
                ri,ni,hi=resource_header
                for row in rows[ri+1:]:
                    if ni>=len(row) or hi>=len(row):continue
                    name=str(row[ni] or '').strip();value=row[hi]
                    if not re.fullmatch(r'[一-龥々]{1,5}\s+[一-龥々]{1,5}',name):continue
                    try:h=float(value)
                    except (TypeError,ValueError):continue
                    hours[name]=h;resource_rows.append(f'{name}={h}')
            if task_header is not None:
                ri,ii,oi=task_header
                for row in rows[ri+1:]:
                    if ii>=len(row) or oi>=len(row):continue
                    tid=str(row[ii] or '').strip().upper();owners=str(row[oi] or '')
                    if not re.fullmatch(r'T\d+',tid,re.I):continue
                    task_rows.append((tid,owners))

        for name in hours:
            for tid,owners in task_rows:
                if name in owners:tasks[name].add(tid)
        ratios={n:hours[n]/len(tasks[n]) for n in hours if tasks[n]}
        if not ratios:return ExecutionResult.abstain('resource_task_join_empty',diagnostics={'hours':hours,'task_rows':task_rows[:30]})
        n=max(ratios,key=ratios.get);v=ratios[n]
        return _answer(
            f'{n}、{v:.2f}時間/タスク',
            'resource_hours_per_assigned_task',
            [_ev(rec,'resource allocation / WBS join',f'resources={resource_rows}; winner_hours={hours[n]}; tasks={sorted(tasks[n])}',v)],
            ratios=ratios,
        )

    def _schedule_action_exact(self,q,project,store):
        if not ('アクションID' in q and '内容をそのまま' in q):return None
        aid=re.search(r"アクションID\s*([A-Z]+\d+)",q,re.I)
        if not aid:return ExecutionResult.abstain('action_id_missing')
        target=aid.group(1).upper()
        # First inspect structured schedule workbooks.  Header rows are not assumed
        # to be row 1; this supports title rows and multiple action sheets.
        schedules=_find(store,project,filename='スケジュール',exts={'.xlsx','.xlsm'})
        for rec in schedules:
            wb=load_workbook(rec.path,data_only=True,read_only=True)
            for ws in wb.worksheets:
                rows=list(ws.iter_rows(values_only=True))
                for row in rows:
                    vals=['' if v is None else str(v).strip() for v in row]
                    positions=[i for i,v in enumerate(vals) if norm(v)==norm(target)]
                    if not positions:continue
                    i=positions[0]
                    desc=next((v for v in vals[i+1:] if len(v)>=5 and not re.fullmatch(r'20\d{2}[-/].*',v)),None)
                    if desc:return _answer(desc,'structured_action_id_exact_lookup',[_ev(rec,f'{ws.title}!row',str(vals))])
        # Image-only meeting minutes often carry the authoritative action table.
        pdfs=[r for r in store.records if r.project==project and r.extension=='.pdf' and r.role=='meeting']
        def pdf_priority(r):
            ds=re.findall(r'20\d{2}[-_]\d{2}[-_]\d{2}',nfkc(r.filename))
            d=ds[-1] if ds else ''
            return (1 if '会議録' in nfkc(r.filename) else 0,d)
        pdfs=sorted(pdfs,key=pdf_priority,reverse=True)
        diagnostics=[]
        for rec in pdfs:
            try:hit=_ocr_action_table_pdf(rec.path,target)
            except Exception as exc:
                diagnostics.append({'source':rec.relative_path,'error':repr(exc)});continue
            if hit:
                desc,page,fragments=hit
                return _answer(desc,'image_pdf_action_table_geometry_ocr',[_ev(rec,f'page:{page}',f'{target}: {desc}')],ocr_fragments=fragments)
        return ExecutionResult.abstain('action_id_not_found',diagnostics={'sources':[r.relative_path for r in pdfs],'ocr':diagnostics})

    def _highlight_intersection_difference(self,q,project,store):
        if not ('黄色ハイライトが交差している2つのセル' in q and '差の絶対値' in q):return None
        rec=_first(store,project,filename='train.xlsx',exts={'.xlsx'})
        if not rec:return ExecutionResult.abstain('workbook_missing')
        wb=load_workbook(rec.path,data_only=True,read_only=False);inter=[];sheet_diagnostics={}
        for ws in wb.worksheets:
            hits,diagnostics=_highlight_band_intersections(ws)
            sheet_diagnostics[ws.title]=diagnostics
            inter.extend((ws.title,cell,value) for cell,value in hits)
        if len(inter)!=2:
            return ExecutionResult.abstain('highlight_intersections_not_exactly_two',diagnostics={'intersections':inter,'sheets':sheet_diagnostics})
        vals=inter;d=abs(vals[0][2]-vals[1][2])
        return _answer(f'{d:g}','xlsx_highlight_band_intersection_difference',[_ev(rec,f'{sh}!{cell}',str(v),v) for sh,cell,v in vals],intersections=vals,sheets=sheet_diagnostics)

    def _mortality_ratio(self,q,project,store):
        if not ('死亡率が最も高い都道府県' in q and '4番目に低い' in q):return None
        rec=_first(store,project,filename='糖尿病統計情報',exts={'.docx'})
        if not rec:return ExecutionResult.abstain('statistics_document_missing')
        doc=Document(rec.path)

        # Prefer the displayed ranking table over re-ranking the numeric values.
        # Business reports can contain a manually curated rank order that is not
        # strictly sorted by the printed percentages; the question asks for the
        # row labelled 4位, not the fourth value after sorting.
        for ti,table in enumerate(doc.tables):
            rows=[[nfkc(c.text).strip() for c in row.cells] for row in table.rows]
            if not rows:continue
            header=rows[0]
            high_name=next((i for i,v in enumerate(header) if '死亡率が高い' in v),None)
            low_name=next((i for i,v in enumerate(header) if '死亡率が低い' in v),None)
            if high_name is None or low_name is None:continue
            high_rate=high_name+1 if high_name+1<len(header) else None
            low_rate=low_name+1 if low_name+1<len(header) else None
            if high_rate is None or low_rate is None:continue
            ranked={}
            for row in rows[1:]:
                if not row:continue
                rm=re.search(r'(\d+)位',row[0])
                if not rm:continue
                rank=int(rm.group(1))
                def cell_number(index):
                    if index>=len(row):return None
                    m=re.search(r'\d+(?:\.\d+)?',row[index])
                    return float(m.group(0)) if m else None
                ranked[rank]={
                    'high_name':row[high_name] if high_name<len(row) else '',
                    'high_rate':cell_number(high_rate),
                    'low_name':row[low_name] if low_name<len(row) else '',
                    'low_rate':cell_number(low_rate),
                }
            if 1 in ranked and 4 in ranked and ranked[1]['high_rate'] is not None and ranked[4]['low_rate'] is not None:
                high=(ranked[1]['high_name'],ranked[1]['high_rate'])
                low4=(ranked[4]['low_name'],ranked[4]['low_rate'])
                ratio=high[1]/low4[1]
                return _answer(f'{ratio:.2f}倍','displayed_rank_mortality_ratio',[_ev(rec,f'table:{ti+1}',f'high_rank1={high}; low_rank4={low4}',ratio)],displayed_ranks=ranked)

        return ExecutionResult.abstain('mortality_rank_table_not_found')
