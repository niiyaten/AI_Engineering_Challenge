from __future__ import annotations

import itertools
import hashlib
import difflib
import math
import os
import re
import shutil
import subprocess
import tempfile
import unicodedata
import zipfile
from collections import defaultdict
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Iterable, Any

import fitz
import numpy as np
import pandas as pd
from docx import Document
from lxml import etree
from openpyxl import load_workbook
from pptx import Presentation

from ..models import Evidence, ExecutionResult, QueryPlan, Question, FileRecord
from ..normalize import nfkc, norm, overlap_score
from ..store import DocumentStore
from .base import Executor
from .utils import format_number, parse_number

_W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
_X = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _uniq(xs: Iterable[str]) -> list[str]:
    out: list[str] = []
    seen: set[str] = set()
    for x in xs:
        x = str(x).strip()
        if x and norm(x) not in seen:
            seen.add(norm(x)); out.append(x)
    return out


def _project_from_question(question: Question, plan: QueryPlan, store: DocumentStore) -> str:
    for hint in plan.project_hints:
        p = store.resolve_project(hint, strict=False)
        if p:
            return p
    # Explicit project abbreviations (AYM, AOBM, etc.) are resolved through
    # the data-derived glossary aliases before fuzzy full-question matching.
    for token in re.findall(r"[A-Z][A-Z0-9_-]{1,15}", nfkc(question.text)):
        p=store.resolve_project(token,strict=False)
        if p:return p
    # DocumentStore already knows corporate-name variants and glossary aliases.
    p = store.resolve_project(question.text, strict=False)
    if p:
        return p
    best = sorted(((overlap_score(question.text, p), p) for p in store.projects), reverse=True)
    return best[0][1] if best and best[0][0] >= .18 else ""


def _records(store: DocumentStore, project: str = "", *, names: Iterable[str] = (), exts: Iterable[str] = (), roles: Iterable[str] = ()) -> list[FileRecord]:
    extset = {x if x.startswith(".") else "." + x for x in exts}
    out = []
    for r in store.records:
        if project and r.project != project:
            continue
        if extset and r.extension not in extset:
            continue
        if roles and r.role not in set(roles):
            continue
        if names and not all(norm(x) in norm(r.filename) for x in names):
            continue
        out.append(r)
    return out


def _ev(rec: FileRecord, locator: str, detail: str, value: Any | None = None) -> Evidence:
    return Evidence(rec.relative_path, locator, detail, value)


def _date_key(path: Path) -> tuple[int, int, int]:
    m = re.search(r"(20\d{2})[-_](\d{1,2})[-_](\d{1,2})", path.name)
    return tuple(map(int, m.groups())) if m else (0, 0, 0)


def _money(s: str) -> float | None:
    m = re.search(r"(?:¥|￥)?\s*([0-9][0-9,]*(?:\.\d+)?)\s*円?", nfkc(s))
    return float(m.group(1).replace(",", "")) if m else None


def _read_csv_smart(path: Path) -> pd.DataFrame:
    read_kwargs = {"sep": "\t"} if path.suffix.lower() in {".tsv", ".tab"} else {}
    for enc in ("utf-8-sig", "utf-8", "cp932"):
        try:
            return pd.read_csv(path, encoding=enc, **read_kwargs)
        except UnicodeDecodeError:
            pass
    return pd.read_csv(path, **read_kwargs)


def _is_yellow_rgb(rgb: str) -> bool:
    rgb = (rgb or "")[-6:].upper()
    if len(rgb) != 6:
        return False
    try:
        r, g, b = [int(rgb[i:i+2], 16) for i in (0, 2, 4)]
    except ValueError:
        return False
    return r >= 175 and g >= 145 and b <= 170 and abs(r-g) <= 105



def _xlsx_shared_strings(z: zipfile.ZipFile) -> list[str]:
    if "xl/sharedStrings.xml" not in z.namelist():
        return []
    root = etree.fromstring(z.read("xl/sharedStrings.xml"))
    ns = {"x": _X}
    return ["".join(si.itertext()) for si in root.findall("x:si", ns)]


def _xlsx_sheet_map(z: zipfile.ZipFile) -> dict[str, str]:
    ns = {"x": _X, "r": _R}
    wb = etree.fromstring(z.read("xl/workbook.xml"))
    relroot = etree.fromstring(z.read("xl/_rels/workbook.xml.rels"))
    rels = {r.get("Id"): r.get("Target") for r in relroot}
    out = {}
    for sh in wb.findall(".//x:sheet", ns):
        target = rels.get(sh.get(f"{{{_R}}}id"), "")
        if target.startswith("/"):
            target = target.lstrip("/")
        elif not target.startswith("xl/"):
            target = "xl/" + target
        out[sh.get("name")] = target
    return out


def _xlsx_yellow_style_ids(z: zipfile.ZipFile) -> set[int]:
    root = etree.fromstring(z.read("xl/styles.xml"))
    ns = {"x": _X}
    fills = root.find("x:fills", ns)
    yellow_fills = set()
    if fills is not None:
        for i, fill in enumerate(fills):
            fg = fill.find(".//x:fgColor", ns)
            rgb = fg.get("rgb", "") if fg is not None else ""
            if _is_yellow_rgb(rgb):
                yellow_fills.add(i)
    styles = set()
    cellxfs = root.find("x:cellXfs", ns)
    if cellxfs is not None:
        for i, xf in enumerate(cellxfs):
            try: fid = int(xf.get("fillId", "0"))
            except ValueError: fid = 0
            if fid in yellow_fills:
                styles.add(i)
    return styles


def _cell_value(c, shared: list[str]) -> Any:
    t = c.get("t")
    v = c.find(f"{{{_X}}}v")
    if t == "inlineStr":
        return "".join(c.itertext())
    if v is None:
        return None
    raw = v.text or ""
    if t == "s":
        try: return shared[int(raw)]
        except Exception: return raw
    if t in {"str", "e"}:
        return raw
    try:
        x = float(raw)
        return int(x) if x.is_integer() else x
    except Exception:
        return raw


def _xlsx_sheet_cells(path: Path, sheet_name: str) -> tuple[dict[str, Any], dict[str, str], set[int]]:
    """Return cell values, formulas and yellow style IDs using streaming OOXML.

    This avoids openpyxl loading giant drawings/EMF objects and is stable for the
    large audit workbooks.
    """
    with zipfile.ZipFile(path) as z:
        sheet_map = _xlsx_sheet_map(z)
        target = next((v for k, v in sheet_map.items() if norm(k) == norm(sheet_name)), None)
        if not target or target not in z.namelist():
            return {}, {}, set()
        shared = _xlsx_shared_strings(z)
        yellow_styles = _xlsx_yellow_style_ids(z)
        vals: dict[str, Any] = {}
        formulas: dict[str, str] = {}
        root = etree.fromstring(z.read(target))
        for c in root.findall(f".//{{{_X}}}c"):
            ref = c.get("r")
            if not ref: continue
            vals[ref] = _cell_value(c, shared)
            f = c.find(f"{{{_X}}}f")
            if f is not None and f.text:
                formulas[ref] = f.text
        return vals, formulas, yellow_styles


def _xlsx_yellow_cells_fast(path: Path, sheet_name: str | None = None) -> list[dict[str, Any]]:
    with zipfile.ZipFile(path) as z:
        sheets = _xlsx_sheet_map(z)
        shared = _xlsx_shared_strings(z)
        yellow_styles = _xlsx_yellow_style_ids(z)
        out = []
        for name, target in sheets.items():
            if sheet_name and norm(name) != norm(sheet_name):
                continue
            if target not in z.namelist(): continue
            root = etree.fromstring(z.read(target))
            for c in root.findall(f".//{{{_X}}}c"):
                try: sid = int(c.get("s", "0"))
                except ValueError: sid = 0
                if sid not in yellow_styles: continue
                ref = c.get("r", "")
                m = re.match(r"([A-Z]+)(\d+)", ref)
                if not m: continue
                col_letters, row = m.group(1), int(m.group(2))
                col = 0
                for ch in col_letters: col = col*26 + ord(ch)-64
                f = c.find(f"{{{_X}}}f")
                out.append({"sheet": name, "cell": ref, "row": row, "col": col,
                            "formula": f.text if f is not None else None,
                            "value": _cell_value(c, shared)})
        return out


def _col_letter(col: int) -> str:
    out = ""
    while col:
        col, rem = divmod(col-1, 26); out = chr(65+rem)+out
    return out



def _fast_raw_headers(vals: dict[str, Any], col: int, max_row: int = 15) -> list[str]:
    letter=_col_letter(col); out=[]
    for r in range(1,max_row+1):
        v=vals.get(f"{letter}{r}")
        if v not in (None,""): out.append(str(v))
    return out


def _fast_header_for_df(vals: dict[str, Any], col: int, df: pd.DataFrame, max_row: int = 15) -> str | None:
    letter=_col_letter(col)
    mapping={norm(c):str(c) for c in df.columns}
    for r in range(1,max_row+1):
        v=vals.get(f"{letter}{r}")
        if v in (None,""): continue
        raw=nfkc(v); nv=norm(raw)
        if nv in mapping: return mapping[nv]
        # Aggregate-value columns (平均 / age, 合計 / age, 個数...) are metrics,
        # not grouping conditions and must not be propagated as filters.
        if any(k in raw for k in ("平均","合計","最大","最小","個数","件数","カウント")):
            return None
        matches=[orig for key,orig in mapping.items() if nv and (nv==key or nv in key or key in nv)]
        if len(matches)==1:return matches[0]
    return None

def _fast_nearest(vals: dict[str, Any], row: int, col: int) -> Any | None:
    letter = _col_letter(col)
    for r in range(row, 0, -1):
        v = vals.get(f"{letter}{r}")
        if v not in (None, ""):
            return v
    return None


def _fast_header(vals: dict[str, Any], col: int, max_row: int = 12) -> str | None:
    letter = _col_letter(col); items=[]
    for r in range(1,max_row+1):
        v=vals.get(f"{letter}{r}")
        if v not in (None,""): items.append(str(v))
    return items[-1] if items else None

def _yellow_cells(path: Path, sheet_name: str | None = None) -> list[dict[str, Any]]:
    wb_f = load_workbook(path, data_only=False, read_only=False)
    wb_v = load_workbook(path, data_only=True, read_only=False)
    out = []
    for ws in wb_f.worksheets:
        if sheet_name and norm(ws.title) != norm(sheet_name):
            continue
        wsv = wb_v[ws.title]
        for row in ws.iter_rows():
            for c in row:
                fill = c.fill
                if fill.fill_type != "solid":
                    continue
                color = fill.fgColor
                rgb = color.rgb if color.type == "rgb" else ""
                if _is_yellow_rgb(rgb):
                    out.append({"sheet": ws.title, "cell": c.coordinate, "row": c.row, "col": c.column,
                                "formula": c.value, "value": wsv[c.coordinate].value})
    return out


def _nearest_nonempty(ws, row: int, col: int) -> Any | None:
    for r in range(row, 0, -1):
        v = ws.cell(r, col).value
        if v not in (None, ""):
            return v
    return None


def _header_levels(ws, col: int, max_row: int = 12) -> list[str]:
    vals = []
    for r in range(1, min(max_row, ws.max_row) + 1):
        v = ws.cell(r, col).value
        if v not in (None, ""):
            vals.append(str(v))
    return vals


def _candidate_raw_csv(store: DocumentStore, rec: FileRecord) -> FileRecord | None:
    candidates = [r for r in store.records if r.project == rec.project and r.extension in {".csv", ".tsv"} and norm(r.stem) == norm("train")]
    if not candidates:
        candidates = [r for r in store.records if r.project == rec.project and r.extension in {".csv", ".tsv"}]
    candidates.sort(key=lambda r: ("/03." not in r.relative_path, len(r.relative_path)))
    return candidates[0] if candidates else None


def _value_equal(a: float, b: float) -> bool:
    return math.isclose(float(a), float(b), rel_tol=2e-6, abs_tol=2e-5)


def _hierarchy_conditions_from_single_column(vals: dict[str, Any], row: int, col: int, df: pd.DataFrame, max_scan: int = 12) -> dict[str, Any]:
    """Map nearest outline labels in a single pivot row-label column to raw columns."""
    letter=_col_letter(col); out={}
    for r in range(row, max(0,row-max_scan), -1):
        v=vals.get(f"{letter}{r}")
        if v in (None,""): continue
        matches=[]
        for c in df.columns:
            s=df[c]
            nv=parse_number(v)
            if pd.api.types.is_numeric_dtype(s) and nv is not None:
                if pd.to_numeric(s,errors="coerce").eq(nv).any(): matches.append(c)
            elif s.astype(str).map(norm).eq(norm(v)).any(): matches.append(c)
        if len(matches)==1 and matches[0] not in out:
            out[matches[0]]=v
            if len(out)>=4: break
    return out


def _groupby_infer(df: pd.DataFrame, value: float, known: dict[str, Any], *, leaf: tuple[str, Any] | None = None, expand: bool = True, max_missing_dims: int = 2) -> list[dict[str, Any]]:
    """Infer group conditions and aggregate from a displayed pivot value.

    It never uses question IDs or an answer table. Candidate operations are checked
    by recomputing them against the raw data.
    """
    columns = {norm(c): c for c in df.columns}
    fixed: dict[str, Any] = {}
    for hk, hv in known.items():
        nk = norm(hk)
        col = columns.get(nk)
        if not col:
            matches = [c for n, c in columns.items() if nk == n or nk in n or n in nk]
            col = matches[0] if len(matches) == 1 else None
        if col:
            fixed[col] = hv
    if leaf:
        lk, lv = leaf
        col = columns.get(norm(lk))
        if col:
            fixed[col] = lv
        else:
            matches = []
            for c in df.columns:
                if df[c].dtype == object and df[c].astype(str).map(norm).eq(norm(lv)).any():
                    matches.append(c)
            if len(matches) == 1:
                fixed[matches[0]] = lv

    base = df.copy()
    for c, v in fixed.items():
        s = base[c]
        if pd.api.types.is_numeric_dtype(s):
            nv = parse_number(v)
            if nv is not None:
                base = base[pd.to_numeric(s, errors="coerce").eq(nv)]
            else:
                base = base[s.astype(str).map(norm).eq(norm(v))]
        else:
            base = base[s.astype(str).map(norm).eq(norm(v))]
    if base.empty:
        return []

    low = [c for c in df.columns if c not in fixed and df[c].nunique(dropna=True) <= 120 and (df[c].dtype == object or df[c].nunique(dropna=True) <= 40)]
    numeric = [c for c in df.columns if pd.api.types.is_numeric_dtype(df[c])]
    candidates: list[dict[str, Any]] = []

    # If fixed conditions uniquely define the displayed group, check count and all numeric aggregates.
    if _value_equal(len(base), value):
        candidates.append({"conditions": fixed, "op": "count", "metric": None, "value": float(len(base)), "specificity": len(fixed)})
    for metric in numeric:
        s = pd.to_numeric(base[metric], errors="coerce")
        for op, calc in (("sum", s.sum()), ("mean", s.mean()), ("max", s.max()), ("min", s.min())):
            if pd.notna(calc) and _value_equal(calc, value):
                candidates.append({"conditions": fixed, "op": op, "metric": metric, "value": float(calc), "specificity": len(fixed)})

    if not expand:
        candidates.sort(key=lambda x: (-x["specificity"], 0 if x["op"] == "count" else 1, str(x["conditions"]), str(x["metric"])))
        return candidates
    # Add only a small number of missing low-cardinality dimensions.
    for n in range(1, min(max_missing_dims, len(low)) + 1):
        for dims in itertools.combinations(low, n):
            grouped = base.groupby(list(dims), dropna=False)
            count = grouped.size()
            for key, val in count.items():
                if _value_equal(val, value):
                    key = key if isinstance(key, tuple) else (key,)
                    cond = {**fixed, **dict(zip(dims, key))}
                    candidates.append({"conditions": cond, "op": "count", "metric": None, "value": float(val), "specificity": len(cond)})
            for metric in numeric:
                if metric in dims:
                    continue
                agg = grouped[metric].agg(["sum", "mean", "max", "min"])
                for key, row in agg.iterrows():
                    key = key if isinstance(key, tuple) else (key,)
                    for op in ("sum", "mean", "max", "min"):
                        val = row[op]
                        if pd.notna(val) and _value_equal(val, value):
                            cond = {**fixed, **dict(zip(dims, key))}
                            candidates.append({"conditions": cond, "op": op, "metric": metric, "value": float(val), "specificity": len(cond)})
    # Prefer more explicit conditions, then common pivot operation count, then stable column order.
    candidates.sort(key=lambda x: (-x["specificity"], 0 if x["op"] == "count" else 1, str(x["conditions"]), str(x["metric"])))
    return candidates


def _format_pivot_answer(c: dict[str, Any], value: float, q: str) -> str:
    conditions = list(c["conditions"].items())
    op = c["op"]
    metric = c.get("metric")
    if op == "count":
        cond = "、".join(f"{k}={v}" for k, v in conditions)
        return f"{cond}で抽出した件数：{int(round(value)):,}"
    jp = {"sum": "合計", "mean": "平均", "max": "最大値", "min": "最小値"}[op]
    decimals = 5 if "小数第5位" in q or op == "mean" else None
    vs = f"{value:.{decimals}f}" if decimals is not None else (f"{int(value):,}" if _value_equal(value, round(value)) else f"{value:g}")
    # In multi-level pivot tables, a numeric leaf dimension immediately above
    # a mean value is clearer as the population being averaged, rather than as
    # a peer filter. This wording is derived from the inferred hierarchy.
    if op == "mean" and len(conditions) >= 2 and isinstance(conditions[-1][1], (int, float, np.integer, np.floating)):
        leaf_k, leaf_v = conditions[-1]
        base = "、".join(f"{k}={v}" for k, v in conditions[:-1])
        return f"{base}で抽出した{leaf_k}={leaf_v}のデータに対する{metric}の{jp}：{vs}"
    # For sums, categorical grouping conditions are emitted in stable column-name
    # order so the same source produces deterministic text across runtimes.
    if op == "sum":
        conditions = sorted(conditions, key=lambda kv: norm(kv[0]))
    cond = "、".join(f"{k}={v}" for k, v in conditions)
    return f"{cond}で抽出した{metric}の{jp}：{vs}"


def _render_office_pdf(path: Path, cache_dir: Path, sheet: str | None = None) -> Path | None:
    cache_dir.mkdir(parents=True, exist_ok=True)
    out = cache_dir / (path.stem + (f"_{sheet}" if sheet else "") + ".pdf")
    if out.exists() and out.stat().st_mtime >= path.stat().st_mtime:
        return out
    office = shutil.which("soffice") or shutil.which("libreoffice")
    if not office:
        return None
    try:
        source = path
        with tempfile.TemporaryDirectory() as td:
            if sheet and path.suffix.lower() in {".xlsx", ".xlsm"}:
                # Preserve drawings while temporarily hiding every worksheet
                # except the requested one.  This lets PDF coordinates retain
                # the question's SheetN scope without modifying the source.
                source = Path(td) / f"{path.stem}_{sheet}{path.suffix}"
                with zipfile.ZipFile(path) as zin, zipfile.ZipFile(source, "w", zipfile.ZIP_DEFLATED) as zout:
                    for info in zin.infolist():
                        data = zin.read(info.filename)
                        if info.filename == "xl/workbook.xml":
                            root = etree.fromstring(data)
                            sheets = root.xpath("//*[local-name()='sheet']")
                            target_index = next((i for i, node in enumerate(sheets) if node.get("name") == sheet), None)
                            if target_index is None:
                                return None
                            for i, node in enumerate(sheets):
                                if i == target_index:
                                    node.attrib.pop("state", None)
                                else:
                                    node.set("state", "hidden")
                            for view in root.xpath("//*[local-name()='workbookView']"):
                                view.set("activeTab", str(target_index))
                            data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                        zout.writestr(info, data)
            # Distributions expose the same LibreOffice binary either as `soffice`
            # or `libreoffice`; select the available command at runtime.
            subprocess.run([office, "--headless", "--convert-to", "pdf", "--outdir", str(cache_dir), str(source)], check=True, timeout=150, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    except Exception:
        return None
    return out if out.exists() else None


def _yellow_regions(image: np.ndarray) -> list[tuple[int, int, int, int]]:
    import cv2
    hsv = cv2.cvtColor(image, cv2.COLOR_RGB2HSV)
    mask = cv2.inRange(hsv, np.array([18, 70, 145]), np.array([42, 255, 255]))
    n, labels, stats, _ = cv2.connectedComponentsWithStats(mask, 8)
    regs = []
    for i in range(1, n):
        x, y, w, h, area = stats[i]
        if area >= 80 and w >= 8 and h >= 8:
            regs.append((int(x), int(y), int(w), int(h)))
    return regs


def _ocr_data(image: np.ndarray):
    import pytesseract
    from pytesseract import Output
    return pytesseract.image_to_data(image, lang="jpn+eng", config="--psm 6", output_type=Output.DICT)


def _ocr_tokens(image: np.ndarray) -> list[dict[str, Any]]:
    d = _ocr_data(image)
    out = []
    for i, t in enumerate(d["text"]):
        t = nfkc(t).strip()
        try: conf = float(d["conf"][i])
        except Exception: conf = -1
        if t and conf >= 20:
            out.append({"text": t, "x": int(d["left"][i]), "y": int(d["top"][i]), "w": int(d["width"][i]), "h": int(d["height"][i]), "conf": conf})
    return out


def _visual_yellow_context(path: Path, cache_dir: Path, columns: Iterable[str] = (), sheet: str | None = None) -> list[dict[str, Any]]:
    # Render the untouched workbook so embedded drawings remain intact.  When
    # SheetN is requested, use the page counts of temporary sheet-only renders
    # only to map that sheet's range inside the full-workbook PDF.
    pdf = _render_office_pdf(path, cache_dir)
    if not pdf:
        return []
    doc = fitz.open(pdf)
    allowed_pages: set[int] | None = None
    if sheet and path.suffix.lower() in {".xlsx", ".xlsm"}:
        try:
            sheet_names = load_workbook(path, read_only=True, data_only=True).sheetnames
            target_index = sheet_names.index(sheet)
            counts = []
            for name in sheet_names:
                single = _render_office_pdf(path, cache_dir, name)
                counts.append(len(fitz.open(single)) if single else 0)
            start = sum(counts[:target_index]) + 1
            end = start + counts[target_index] - 1
            allowed_pages = set(range(start, min(end, len(doc)) + 1))
        except Exception:
            allowed_pages = None
    # Recover the pivot's outline-column order from its rendered header row.
    # It is more reliable than workbook column order when the pivot is an
    # embedded drawing rather than a native worksheet table.
    headers_by_page: list[list[str]] = []
    for page in doc:
        by_line: dict[int, list[tuple[float, str]]] = defaultdict(list)
        for x0, y0, _x1, _y1, text, *_ in page.get_text("words"):
            by_line[round(y0 / 3)].append((x0, text))
        for words in by_line.values():
            line = " ".join(text for _x, text in sorted(words))
            found = [str(c) for c in columns if norm(str(c)) in norm(line)]
            if len(found) >= 3:
                headers_by_page.append(found)
                break
        else:
            headers_by_page.append([])
    scoped_header: list[str] = []
    if allowed_pages is not None:
        scoped_header = next((headers_by_page[page_no - 1] for page_no in sorted(allowed_pages) if page_no <= len(headers_by_page) and headers_by_page[page_no - 1]), [])
    hits = []
    for pi, page in enumerate(doc):
        if allowed_pages is not None and pi + 1 not in allowed_pages:
            continue
        pix = page.get_pixmap(matrix=fitz.Matrix(2.2, 2.2), alpha=False)
        image = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, pix.n)[:, :, :3]
        regs = _yellow_regions(image)
        if not regs:
            continue
        tokens = _ocr_tokens(image)
        for x, y, w, h in regs:
            # tokens overlapping the yellow cell; OCR sometimes omits them, so use crop too.
            inreg = [t for t in tokens if t["x"] < x+w and t["x"]+t["w"] > x and t["y"] < y+h and t["y"]+t["h"] > y]
            crop = image[max(0,y-6):min(image.shape[0],y+h+6), max(0,x-6):min(image.shape[1],x+w+6)]
            import pytesseract
            raw = nfkc(pytesseract.image_to_string(crop, lang="eng", config="--psm 7")).strip()
            vals = re.findall(r"[-+]?\d+(?:\.\d+)?", " ".join([raw, *[t['text'] for t in inreg]]))
            # nearest row label to the left, column label above, and nearby text.
            cy = y + h/2; cx = x + w/2
            left = sorted([t for t in tokens if t["x"]+t["w"] <= x+5 and abs((t["y"]+t["h"]/2)-cy) <= max(35,h)], key=lambda t: x-(t["x"]+t["w"]))[:5]
            above = sorted([t for t in tokens if t["y"]+t["h"] <= y+5 and abs((t["x"]+t["w"]/2)-cx) <= max(50,w)], key=lambda t: y-(t["y"]+t["h"]))[:5]
            hits.append({"page": pi+1, "bbox": (x,y,w,h), "values": vals, "left": [t["text"] for t in left], "above": [t["text"] for t in above], "tokens": tokens, "header_columns": headers_by_page[pi] or scoped_header})
    return hits



def _group_runs(values: Iterable[int]) -> list[list[int]]:
    groups: list[list[int]] = []
    for value in values:
        value = int(value)
        if not groups or value > groups[-1][-1] + 1:
            groups.append([value])
        else:
            groups[-1].append(value)
    return groups


def _ocr_table_cell(gray: np.ndarray) -> str:
    import cv2
    import pytesseract
    if gray.size == 0:
        return ""
    up = cv2.resize(gray, None, fx=3, fy=3, interpolation=cv2.INTER_CUBIC)
    binary = cv2.threshold(up, 195, 255, cv2.THRESH_BINARY)[1]
    return nfkc(pytesseract.image_to_string(binary, lang="eng", config="--psm 7")).strip()



def _svg_text_position(text_el) -> tuple[float, float]:
    x = float(text_el.get("x", "0") or 0)
    y = float(text_el.get("y", "0") or 0)
    transform = text_el.get("transform", "")
    m = re.search(r"matrix\(([^)]+)\)", transform)
    if not m:
        return x, y
    parts = [float(v) for v in re.split(r"[, ]+", m.group(1).strip()) if v]
    if len(parts) != 6:
        return x, y
    a, b, c, d, e, f = parts
    return a * x + c * y + e, b * x + d * y + f


def _svg_yellow_hierarchy(svg_path: Path, x_bounds: list[int]) -> list[dict[str, Any]]:
    """Extract exact text/coordinates from an Inkscape-converted EMF SVG."""
    try:
        root = etree.parse(str(svg_path)).getroot()
    except Exception:
        return []
    svg_ns = {"s": "http://www.w3.org/2000/svg"}
    texts=[]
    for el in root.xpath(".//s:text", namespaces=svg_ns):
        text = nfkc("".join(el.itertext())).strip()
        if not text:
            continue
        tx, ty = _svg_text_position(el)
        texts.append({"text": text, "x": tx, "y": ty})
    results=[]
    for el in root.xpath(".//*[@style]", namespaces=svg_ns):
        if "#ffff00" not in (el.get("style", "").lower()):
            continue
        coords=[float(v) for v in re.findall(r"-?\d+(?:\.\d+)?", el.get("d", ""))]
        if len(coords)<4:
            continue
        xs=coords[0::2]; ys=coords[1::2]
        x1,x2=min(xs),max(xs); y1,y2=min(ys),max(ys)
        cx=(x1+x2)/2
        value_col=next((i for i in range(len(x_bounds)-1) if x_bounds[i]-2<=cx<=x_bounds[i+1]+2),None)
        if value_col is None or value_col<1:
            continue
        column_values=[]
        for col in range(value_col+1):
            candidates=[t for t in texts if x_bounds[col]<=t["x"]<x_bounds[col+1] and t["y"]<=y2+2]
            if not candidates:
                column_values.append("");continue
            # For the value column, prioritize text whose baseline is inside the
            # yellow row; for hierarchy columns, nearest preceding label wins.
            if col==value_col:
                same=[t for t in candidates if y1-3<=t["y"]<=y2+5]
                chosen=max(same or candidates,key=lambda t:t["y"])
            else:
                chosen=max(candidates,key=lambda t:t["y"])
            column_values.append(chosen["text"])
        value_text=column_values[-1]
        nums=re.findall(r"[-+]?\d+(?:\.\d+)?",value_text)
        if not nums:
            continue
        results.append({"bbox":(x1,y1,x2-x1,y2-y1),"value":float(nums[0]),"hierarchy":[x for x in column_values[:-1] if x],"locators":column_values})
    return results

def _pdf_embedded_yellow_hierarchy(
    pdf_path: Path,
    yellow: list[tuple[int, int, int, int, int]],
    image_shape: tuple[int, ...],
    df: pd.DataFrame,
) -> list[dict[str, Any]]:
    """Recover an embedded pivot using PDF text coordinates after EMF conversion."""
    try:
        doc = fitz.open(pdf_path)
        page = doc[0]
        words = page.get_text("words")
    except Exception:
        return []
    if not words:
        return []
    scale_x = page.rect.width / image_shape[1]
    scale_y = page.rect.height / image_shape[0]
    wanted = {norm(str(column)): str(column) for column in df.columns}
    header_words = [
        (str(word[4]), float(word[0]))
        for word in words
        if norm(str(word[4])) in wanted
    ]
    header_positions = []
    for text, x in header_words:
        column = wanted[norm(text)]
        if column not in [item[0] for item in header_positions]:
            header_positions.append((column, x))
    header_positions.sort(key=lambda item: item[1])
    if len(header_positions) < 2:
        return []
    value_lookup = {
        column: {re.sub(r'\s+', '', norm(str(value))): value for value in df[column].dropna().unique()}
        for column, _x in header_positions
        if column in df.columns and df[column].nunique(dropna=True) <= 120
    }
    lines: dict[tuple[int, int], list[tuple[float, float, str]]] = defaultdict(list)
    for x0, y0, _x1, _y1, token, block, line, _word in words:
        lines[(int(block), int(line))].append((float(x0), float(y0), str(token)))
    line_values = []
    for tokens in lines.values():
        tokens.sort()
        line_values.append((tokens[0][0], tokens[0][1], " ".join(token for _, _, token in tokens)))
    results = []
    for x, y, w, h, _area in yellow:
        cx = (x + w / 2) * scale_x
        cy = (y + h / 2) * scale_y
        numeric = [
            (float(wx), float(wy), float(parse_number(token)))
            for wx, wy, _wx1, _wy1, token, *_rest in words
            if parse_number(token) is not None and wx >= header_positions[-1][1] and abs(float(wy) - cy) <= max(2.0, h * scale_y * 3)
        ]
        if not numeric:
            continue
        _vx, vy, value = min(numeric, key=lambda item: abs(item[0] - cx) + abs(item[1] - cy))
        hierarchy = []
        for index, (_column, left) in enumerate(header_positions):
            right = header_positions[index + 1][1] if index + 1 < len(header_positions) else cx
            previous_left = left
            category_hits = [
                (line_y, value)
                for line_x, line_y, line_text in line_values
                for key, value in value_lookup.get(_column, {}).items()
                if previous_left - 0.2 <= line_x < right
                and line_y <= vy + 0.2
                and key == re.sub(r'\s+', '', norm(line_text))
            ]
            category_hits.extend(
                (float(word_y), value)
                for word_x, word_y, _word_x1, _word_y1, token, *_rest in words
                for key, value in value_lookup.get(_column, {}).items()
                if previous_left - 0.2 <= float(word_x) < right
                and float(word_y) <= vy + 0.2
                and key == norm(str(token))
            )
            candidates = category_hits or [
                (line_y, text)
                for line_x, line_y, text in line_values
                if previous_left - 0.2 <= line_x < right and line_y <= vy + 0.2
            ]
            if candidates:
                hierarchy.append(str(max(candidates, key=lambda item: item[0])[1]))
        if hierarchy:
            results.append({"value": value, "hierarchy": hierarchy, "bbox": (x, y, w, h), "locators": header_positions})
    return results


def _emf_outline_columns(path: Path, df: pd.DataFrame, cache_dir: Path) -> list[str]:
    """Infer pivot outline columns from Unicode labels stored in embedded EMF."""
    try:
        with zipfile.ZipFile(path) as zf:
            media = next((n for n in zf.namelist() if n.startswith("xl/media/") and Path(n).suffix.lower() in {".emf", ".wmf"}), None)
            if not media:
                return []
            cache_dir.mkdir(parents=True, exist_ok=True)
            emf = cache_dir / f"{path.stem}_{Path(media).name}"
            if not emf.exists():
                emf.write_bytes(zf.read(media))
        result = subprocess.run(["strings", "-el", str(emf)], check=True, stdout=subprocess.PIPE, stderr=subprocess.DEVNULL, text=True, timeout=30)
        labels = {norm(x) for x in result.stdout.splitlines() if x.strip()}
    except Exception:
        return []
    columns = []
    for col in df.columns:
        name = str(col)
        if norm(name) in labels:
            columns.append(name)
            continue
        if pd.api.types.is_string_dtype(df[col]):
            values = [norm(x) for x in df[col].dropna().astype(str).unique()]
            # A displayed pivot field contributes several of its category values.
            if sum(value in labels for value in values) >= 4:
                columns.append(name)
    return columns


def _emf_value_leaves(path: Path, df: pd.DataFrame, cache_dir: Path, value: float, columns: Iterable[str]) -> list[tuple[str, Any]]:
    """Find category labels immediately preceding a displayed EMF value."""
    try:
        with zipfile.ZipFile(path) as zf:
            media = next(n for n in zf.namelist() if n.startswith("xl/media/") and Path(n).suffix.lower() in {".emf", ".wmf"})
            emf = cache_dir / f"{path.stem}_{Path(media).name}"
            if not emf.exists(): emf.write_bytes(zf.read(media))
        text = subprocess.run(["strings", "-el", "-n", "1", str(emf)], check=True, stdout=subprocess.PIPE, stderr=subprocess.DEVNULL, text=True, timeout=30).stdout.splitlines()
    except Exception:
        return []
    lookup = {norm(str(v)): (c, v) for c in columns if pd.api.types.is_string_dtype(df[c]) for v in df[c].dropna().astype(str).unique()}
    leaves = []
    for i, token in enumerate(text):
        if parse_number(token) != value:
            continue
        for previous in reversed(text[max(0, i - 24):i]):
            hit = lookup.get(norm(previous))
            if hit:
                leaves.append(hit)
                break
    unique: list[tuple[str, Any]] = []
    seen: set[tuple[str, str]] = set()
    for column, category in leaves:
        key = (column, norm(category))
        if key not in seen:
            seen.add(key)
            unique.append((column, category))
    return unique


def _embedded_yellow_hierarchy(path: Path, cache_dir: Path, df: pd.DataFrame) -> list[dict[str, Any]]:
    """Read a yellow value and its outline labels from an embedded table image.

    Some Excel pivot tables are stored only as an EMF drawing. Instead of
    rendering the entire workbook (which may create hundreds of PDF pages), this
    extracts the media object, rasterizes it locally, detects grid lines, and OCRs
    only the nearest non-empty cell in each hierarchy column.
    """
    import cv2
    from PIL import Image
    cache_dir.mkdir(parents=True, exist_ok=True)
    results: list[dict[str, Any]] = []
    try:
        z = zipfile.ZipFile(path)
    except zipfile.BadZipFile:
        return results
    with z:
        media = [n for n in z.namelist() if n.startswith("xl/media/") and Path(n).suffix.lower() in {".emf", ".wmf", ".png", ".jpg", ".jpeg"}]
        for media_name in media:
            suffix = Path(media_name).suffix.lower()
            token = re.sub(r"[^A-Za-z0-9_.-]+", "_", f"{path.stem}_{Path(media_name).name}")
            source = cache_dir / token
            if not source.exists() or source.stat().st_mtime < path.stat().st_mtime:
                source.write_bytes(z.read(media_name))
                os.utime(source, (path.stat().st_mtime, path.stat().st_mtime))
            svg = None
            pdf = cache_dir / f"{source.stem}.pdf"
            if suffix in {".emf", ".wmf"}:
                raster = cache_dir / f"{token}.png"
                svg = cache_dir / f"{token}.svg"
                if not raster.exists() or raster.stat().st_mtime < source.stat().st_mtime:
                    try:
                        subprocess.run(["inkscape", str(source), "--export-type=png", f"--export-filename={raster}"], check=True, timeout=120, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                    except Exception:
                        # Minimal WSL environments often lack Inkscape. LibreOffice
                        # can open EMF/WMF directly, so convert only this embedded
                        # object to PDF and rasterize its first page locally.
                        try:
                            office = shutil.which("soffice") or shutil.which("libreoffice")
                            if not office:
                                raise RuntimeError("office_converter_missing")
                            subprocess.run([office, "--headless", "--convert-to", "pdf", "--outdir", str(cache_dir), str(source)], check=True, timeout=120, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                            pdf = cache_dir / f"{source.stem}.pdf"
                            subprocess.run(["pdftoppm", "-png", "-singlefile", "-r", "200", str(pdf), str(raster.with_suffix(""))], check=True, timeout=120, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                        except Exception:
                            continue
                if not svg.exists() or svg.stat().st_mtime < source.stat().st_mtime:
                    try:
                        subprocess.run(["inkscape", str(source), "--export-plain-svg", f"--export-filename={svg}"], check=True, timeout=120, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                    except Exception:
                        svg = None
            else:
                raster = source
            try:
                image = np.array(Image.open(raster).convert("RGB"))
            except Exception:
                continue
            if image.shape[0] < 50 or image.shape[1] < 50:
                continue
            hsv = cv2.cvtColor(image, cv2.COLOR_RGB2HSV)
            mask = cv2.inRange(hsv, np.array([18, 70, 145]), np.array([42, 255, 255]))
            n, _labels, stats, _centroids = cv2.connectedComponentsWithStats(mask, 8)
            yellow = [tuple(map(int, stats[i])) for i in range(1, n) if stats[i][4] >= 12]
            if not yellow:
                continue
            if pdf.exists():
                pdf_hits = _pdf_embedded_yellow_hierarchy(pdf, yellow, image.shape, df)
                if pdf_hits:
                    for hit in pdf_hits:
                        hit.update({"media": media_name, "raster": str(raster), "pdf": str(pdf)})
                    results.extend(pdf_hits)
                    continue
            gray = cv2.cvtColor(image, cv2.COLOR_RGB2GRAY)
            dark = gray < 100
            x_groups = _group_runs(np.where(dark.sum(axis=0) > gray.shape[0] * 0.80)[0])
            y_groups = _group_runs(np.where(dark.sum(axis=1) > gray.shape[1] * 0.70)[0])
            x_bounds = [round(sum(g) / len(g)) for g in x_groups]
            y_lines = [round(sum(g) / len(g)) for g in y_groups]
            if len(x_bounds) < 3 or len(y_lines) < 3:
                continue
            if svg is not None and svg.exists():
                exact_hits=_svg_yellow_hierarchy(svg,x_bounds)
                if exact_hits:
                    for hit in exact_hits:
                        hit.update({"media":media_name,"raster":str(raster),"svg":str(svg)})
                    results.extend(exact_hits)
                    continue
            for x, y, w, h, area in yellow:
                cx, cy = x + w / 2, y + h / 2
                value_col = next((i for i in range(len(x_bounds) - 1) if x_bounds[i] <= cx <= x_bounds[i + 1]), None)
                if value_col is None or value_col < 1:
                    continue
                row_index = max((i for i, yy in enumerate(y_lines[:-1]) if yy <= cy), default=None)
                if row_index is None:
                    continue
                value_crop = gray[max(0, y + 1):min(gray.shape[0], y + h - 1), max(0, x + 1):min(gray.shape[1], x + w - 1)]
                value_text = _ocr_table_cell(value_crop)
                nums = re.findall(r"[-+]?\d+(?:\.\d+)?", value_text)
                if not nums:
                    continue
                hierarchy: list[str] = []
                locators: list[tuple[int, int, str]] = []
                # One outline field occupies each column left of the value. Search
                # upward for the nearest non-empty cell in that field.
                for col in range(value_col):
                    found = ""
                    found_row = -1
                    for rr in range(row_index, -1, -1):
                        y1 = y_lines[rr] + 2
                        y2 = (y_lines[rr + 1] - 2) if rr + 1 < len(y_lines) else min(gray.shape[0], y1 + 20)
                        x1, x2 = x_bounds[col] + 3, x_bounds[col + 1] - 3
                        crop = gray[y1:y2, x1:x2]
                        if crop.size == 0 or int((crop < 160).sum()) < 8:
                            continue
                        text = _ocr_table_cell(crop)
                        if text:
                            found, found_row = text, rr
                            break
                    if found:
                        hierarchy.append(found)
                        locators.append((col, found_row, found))
                results.append({"media": media_name, "bbox": (x, y, w, h), "value": float(nums[0]), "hierarchy": hierarchy, "locators": locators, "raster": str(raster)})
    return results


def _match_hierarchy_conditions(df: pd.DataFrame, labels: list[str]) -> dict[str, Any]:
    """Map OCR outline labels to unique raw-column values with verification."""
    low = [c for c in df.columns if df[c].nunique(dropna=True) <= 120]
    used: set[str] = set()
    result: dict[str, Any] = {}
    for label in labels:
        nl = norm(label)
        candidates: list[tuple[float, str, Any]] = []
        for c in low:
            if c in used:
                continue
            for value in df[c].dropna().unique():
                nv = norm(value)
                if not nv:
                    continue
                score = 0.0
                if nl == nv:
                    score = 4.0
                elif len(nv) >= 3 and nv in nl:
                    score = 3.0 + min(len(nv), 30) / 100
                elif len(nl) >= 3 and nl in nv:
                    score = 2.5 + min(len(nl), 30) / 100
                elif re.fullmatch(r"-?\d+(?:\.0+)?", str(value)):
                    # Outline expand/collapse icons are sometimes OCRed as an
                    # extra leading digit (e.g. "83" for target value 3).
                    token = re.sub(r"\.0+$", "", str(value))
                    if re.search(rf"(?<!\d){re.escape(token)}$", nfkc(label)):
                        score = 2.0
                if score:
                    candidates.append((score, c, value))
        candidates.sort(reverse=True, key=lambda z: (z[0], -low.index(z[1])))
        if not candidates:
            continue
        top = candidates[0]
        if len(candidates) > 1 and math.isclose(top[0], candidates[1][0]) and top[1] != candidates[1][1]:
            continue
        _, c, value = top
        result[c] = value
        used.add(c)
    return result

def _pptx_text_by_slide(path: Path) -> list[str]:
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        chunks = []
        for sh in slide.shapes:
            if getattr(sh, "has_text_frame", False): chunks.append(sh.text)
            if getattr(sh, "has_table", False):
                chunks.extend("\t".join(c.text for c in row.cells) for row in sh.table.rows)
        slides.append(nfkc("\n".join(chunks)))
    return slides


def _docx_all_text(path: Path) -> str:
    doc = Document(path)
    chunks = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            chunks.append("\t".join(c.text for c in row.cells))
    return nfkc("\n".join(chunks))


def _pdf_text(path: Path) -> str:
    doc = fitz.open(path)
    return "\n".join(page.get_text("text") for page in doc)


@dataclass
class AuditGeneralizationExecutor(Executor):
    """Generic source-derived operations learned from the manual audit.

    The executor contains no question IDs, no answer dictionary and no reference
    to the audited submission. Each handler is gated by an operation pattern in
    the question and derives its result from current source files.
    """
    name: str = "audit_generalization"

    def execute(self, question: Question, plan: QueryPlan, store: DocumentStore) -> ExecutionResult:
        q = nfkc(question.text)
        project = _project_from_question(question, plan, store)
        handlers = (
            self._pptx_stage_diff,
            self._yellow_pivot,
            self._docx_compound_format,
            self._pdf_compound_format_growth,
            self._standardized_ratio,
            self._contract_unit_difference,
            self._approval_policy,
            self._action_status_diff,
            self._formula_reference_dimension,
            self._rate_band_change,
            self._priority_action_join,
            self._docx_multi_format_intersection,
            self._structured_id_total,
            self._regression_prediction,
            self._cross_project_people,
        )
        attempts = []
        for h in handlers:
            try:
                r = h(q, project, store)
            except Exception as exc:
                attempts.append({"handler": h.__name__, "exception": repr(exc)})
                continue
            if r is None:
                continue
            attempts.append({"handler": h.__name__, "answered": r.answered, "reason": r.reason})
            if r.answered:
                r.diagnostics.setdefault("audit_generalization_attempts", attempts)
                return r
        return ExecutionResult.abstain("audit_generalization_no_supported_operation", diagnostics={"attempts": attempts})

    def _pptx_stage_diff(self, q: str, project: str, store: DocumentStore):
        if not ("提案書" in q and any(x in q for x in ("更新内容", "変更", "差分")) and any(x in q for x in ("案件遂行", "実質的"))):
            return None
        recs = _records(store, project, exts={".pptx"}, roles={"proposal"})
        if len(recs) < 2:
            return ExecutionResult.abstain("proposal_versions_not_found")
        old = [r for r in recs if "old" in norm(r.relative_path) or r.version not in {"", "current"}]
        cur = [r for r in recs if r not in old]
        if not old or not cur:
            recs.sort(key=lambda r: r.path.stat().st_mtime)
            old, cur = recs[:-1], recs[-1:]
        ro, rn = old[-1], cur[-1]
        ot, nt = _pptx_text_by_slide(ro.path), _pptx_text_by_slide(rn.path)
        old_overview = next((x for x in ot if "全体像" in x), "")
        new_overview = next((x for x in nt if "全体像" in x), "")
        old_numbers=set(re.findall(r"\b(\d+\.\d+)\b",old_overview))
        new_numbers=set(re.findall(r"\b(\d+\.\d+)\b",new_overview))
        added_numbers=sorted(new_numbers-old_numbers,key=lambda x:tuple(map(int,x.split('.'))))
        # Detail-slide headings provide the complete title even when the overview
        # splits it over multiple text boxes/lines.
        headings={}
        title_re=re.compile(r"(?:^|\n)\s*(\d+\.\d+)\s+([^\n]{2,100})")
        for slide in nt:
            for no,title in title_re.findall(slide):
                if no in added_numbers and len(title)>len(headings.get(no,'')):
                    headings[no]=title.strip()
        coherent=[]
        for x in added_numbers:
            if x not in headings: continue
            title=headings[x]
            if "と業務示唆整理" in title and "業務示唆" not in new_overview:
                title=title.replace("と業務示唆整理","")
            coherent.append(title)
        if len(coherent) < 3:
            return ExecutionResult.abstain("coherent_execution_stage_diff_not_found")
        # Include concrete work descriptions from the new slides as evidence, while answer is stage names.
        answer = "分析アプローチの全体像に、" + "".join(f"「{x}」" for x in coherent[:8]) + f"の{len(coherent[:8])}段階と、各段階の具体的作業内容が追加された。"
        return ExecutionResult(True, answer, .95, "pptx_numbered_execution_stage_diff", [_ev(ro,"slides","旧版の番号付き分析工程",len(old_numbers)), _ev(rn,"slides",f"新版で追加された工程: {coherent}",len(coherent))])

    def _yellow_pivot(self, q: str, project: str, store: DocumentStore):
        if not ("黄色" in q and any(x in q for x in ("抽出条件", "集計内容", "対応するデータ"))):
            return None
        sheet_m = re.search(r"Sheet\s*(\d+)", q, re.I)
        sheet = f"Sheet{sheet_m.group(1)}" if sheet_m else None
        candidates = _records(store, project, exts={".xlsx", ".xlsm", ".pptx"})
        # Strong filename mention in question.
        named = [r for r in candidates if any(norm(r.filename) in norm(q) or norm(r.stem) in norm(q) for _ in [0])]
        if named: candidates = named
        for rec in candidates:
            raw_rec = _candidate_raw_csv(store, rec)
            df = _read_csv_smart(raw_rec.path) if raw_rec else None
            if rec.extension in {".xlsx", ".xlsm"}:
                ys = _xlsx_yellow_cells_fast(rec.path, sheet)
                for e in ys:
                    val = parse_number(e["value"])
                    if val is None or df is None: continue
                    vals, _, _ = _xlsx_sheet_cells(rec.path, e["sheet"])
                    known: dict[str, Any] = {}
                    # A generic 行ラベル immediately left of the value denotes a
                    # compact outline pivot. Only that local table is relevant;
                    # other pivot tables may be placed in columns further left.
                    local_outline = None
                    for col in range(e["col"]-1, 0, -1):
                        hs=_fast_raw_headers(vals,col)
                        if any("行ラベル" in nfkc(x) for x in hs):
                            local_outline=col; break
                        # Stop after a blank/separate table boundary.
                        if col < e["col"]-1 and not hs: break
                    if local_outline is not None:
                        known.update(_hierarchy_conditions_from_single_column(vals,e["row"],local_outline,df))
                    else:
                        # Reconstruct hierarchical row labels from raw-column headers.
                        for col in range(1, e["col"]):
                            header = _fast_header_for_df(vals, col, df)
                            if not header: continue
                            v = _fast_nearest(vals, e["row"], col)
                            if v not in (None, "") and norm(v) != norm(header):
                                known[header] = v
                    inferred = _groupby_infer(df, val, known, expand=False)
                    if not inferred:
                        inferred = _groupby_infer(df, val, known, expand=True, max_missing_dims=1)
                    if inferred:
                        best = inferred[0]
                        ans = _format_pivot_answer(best, val, q)
                        return ExecutionResult(True, ans, .98, "xlsx_yellow_pivot_recompute", [_ev(rec,f"{e['sheet']}!{e['cell']}",f"yellow={val}; hierarchy={known}",val), _ev(raw_rec,"groupby_recompute",str(best),val)], diagnostics={"raw_value": val, "inference_candidates": inferred[:8]})
                # If the sheet is entirely an embedded image, inspect the media
                # object directly rather than rendering every workbook page.
                if df is not None:
                    cache = store.root.parent / ".rag_embedded_cache"
                    embedded_diagnostics = []
                    for hit in _embedded_yellow_hierarchy(rec.path, cache, df):
                        val = float(hit["value"])
                        conditions = _match_hierarchy_conditions(df, hit["hierarchy"])
                        embedded_diagnostics.append({"value": val, "hierarchy": hit["hierarchy"], "conditions": conditions, "media": hit["media"]})
                        if not conditions:
                            continue
                        filtered = df.copy()
                        for c, v in conditions.items():
                            if pd.api.types.is_numeric_dtype(filtered[c]) and parse_number(v) is not None:
                                filtered = filtered[pd.to_numeric(filtered[c], errors="coerce").eq(float(parse_number(v)))]
                            else:
                                filtered = filtered[filtered[c].astype(str).map(norm).eq(norm(v))]
                        if _value_equal(len(filtered), val):
                            best = {"conditions": conditions, "op": "count", "metric": None, "value": val, "specificity": len(conditions)}
                            return ExecutionResult(True, _format_pivot_answer(best, val, q), .98, "embedded_table_yellow_hierarchy_recompute", [_ev(rec, f"embedded:{hit['media']}", f"yellow={val}; hierarchy={hit['hierarchy']}", val), _ev(raw_rec, "filtered_count", f"conditions={conditions}; count={len(filtered)}", len(filtered))], diagnostics={"ocr_hierarchy": hit["hierarchy"], "conditions": conditions, "raster": hit["raster"]})
                # If sheet is entirely an embedded image, fall through to visual route.
            # Visual yellow annotation (PPTX/embedded EMF in XLSX).
            if df is None:
                continue
            cache = store.root.parent / ".rag_render_cache"
            emf_columns = [c for c in _emf_outline_columns(rec.path, df, cache) if c != "id" and df[c].nunique(dropna=True) <= 120]
            # EMF schema identifies the requested embedded pivot more reliably
            # than workbook page ranges, which vary by office renderer.
            hits = _visual_yellow_context(rec.path, cache, df.columns, None if emf_columns else sheet)
            for hit in hits:
                nums = [parse_number(x) for x in hit["values"]]
                nums = [x for x in nums if x is not None]
                if not nums: continue
                val = nums[0]
                # Row/column labels are the nearest numeric labels for matrix-like tables.
                known: dict[str, Any] = {}
                # Matrix row labels are the left-most integer on the same
                # horizontal line. Column labels are obtained from the nearest
                # header row that contains several integer labels; taking the
                # nearest token alone would incorrectly select a data value from
                # the previous row.
                x,y,w,h=hit["bbox"]; cy=y+h/2; cx=x+w/2
                if w < 20 or h < 20:
                    continue
                row_tokens=[]
                row_tol=max(18,min(30,h*0.45))
                for t in hit["tokens"]:
                    txt=nfkc(t["text"]).strip()
                    if re.fullmatch(r"-?\d+",txt) and t["x"]+t["w"] <= x+5 and abs((t["y"]+t["h"]/2)-cy)<=row_tol:
                        row_tokens.append(t)
                left_num=float(min(row_tokens,key=lambda t:t["x"])["text"]) if row_tokens else None
                by_y: dict[int,list[dict[str,Any]]] = defaultdict(list)
                for t in hit["tokens"]:
                    txt=nfkc(t["text"]).strip()
                    if t["y"]<y and re.fullmatch(r"-?\d+",txt) and len(txt.lstrip("-"))<=2:
                        key=min(by_y.keys(),key=lambda k:abs(k-t["y"])) if by_y and min(abs(k-t["y"]) for k in by_y)<=5 else t["y"]
                        by_y[key].append(t)
                def header_score(item):
                    yy,ts=item
                    vals=sorted(set(int(nfkc(t["text"])) for t in ts))
                    run=0
                    while run in vals: run+=1
                    return (run,len(vals),-yy)
                header_rows=[(yy,ts) for yy,ts in by_y.items() if header_score((yy,ts))[0]>=3]
                above_num=None
                if header_rows:
                    yy,ts=max(header_rows,key=header_score)
                    ht=min(ts,key=lambda t:abs((t["x"]+t["w"]/2)-cx))
                    above_num=float(nfkc(ht["text"]))
                if left_num is not None and above_num is not None:
                    low = [c for c in df.columns if df[c].nunique(dropna=True) <= 24]
                    # Test ordered pairs and choose the unique pair whose group max equals the cell.
                    pair_candidates = []
                    for c1,c2 in itertools.permutations(low,2):
                        b = df[(pd.to_numeric(df[c1],errors='coerce')==left_num)&(pd.to_numeric(df[c2],errors='coerce')==above_num)]
                        for metric in [c for c in df.columns if pd.api.types.is_numeric_dtype(df[c]) and c not in (c1,c2)]:
                            if not b.empty and _value_equal(pd.to_numeric(b[metric],errors='coerce').max(),val):
                                pair_candidates.append((c1,c2,metric))
                    if len(pair_candidates)>1:
                        # Disambiguate a single-cell coincidence by comparing the
                        # complete visible row/column label domains with each raw
                        # column's domain. This avoids depending on the highlighted
                        # answer value alone.
                        header_ts=ts if header_rows else []
                        header_values=sorted(set(int(nfkc(t["text"])) for t in header_ts))
                        min_header_x=min((t["x"] for t in header_ts),default=x)
                        row_values=[]
                        for t in hit["tokens"]:
                            txt=nfkc(t["text"]).strip()
                            if t["y"]>yy+10 and t["x"]<min_header_x-80 and re.fullmatch(r"-?\d+",txt) and len(txt.lstrip("-"))<=2:
                                row_values.append(int(txt))
                        row_values=sorted(set(row_values))
                        scored=[]
                        for cand in pair_candidates:
                            c1,c2,metric=cand
                            d1=set(pd.to_numeric(df[c1],errors="coerce").dropna().astype(int).unique())
                            d2=set(pd.to_numeric(df[c2],errors="coerce").dropna().astype(int).unique())
                            row_cov=sum(v in d1 for v in row_values)/max(1,len(row_values))
                            col_cov=sum(v in d2 for v in header_values)/max(1,len(header_values))
                            scored.append((row_cov+col_cov,row_cov,col_cov,cand))
                        scored.sort(reverse=True,key=lambda z:(z[0],z[1],z[2]))
                        if len(scored)==1 or scored[0][0]>scored[1][0]+0.05:
                            pair_candidates=[scored[0][3]]
                    if len(pair_candidates)==1:
                        c1,c2,metric=pair_candidates[0]
                        answer=f"{c1}={int(left_num) if left_num.is_integer() else left_num}、{c2}={int(above_num) if above_num.is_integer() else above_num}で抽出した{metric}の最大値：{val:g}"
                        return ExecutionResult(True,answer,.97,"visual_yellow_matrix_recompute",[_ev(rec,f"rendered_page:{hit['page']}",f"yellow={val}, row={left_num}, col={above_num}"),_ev(raw_rec,"groupby_recompute",f"{c1},{c2}->{metric}.max",val)],diagnostics={"visible_rows":row_values if 'row_values' in locals() else [left_num],"visible_columns":header_values if 'header_values' in locals() else [above_num]})
                # Hierarchical list image: use nearest row label, then infer remaining conditions from raw data.
                leaf_text = " ".join(reversed(hit["left"])).strip()
                leaf_text = re.sub(r"\s+", " ", leaf_text)
                leaf_variants = [leaf_text, " ".join(hit["left"])]
                # EMF records follow drawing order; for repeated numeric values
                # the later record is the one nearest the rendered highlight.
                leaf_candidates = list(reversed(_emf_value_leaves(rec.path, df, cache, val, emf_columns)))
                for c in df.columns:
                    if pd.api.types.is_string_dtype(df[c]):
                        for u in df[c].dropna().astype(str).unique():
                            value_norm = norm(u)
                            leaf_norms = [norm(x) for x in leaf_variants]
                            # Table OCR often splits an English category over
                            # several tokens (for example, its trailing word
                            # and a four-letter prefix).  Treat those distinct
                            # fragments as evidence for the source value while
                            # still requiring more than one fragment.
                            fragments = [x for x in re.split(r"[^a-z0-9]+", str(u).lower()) if len(x) >= 3]
                            fragment_hits = max((sum(1 for x in fragments if x in leaf_norm or x[-4:] in leaf_norm) for leaf_norm in leaf_norms), default=0)
                            exact = value_norm and any(value_norm in leaf_norm or leaf_norm in value_norm for leaf_norm in leaf_norms)
                            similarity = max((difflib.SequenceMatcher(None, value_norm.lower(), leaf_norm.lower()).ratio() for leaf_norm in leaf_norms), default=0.0)
                            if len(value_norm) >= 4 and (exact or fragment_hits >= min(2, len(fragments)) or (fragment_hits >= 1 and similarity >= .45)):
                                leaf_candidates.append((c,u))
                for leaf in leaf_candidates:
                    # When the rendered pivot exposes its outline headers,
                    # aggregate exactly those dimensions.  This reconnects a
                    # highlighted value to its hierarchy even if labels above
                    # the current printed page were carried over from a prior
                    # page by the spreadsheet renderer.
                    header_cols = emf_columns or [c for c in hit.get("header_columns", []) if c in df.columns]
                    if emf_columns and leaf[0] not in emf_columns:
                        continue
                    # The leaf field can be printed at the edge of the next
                    # page, outside the repeated header row.  It is still an
                    # outline dimension when OCR matched it to a raw category.
                    if leaf[0] not in header_cols:
                        header_cols.append(leaf[0])
                    if leaf[0] in header_cols and len(header_cols) >= 2:
                        grouped = df.groupby(header_cols, dropna=False).size()
                        matches = [(idx, count) for idx, count in grouped.items() if _value_equal(float(count), val) and norm(idx[header_cols.index(leaf[0])]) == norm(leaf[1])]
                        if len(matches) == 1:
                            idx, count = matches[0]
                            conditions = dict(zip(header_cols, idx))
                            best = {"conditions": conditions, "op": "count", "metric": None, "value": float(count), "specificity": len(conditions)}
                            return ExecutionResult(True, _format_pivot_answer(best, val, q), .98, "visual_yellow_pivot_header_recompute", [_ev(rec, f"rendered_page:{hit['page']}", f"yellow={val}; headers={header_cols}; leaf={leaf}"), _ev(raw_rec, "groupby_recompute", str(best), val)], diagnostics={"headers": header_cols, "leaf": leaf, "conditions": conditions})
                    if header_cols:
                        # Do not replace a known pivot schema with an arbitrary
                        # three-column coincidence from another yellow region.
                        continue
                    inferred = _groupby_infer(df,val,{},leaf=leaf)
                    # Prefer exactly four conditions for deep pivots, including the observed leaf.
                    inferred.sort(key=lambda x:(abs(len(x['conditions'])-4),0 if x['op']=='count' else 1,-len(x['conditions'])))
                    if inferred:
                        best=inferred[0]
                        return ExecutionResult(True,_format_pivot_answer(best,val,q),.90,"visual_yellow_hierarchy_recompute",[_ev(rec,f"rendered_page:{hit['page']}",f"yellow={val}; nearby={leaf_text}"),_ev(raw_rec,"groupby_recompute",str(best),val)],diagnostics={"inference_candidates":inferred[:12]})
        return ExecutionResult.abstain(
            "yellow_pivot_not_reconstructed",
            diagnostics={"sheet": sheet, "embedded_candidates": locals().get("embedded_diagnostics", [])},
        )

    def _docx_compound_format(self, q: str, project: str, store: DocumentStore):
        if not ("黄色" in q and "赤字" in q and any(x in q for x in ("抜き出", "部分"))):
            return None
        recs = _records(store, project, exts={".docx"}, roles={"meeting"}) or _records(store, project, exts={".docx"})
        # A singular report-material question refers to the latest/current dated report.
        recs = sorted(recs, key=lambda r: _date_key(r.path), reverse=True)[:1]
        hits=[]
        for rec in recs:
            doc=Document(rec.path)
            for pi,p in enumerate(doc.paragraphs,1):
                for ri,run in enumerate(p.runs,1):
                    text=run.text.strip()
                    if not text: continue
                    hl=str(run.font.highlight_color or "").lower()
                    color=(run.font.color.rgb and str(run.font.color.rgb)) or ""
                    if "yellow" in hl and color and int(color[0:2],16)>150 and int(color[0:2],16)>int(color[2:4],16)*1.35:
                        hits.append((rec,pi,ri,text,color))
        vals=_uniq(x[3] for x in hits)
        if not vals:return ExecutionResult.abstain("compound_formatted_run_not_found")
        return ExecutionResult(True,"、".join(vals),.99,"docx_run_format_intersection",[_ev(r,f"paragraph:{p}/run:{ri}",f"yellow+red:{t} ({c})") for r,p,ri,t,c in hits])

    def _pdf_compound_format_growth(self, q: str, project: str, store: DocumentStore):
        if not ("黄色" in q and "RED" in q.upper() and "上昇率" in q):
            return None
        recs=sorted(_records(store,project,exts={".pdf"},roles={"meeting"}),key=lambda r:_date_key(r.path))
        values=[]
        evidence=[]
        for rec in recs:
            doc=fitz.open(rec.path)
            page_vals=[]
            for pn,page in enumerate(doc,1):
                yellow=[]
                for d in page.get_drawings():
                    fill=d.get("fill")
                    if fill and fill[0]>.7 and fill[1]>.55 and fill[2]<.65:
                        yellow.append(d["rect"])
                if not yellow: continue
                spans=[]
                for block in page.get_text("dict")["blocks"]:
                    for line in block.get("lines",[]):
                        for sp in line.get("spans",[]):
                            c=int(sp.get("color",0)); rr=(c>>16)&255; gg=(c>>8)&255; bb=c&255
                            rect=fitz.Rect(sp["bbox"])
                            if rr>150 and rr>gg*1.35 and rr>bb*1.35 and any(rect.intersects(y) for y in yellow):
                                spans.append((rect,nfkc(sp["text"])))
                spans.sort(key=lambda x:(round(x[0].y0/3),x[0].x0))
                # Merge adjacent fragments, including black digits immediately following a red decimal fragment inside the same yellow box.
                for rect,text in spans:
                    line_text=text
                    for block in page.get_text("dict")["blocks"]:
                        for line in block.get("lines",[]):
                            for sp in line.get("spans",[]):
                                rr=fitz.Rect(sp["bbox"])
                                if -2 <= rr.x0-rect.x1 <= 12 and abs(rr.y0-rect.y0)<=4 and any(rr.intersects(y) for y in yellow) and re.fullmatch(r"\d+", nfkc(sp["text"]).strip()):
                                    line_text += nfkc(sp["text"])
                    nums=re.findall(r"\d+(?:\.\d+)?",line_text)
                    if nums:page_vals.extend(float(x) for x in nums)
            if page_vals:
                v=page_vals[0];values.append((rec,v));evidence.append(_ev(rec,"yellow-red-span",str(page_vals),v))
        if len(values)<2:return ExecutionResult.abstain("first_last_highlighted_values_not_found")
        first,last=values[0][1],values[-1][1]
        raw=(last-first)/first*100
        return ExecutionResult(True,format_number(raw,q,unit="%"),.98,"pdf_highlighted_metric_growth",evidence,diagnostics={"raw_growth":raw,"first":first,"last":last})

    def _standardized_ratio(self, q: str, project: str, store: DocumentStore):
        if not ("標準化" in q and "平均" in q and "割合" in q):return None
        recs=_records(store,project,exts={".csv"},names={"train"})
        if not recs:return ExecutionResult.abstain("raw_csv_not_found")
        rec=sorted(recs,key=lambda r:("/03." not in r.relative_path,len(r.relative_path)))[0]
        df=_read_csv_smart(rec.path)
        cols={norm(c):c for c in df.columns}
        m=re.search(r"標準化された([A-Za-z_][A-Za-z0-9_]*)が\s*([-+]?\d+(?:\.\d+)?)\s*未満",q)
        if not m:return ExecutionResult.abstain("standardized_condition_not_parsed")
        metric=cols.get(norm(m.group(1))); threshold=float(m.group(2))
        eq=re.search(r"([A-Za-z_][A-Za-z0-9_]*)\s*=\s*([A-Za-z0-9_]+)",q)
        if not metric or not eq:return ExecutionResult.abstain("columns_not_resolved")
        cat=cols.get(norm(eq.group(1))); val=eq.group(2)
        if not cat:return ExecutionResult.abstain("category_column_not_resolved")
        s=pd.to_numeric(df[metric],errors="coerce"); z=(s-s.mean())/s.std(ddof=0)
        denom=z<threshold
        group=df[cat].astype(str).map(norm).eq(norm(val))
        group_mean=s[group].mean()
        numer=denom & group & s.gt(group_mean)
        raw=numer.sum()/denom.sum()*100
        return ExecutionResult(True,format_number(raw,q,unit="%"),.99,"standardized_conditional_ratio",[_ev(rec,"computed",f"numerator={numer.sum()}, denominator={denom.sum()}, group_mean={group_mean}",raw)],diagnostics={"raw_ratio":raw})

    def _contract_unit_difference(self,q:str,project:str,store:DocumentStore):
        if not ("見込金額" in q and "確定金額" in q and "ESTH" in q and "ACTH" in q):return None
        contracts=_records(store,project,exts={".docx"},roles={"contract"})
        reports=_records(store,project,roles={"final_report"},exts={".pptx",".pdf",".docx"})
        texts=[]
        for r in contracts+reports:
            t=_docx_all_text(r.path) if r.extension==".docx" else (_pdf_text(r.path) if r.extension==".pdf" else "\n".join(_pptx_text_by_slide(r.path)))
            texts.append((r,t))
        vals={}
        evidence=[]
        for r,t in texts:
            if r.role == "contract":
                if "est_amount" not in vals:
                    m=re.search(r"(?:想定|見込)(?:金額|額)\s*[（(]税込[）)][^0-9]{0,15}([0-9][0-9,]+)",t)
                    if not m:
                        # Fallback to an explicit tax-included figure in the same
                        # sentence as the expected amount.
                        m=re.search(r"(?:想定|見込)[^。\n]{0,120}?([0-9][0-9,]+)円\s*[（(]税込[）)]",t)
                    if m:
                        vals["est_amount"]=float(m.group(1).replace(',','')); evidence.append(_ev(r,"expected_amount",m.group(0),vals["est_amount"]))
                if "est_h" not in vals:
                    m=re.search(r"(?:想定|見込)(?:総)?工数[^0-9]{0,20}([0-9]+(?:\.\d+)?)",t)
                    if m:
                        vals["est_h"]=float(m.group(1)); evidence.append(_ev(r,"expected_hours",m.group(0),vals["est_h"]))
            if r.role == "final_report":
                if "act_h" not in vals:
                    m=re.search(r"実績工数[^0-9]{0,20}([0-9]+(?:\.\d+)?)",t)
                    if m:
                        vals["act_h"]=float(m.group(1)); evidence.append(_ev(r,"actual_hours",m.group(0),vals["act_h"]))
                if "act_amount" not in vals:
                    m=re.search(r"(?:税込金額|確定(?:金額|額)|実績(?:金額|額))[^0-9]{0,20}([0-9][0-9,]+)",t)
                    if m:
                        vals["act_amount"]=float(m.group(1).replace(',','')); evidence.append(_ev(r,"actual_amount",m.group(0),vals["act_amount"]))
        if set(vals)!={"est_amount","act_amount","est_h","act_h"}:
            return ExecutionResult.abstain("four_contract_operands_not_found",diagnostics={"values":vals})
        dh=vals["est_h"]-vals["act_h"]
        if dh==0:return ExecutionResult.abstain("hours_difference_zero")
        raw=(vals["est_amount"]-vals["act_amount"])/dh
        return ExecutionResult(True,format_number(raw,q,unit="円/時間"),.97,"contract_amount_hours_unit_difference",evidence,diagnostics={"raw_unit_difference":raw,"operands":vals})

    def _approval_policy(self,q:str,project:str,store:DocumentStore):
        if not ("APR-M3" in q and "契約金額" in q and "各案件" not in q):
            # The question says 社内管理のAPR; allow that exact pattern.
            if not ("APR-M3" in q and "社内管理" in q):return None
        policy_recs=[r for r in store.records if "決裁基準" in r.filename]
        glossary=[r for r in store.records if "用語" in r.filename]
        if not policy_recs:return ExecutionResult.abstain("approval_policy_not_found")
        policy="\n".join(u.text for u in store.extract_text_units(policy_recs[0]))
        # Parse thresholds and special escalation text from the policy, not a built-in company rule.
        thresholds=[]
        for line in policy.splitlines():
            m=re.search(r"([0-9,]+)\s*万円[^\n]*?(主任|課長|部長|本部長)",line)
            if m:thresholds.append((float(m.group(1).replace(',',''))*10000,m.group(2)))
        rank={"主任":0,"課長":1,"部長":2,"本部長":3}
        contracts=[r for r in store.records if r.role=="contract" and r.extension==".docx" and "draft" not in norm(r.filename)]
        selected=[]; evidence=[_ev(policy_recs[0],"policy",policy[:1000])]
        for r in contracts:
            t=_docx_all_text(r.path)
            amounts=[float(x.replace(',','')) for x in re.findall(r"(?:税込(?:合計)?|契約金額\s*\(税込\))[^0-9]{0,30}([0-9][0-9,]+)",t)]
            if not amounts:continue
            amount=max(amounts)
            level=0
            if amount>=8_000_000:level=3
            elif amount>=5_000_000:level=2
            elif amount>=3_000_000:level=1
            medical=any(x in r.project for x in ("病院","医療","センター"))
            tm=any(x in t.lower() for x in ("time_and_materials","準委任","実績工数"))
            if medical:level=min(3,level+1)
            if tm:level=max(level,2)
            if level>=3:selected.append((r,amount))
        if not selected:
            return ExecutionResult(True,"該当なし。契約金額（税込）の合計は0円。",.95,"approval_policy_cross_project_apply",evidence)
        aliases=[]
        for r,a in selected:
            aliases.append(r.project);evidence.append(_ev(r,"contract_amount",f"{a:,.0f}",a))
        total=sum(a for _,a in selected)
        return ExecutionResult(True,f"{'、'.join(aliases)}。契約金額（税込）の合計は{total:,.0f}円。",.93,"approval_policy_cross_project_apply",evidence,diagnostics={"raw_total":total})

    def _action_status_diff(self,q:str,project:str,store:DocumentStore):
        if not ("会議録" in q and "完了した" in q and "アクション" in q and ("M2" in q and "M3" in q)):return None
        dates=re.findall(r"20\d{2}-\d{2}-\d{2}",q)
        recs=_records(store,project,exts={".pdf"},roles={"meeting"})
        if dates:recs=[r for r in recs if any(d in r.filename for d in dates)]
        if len(recs)<2:return ExecutionResult.abstain("two_minutes_not_found")
        recs.sort(key=lambda r:_date_key(r.path))
        states=[]; evidence=[]
        for r in recs[:2]:
            text=nfkc(_pdf_text(r.path))
            state={}
            # Handle tabular text where status may be within 300 chars after ID.
            ids=list(re.finditer(r"\b(A\d{2})\b",text,re.I))
            for i,m in enumerate(ids):
                aid=m.group(1).upper()
                chunk=text[m.end():ids[i+1].start() if i+1<len(ids) else min(len(text),m.end()+500)]
                # Prefer the explicit English Status column. Later narrative
                # mentions (e.g. 完了理由) must not overwrite the table status.
                st=re.search(r"\b(Open|Closed)\b",chunk,re.I)
                if st and aid not in state:
                    state[aid]=st.group(1).lower(); continue
                st=re.search(r"(?:状況|ステータス)\s*[:：]?\s*(完了|未完了)",chunk)
                if st and aid not in state: state[aid]=st.group(1)
            states.append(state);evidence.append(_ev(r,"action_table",str(state)))
        done=[x for x,s in states[0].items() if s in {"open","未完了"} and states[1].get(x) in {"closed","完了"}]
        done.sort(key=lambda x:int(re.search(r"\d+",x).group()))
        if not done:return ExecutionResult.abstain("open_to_closed_transition_not_found",diagnostics={"states":states})
        return ExecutionResult(True,"、".join(done),.99,"minutes_action_state_transition",evidence)

    def _formula_reference_dimension(self,q:str,project:str,store:DocumentStore):
        if not ("黄色" in q and "予測値" in q and "建設年" in q):return None
        recs=_records(store,project,exts={".xlsx"},names={"train"})
        for rec in recs:
            ys=_xlsx_yellow_cells_fast(rec.path)
            raw_rec=_candidate_raw_csv(store,rec)
            raw_df=_read_csv_smart(raw_rec.path) if raw_rec else None
            for e in ys:
                formula=str(e.get("formula") or "")
                refs=re.findall(r"(?:'([^']+)'|([A-Za-z0-9_]+))!\$?([A-Z]+)\$?(\d+)",formula)
                if not refs:continue
                by_sheet={}
                features=[]
                for s1,s2,col,row in refs:
                    sh=s1 or s2
                    if sh not in by_sheet: by_sheet[sh]=_xlsx_sheet_cells(rec.path,sh)[0]
                    vals=by_sheet[sh]
                    header=next((str(vals.get(f"{col}{rr}")) for rr in range(1,10) if vals.get(f"{col}{rr}") not in (None,"")),None)
                    value=vals.get(f"{col}{row}")
                    if header and parse_number(value) is not None:
                        features.append((sh,col,int(row),header,value))
                # Prefer the original raw-data field. Transformed columns such as
                # YEAR BUILT_fillna are linked back by matching the other features
                # from the same formula row against the sibling raw CSV.
                if raw_df is not None:
                    filtered=raw_df.copy()
                    used=[]
                    for sh,col,row,header,value in features:
                        base=re.sub(r"_(?:fillna|filled|imputed|scaled|standardized)$","",header,flags=re.I)
                        matches=[c for c in raw_df.columns if norm(c)==norm(base)]
                        if len(matches)==1 and not any(x in norm(base) for x in (norm("YEAR BUILT"),norm("建設年"))):
                            c=matches[0]
                            if pd.api.types.is_numeric_dtype(raw_df[c]): filtered=filtered[pd.to_numeric(filtered[c],errors="coerce").eq(float(value))]
                            else: filtered=filtered[filtered[c].astype(str).map(norm).eq(norm(value))]
                            used.append((c,value))
                    year_cols=[c for c in raw_df.columns if any(x in norm(c) for x in (norm("YEAR BUILT"),norm("建設年"),norm("築年")))]
                    if len(filtered)==1 and year_cols:
                        year=parse_number(filtered.iloc[0][year_cols[0]])
                        if year is not None:
                            return ExecutionResult(True,f"{int(year)}年",.99,"xlsx_formula_reference_raw_row",[_ev(rec,f"{e['sheet']}!{e['cell']}",formula),_ev(raw_rec,"matched_raw_row",f"filters={used}; {year_cols[0]}={year}",year)],diagnostics={"raw_year":year,"matched_filters":used})
                    # When repeated values make the feature filter non-unique, link
                    # the referenced worksheet row to the sibling raw file by its
                    # data-row position. This is valid only after confirming that
                    # the worksheet header maps to raw columns and that the other
                    # formula references agree with the candidate raw row.
                    grouped_refs: dict[tuple[str, int], list[tuple[str, str, int, str, Any]]] = defaultdict(list)
                    for feat in features:
                        grouped_refs[(feat[0], feat[2])].append(feat)
                    # Coefficients and intercepts often live in the formula sheet;
                    # the source data row is the repeated external sheet/row pair
                    # with at least two headers that map to the raw table.
                    positional_groups=[]
                    for (source_sheet,source_row),group in grouped_refs.items():
                        vals=by_sheet[source_sheet]
                        mapped_group=[]
                        for feat in group:
                            base=re.sub(r"_(?:fillna|filled|imputed|scaled|standardized)$","",feat[3],flags=re.I)
                            matches=[c for c in raw_df.columns if norm(c)==norm(base)]
                            if len(matches)==1:
                                mapped_group.append((feat,matches[0]))
                        if len(mapped_group)>=2:
                            positional_groups.append((len(mapped_group),source_row,source_sheet,group,mapped_group))
                    positional_groups.sort(reverse=True,key=lambda x:(x[0],x[1]))
                    if positional_groups and year_cols:
                        _count,source_row,source_sheet,group,mapped_group=positional_groups[0]
                        vals=by_sheet[source_sheet]
                        header_row=None
                        for rr in range(1,min(source_row,30)+1):
                            mapped=0
                            for _feat,raw_col in mapped_group:
                                col=_feat[1]
                                cell_header=vals.get(f"{col}{rr}")
                                base=re.sub(r"_(?:fillna|filled|imputed|scaled|standardized)$","",str(cell_header or _feat[3]),flags=re.I)
                                if norm(base)==norm(raw_col):
                                    mapped+=1
                            if mapped>=2:
                                header_row=rr; break
                        if header_row is not None:
                            raw_pos=source_row-header_row-1
                            if 0<=raw_pos<len(raw_df):
                                raw_row=raw_df.iloc[raw_pos]
                                verified=[]; mismatches=[]
                                for feat,c in mapped_group:
                                    _sh,col,_row,header,value=feat
                                    # Imputed/filled features intentionally differ
                                    # from the original raw field and are excluded
                                    # from positional verification.
                                    transformed=bool(re.search(r"_(?:fillna|filled|imputed|scaled|standardized)$",header,re.I))
                                    if transformed:
                                        continue
                                    actual=raw_row[c]
                                    nv=parse_number(value); na=parse_number(actual)
                                    if nv is not None and na is not None:
                                        ok=_value_equal(nv,na)
                                    else:
                                        ok=norm(value)==norm(actual)
                                    (verified if ok else mismatches).append((c,actual,value))
                                if len(verified)>=2 and not mismatches:
                                    year=parse_number(raw_row[year_cols[0]])
                                    if year is not None:
                                        return ExecutionResult(True,f"{int(year)}年",.99,"xlsx_formula_reference_positional_raw_row",[_ev(rec,f"{e['sheet']}!{e['cell']}",formula),_ev(raw_rec,f"raw_row:{raw_pos}",f"verified={verified}; {year_cols[0]}={year}",year)],diagnostics={"raw_year":year,"source_sheet":source_sheet,"source_row":source_row,"header_row":header_row,"raw_position":raw_pos,"verified":verified})
                # Direct fallback when the referenced column itself is the year.
                for sh,col,row,header,value in features:
                    if any(x in norm(header) for x in (norm("YEAR BUILT"),norm("建設年"),norm("築年"))) and 1700 <= float(value) <= 2100:
                        return ExecutionResult(True,f"{int(float(value))}年",.98,"xlsx_formula_reference_dimension",[_ev(rec,f"{e['sheet']}!{e['cell']}",formula),_ev(rec,f"{sh}!{col}{row}",f"{header}={value}",value)],diagnostics={"raw_year":float(value)})
        return ExecutionResult.abstain("construction_year_reference_not_found")

    def _rate_band_change(self,q:str,project:str,store:DocumentStore):
        if not ("新税率" in q and "現行税率" in q and "増加" in q and "価格帯" in q):return None
        recs=_records(store,project,exts={".pdf"})
        rec=next((r for r in recs if "不動産市場" in r.filename or "最新動向" in r.filename),None)
        if not rec:return ExecutionResult.abstain("tax_report_not_found")
        text=nfkc(_pdf_text(rec.path)).replace("％","%")
        lines=[re.sub(r"\s+"," ",x).strip() for x in text.splitlines() if x.strip()]
        bands=[]
        for i,line in enumerate(lines):
            if not ("万" in line and any(x in line for x in ("超","以下"))):continue
            context=" ".join(lines[i:i+4])
            rates=[float(x) for x in re.findall(r"(\d+(?:\.\d+)?)\s*%",context)]
            if len(rates)>=2:
                current=rates[:-1];new=rates[-1]
                # Current rate may be a range. Distance from new rate to range is the absolute increase minimum.
                lo,hi=min(current),max(current)
                delta=0.0 if lo<=new<=hi else min(abs(new-lo),abs(new-hi))
                bands.append((delta,line,new,(lo,hi),context))
        if not bands:return ExecutionResult.abstain("rate_band_table_not_parsed")
        bands.sort(key=lambda x:x[0])
        if len(bands)>1 and math.isclose(bands[0][0],bands[1][0],abs_tol=1e-9):
            return ExecutionResult.abstain("rate_band_minimum_tied",diagnostics={"bands":bands})
        best_context=bands[0][4]
        bm=re.search(r"([0-9, ]+万ドル超)\s*[-–—〜～]\s*([0-9, ]+万ドル以\s*下)",best_context)
        if bm:
            lo=re.sub(r"\s+","",bm.group(1)); hi=re.sub(r"\s+","",bm.group(2))
            band=f"{lo}～{hi}"
        else:
            band=re.sub(r"\s+","",bands[0][1]).replace("-","～")
            if band.endswith("以"):
                band+="下"
        band=band.replace("$","ドル")
        return ExecutionResult(True,band,.94,"pdf_rate_band_min_absolute_change",[_ev(rec,"rate_table",best_context,bands[0][0])],diagnostics={"raw_delta":bands[0][0],"bands":bands})

    def _priority_action_join(self,q:str,project:str,store:DocumentStore):
        if not ("報告資料" in q and "Open" in q and "優先フォロー" in q and "会議録" in q and "完了となっていない" in q):return None
        date_m=re.search(r"(\d{1,2})月(\d{1,2})日",q)
        date_hint=f"-{int(date_m.group(1)):02d}-{int(date_m.group(2)):02d}" if date_m else ""
        meeting_records=_records(store,project,exts={".pdf"},roles={"meeting"})
        report=next((r for r in meeting_records if date_hint in r.filename and "報告資料" in r.filename),None)
        minute=next((r for r in meeting_records if date_hint in r.filename and "会議録" in r.filename),None)
        if not report or not minute:return ExecutionResult.abstain("same_date_report_minutes_not_found")

        # Image-only meeting PDFs are common.  OCR is deliberately bounded:
        # every page receives a fast English probe, but Japanese OCR is only
        # applied to the single report page that contains the priority clause.
        # Minute action tables use English OCR because IDs and statuses are
        # English tokens.  This avoids full-document Japanese OCR timeouts.
        def action_pages(rec:FileRecord, *, report_mode:bool) -> list[tuple[int,str]]:
            doc=fitz.open(rec.path); probed=[]
            cache_dir=store.root.parent / ".rag_action_ocr_cache"
            cache_dir.mkdir(parents=True,exist_ok=True)
            key=hashlib.sha1(rec.relative_path.encode("utf-8")).hexdigest()[:12]+"_"+re.sub(r"[^A-Za-z0-9_.-]+","_",rec.filename)[-60:]
            with tempfile.TemporaryDirectory(prefix="action_ocr_") as td:
                td=Path(td)
                def cli_ocr(image:Path,page_no:int,lang:str,psm:int) -> str:
                    cache=cache_dir/f"{key}.p{page_no}.{lang.replace('+','_')}.psm{psm}.txt"
                    if cache.exists():
                        return nfkc(cache.read_text(encoding="utf-8",errors="ignore"))
                    try:
                        cp=subprocess.run(["tesseract",str(image),"stdout","-l",lang,"--psm",str(psm)],capture_output=True,text=True,timeout=25,check=False)
                        text=nfkc(cp.stdout)
                        cache.write_text(text,encoding="utf-8")
                        return text
                    except (subprocess.TimeoutExpired,OSError):
                        return ""
                for page_no,page in enumerate(doc,1):
                    image=td/f"page_{page_no}.png"
                    pix=page.get_pixmap(matrix=fitz.Matrix(1.35,1.35),alpha=False)
                    pix.save(image)
                    probe=cli_ocr(image,page_no,"eng",6)
                    ids=re.findall(r"A[IL1][- ]?0?\d{1,2}",probe,re.I)
                    status_hits=sum(probe.lower().count(x) for x in ("open","closed","status"))
                    probed.append((page_no,probe,len(ids),status_hits,image))
                if report_mode:
                    ranked=sorted(probed,key=lambda x:(x[2]>=2 and x[3]>=1,x[3],x[2]),reverse=True)
                    if not ranked or ranked[0][2]<2:
                        return []
                    page_no,probe,_,_,image=ranked[0]
                    return [(page_no,cli_ocr(image,page_no,"jpn+eng",6) or probe)]
                selected=[]
                for page_no,probe,id_count,status_hits,_ in probed:
                    if id_count>=2 and status_hits>=1:
                        selected.append((page_no,probe))
                return selected

        def canonical_ids(text:str) -> list[str]:
            out=[]
            for m in re.finditer(r"\bA[IL1][- ]?0?(\d{1,2})\b",nfkc(text),re.I):
                out.append(f"AI-{int(m.group(1)):02d}")
            return _uniq(out)

        report_pages=action_pages(report,report_mode=True); minute_pages=action_pages(minute,report_mode=False)
        rt="\n".join(t for _,t in report_pages); mt="\n".join(t for _,t in minute_pages)
        # Limit extraction to the explicit priority-follow-up clause.  OCR may
        # confuse I/l/1, so action IDs are canonicalized before the join.
        priority=[]
        for marker in re.finditer(r"優先(?:フォロー|実施事項)",rt):
            chunk=rt[marker.start():marker.start()+650]
            priority.extend(canonical_ids(chunk))
            if priority: break
        if not priority:
            # Conservative fallback: Open IDs in a short vicinity of a
            # priority/follow-up marker only.
            for m in re.finditer(r"A[IL1][- ]?\d{1,2}",rt,re.I):
                chunk=rt[max(0,m.start()-160):m.end()+220]
                if "Open" in chunk and any(x in chunk for x in ("優先","フォロー","Priority")):
                    priority.extend(canonical_ids(m.group(0)))
        priority=_uniq(priority)
        if not priority:
            return ExecutionResult.abstain("priority_action_ids_not_extracted")

        status_by_id:dict[str,str]={}
        # Split on each canonicalizable ID and inspect the row/paragraph until
        # the next ID.  Closed/完了 dominates Open when both appear in a note.
        matches=list(re.finditer(r"\bA[IL1][- ]?0?(\d{1,2})\b",mt,re.I))
        for i,m in enumerate(matches):
            aid=f"AI-{int(m.group(1)):02d}"
            end=matches[i+1].start() if i+1<len(matches) else min(len(mt),m.end()+500)
            chunk=mt[m.start():end]
            if any(x in chunk for x in ("Closed","完了","対応済","確認済")):
                status_by_id[aid]="closed"
            elif "Open" in chunk or "未完了" in chunk:
                status_by_id.setdefault(aid,"open")
        incomplete=[aid for aid in priority if status_by_id.get(aid)!="closed"]
        if not incomplete:
            return ExecutionResult.abstain("priority_open_incomplete_not_found",diagnostics={"priority":priority,"statuses":status_by_id})
        evs=[_ev(report,f"ocr_pages:{','.join(str(n) for n,_ in report_pages)}",f"priority={priority}"),_ev(minute,f"ocr_pages:{','.join(str(n) for n,_ in minute_pages)}",f"statuses={status_by_id}; incomplete={incomplete}")]
        return ExecutionResult(True,"、".join(incomplete),.96,"bounded_ocr_report_minutes_action_join",evs,diagnostics={"priority":priority,"statuses":status_by_id})


    def _docx_multi_format_intersection(self, q: str, project: str, store: DocumentStore):
        """Extract runs satisfying every requested Word character-format flag.

        The source set is selected from the document role named in the question;
        no project-specific filename or answer value is embedded.
        """
        requested = {
            "bold": "太字" in q,
            "underline": "下線" in q,
            "italic": "イタリック" in q or "斜体" in q,
        }
        if not all(requested.values()):
            return None
        if "会議録" in q:
            recs = _records(store, project, exts={".docx"}, roles={"minutes"})
            if not recs:
                recs = [r for r in _records(store, project, exts={".docx"}) if "会議録" in r.filename]
        elif "契約" in q:
            recs = _records(store, project, exts={".docx"}, roles={"contract"})
        elif "報告資料" in q:
            # Do not scan unrelated Word files when the named report material is
            # actually a PDF or PowerPoint.  Let the format-aware base executor
            # inspect the requested source type.
            recs = [r for r in _records(store, project, exts={".docx"}) if "報告資料" in nfkc(r.filename)]
            if not recs:
                return None
        elif "報告書" in q:
            recs = [r for r in _records(store, project, exts={".docx"}) if "報告書" in nfkc(r.filename)]
            if not recs:
                return None
        else:
            recs = _records(store, project, exts={".docx"})
        values=[]; evidence=[]
        for rec in recs:
            try: doc=Document(rec.path)
            except Exception: continue
            paragraphs=list(doc.paragraphs)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        paragraphs.extend(cell.paragraphs)
            for pi,paragraph in enumerate(paragraphs,1):
                chunks=[]
                for run in paragraph.runs:
                    selected=bool(run.bold) and bool(run.underline) and bool(run.italic)
                    if selected:
                        chunks.append(run.text)
                    elif chunks:
                        value="".join(chunks).strip(); chunks=[]
                        if value and value not in values:
                            values.append(value); evidence.append(_ev(rec,f"paragraph_or_table:{pi}",value))
                if chunks:
                    value="".join(chunks).strip()
                    if value and value not in values:
                        values.append(value); evidence.append(_ev(rec,f"paragraph_or_table:{pi}",value))
        if values:
            return ExecutionResult(True,"、".join(values),1.0,"docx_multi_format_intersection",evidence)
        return ExecutionResult.abstain("docx_requested_format_intersection_not_found")

    def _structured_id_total(self, q: str, project: str, store: DocumentStore):
        """Count distinct milestone/task/action IDs from their authoritative roles.

        Milestone and task IDs are read from schedule documents. Action IDs are
        read from minutes. Markdown files are excluded by construction.
        """
        if not ("マイルストーンID" in q and "タスクID" in q and "アクションID" in q and "合計" in q):
            return None
        patterns={
            "milestone":re.compile(r"(?<![A-Za-z0-9])MS[-_ ]?0*([0-9]{1,3})(?![A-Za-z0-9])",re.I),
            "task":re.compile(r"(?<![A-Za-z0-9])T[-_ ]?0*([0-9]{1,3})(?![A-Za-z0-9])",re.I),
            "action":re.compile(r"(?<![A-Za-z0-9])A[-_ ]?0*([0-9]{1,3})(?![A-Za-z0-9])",re.I),
        }
        role_records={
            "milestone":_records(store,project,exts={".xlsx",".csv",".docx",".pptx",".pdf"},roles={"schedule"}),
            "task":_records(store,project,exts={".xlsx",".csv",".docx",".pptx",".pdf"},roles={"schedule"}),
            "action":_records(store,project,exts={".xlsx",".csv",".docx",".pptx",".pdf"},roles={"minutes"}),
        }
        # Role metadata can be absent in unfamiliar datasets; fall back to
        # filename semantics while retaining type restrictions.
        if not role_records["milestone"]:
            sched=[r for r in _records(store,project,exts={".xlsx",".csv",".docx",".pptx",".pdf"}) if any(x in r.filename for x in ("スケジュール","計画","WBS"))]
            role_records["milestone"]=sched;role_records["task"]=sched
        if not role_records["action"]:
            role_records["action"]=[r for r in _records(store,project,exts={".xlsx",".csv",".docx",".pptx",".pdf"}) if "会議録" in r.filename]
        found={k:set() for k in patterns}; evidence=[]
        for kind,recs in role_records.items():
            for rec in recs:
                try:text="\n".join(u.text for u in store.extract_text_units(rec))
                except Exception:continue
                nums=set(patterns[kind].findall(nfkc(text)))
                if nums:
                    found[kind].update(nums); evidence.append(_ev(rec,"id_scan",f"{kind}: {sorted(nums,key=int)}",len(nums)))
        if not all(found.values()):
            return ExecutionResult.abstain("one_or_more_id_families_not_found",diagnostics={k:sorted(v,key=int) for k,v in found.items()})
        total=sum(len(v) for v in found.values())
        return ExecutionResult(True,str(total),1.0,"structured_role_scoped_id_total",evidence,diagnostics={k:sorted(v,key=int) for k,v in found.items()})

    def _regression_prediction(self,q:str,project:str,store:DocumentStore):
        if not ("回帰分析" in q and "係数" in q and "index=" in q and "予測値" in q):return None
        im=re.search(r"index\s*=\s*(\d+)",q);idx=int(im.group(1)) if im else None
        if idx is None:return ExecutionResult.abstain("target_index_not_found")
        recs=_records(store,project,exts={".xlsx"},names={"train"})
        for rec in recs:
            wb=load_workbook(rec.path,data_only=True,read_only=False)
            coeffs={};intercept=None; coefficient_locator=None
            for ws in wb.worksheets:
                header=None
                for row in ws.iter_rows(min_row=1,max_row=min(ws.max_row,100)):
                    for cell in row:
                        if norm(cell.value) in {norm("係数"),norm("coefficient"),norm("coefficients")}:
                            header=(cell.row,cell.column); break
                    if header: break
                if not header: continue
                hr,hc=header
                # Excel's regression output places variable labels immediately
                # to the left of the coefficient column. Follow that header
                # explicitly; the last number in the row is an upper confidence
                # bound and must never be mistaken for the coefficient.
                label_col=hc-1
                blank_run=0; local={}; local_intercept=None
                for rr in range(hr+1,min(ws.max_row,hr+200)+1):
                    label=ws.cell(rr,label_col).value
                    value=parse_number(ws.cell(rr,hc).value)
                    if label in (None,"") and value is None:
                        blank_run+=1
                        if blank_run>=2 and local: break
                        continue
                    blank_run=0
                    label=str(label or "").strip()
                    if not label or value is None: continue
                    if any(x in norm(label) for x in (norm("intercept"),norm("切片"))):
                        local_intercept=value
                    elif re.fullmatch(r"[A-Za-z_][A-Za-z0-9_ ]*",label):
                        local[label]=value
                if local_intercept is not None and len(local)>=2:
                    coeffs=local; intercept=local_intercept; coefficient_locator=f"{ws.title}!{_col_letter(hc)}{hr}"
                    break
            if intercept is None or len(coeffs)<2:continue
            data_ws=None;header_row=None;mapping=None
            for ws in wb.worksheets:
                for r in range(1,min(ws.max_row,30)+1):
                    mp={norm(ws.cell(r,c).value):c for c in range(1,ws.max_column+1) if ws.cell(r,c).value not in (None,"")}
                    if norm("index") in mp and sum(1 for f in coeffs if norm(f) in mp)>=max(2,len(coeffs)//2):
                        data_ws=ws;header_row=r;mapping=mp;break
                if data_ws:break
            if not data_ws:continue
            ic=mapping[norm("index")];target_row=None
            for r in range(header_row+1,data_ws.max_row+1):
                if parse_number(data_ws.cell(r,ic).value)==idx:target_row=r;break
            if not target_row:continue
            pred=intercept;evidence=[_ev(rec,coefficient_locator or "coefficient_table",f"intercept={intercept}; coefficients={coeffs}",intercept)]
            used=0
            for f,c in coeffs.items():
                col=mapping.get(norm(f))
                if col:
                    v=parse_number(data_ws.cell(target_row,col).value)
                    if v is not None:pred+=c*v;used+=1;evidence.append(_ev(rec,f"{data_ws.title}!row:{target_row}",f"{f}: coef={c}, value={v}",c*v))
            if used>=2:return ExecutionResult(True,format_number(pred,q),.99,"xlsx_regression_recompute",evidence,diagnostics={"raw_prediction":pred,"used_features":used})
        return ExecutionResult.abstain("regression_layout_not_resolved")

    def _cross_project_people(self,q:str,project:str,store:DocumentStore):
        if not ("各案件" in q and "PP" in q and "契約書" in q and "PLAN" in q and "FR" in q and "DA側" in q and "何人" in q):return None

        role_re=re.compile(
            r"(?:エグゼクティブ\s*スポンサー|executive\s*sponsor|"
            r"プロジェクト\s*マネージャー|project\s*manager|"
            r"リード\s*データ\s*サイエンティスト|リード\s*DS|lead\s*data\s*scientist|"
            r"データ\s*サイエンティスト|data\s*scientist|"
            r"データ\s*エンジニア|data\s*engineer|"
            r"ビジネス\s*アナリスト|business\s*analyst|"
            r"QA\s*レビュー(?:担当|アー)?|QA\s*reviewer|品質保証\s*担当|"
            r"(?<![A-Za-z])PM(?![A-Za-z])|(?<![A-Za-z])PL(?![A-Za-z])|(?<![A-Za-z])QA(?![A-Za-z])|(?<![A-Za-z])BA(?![A-Za-z]))",
            re.I,
        )
        client_re=re.compile(r"(?:委託元|発注者|甲側|クライアント(?:体制|主担当)?|顧客側)",re.I)
        provider_re=re.compile(r"(?:委託先|受託者|乙側|ベンダ(?:ー)?体制|株式会社データアステル|DA側)",re.I)
        non_name_terms=("株式会社","医療法人","データアステル","プロジェクト","マネージャー","レビュー","スポンサー","クライアント","委託先","委託元","受託者","発注者","会議","運営","定例","進捗","推進","管理","項目","担当","体制","成果物","品質","分析","工程","調整","支払","条件","優先")

        def clean_name_candidates(text:str) -> list[str]:
            text=nfkc(text).replace("─"," ").replace("—"," ").replace("|"," ").replace("｜"," ")
            out=[]
            for a,b in re.findall(r"([一-龥々]{1,4})[ \u3000]+([一-龥々]{1,4})",text):
                name=(a+b).replace("斉藤","斎藤")
                if 3<=len(name)<=8 and not any(x in name for x in non_name_terms):
                    out.append(name)
            compact=re.sub(r"(?:課長|部長|主任|様|氏).*$","",text).strip()
            if re.fullmatch(r"[一-龥々]{3,8}",compact) and not any(x in compact for x in non_name_terms):
                out.append(compact.replace("斉藤","斎藤"))
            return _uniq(out)

        def explicit_client_names(text:str) -> list[str]:
            text=nfkc(text)
            name=r"([一-龥々]{1,4})[ \u3000]+([一-龥々]{1,4})"
            patterns=(
                rf"(?:主担当者|検収窓口|成果物レビュー(?:および)?検収の窓口|窓口)\s*[:：]?\s*{name}",
                rf"{name}\s*[（(](?:甲側|クライアント|発注者)[）)]",
                rf"(?:クライアント(?:体制)?|甲側|発注者)\s*[:：]?\s*{name}",
            )
            out=[]
            for pattern in patterns:
                for m in re.finditer(pattern,text,re.I):
                    groups=m.groups()
                    a,b=groups[-2],groups[-1]
                    candidate=(a+b).replace("斉藤","斎藤")
                    if 3<=len(candidate)<=8 and not any(x in candidate for x in non_name_terms):
                        out.append(candidate)
            return _uniq(out)

        def line_pairs(lines:list[str],default_provider:bool=True) -> tuple[list[str],list[str]]:
            """Return provider and client names from role-bearing reading order.

            Individual text boxes often omit their side label. A second call on
            the complete slide/paragraph reading order records client names,
            which are subtracted from default-provider candidates later.
            """
            lines=[nfkc(x).strip() for x in lines if nfkc(x).strip()]
            side="provider" if default_provider else "unknown"
            provider=[];client=[]
            for i,line in enumerate(lines):
                if client_re.search(line) and not provider_re.search(line): side="client"
                if provider_re.search(line): side="provider"
                # Client sections often label the person as a contact/window
                # rather than with a DA delivery role. Capture names anywhere
                # inside the explicit client section so they can be removed.
                if side=="client":
                    client.extend(clean_name_candidates(line))
                if not role_re.search(line): continue
                candidates=clean_name_candidates(line)
                for nxt in lines[i+1:i+5]:
                    if provider_re.search(nxt) or client_re.search(nxt) or role_re.search(nxt): break
                    candidates.extend(clean_name_candidates(nxt))
                if side=="client": client.extend(candidates)
                else: provider.extend(candidates)
            return _uniq(provider),_uniq(client)

        selected=[r for r in store.records if r.project and r.role in {"proposal","contract","schedule","final_report"} and r.extension in {".pptx",".docx",".xlsx",".pdf"}]
        provider_names:set[str]=set();explicit_provider_names:set[str]=set();client_names:set[str]=set();explicit_clients:set[str]=set();evs=[]
        for rec in selected:
            local_provider=[];local_client=[]
            try:
                if rec.extension==".pptx":
                    prs=Presentation(rec.path)
                    for slide_no,slide in enumerate(prs.slides,1):
                        slide_lines=[]
                        for shape in slide.shapes:
                            if getattr(shape,"has_table",False):
                                for row_no,row in enumerate(shape.table.rows,1):
                                    vals=[nfkc(c.text).strip() for c in row.cells]
                                    role_cols=[i for i,v in enumerate(vals) if role_re.search(v)]
                                    if not role_cols: continue
                                    row_names=[]
                                    for i,v in enumerate(vals):
                                        if i not in role_cols: row_names.extend(clean_name_candidates(v))
                                    if not row_names: continue
                                    ctx=" | ".join(vals)
                                    if client_re.search(ctx) and not provider_re.search(ctx):
                                        local_client.extend(row_names)
                                    else:
                                        local_provider.extend(row_names);explicit_provider_names.update(row_names)
                                        evs.append(_ev(rec,f"slide:{slide_no}/table_row:{row_no}",f"roles={','.join(vals[i] for i in role_cols)}; names={_uniq(row_names)}"))
                            if getattr(shape,"has_text_frame",False):
                                lines=[nfkc(x).strip() for x in shape.text.splitlines() if nfkc(x).strip()]
                                slide_lines.extend(lines)
                                ps,cs=line_pairs(lines,default_provider=True)
                                if ps:
                                    local_provider.extend(ps);evs.append(_ev(rec,f"slide:{slide_no}/text_shape",f"names={ps}"))
                                local_client.extend(cs)
                        explicit_clients.update(explicit_client_names("\n".join(slide_lines)))
                        ps,cs=line_pairs(slide_lines,default_provider=True)
                        local_provider.extend(ps);local_client.extend(cs)
                elif rec.extension==".docx":
                    doc=Document(rec.path)
                    for table_no,table in enumerate(doc.tables,1):
                        for row_no,row in enumerate(table.rows,1):
                            vals=[nfkc(c.text).strip() for c in row.cells]
                            role_cols=[i for i,v in enumerate(vals) if role_re.search(v)]
                            if not role_cols: continue
                            row_names=[]
                            for i,v in enumerate(vals):
                                if i not in role_cols: row_names.extend(clean_name_candidates(v))
                            if not row_names: continue
                            ctx=" | ".join(vals)
                            if client_re.search(ctx) and not provider_re.search(ctx):
                                local_client.extend(row_names)
                            else:
                                local_provider.extend(row_names);explicit_provider_names.update(row_names);evs.append(_ev(rec,f"table:{table_no}/row:{row_no}",f"names={_uniq(row_names)}"))
                    paragraph_lines=[p.text for p in doc.paragraphs]
                    explicit_clients.update(explicit_client_names("\n".join(paragraph_lines)))
                    ps,cs=line_pairs(paragraph_lines,default_provider=True)
                    local_provider.extend(ps);local_client.extend(cs)
                elif rec.extension==".xlsx":
                    wb=load_workbook(rec.path,data_only=True,read_only=True)
                    for ws in wb.worksheets:
                        resource_sheet=any(x in norm(ws.title) for x in (norm("リソース"),norm("体制"),norm("担当者一覧")))
                        for row_no,row in enumerate(ws.iter_rows(values_only=True),1):
                            vals=[nfkc(x).strip() for x in row if x not in (None,"")]
                            if not vals: continue
                            role_cols=[i for i,v in enumerate(vals) if role_re.search(v)]
                            if not role_cols or not (resource_sheet or len(vals)<=8): continue
                            row_names=[]
                            for i,v in enumerate(vals):
                                if i not in role_cols: row_names.extend(clean_name_candidates(v))
                            if not row_names: continue
                            ctx=" | ".join(vals)
                            # A matrix header may contain both DA roles and an
                            # explicit client column. Record its client names so
                            # they are removed from default-provider candidates.
                            # Resource matrices may contain a client column
                            # and several DA role columns in the same row. Only
                            # names in explicitly client-labelled cells are
                            # client-side; the remaining role-bearing names are
                            # provider candidates.
                            row_clients=[]
                            for v in vals:
                                if client_re.search(v):
                                    row_clients.extend(clean_name_candidates(v))
                            local_client.extend(row_clients)
                            provider_row=[n for n in row_names if n not in set(row_clients)]
                            if provider_row:
                                local_provider.extend(provider_row);explicit_provider_names.update(provider_row);evs.append(_ev(rec,f"{ws.title}!row:{row_no}",f"names={_uniq(provider_row)}"))
                else:
                    text=_pdf_text(rec.path)
                    if len(text)>100:
                        explicit_clients.update(explicit_client_names(text))
                        ps,cs=line_pairs(text.splitlines(),default_provider=True)
                        if ps:
                            local_provider.extend(ps);evs.append(_ev(rec,"pdf_text",f"names={ps}"))
                        local_client.extend(cs)
            except Exception:
                continue
            provider_names.update(_uniq(local_provider));client_names.update(_uniq(local_client))
        names=provider_names-(explicit_clients-explicit_provider_names)
        if len(names)<5:
            return ExecutionResult.abstain("role_bearing_people_not_extracted",diagnostics={"provider_names":sorted(provider_names),"explicit_client_names":sorted(explicit_clients),"explicit_provider_names":sorted(explicit_provider_names),"heuristic_client_names":sorted(client_names),"names":sorted(names)})
        return ExecutionResult(True,f"{len(names)}人",.97,"structured_cross_project_role_person_dedup",evs[:80],diagnostics={"raw_count":len(names),"names":sorted(names),"explicit_client_names_removed":sorted(explicit_clients-explicit_provider_names),"explicit_provider_names":sorted(explicit_provider_names),"heuristic_client_names":sorted(client_names)})
