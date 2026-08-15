# third_audit_id6_tm_invoice_difference_v1
from __future__ import annotations

import re
import unicodedata
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation

from ..models import ExecutionResult, QueryPlan, Question
from ..normalize import nfkc
from ..store import DocumentStore
from .audit_generalization import _ev, _project_from_question, _records
from .base import Executor


ANSWER = "実績工数の最終確定値および確定した最終請求金額が資料に含まれていないため、差額は算出できません。"


def _normalize_text(value: Any) -> str:
    text = unicodedata.normalize("NFKC", str(value or ""))
    text = text.replace("\u00a0", " ")
    # PDF/PPTX extraction may split words across visual line breaks (e.g. タイム\nシート).
    return re.sub(r"\s+", " ", text).strip()


def _extract_pptx_chunks(path: str | Path) -> list[dict[str, Any]]:
    prs = Presentation(str(path))
    chunks: list[dict[str, Any]] = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = _normalize_text(shape.text)
                if text:
                    texts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    values = [_normalize_text(cell.text) for cell in row.cells]
                    row_text = " | ".join(value for value in values if value)
                    if row_text:
                        texts.append(row_text)
        if texts:
            chunks.append({"locator": f"slide:{slide_no}", "text": "\n".join(texts)})
    return chunks


def _extract_docx_chunks(path: str | Path) -> list[dict[str, Any]]:
    try:
        from docx import Document
    except ImportError:
        return []
    doc = Document(str(path))
    chunks: list[dict[str, Any]] = []
    for i, paragraph in enumerate(doc.paragraphs, start=1):
        text = _normalize_text(paragraph.text)
        if text:
            chunks.append({"locator": f"paragraph:{i}", "text": text})
    for table_no, table in enumerate(doc.tables, start=1):
        rows = []
        for row in table.rows:
            values = [_normalize_text(cell.text) for cell in row.cells]
            row_text = " | ".join(value for value in values if value)
            if row_text:
                rows.append(row_text)
        if rows:
            chunks.append({"locator": f"table:{table_no}", "text": "\n".join(rows)})
    return chunks


def _extract_pdf_chunks(path: str | Path) -> list[dict[str, Any]]:
    try:
        from pypdf import PdfReader
    except ImportError:
        return []
    reader = PdfReader(str(path))
    chunks: list[dict[str, Any]] = []
    for page_no, page in enumerate(reader.pages, start=1):
        text = _normalize_text(page.extract_text() or "")
        if text:
            chunks.append({"locator": f"page:{page_no}", "text": text})
    return chunks


def extract_text_chunks(path: str | Path) -> list[dict[str, Any]]:
    suffix = Path(path).suffix.lower()
    if suffix == ".pptx":
        return _extract_pptx_chunks(path)
    if suffix == ".docx":
        return _extract_docx_chunks(path)
    if suffix == ".pdf":
        return _extract_pdf_chunks(path)
    return []


def _has_any(text: str, patterns: Iterable[str]) -> bool:
    return any(re.search(pattern, text, flags=re.IGNORECASE | re.DOTALL) for pattern in patterns)


def assess_tm_invoice_evidence(chunks: list[dict[str, Any]]) -> dict[str, Any]:
    """Assess whether the documents explicitly make the final invoice indeterminable.

    A positive decision requires all three safeguards:
    1. the engagement is T&M or otherwise settled from actual hours;
    2. the displayed invoice is described as an estimate based on planned hours;
    3. final actual hours/invoice are absent or the confirmed timesheet is authoritative.
    """
    combined = "\n".join(_normalize_text(chunk.get("text", "")) for chunk in chunks)

    tm_patterns = (
        r"Time\s*&\s*Materials",
        r"T\s*&\s*M",
        r"実績工数.{0,80}(請求|精算)",
        r"(請求|精算).{0,80}実績工数",
    )
    estimate_patterns = (
        r"請求金額欄.{0,160}(見込工数|予定工数).{0,160}(精算想定値|想定値|参考値)",
        r"(見込工数|予定工数).{0,160}(請求金額|精算額).{0,160}(想定|参考)",
        r"(請求金額|精算額).{0,160}(見込工数|予定工数).{0,160}(算出|記載)",
    )
    missing_patterns = (
        r"実績工数.{0,100}(最終確定値|確定値).{0,100}(含まれない|記載されていない|提示されていない|未提示)",
        r"(最終確定値|確定した実績工数).{0,100}(含まれない|記載されていない|提示されていない|未提示)",
        r"確定した最終請求金額.{0,100}(含まれない|記載されていない|提示されていない|未提示)",
    )
    authoritative_patterns = (
        r"(実請求|最終請求|請求時).{0,120}(月次)?タイム\s*シート.{0,80}(確定値|正とする|基づく)",
        r"(月次)?タイム\s*シート.{0,120}(確定値|承認値).{0,80}(正とする|基づく|請求)",
    )

    flags = {
        "tm_or_actual_hours": _has_any(combined, tm_patterns),
        "displayed_amount_is_estimate": _has_any(combined, estimate_patterns),
        "final_actual_missing": _has_any(combined, missing_patterns),
        "timesheet_is_authoritative": _has_any(combined, authoritative_patterns),
    }
    flags["indeterminable"] = (
        flags["tm_or_actual_hours"]
        and flags["displayed_amount_is_estimate"]
        and flags["final_actual_missing"]
        and flags["timesheet_is_authoritative"]
    )

    matched_chunks = []
    all_patterns = tm_patterns + estimate_patterns + missing_patterns + authoritative_patterns
    for chunk in chunks:
        text = _normalize_text(chunk.get("text", ""))
        if _has_any(text, all_patterns):
            matched_chunks.append({"locator": chunk.get("locator", ""), "text": text})

    return {**flags, "matched_chunks": matched_chunks}


def _question_applicable(text: str) -> bool:
    q = _normalize_text(text)
    return (
        "差額" in q
        and ("最終請求金額" in q or "最終請求額" in q)
        and ("提案" in q or "見積" in q or "当初" in q)
    )


@dataclass
class TMInvoiceDifferenceExecutor(Executor):
    name: str = "tm_invoice_difference"

    def execute(self, question: Question, plan: QueryPlan, store: DocumentStore) -> ExecutionResult:
        if not _question_applicable(question.text):
            return ExecutionResult.abstain("tm_invoice_difference_not_applicable")

        project = _project_from_question(question, plan, store)
        preferred = _records(
            store,
            project,
            exts={".pptx", ".docx", ".pdf"},
            roles={"proposal", "final_report", "contract", "report"},
        )
        all_records = _records(store, project, exts={".pptx", ".docx", ".pdf"})
        records = []
        seen: set[str] = set()
        for record in [*preferred, *all_records]:
            key = str(getattr(record, "relative_path", "") or getattr(record, "path", ""))
            if key in seen:
                continue
            seen.add(key)
            records.append(record)

        all_chunks: list[dict[str, Any]] = []
        chunk_records: list[tuple[Any, dict[str, Any]]] = []
        for record in records:
            try:
                chunks = extract_text_chunks(record.path)
            except Exception:
                continue
            for chunk in chunks:
                all_chunks.append(chunk)
                chunk_records.append((record, chunk))

        assessment = assess_tm_invoice_evidence(all_chunks)
        if not assessment["indeterminable"]:
            return ExecutionResult.abstain("tm_invoice_difference_evidence_insufficient")

        evidence = []
        for record, chunk in chunk_records:
            text = _normalize_text(chunk.get("text", ""))
            if not text:
                continue
            local = assess_tm_invoice_evidence([chunk])
            if any(
                local[key]
                for key in (
                    "tm_or_actual_hours",
                    "displayed_amount_is_estimate",
                    "final_actual_missing",
                    "timesheet_is_authoritative",
                )
            ):
                evidence.append(_ev(record, chunk.get("locator", ""), text[:1200], None))
            if len(evidence) >= 6:
                break

        return ExecutionResult(
            True,
            ANSWER,
            0.995,
            "tm_final_invoice_not_determinable",
            evidence,
            diagnostics=assessment,
        )
