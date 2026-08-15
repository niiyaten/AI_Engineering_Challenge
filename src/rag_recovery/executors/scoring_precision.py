# third_audit_scoring_precision_v2
from __future__ import annotations

import re
import unicodedata
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation

from ..models import ExecutionResult, QueryPlan, Question
from ..normalize import nfkc, norm
from ..store import DocumentStore
from .audit_generalization import _ev, _project_from_question, _records
from .base import Executor


def _norm_token(value: Any) -> str:
    text = unicodedata.normalize("NFKC", str(value or ""))
    return re.sub(r"[^0-9a-z一-龥ぁ-んァ-ヶ]+", "", text.lower())


def _float_or_none(value: str) -> float | None:
    try:
        return float(unicodedata.normalize("NFKC", value).replace(",", "").strip())
    except (TypeError, ValueError):
        return None


def _find_header_index(headers: list[str], aliases: tuple[str, ...]) -> int | None:
    normalized = [_norm_token(x) for x in headers]
    alias_tokens = [_norm_token(x) for x in aliases]
    for i, token in enumerate(normalized):
        if token in alias_tokens:
            return i
    for i, token in enumerate(normalized):
        if any(alias and (alias in token or token in alias) for alias in alias_tokens):
            return i
    return None


def extract_displayed_ranked_metric(
    pptx_path: str | Path,
    *,
    reference_model: str,
    rank_metric: str = "F1",
    output_metric: str = "Accuracy",
) -> tuple[str, dict[str, Any]] | None:
    """Return the exact displayed cell text for the model ranked after a reference.

    The key behavior is preserving the text shown in the named PowerPoint instead
    of replacing it with a higher-precision value from an attached CSV.
    """
    prs = Presentation(str(pptx_path))
    candidates: list[tuple[float, int, int, str, str, list[list[str]]]] = []

    for slide_no, slide in enumerate(prs.slides, start=1):
        for table_no, shape in enumerate((s for s in slide.shapes if getattr(s, "has_table", False)), start=1):
            rows = [[nfkc(cell.text).strip() for cell in row.cells] for row in shape.table.rows]
            if len(rows) < 2:
                continue
            for header_row in range(min(3, len(rows) - 1)):
                headers = rows[header_row]
                model_col = _find_header_index(headers, ("model", "モデル", "algorithm", "アルゴリズム"))
                rank_col = _find_header_index(headers, (rank_metric, "F1 score", "F1スコア", "primary_value"))
                out_col = _find_header_index(headers, (output_metric, "accuracy", "正解率"))
                if None in (model_col, rank_col, out_col):
                    continue

                parsed: list[tuple[float, str, str, int]] = []
                for row_index, row in enumerate(rows[header_row + 1 :], start=header_row + 2):
                    width = max(model_col, rank_col, out_col)
                    if width >= len(row):
                        continue
                    model = row[model_col].strip()
                    rank_value = _float_or_none(row[rank_col])
                    displayed = row[out_col].strip()
                    if model and rank_value is not None and displayed:
                        parsed.append((rank_value, model, displayed, row_index))
                if len(parsed) < 2:
                    continue

                ranked = sorted(parsed, key=lambda x: (-x[0], _norm_token(x[1])))
                ref_index = next(
                    (i for i, item in enumerate(ranked) if _norm_token(item[1]) == _norm_token(reference_model)),
                    None,
                )
                if ref_index is None or ref_index + 1 >= len(ranked):
                    continue
                target = ranked[ref_index + 1]
                candidates.append(
                    (
                        ranked[ref_index][0],
                        slide_no,
                        table_no,
                        target[2],
                        target[1],
                        rows,
                    )
                )

    if not candidates:
        return None

    _, slide_no, table_no, displayed, target_model, rows = max(candidates, key=lambda x: x[0])
    return displayed, {
        "slide": slide_no,
        "table": table_no,
        "reference_model": reference_model,
        "target_model": target_model,
        "rank_metric": rank_metric,
        "output_metric": output_metric,
        "displayed_value": displayed,
        "table_rows": rows,
    }


def _explicit_pptx_name(question: str) -> str:
    matches = re.findall(r"([^\s、。/\\]+?\.pptx)", nfkc(question), flags=re.IGNORECASE)
    return matches[-1] if matches else ""


@dataclass
class ScoringPrecisionExecutor(Executor):
    name: str = "scoring_precision"

    def execute(self, question: Question, plan: QueryPlan, store: DocumentStore) -> ExecutionResult:
        q = nfkc(question.text)
        applicable = (
            ".pptx" in q.lower()
            and "F1" in q.upper()
            and "Accuracy" in q
            and ("次ぐ" in q or "次" in q)
        )
        if not applicable:
            return ExecutionResult.abstain("scoring_precision_not_applicable")

        ref = re.search(r"([A-Za-z][A-Za-z0-9_-]+)\s*に次ぐ", q)
        if not ref:
            return ExecutionResult.abstain("reference_model_not_found")
        reference_model = ref.group(1)

        project = _project_from_question(question, plan, store)
        records = _records(store, project, exts={".pptx"}, roles={"final_report"})
        if not records:
            records = _records(store, project, exts={".pptx"})

        explicit = _explicit_pptx_name(q)
        if explicit:
            explicit_norm = norm(explicit)
            records = sorted(
                records,
                key=lambda r: (explicit_norm not in norm(r.filename), "old" in norm(r.relative_path), r.relative_path),
            )
        else:
            records = sorted(records, key=lambda r: ("old" in norm(r.relative_path), r.relative_path))

        for rec in records:
            result = extract_displayed_ranked_metric(
                rec.path,
                reference_model=reference_model,
                rank_metric="F1",
                output_metric="Accuracy",
            )
            if not result:
                continue
            displayed, diag = result
            numeric = _float_or_none(displayed)
            locator = f"slide:{diag['slide']}/table:{diag['table']}"
            detail = (
                f"rank=F1; reference={reference_model}; next_model={diag['target_model']}; "
                f"displayed Accuracy={displayed}"
            )
            return ExecutionResult(
                True,
                displayed,
                0.995,
                "pptx_displayed_ranked_metric",
                [_ev(rec, locator, detail, numeric)],
                diagnostics=diag,
            )

        return ExecutionResult.abstain("pptx_displayed_ranked_metric_not_found")
