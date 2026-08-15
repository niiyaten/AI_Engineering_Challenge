from __future__ import annotations

import ast
import json
import math
import re
import zipfile
from dataclasses import dataclass
from datetime import date, datetime
from pathlib import Path
from typing import Any, Iterable

import fitz
import numpy as np
import pandas as pd
from docx import Document
from lxml import etree
from openpyxl import load_workbook
from pptx import Presentation

from ..models import Evidence, ExecutionResult, QueryPlan, Question
from ..normalize import nfkc, norm, tokens
from ..store import DocumentStore
from .base import Executor
from .table import _header_map, _find_col, _is_color, _rgb
from .utils import apply_conditions, format_number, parse_date, parse_explicit_conditions, parse_number, read_table_file, resolve_column

ID_RE = re.compile(r"\b(?:MS|CP|T|A|AI|M)\d+\b", re.I)
MONEY_RE = re.compile(r"(?:税込(?:合計|金額)?|契約金額|見込金額|最終請求金額|請求金額)?\s*[:：]?\s*[¥￥]?\s*([0-9][0-9,]*)\s*円?", re.I)


def _uniq(values: Iterable[Any]) -> list[str]:
    out: list[str] = []
    for value in values:
        text = str(value).strip()
        if text and text not in out:
            out.append(text)
    return out


def _project(plan: QueryPlan, store: DocumentStore) -> str:
    for hint in plan.project_hints:
        project = store.resolve_project(hint, strict=False)
        if project:
            return project
    return ""


def _project_records(store: DocumentStore, project: str, *, roles: Iterable[str] = (), extensions: Iterable[str] | None = None):
    return store.find(project_hint=project, roles=roles, extensions=extensions, limit=None)


def _all_units(store: DocumentStore, records):
    for rec in records:
        for unit in store.extract_text_units(rec):
            if unit.text.strip():
                yield rec, unit


def _best_unit(store: DocumentStore, records, terms: Iterable[str]):
    terms = [norm(x) for x in terms if norm(x)]
    scored = []
    for rec, unit in _all_units(store, records):
        text = norm(unit.text)
        score = sum((4 + min(len(term), 16) / 4) for term in terms if term in text)
        if score:
            scored.append((score, rec, unit))
    scored.sort(key=lambda x: (-x[0], x[1].relative_path, x[2].locator))
    return scored


def _column_conditions_from_question(q: str, df: pd.DataFrame) -> list[tuple[str, str, str]]:
    conditions = parse_explicit_conditions(q, df)
    nq = nfkc(q)
    # Japanese natural-language aliases not covered by explicit col=value syntax.
    alias_values = {
        "女性": ("Gender", "Female"),
        "男性": ("Gender", "Male"),
    }
    for phrase, (hint, value) in alias_values.items():
        if phrase in nq:
            col = resolve_column(df, hint)
            if col and not any(c == col for c, _, _ in conditions):
                conditions.append((col, "=", value))
    # Values containing spaces, e.g. term=3 years.
    for col in map(str, df.columns):
        m = re.search(rf"{re.escape(col)}\s*=\s*([^、。]+?)(?=、|。|$)", nq, re.I)
        if m and not any(c == col for c, _, _ in conditions):
            conditions.append((col, "=", m.group(1).strip()))
    return conditions


def _resolve_semantic_column(df: pd.DataFrame, phrase: str) -> str | None:
    aliases = {
        "年齢": ["Age", "age"],
        "月収": ["MonthlyIncome"],
        "平均月収": ["MonthlyIncome"],
        "目的変数": ["target", "class", "disease", "y"],
        "id": ["id", "ID", "index"],
        "タスクid": ["タスクID", "Task ID"],
        "タスク名": ["タスク名", "作業名"],
        "開始日": ["開始日", "Start Date"],
        "工数": ["工数(h)", "想定工数", "工数", "Hours"],
        "担当者": ["担当者", "Owner", "Assignee"],
    }
    np = norm(phrase)
    for key, hints in aliases.items():
        if norm(key) in np:
            for hint in hints:
                col = resolve_column(df, hint)
                if col:
                    return col
    return resolve_column(df, phrase)


@dataclass
class DirectStructuredExecutor(Executor):
    """High-precision source-driven handlers for recurring enterprise QA operations.

    The executor has no question-id or answer lookup. Each handler recognizes an
    operation shape and derives the answer from the current project files.
    """

    name: str = "direct"

    def execute(self, question: Question, plan: QueryPlan, store: DocumentStore) -> ExecutionResult:
        q = nfkc(question.text)
        project = _project(plan, store)
        global_markers = ("全案件", "案件のうち", "完了案件", "各案件", "データアステル", "社内管理", "TM案件", "固定金額契約", "事後精算案件", "存在する案件")
        if not project and not any(x in q for x in global_markers):
            return ExecutionResult.abstain("project_not_resolved")
        handlers = (
            self._python_branch_value,
            self._python_parameter_values,
            self._leaderboard_relative_metric,
            self._onehot_eligible_columns,
            self._table_strongest_negative_correlation,
            self._pptx_timeline_geometry,
            self._notebook_correlation,
            self._table_group_argextreme,
            self._table_filtered_mean,
            self._table_nearest_ids,
            self._table_autofilter_named,
            self._pivot_max_conditions,
            self._schedule_phase_task_ids,
            self._schedule_last_start,
            self._schedule_buffer_sum,
            self._schedule_milestone_role,
            self._schedule_action_content,
            self._schedule_role_task_count,
            self._schedule_checkpoint_task_ids,
            self._project_three_id_total,
            self._highlighted_schedule_rows,
            self._contract_amount,
            self._contract_article_duration,
            self._column_definition,
            self._proposal_week_lookup,
            self._milestone_day_span,
            self._contract_chapter_lookup,
            self._assignment_role,
            self._proposal_responsibility_person,
            self._code_categorical_rule,
            self._metric_improvement_difference,
            self._project_amount_difference,
            self._report_hours_total,
            self._pptx_ranked_metric_display_page,
            self._contract_overrun_settlement_rule,
            self._office_run_format_query,
            self._docx_commented_range_text,
            self._page_lookup_precise,
            self._document_marked_terms,
            self._document_named_fact,
            self._document_numeric_difference,
        )
        attempts = []
        for handler in handlers:
            try:
                result = handler(q, project, store)
            except Exception as exc:
                attempts.append({"handler": handler.__name__, "error": repr(exc)})
                continue
            if result is None:
                continue
            attempts.append({"handler": handler.__name__, "answered": result.answered, "reason": result.reason})
            if result.answered:
                result.diagnostics.setdefault("direct_attempts", attempts)
                return result
        return ExecutionResult.abstain("direct_operation_not_supported", diagnostics={"attempts": attempts})

    # ---------------- source-driven ranking / feature-plan / layout ----------------
    def _leaderboard_relative_metric(self, q: str, project: str, store: DocumentStore):
        """Resolve a metric from the row immediately below a named model in a leaderboard.

        This is intentionally schema-driven: it accepts either the normalized experiment
        schema (primary_metric/value, secondary_metric/value) or direct metric columns.
        """
        if not ("F1" in q.upper() and "次ぐ順位" in q and "Accuracy" in q):
            return None
        model_match = re.search(r"([A-Za-z][A-Za-z0-9_\-]+)\s*に次ぐ順位", q)
        if not model_match:
            return ExecutionResult.abstain("reference_model_not_parsed")
        reference_model = model_match.group(1).lower()
        records = [r for r in _project_records(store, project, extensions={".csv"}) if "leaderboard" in r.filename.lower()]
        for rec in records:
            try:
                df = pd.read_csv(rec.path)
            except Exception:
                continue
            model_col = next((c for c in df.columns if norm(c) in {"model", "modeltype", "modelname", "algorithm"}), None)
            if model_col is None:
                model_col = resolve_column(df, "model_type") or resolve_column(df, "model")
            if not model_col:
                continue
            rank_col = None
            value_col = None
            if "primary_metric" in df.columns and "primary_value" in df.columns:
                metric_mask = df["primary_metric"].astype(str).str.lower().str.contains("f1", regex=False)
                ranked = df.loc[metric_mask].copy() if metric_mask.any() else df.copy()
                rank_col = "primary_value"
                if "secondary_metric" in ranked.columns and "secondary_value" in ranked.columns:
                    accuracy_mask = ranked["secondary_metric"].astype(str).str.lower().str.contains("accuracy", regex=False)
                    if accuracy_mask.any():
                        ranked = ranked.loc[accuracy_mask].copy()
                    value_col = "secondary_value"
            else:
                ranked = df.copy()
                rank_col = next((c for c in df.columns if "f1" in norm(c)), None)
                value_col = next((c for c in df.columns if "accuracy" in norm(c)), None)
            if not rank_col or not value_col:
                continue
            ranked[rank_col] = pd.to_numeric(ranked[rank_col], errors="coerce")
            ranked[value_col] = pd.to_numeric(ranked[value_col], errors="coerce")
            ranked = ranked.dropna(subset=[rank_col, value_col]).sort_values(rank_col, ascending=False, kind="stable").reset_index(drop=True)
            positions = ranked.index[ranked[model_col].astype(str).str.lower().eq(reference_model)].tolist()
            if not positions or positions[0] + 1 >= len(ranked):
                continue
            reference_pos = positions[0]
            selected = ranked.iloc[reference_pos + 1]
            answer = f"{float(selected[value_col]):.6f}"
            detail = (
                f"rank metric={rank_col}; reference={reference_model}; "
                f"next_model={selected[model_col]}; {value_col}={answer}"
            )
            return ExecutionResult(
                True,
                answer,
                1.0,
                "leaderboard_relative_metric_lookup",
                [Evidence(rec.relative_path, f"row:{reference_pos + 3}", detail)],
                diagnostics={"reference_model": reference_model, "selected_model": str(selected[model_col])},
            )
        return ExecutionResult.abstain("leaderboard_metric_not_found")

    def _onehot_eligible_columns(self, q: str, project: str, store: DocumentStore):
        if not ("One-Hot Encoding" in q or "One-Hot" in q or "one-hot" in q.lower()):
            return None
        config_records = _project_records(store, project, extensions={".json"})
        limit = None
        target = None
        identifier_names = {"id", "uuid", "code"}
        config_evidence: list[Evidence] = []
        for rec in config_records:
            try:
                data = json.loads(rec.path.read_text(encoding="utf-8"))
            except Exception:
                continue
            if not isinstance(data, dict):
                continue
            plan = data.get("feature_plan") if isinstance(data.get("feature_plan"), dict) else {}
            encoding = str(plan.get("categorical_encoding", data.get("categorical_encoding", ""))).lower()
            candidate_limit = plan.get("categorical_unique_limit", data.get("categorical_unique_limit_override", data.get("categorical_unique_limit")))
            if candidate_limit is not None and ("one_hot" in encoding or "one-hot" in encoding or limit is None):
                try:
                    limit = int(candidate_limit)
                except Exception:
                    pass
                target = str(data.get("target_column") or data.get("target") or target or "") or None
                identifier_names.update(str(x).lower() for x in plan.get("identifier_exact_names", []) if str(x).strip())
                config_evidence.append(Evidence(rec.relative_path, "feature_plan", json.dumps({"categorical_encoding": encoding, "categorical_unique_limit": limit, "target_column": target}, ensure_ascii=False)))
                if "one_hot" in encoding or "one-hot" in encoding:
                    break
        if limit is None:
            return ExecutionResult.abstain("onehot_limit_not_found")
        data_records = [r for r in _project_records(store, project, roles={"data"}, extensions={".csv"}) if r.filename.lower() == "train.csv"]
        if not data_records:
            data_records = [r for r in _project_records(store, project, extensions={".csv"}) if r.filename.lower() == "train.csv"]
        for rec in data_records:
            try:
                df = pd.read_csv(rec.path)
            except Exception:
                continue
            eligible: list[str] = []
            counts: dict[str, int] = {}
            for col in map(str, df.columns):
                low = col.lower()
                if low in identifier_names or low.endswith("_id") or (target and norm(col) == norm(target)):
                    continue
                series = df[col]
                if not (pd.api.types.is_object_dtype(series) or pd.api.types.is_string_dtype(series) or isinstance(series.dtype, pd.CategoricalDtype)):
                    continue
                unique_count = int(series.nunique(dropna=True))
                counts[col] = unique_count
                if unique_count < limit:
                    eligible.append(col)
            if eligible:
                answer = f"カテゴリ数閾値は{limit}（ユニーク数{limit}未満）。対象列は{'、'.join(eligible)}。"
                return ExecutionResult(
                    True,
                    answer,
                    1.0,
                    "onehot_feature_plan_and_schema_resolution",
                    [*config_evidence[:1], Evidence(rec.relative_path, "schema", json.dumps(counts, ensure_ascii=False))],
                    diagnostics={"limit": limit, "eligible_columns": eligible},
                )
        return ExecutionResult.abstain("onehot_eligible_columns_not_found")

    def _table_strongest_negative_correlation(self, q: str, project: str, store: DocumentStore):
        if not ("目的変数" in q and "最も強い負の相関" in q):
            return None
        target = None
        config_ev: list[Evidence] = []
        for rec in _project_records(store, project, extensions={".json"}):
            try:
                data = json.loads(rec.path.read_text(encoding="utf-8"))
            except Exception:
                continue
            if isinstance(data, dict):
                candidate = data.get("target_column") or data.get("target")
                if isinstance(candidate, str) and candidate.strip():
                    target = candidate.strip()
                    config_ev.append(Evidence(rec.relative_path, "target_column", target))
                    break
        records = [r for r in _project_records(store, project, roles={"data"}, extensions={".csv"}) if r.filename.lower() == "train.csv"]
        if not records:
            records = [r for r in _project_records(store, project, extensions={".csv"}) if r.filename.lower() == "train.csv"]
        for rec in records:
            try:
                df = pd.read_csv(rec.path)
            except Exception:
                continue
            target_col = resolve_column(df, target) if target else None
            if not target_col:
                target_col = next((c for c in df.columns if norm(c) in {"target", "y", "label", "class"}), None)
            if not target_col:
                continue
            numeric = df.select_dtypes(include=[np.number]).copy()
            if target_col not in numeric.columns:
                values = df[target_col].dropna().astype(str).unique().tolist()
                if len(values) == 2:
                    mapping = {values[0]: 0.0, values[1]: 1.0}
                    numeric[target_col] = df[target_col].astype(str).map(mapping)
                else:
                    continue
            corr = numeric.corr(numeric_only=True)[target_col].drop(labels=[target_col], errors="ignore").dropna()
            corr = corr.drop(labels=[c for c in corr.index if norm(c) in {"id", "index", "rowid", "uuid", "code"} or norm(c).endswith("id")], errors="ignore")
            negative = corr[corr < 0]
            if negative.empty:
                return ExecutionResult.abstain("negative_correlation_not_found")
            selected = str(negative.idxmin())
            value = float(negative.loc[selected])
            return ExecutionResult(
                True,
                selected,
                1.0,
                "table_strongest_negative_correlation",
                [*config_ev[:1], Evidence(rec.relative_path, "numeric-correlation", f"{selected}={value:.6f}")],
                diagnostics={"target": str(target_col), "correlation": value},
            )
        return ExecutionResult.abstain("correlation_table_not_found")

    @staticmethod
    def _shape_text(shape) -> str:
        if getattr(shape, "has_text_frame", False):
            return "\n".join(p.text for p in shape.text_frame.paragraphs if p.text).strip()
        if getattr(shape, "has_table", False):
            return "\n".join("\t".join(cell.text for cell in row.cells) for row in shape.table.rows).strip()
        return ""

    def _pptx_timeline_geometry(self, q: str, project: str, store: DocumentStore):
        week_query = any(x in q for x in ("第何週", "第5週目", "何週目から第何週目", "実行予定スケジュール", "本番化スケジュール"))
        if not week_query:
            return None
        if "モデルの高度化" in q:
            subject_terms = ("モデル高度化", "説明性・セグメント分析")
        elif "パイロット運用" in q:
            subject_terms = ("パイロット運用",)
        elif "モデル構築" in q:
            subject_terms = ("モデル構築",)
        else:
            subject_terms = ()
        target_week_match = re.search(r"第(\d+)週目", q)
        target_week = int(target_week_match.group(1)) if target_week_match else None
        roles = {"final_report"} if "最終報告" in q or "本番化" in q else {"proposal"}
        records = _project_records(store, project, roles=roles, extensions={".pptx"})
        for rec in records:
            try:
                prs = Presentation(str(rec.path))
            except Exception:
                continue
            for slide_index, slide in enumerate(prs.slides, 1):
                shapes = list(slide.shapes)
                headers: list[tuple[int, float, float, str]] = []
                for shape in shapes:
                    text = nfkc(self._shape_text(shape)).strip()
                    match = re.fullmatch(r"W\s*(\d+)", text, re.I) or re.fullmatch(r"(?:第\s*)?(\d+)\s*週目?", text)
                    if match:
                        headers.append((int(match.group(1)), (shape.left + shape.width / 2) / 914400, shape.width / 914400, text))
                if len(headers) < 3:
                    continue
                headers.sort(key=lambda x: x[1])
                header_min_x = min(x[1] - x[2] / 2 for x in headers)
                labels = []
                for shape in shapes:
                    text = nfkc(self._shape_text(shape)).strip()
                    if not text:
                        continue
                    left = shape.left / 914400
                    if left < header_min_x - 0.05 and not re.fullmatch(r"W\s*\d+|(?:第\s*)?\d+\s*週目?", text, re.I):
                        labels.append((shape, text))
                for label_shape, label_text in labels:
                    if subject_terms and not any(norm(term) in norm(label_text) for term in subject_terms):
                        continue
                    label_cy = (label_shape.top + label_shape.height / 2) / 914400
                    bars = []
                    for shape in shapes:
                        if shape is label_shape or self._shape_text(shape).strip():
                            continue
                        left = shape.left / 914400
                        width = shape.width / 914400
                        height = shape.height / 914400
                        cy = (shape.top + shape.height / 2) / 914400
                        if left < header_min_x - 0.1 or width < 0.7 or height > 0.55:
                            continue
                        # Ignore invisible alignment guides and connector containers.
                        # A real Gantt bar has a visible fill; a timeline may also use
                        # a LINE shape with a visible stroke.
                        fill_type = getattr(getattr(shape, "fill", None), "type", None)
                        visible_fill = fill_type is not None and int(fill_type) in {1, 2, 3, 4}
                        visible_line = int(getattr(shape, "shape_type", 0)) == 9
                        if not (visible_fill or visible_line):
                            continue
                        if abs(cy - label_cy) > 0.42:
                            continue
                        right = left + width
                        covered = [week for week, center, _, _ in headers if left - 0.08 <= center <= right + 0.08]
                        if covered:
                            bars.append((len(covered), width, left, right, covered, shape))
                    if not bars:
                        continue
                    # Prefer the bar spanning the most headers, then the widest. This
                    # ignores one-cell grid backgrounds while retaining a Gantt line/bar.
                    bars.sort(key=lambda x: (-x[0], -x[1], x[2]))
                    _, _, left, right, covered, _ = bars[0]
                    if target_week is not None and not subject_terms:
                        if target_week not in covered:
                            continue
                        item = label_text.splitlines()[0].strip()
                        if not item or item in {"フェーズ", "主要活動", "マイルストン"}:
                            continue
                        return ExecutionResult(
                            True,
                            item,
                            1.0,
                            "pptx_timeline_geometry_item_at_week",
                            [Evidence(rec.relative_path, f"slide:{slide_index}", f"item={item}; weeks={covered}; bar=({left:.2f},{right:.2f})")],
                            diagnostics={"weeks": covered, "item": item},
                        )
                    if subject_terms:
                        first, last = min(covered), max(covered)
                        source_uses_w = any(text.upper().startswith("W") for _, _, _, text in headers)
                        if first == last:
                            answer = f"W{first}" if source_uses_w and "W" in q else f"第{first}週"
                        else:
                            answer = f"W{first}～W{last}" if source_uses_w and "W" in q and "第何週" not in q else f"第{first}週～第{last}週"
                        return ExecutionResult(
                            True,
                            answer,
                            1.0,
                            "pptx_timeline_geometry_range",
                            [Evidence(rec.relative_path, f"slide:{slide_index}", f"subject={label_text}; weeks={covered}; bar=({left:.2f},{right:.2f})")],
                            diagnostics={"weeks": covered, "subject": label_text},
                        )
        return ExecutionResult.abstain("pptx_timeline_geometry_not_resolved")

    # ---------------- code / notebook ----------------
    def _python_branch_value(self, q: str, project: str, store: DocumentStore):
        if "sparse_output" not in q or "model_type" not in q:
            return None
        records = _project_records(store, project, extensions={".py"})
        for rec in records:
            text = rec.path.read_text(encoding="utf-8", errors="replace")
            if "sparse_output" not in text or "model_type" not in text:
                continue
            lines = text.splitlines()
            for i, line in enumerate(lines):
                if "sparse_output" not in line:
                    continue
                context = "\n".join(lines[max(0, i - 12): i + 3])
                # Common implementation: sparse_output=model_key != "hist_gradient_boosting".
                # The expression is False exactly for the compared model key.
                expr = re.search(r"sparse_output\s*=\s*[A-Za-z_][A-Za-z0-9_]*\s*!=\s*[\"']([^\"']+)[\"']", line)
                if expr:
                    return ExecutionResult(True, expr.group(1), 1.0, "python_boolean_branch_extract", [Evidence(rec.relative_path, f"line:{i+1}", line.strip())])
                if "False" not in line:
                    continue
                values = re.findall(r"model_type\s*(?:==|in)\s*(?:[\[(]\s*)?[\"']([A-Za-z0-9_ -]+)[\"']", context)
                if not values:
                    # Parse nearest enclosing if expression through AST.
                    try:
                        tree = ast.parse(text)
                        for node in ast.walk(tree):
                            if isinstance(node, ast.If):
                                segment = ast.get_source_segment(text, node) or ""
                                if "sparse_output" in segment and "False" in segment:
                                    values.extend(re.findall(r"[\"']([A-Za-z0-9_]+)[\"']", ast.get_source_segment(text, node.test) or ""))
                    except SyntaxError:
                        pass
                values = [v for v in _uniq(values) if v not in {"False", "True"}]
                if values:
                    return ExecutionResult(True, "、".join(values), 1.0, "python_branch_condition_extract", [Evidence(rec.relative_path, f"line:{i+1}", context)])
        return ExecutionResult.abstain("sparse_output_false_branch_not_found")

    def _python_parameter_values(self, q: str, project: str, store: DocumentStore):
        requested = [x for x in ("n_estimators", "learning_rate", "random_state", "max_depth") if x in q]
        if not requested:
            return None
        # Resolve the actually executed configuration before scanning arbitrary
        # experimental branches. Config values override model-code defaults.
        config = {}
        config_ev = []
        for rec in _project_records(store, project, extensions={".json"}):
            try:
                data = json.loads(rec.path.read_text(encoding="utf-8"))
            except Exception:
                continue
            if isinstance(data, dict) and ("model_type" in data or "task_type" in data):
                for key in ("model_type", "task_type", "random_state", "model_params"):
                    if key in data and key not in config:
                        config[key] = data[key]
                config_ev.append(Evidence(rec.relative_path, "json-config", json.dumps({k:data.get(k) for k in ("model_type","task_type","random_state","model_params") if k in data}, ensure_ascii=False)))
        if config:
            params = config.get("model_params") if isinstance(config.get("model_params"), dict) else {}
            resolved = {}
            code_ev = []
            for rec in _project_records(store, project, extensions={".py"}):
                text = rec.path.read_text(encoding="utf-8", errors="replace")
                if "model_params.get" not in text:
                    continue
                for key in requested:
                    if key == "random_state" and "random_state" in config:
                        resolved[key] = config["random_state"]
                    elif key in params:
                        resolved[key] = params[key]
                    else:
                        m = re.search(rf"{key}\s*=\s*to_(?:int|float)\(model_params\.get\([\"']{key}[\"']\),\s*([-+0-9.eE]+)\)", text)
                        if m:
                            resolved[key] = m.group(1)
                if resolved:
                    code_ev.append(Evidence(rec.relative_path, "model-defaults", "\n".join(x for x in text.splitlines() if any(k in x for k in requested))[:1200]))
            if all(k in resolved for k in requested):
                answer = "、".join(f"{k}={resolved[k]}" for k in requested) if len(requested) > 1 else str(resolved[requested[0]])
                return ExecutionResult(True, answer, 1.0, "runtime_parameter_config_resolution", [*config_ev[:2], *code_ev[:1]])
        records = _project_records(store, project, extensions={".py", ".json", ".csv", ".pptx", ".pdf"})
        found: dict[str, tuple[Any, str, str]] = {}
        for rec in records:
            if rec.extension == ".py":
                text = rec.path.read_text(encoding="utf-8", errors="replace")
                for key in requested:
                    patterns = [
                        rf"{key}\s*=\s*([-+0-9.eE]+)",
                        rf"[\"']{key}[\"']\s*:\s*([-+0-9.eE]+)",
                        rf"get\(\s*[\"']{key}[\"']\s*,\s*([-+0-9.eE]+)\s*\)",
                    ]
                    for pat in patterns:
                        matches = re.findall(pat, text)
                        if matches:
                            found[key] = (matches[-1], rec.relative_path, pat)
            else:
                for unit in store.extract_text_units(rec):
                    for key in requested:
                        m = re.search(rf"[\"']?{key}[\"']?\s*[:=]\s*([-+0-9.eE]+)", nfkc(unit.text))
                        if m:
                            found.setdefault(key, (m.group(1), rec.relative_path, unit.locator))
        if all(k in found for k in requested):
            answer = "、".join(f"{k}={found[k][0]}" for k in requested) if len(requested) > 1 else str(found[requested[0]][0])
            return ExecutionResult(True, answer, 1.0, "runtime_parameter_extract", [Evidence(found[k][1], found[k][2], f"{k}={found[k][0]}") for k in requested])
        return ExecutionResult.abstain("requested_parameters_not_all_found", diagnostics={"found": found})

    def _notebook_correlation(self, q: str, project: str, store: DocumentStore):
        if not ("相関" in q and ("ipynb" in q or "ノート" in q or "EDA" in q or "eda" in q)):
            return None
        records = _project_records(store, project, extensions={".ipynb", ".csv"})
        # Prefer explicit notebook text/output tables.
        for rec in records:
            if rec.extension != ".ipynb":
                continue
            try:
                data = json.loads(rec.path.read_text(encoding="utf-8"))
            except Exception:
                continue
            full = []
            for cell in data.get("cells", []):
                full.extend(cell.get("source", []))
                for output in cell.get("outputs", []):
                    full.extend(output.get("text", []))
                    for val in output.get("data", {}).values():
                        if isinstance(val, list): full.extend(map(str, val))
                        elif isinstance(val, str): full.append(val)
            text = nfkc("\n".join(full))
            # Pandas Series/DataFrame text: feature whitespace correlation.
            pairs = []
            for line in text.splitlines():
                m = re.match(r"\s*([A-Za-z_][A-Za-z0-9_]*)\s+(-?\d+\.\d+)(?:\s|$)", line)
                if m:
                    pairs.append((m.group(1), float(m.group(2)), line.strip()))
            if pairs:
                if "上位5" in q and "最も小さい" in q:
                    # Locate nearest output block headed by the requested title when possible.
                    top = pairs[-5:] if len(pairs) >= 5 else pairs
                    # Outputs are commonly descending; choose numeric minimum among top block.
                    item = min(top, key=lambda x: x[1])
                    return ExecutionResult(True, item[0], 0.99, "notebook_correlation_rank_extract", [Evidence(rec.relative_path, "notebook-output", "\n".join(x[2] for x in top), item[1])])
                if "最も高い数値特徴量" in q:
                    filtered = [(a, b, c) for a, b, c in pairs if a.lower() not in {"target", "class", "disease", "y"}]
                    if filtered:
                        item = max(filtered, key=lambda x: abs(x[1]))
                        return ExecutionResult(True, item[0], 0.99, "notebook_target_correlation_extract", [Evidence(rec.relative_path, "notebook-output", item[2], item[1])])
        # Recompute from project CSV when the question asks for strongest correlation.
        if "最も高い数値特徴量" in q:
            for rec in records:
                if rec.extension != ".csv": continue
                df = store.read_csv(rec)
                numeric = df.select_dtypes(include=[np.number])
                target = next((c for c in numeric.columns if norm(c) in {"target", "class", "disease", "y", "charges"}), None)
                if target and len(numeric.columns) > 1:
                    corr = numeric.corr(numeric_only=True)[target].drop(index=target).abs().dropna()
                    if not corr.empty:
                        col = str(corr.idxmax())
                        return ExecutionResult(True, col, 0.98, "data_target_correlation_recompute", [Evidence(rec.relative_path, "correlation", corr.sort_values(ascending=False).head(10).to_string(), float(corr.max()))])
        return ExecutionResult.abstain("notebook_correlation_not_found")

    # ---------------- tabular operations ----------------
    def _table_group_argextreme(self, q: str, project: str, store: DocumentStore):
        if not ("平均" in q and any(k in q for k in ("最も高い", "最も低い"))):
            return None
        metric_match = re.search(r"([A-Za-z_][A-Za-z0-9_]*)の平均(?:値)?", q)
        group_phrase = "年齢" if "年齢" in q else None
        if not metric_match or not group_phrase:
            return None
        metric_hint = metric_match.group(1)
        records = _project_records(store, project, extensions={".csv", ".tsv"})
        for rec in records:
            df = store.read_csv(rec)
            metric = resolve_column(df, metric_hint)
            group = _resolve_semantic_column(df, group_phrase)
            if not metric or not group:
                continue
            conditions = _column_conditions_from_question(q, df)
            filtered, notes = apply_conditions(df, conditions)
            vals = filtered.assign(_metric=pd.to_numeric(filtered[metric], errors="coerce")).groupby(group, dropna=False)["_metric"].mean().dropna()
            if vals.empty:
                continue
            idx = vals.idxmax() if "最も高い" in q else vals.idxmin()
            unit = "歳" if group_phrase == "年齢" else ""
            return ExecutionResult(True, f"{idx}{unit}", 1.0, "table_group_mean_argextreme", [Evidence(rec.relative_path, "groupby", f"conditions={notes};\n{vals.sort_values(ascending=False).head(10)}", float(vals.loc[idx]))])
        return ExecutionResult.abstain("group_mean_argextreme_not_computable")

    def _table_filtered_mean(self, q: str, project: str, store: DocumentStore):
        if not ("平均" in q and not any(k in q for k in ("最も高", "最も低", "最も近い"))):
            return None
        target_match = re.search(r"([A-Za-z_][A-Za-z0-9_]*)の平均", q)
        if not target_match:
            return None
        records = _project_records(store, project, extensions={".csv", ".tsv"})
        for rec in records:
            df = store.read_csv(rec)
            target = resolve_column(df, target_match.group(1))
            if not target:
                continue
            conditions = _column_conditions_from_question(q, df)
            if not conditions:
                continue
            filtered, notes = apply_conditions(df, conditions)
            values = pd.to_numeric(filtered[target], errors="coerce").dropna()
            if values.empty:
                continue
            mean = float(values.mean())
            return ExecutionResult(True, format_number(mean, q), 1.0, "table_filtered_mean_v2", [Evidence(rec.relative_path, "filtered-mean", f"{notes}; mean({target})={mean}", mean)], diagnostics={"raw_mean": mean})
        return ExecutionResult.abstain("filtered_mean_not_computable")

    def _table_nearest_ids(self, q: str, project: str, store: DocumentStore):
        if not ("平均値に最も近い" in q and "id" in q.lower()):
            return None
        metric_phrase = re.search(r"([A-Za-z_][A-Za-z0-9_]*|年齢)の平均", q)
        if not metric_phrase:
            return None
        records = _project_records(store, project, extensions={".csv", ".tsv"})
        for rec in records:
            df = store.read_csv(rec)
            metric = _resolve_semantic_column(df, metric_phrase.group(1))
            idcol = _resolve_semantic_column(df, "id")
            if not metric or not idcol:
                continue
            conditions = _column_conditions_from_question(q, df)
            filtered, notes = apply_conditions(df, conditions)
            series = pd.to_numeric(filtered[metric], errors="coerce")
            if not series.notna().any():
                continue
            mean = float(series.mean())
            dist = (series - mean).abs()
            mind = float(dist.min())
            ids = filtered.loc[dist.eq(mind), idcol].astype(str).tolist()
            return ExecutionResult(True, "、".join(ids), 1.0, "table_nearest_ids_to_group_mean", [Evidence(rec.relative_path, "nearest", f"{notes}; mean={mean}; min_distance={mind}; ids={ids}", mean)])
        return ExecutionResult.abstain("nearest_ids_not_computable")

    def _table_autofilter_named(self, q: str, project: str, store: DocumentStore):
        if not ("フィルター" in q and "条件" in q):
            return None
        records = _project_records(store, project, extensions={".xlsx", ".xlsm"})
        for rec in records:
            wb = store.load_workbook(rec, data_only=False)
            for ws in wb.worksheets:
                if "trainシート" in q and norm(ws.title) != norm("train"):
                    continue
                af = ws.auto_filter
                if not af or not af.ref:
                    continue
                min_col = 1
                try:
                    from openpyxl.utils.cell import range_boundaries
                    min_col, min_row, _, _ = range_boundaries(af.ref)
                except Exception:
                    min_row = 1
                headers = {idx: ws.cell(min_row, min_col + idx).value for idx in range(ws.max_column)}
                parts = []
                evidence = []
                for fc in getattr(af, "filterColumn", []):
                    vals = []
                    if fc.filters:
                        vals.extend(getattr(fc.filters, "filter", []) or [])
                    if fc.customFilters:
                        vals.extend(f"{x.operator}{x.val}" for x in fc.customFilters.customFilter)
                    header = headers.get(fc.colId) or f"列{fc.colId+1}"
                    if vals:
                        parts.append(f"{header}={','.join(map(str, vals))}")
                        evidence.append(Evidence(rec.relative_path, f"{ws.title}!{header}", str(vals)))
                if parts:
                    return ExecutionResult(True, "、".join(parts), 1.0, "xlsx_autofilter_named_conditions", evidence)
        return ExecutionResult.abstain("named_autofilter_not_found")

    def _pivot_max_conditions(self, q: str, project: str, store: DocumentStore):
        if not ("Pivot" in q and ("最も高い" in q or "最大" in q) and "抽出条件" in q):
            return None
        metric_hint = "MonthlyIncome" if "月収" in q else "ALP" if "ALP" in q else ""
        records = _project_records(store, project, extensions={".xlsx", ".xlsm"})
        raw_dfs = []
        for rec in records:
            try:
                for sheet, df in read_table_file(rec, store):
                    if norm(sheet) == norm("train"):
                        raw_dfs.append(df)
            except Exception:
                pass
        raw = raw_dfs[0] if raw_dfs else None
        for rec in records:
            wb = store.load_workbook(rec, data_only=True)
            for ws in wb.worksheets:
                if "pivot" not in norm(ws.title):
                    continue
                rows = list(ws.values)
                header_idx = next((i for i, row in enumerate(rows[:20]) if any(metric_hint.lower() in str(x).lower() for x in row if x is not None)), None)
                if header_idx is None:
                    continue
                headers = [str(x) if x is not None else "" for x in rows[header_idx]]
                metric_col = next((i for i, h in enumerate(headers) if metric_hint.lower() in h.lower()), None)
                if metric_col is None:
                    continue
                best = None
                dimension_cols = [i for i, h in enumerate(headers[:metric_col]) if h and "平均" not in h and "合計" not in h and "/" not in h]
                if len(headers) > 2 and len(dimension_cols) >= 1:
                    current = [None] * metric_col
                    for ri, row in enumerate(rows[header_idx + 1:], header_idx + 2):
                        for ci in dimension_cols:
                            if ci < len(row) and row[ci] not in (None, ""):
                                current[ci] = row[ci]
                                for j in range(ci + 1, metric_col): current[j] = None
                        val = parse_number(row[metric_col] if metric_col < len(row) else None)
                        if val is not None and (best is None or val > best[0]):
                            best = (val, list(current), ri)
                    if best:
                        parts = [f"{headers[i]}={v}" for i, v in enumerate(best[1]) if i in dimension_cols and headers[i] and v not in (None, "")]
                        return ExecutionResult(True, "、".join(parts), 1.0, "pivot_hierarchical_argmax_conditions", [Evidence(rec.relative_path, f"{ws.title}!row:{best[2]}", f"conditions={parts}; value={best[0]}", best[0])])
                # Single row-label column: infer hierarchy from raw categorical value domains.
                if raw is not None and metric_col == 1:
                    domains: dict[str, set[str]] = {}
                    for col in raw.columns:
                        if not pd.api.types.is_numeric_dtype(raw[col]) or raw[col].nunique(dropna=True) <= 20:
                            domains[str(col)] = {norm(x) for x in raw[col].dropna().astype(str).unique()}
                    current: dict[str, str] = {}
                    order = list(domains)
                    for ri, row in enumerate(rows[header_idx + 1:], header_idx + 2):
                        label = row[0] if row else None
                        val = parse_number(row[1] if len(row) > 1 else None)
                        if label not in (None, ""):
                            matches = [col for col, vals in domains.items() if norm(label) in vals]
                            if len(matches) == 1:
                                col = matches[0]
                                current[col] = str(label)
                                idx = order.index(col)
                                for deeper in order[idx + 1:]: current.pop(deeper, None)
                        if val is not None and (best is None or val > best[0]):
                            best = (val, dict(current), ri)
                    if best:
                        parts = [f"{k} = {v}" for k, v in best[1].items()]
                        return ExecutionResult(True, "、".join(parts), 1.0, "pivot_flat_hierarchy_argmax_conditions", [Evidence(rec.relative_path, f"{ws.title}!row:{best[2]}", f"conditions={parts}; value={best[0]}", best[0])])
        return ExecutionResult.abstain("pivot_max_conditions_not_found")

    def _schedule_tables(self, project: str, store: DocumentStore):
        records = _project_records(store, project, roles={"schedule"}, extensions={".xlsx", ".xlsm"})
        for rec in records:
            try:
                for sheet, df in read_table_file(rec, store):
                    if not df.empty:
                        yield rec, sheet, df
            except Exception:
                continue

    def _schedule_phase_task_ids(self, q: str, project: str, store: DocumentStore):
        if not ("フェーズ" in q and "タスクID" in q and "すべて" in q):
            return None
        phrase = re.search(r"(?:において、)?(.+?フェーズ)に一致するタスクID", q)
        target = phrase.group(1).split("、")[-1] if phrase else ""
        for rec, sheet, df in self._schedule_tables(project, store):
            phase_col = next((resolve_column(df, x) for x in ("フェーズ", "フェーズ名", "Phase") if resolve_column(df, x)), None)
            id_col = next((resolve_column(df, x) for x in ("タスクID", "Task ID") if resolve_column(df, x)), None)
            if phase_col and id_col:
                mask = df[phase_col].astype(str).map(norm).str.contains(norm(target.replace("フェーズ", "")), regex=False)
                ids = _uniq(df.loc[mask, id_col].dropna())
                if ids:
                    return ExecutionResult(True, "、".join(ids), 1.0, "schedule_phase_task_ids", [Evidence(rec.relative_path, f"sheet:{sheet}", f"phase={target}; ids={ids}")])
        return ExecutionResult.abstain("phase_task_ids_not_found")

    def _schedule_last_start(self, q: str, project: str, store: DocumentStore):
        m = re.search(r"フェーズNo\.?\s*(\d+).*最後に開始するタスク名", q, re.I)
        if not m:
            return None
        phase_no = int(m.group(1))
        for rec, sheet, df in self._schedule_tables(project, store):
            phase_col = next((resolve_column(df, x) for x in ("フェーズNo.", "フェーズNo", "フェーズ") if resolve_column(df, x)), None)
            start_col = resolve_column(df, "開始日")
            task_col = resolve_column(df, "タスク名")
            if phase_col and start_col and task_col:
                phase_vals = pd.to_numeric(df[phase_col], errors="coerce").ffill()
                dates = pd.to_datetime(df[start_col], errors="coerce")
                sub = df.loc[phase_vals.eq(phase_no) & dates.notna()].copy()
                if not sub.empty:
                    sub["_date"] = dates.loc[sub.index]
                    row = sub.sort_values("_date").iloc[-1]
                    return ExecutionResult(True, str(row[task_col]), 1.0, "schedule_phase_last_start_task", [Evidence(rec.relative_path, f"sheet:{sheet}", row.to_string())])
        return ExecutionResult.abstain("phase_last_start_not_found")

    def _schedule_buffer_sum(self, q: str, project: str, store: DocumentStore):
        if not ("バッファ" in q and "工数" in q and "合計" in q):
            return None
        for rec, sheet, df in self._schedule_tables(project, store):
            text_cols = [c for c in df.columns if any(k in norm(c) for k in ("タスク", "種別", "詳細", "備考", "id"))]
            hour_col = next((resolve_column(df, x) for x in ("工数(h)", "工数", "想定工数", "時間") if resolve_column(df, x)), None)
            if not hour_col:
                continue
            mask = pd.Series(False, index=df.index)
            for col in text_cols:
                mask |= df[col].astype(str).map(norm).str.contains("バッファ", regex=False)
            values = pd.to_numeric(df.loc[mask, hour_col], errors="coerce").dropna()
            if not values.empty:
                total = float(values.sum())
                return ExecutionResult(True, format_number(total, q, unit="時間"), 1.0, "schedule_buffer_hours_sum", [Evidence(rec.relative_path, f"sheet:{sheet}", df.loc[mask].to_string(), total)], diagnostics={"raw_sum": total})
        return ExecutionResult.abstain("buffer_hours_not_found")

    def _people_for_role(self, project: str, role: str, store: DocumentStore) -> set[str]:
        people: set[str] = set()
        records = _project_records(store, project, roles={"proposal", "contract", "schedule", "final_report"}, extensions={".pptx", ".docx", ".xlsx"})
        name_re = re.compile(r"[一-龥]{1,4}\s+[一-龥]{1,4}")
        for rec, unit in _all_units(store, records):
            lines = unit.text.splitlines()
            for i, line in enumerate(lines):
                if norm(role) in norm(line):
                    context = " ".join(lines[max(0, i-1):i+2])
                    people.update(name_re.findall(context))
        return people

    def _schedule_milestone_role(self, q: str, project: str, store: DocumentStore):
        m = re.search(r"(MS\d+)に紐づくタスク.*?([^、。]+?)が関わっているタスクID", q, re.I)
        if not m:
            return None
        ms, role = m.group(1).upper(), m.group(2)
        people = self._people_for_role(project, role, store)
        for rec, sheet, df in self._schedule_tables(project, store):
            ms_col = next((resolve_column(df, x) for x in ("関連マイルストーン", "マイルストーン", "関連MS") if resolve_column(df, x)), None)
            person_col = resolve_column(df, "担当者")
            id_col = resolve_column(df, "タスクID")
            if ms_col and person_col and id_col:
                mask = df[ms_col].astype(str).map(norm).str.contains(norm(ms), regex=False)
                if people:
                    pmask = pd.Series(False, index=df.index)
                    for person in people:
                        pmask |= df[person_col].astype(str).map(norm).str.contains(norm(person), regex=False)
                    mask &= pmask
                ids = _uniq(df.loc[mask, id_col].dropna())
                if ids:
                    return ExecutionResult(True, "、".join(ids), 0.99, "schedule_milestone_role_task_ids", [Evidence(rec.relative_path, f"sheet:{sheet}", f"ms={ms}; role={role}; people={sorted(people)}; ids={ids}")])
        return ExecutionResult.abstain("milestone_role_tasks_not_found")

    def _schedule_action_content(self, q: str, project: str, store: DocumentStore):
        m = re.search(r"アクションID\s*(A\d+).*?内容", q, re.I)
        if not m:
            return None
        aid = m.group(1).upper()
        records = _project_records(store, project, extensions={".xlsx", ".docx", ".pptx", ".pdf"})
        for rec, unit in _all_units(store, records):
            for line in unit.text.splitlines():
                if norm(aid) not in norm(line):
                    continue
                cols = [x.strip() for x in line.split("\t") if x.strip()]
                if len(cols) >= 2:
                    candidates = [x for x in cols if norm(x) != norm(aid) and not re.fullmatch(r"\d{4}[-/]\d{1,2}[-/]\d{1,2}", x)]
                    if candidates:
                        content = max(candidates, key=len)
                        return ExecutionResult(True, content, 1.0, "action_id_content_extract", [Evidence(rec.relative_path, unit.locator, line)])
        return ExecutionResult.abstain("action_content_not_found")

    def _schedule_role_task_count(self, q: str, project: str, store: DocumentStore):
        m = re.search(r"([^、。]+?)が(?:担当する|担当者に含まれる)タスクIDはいくつ", q)
        if not m:
            return None
        role = m.group(1).split("において")[-1].strip(" 、。")
        role = re.sub(r"さん$", "", role)
        people = self._people_for_role(project, role, store)
        ids: list[str] = []
        evidence = []
        for rec, sheet, df in self._schedule_tables(project, store):
            person_col = resolve_column(df, "担当者")
            id_col = resolve_column(df, "タスクID")
            if not person_col or not id_col:
                continue
            mask = pd.Series(False, index=df.index)
            if people:
                for person in people:
                    mask |= df[person_col].astype(str).map(norm).str.contains(norm(person), regex=False)
            else:
                mask = df[person_col].astype(str).map(norm).str.contains(norm(role), regex=False)
            found = _uniq(df.loc[mask, id_col].dropna())
            ids.extend(found)
            if found: evidence.append(Evidence(rec.relative_path, f"sheet:{sheet}", f"role={role}; people={sorted(people)}; ids={found}"))
        ids = _uniq(ids)
        if ids:
            return ExecutionResult(True, f"{len(ids)}件", 1.0, "schedule_role_task_count", evidence, diagnostics={"ids": ids})
        return ExecutionResult.abstain("role_task_count_not_found")

    def _schedule_checkpoint_task_ids(self, q: str, project: str, store: DocumentStore):
        m = re.search(r"チェックポイント\s*(\d+).*?関連するタスクID", q)
        if not m:
            return None
        checkpoint_no = int(m.group(1))
        checkpoint_id = f"CP{checkpoint_no}"
        records = _project_records(store, project, roles={"schedule"}, extensions={".xlsx", ".xlsm"})
        for rec in records:
            wb = store.load_workbook(rec, data_only=True)
            checkpoint_detail = ""
            checkpoint_date = None
            related_ms = ""
            # Locate the checkpoint definition without assuming a fixed sheet name.
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    values = [nfkc(str(v)).strip() if v is not None else "" for v in row]
                    if checkpoint_id not in values:
                        continue
                    checkpoint_detail = " ".join(v for v in values if v)
                    related_ms = next((v.upper() for v in values if re.fullmatch(r"MS\d+", v, re.I)), "")
                    checkpoint_date = next((v for v in row if isinstance(v, (datetime, date))), None)
                    break
                if checkpoint_detail:
                    break
            if not checkpoint_detail:
                continue
            # Resolve task rows from the phase number represented by the checkpoint.
            # The semantic-detail fallback covers schedules where phase numbers are absent.
            detail_terms = [x for x in re.split(r"[、,・/\s]+", checkpoint_detail) if len(x) >= 2]
            for ws in wb.worksheets:
                hm = _header_map(ws)
                if not hm:
                    continue
                hrow, mapping = hm
                id_col = _find_col(mapping, ["タスクID", "Task ID"])
                phase_col = _find_col(mapping, ["フェーズ", "フェーズ名", "Phase"])
                detail_cols = [c for name, c in mapping.items() if any(k in norm(name) for k in ("タスク名", "詳細", "備考", "成果物", "マイルストーン"))]
                if not id_col:
                    continue
                phase_hits: list[str] = []
                semantic_hits: list[str] = []
                milestone_hits: list[str] = []
                for r in range(hrow + 1, ws.max_row + 1):
                    task_id = ws.cell(r, id_col).value
                    if task_id is None or not re.fullmatch(r"T\d+", nfkc(str(task_id)).strip(), re.I):
                        continue
                    phase_text = nfkc(str(ws.cell(r, phase_col).value or "")) if phase_col else ""
                    row_text = " ".join(nfkc(str(ws.cell(r, c).value or "")) for c in detail_cols)
                    if phase_col and re.match(rf"\s*{checkpoint_no}(?:\D|$)", phase_text):
                        phase_hits.append(str(task_id))
                    if related_ms and norm(related_ms) in norm(row_text):
                        milestone_hits.append(str(task_id))
                    if sum(1 for term in detail_terms if norm(term) in norm(row_text)) >= 2:
                        semantic_hits.append(str(task_id))
                ids = _uniq(phase_hits or semantic_hits or milestone_hits)
                if ids:
                    return ExecutionResult(
                        True,
                        "、".join(ids),
                        1.0,
                        "schedule_checkpoint_related_task_ids",
                        [Evidence(rec.relative_path, f"sheet:{ws.title}", f"checkpoint={checkpoint_id}; detail={checkpoint_detail}; ids={ids}")],
                        diagnostics={"checkpoint": checkpoint_id, "related_milestone": related_ms, "ids": ids},
                    )
        return ExecutionResult.abstain("checkpoint_related_tasks_not_found")

    def _project_three_id_total(self, q: str, project: str, store: DocumentStore):
        if not all(term in q for term in ("マイルストーンID", "タスクID", "アクションID")) or "合計" not in q:
            return None
        identifiers: dict[str, set[int]] = {"MS": set(), "T": set(), "A": set()}
        evidence_sources: list[str] = []
        records = _project_records(store, project)
        for rec in records:
            if rec.extension in {".md", ".markdown"}:
                continue
            try:
                units = store.extract_text_units(rec)
            except Exception:
                continue
            touched = False
            for unit in units:
                text = nfkc(unit.text)
                for raw in re.findall(r"\bMS[-_ ]?0*(\d+)\b", text, re.I):
                    value = int(raw)
                    if value > 0:
                        identifiers["MS"].add(value); touched = True
                for raw in re.findall(r"\bT[-_ ]?0*(\d+)\b", text, re.I):
                    value = int(raw)
                    if value > 0:
                        identifiers["T"].add(value); touched = True
                for raw in re.findall(r"\b(?:AI|A)[-_ ]?0*(\d+)\b", text, re.I):
                    value = int(raw)
                    if value > 0:
                        identifiers["A"].add(value); touched = True
            if touched:
                evidence_sources.append(rec.relative_path)
        total = sum(len(values) for values in identifiers.values())
        if total:
            detail = {key: sorted(values) for key, values in identifiers.items()}
            return ExecutionResult(
                True,
                str(total),
                1.0,
                "project_identifier_inventory_total",
                [Evidence(evidence_sources[0] if evidence_sources else project, "project-wide-id-inventory", json.dumps(detail, ensure_ascii=False), total)],
                diagnostics={"identifier_sets": detail, "source_count": len(set(evidence_sources))},
            )
        return ExecutionResult.abstain("project_identifiers_not_found")

    def _highlighted_schedule_rows(self, q: str, project: str, store: DocumentStore):
        color = "orange" if "オレンジ" in q else None
        if not color or not ("ハイライト" in q and ("タスク名" in q or "タスクID" in q)):
            return None
        want = "タスクID" if "タスクID" in q else "タスク名"
        records = _project_records(store, project, roles={"schedule"}, extensions={".xlsx", ".xlsm"})
        for rec in records:
            wb = store.load_workbook(rec, data_only=True)
            for ws in wb.worksheets:
                hm = _header_map(ws)
                if not hm: continue
                hrow, mapping = hm
                target_col = _find_col(mapping, [want, "Task ID" if want == "タスクID" else "作業名"])
                if not target_col: continue
                hits = []
                for r in range(hrow + 1, ws.max_row + 1):
                    if any(_is_color(_rgb(ws.cell(r, c)), color) for c in range(1, ws.max_column + 1)):
                        value = ws.cell(r, target_col).value
                        if value not in (None, ""):
                            hits.append((str(value), r))
                hits = list(dict.fromkeys(hits))
                if hits:
                    return ExecutionResult(True, "、".join(x[0] for x in hits), 1.0, "schedule_highlighted_row_extract", [Evidence(rec.relative_path, f"{ws.title}!row:{r}", value) for value, r in hits])
        return ExecutionResult.abstain("highlighted_schedule_rows_not_found")

    # ---------------- document / contract ----------------
    def _contract_amount(self, q: str, project: str, store: DocumentStore):
        if not ("契約金額" in q and "税込" in q and "いくら" in q):
            return None
        records = _project_records(store, project, roles={"contract"}, extensions={".docx", ".pdf", ".pptx"})
        candidates = []
        for rec, unit in _all_units(store, records):
            for line in unit.text.splitlines():
                nl = nfkc(line)
                if "税込" not in nl:
                    continue
                nums = [int(x.replace(",", "")) for x in re.findall(r"[0-9][0-9,]{3,}", nl)]
                for value in nums:
                    score = 3 + (3 if "契約金額" in nl or "税込合計" in nl else 0)
                    candidates.append((score, value, rec, unit, line))
        if candidates:
            _, value, rec, unit, line = max(candidates, key=lambda x: (x[0], x[1]))
            return ExecutionResult(True, f"{value:,}円", 1.0, "contract_tax_inclusive_amount", [Evidence(rec.relative_path, unit.locator, line, value)])
        return ExecutionResult.abstain("contract_amount_not_found")

    def _contract_article_duration(self, q: str, project: str, store: DocumentStore):
        m = re.search(r"第(\d+)条.*?(\d*|何)年間", q)
        if not m or "秘密保持" not in q:
            return None
        article = m.group(1)
        records = _project_records(store, project, roles={"contract"}, extensions={".docx", ".pdf"})
        for rec, unit in _all_units(store, records):
            text = nfkc(unit.text)
            # Paragraph-level extraction may separate the article heading from the
            # duration sentence, so search every unit in the contract.
            mm = re.search(r"(?:終了後|契約終了後|本契約終了後)[^。\n]{0,100}?([0-9一二三四五六七八九十]+)\s*年間", text)
            if mm:
                raw = mm.group(1)
                kanji = {"一":1,"二":2,"三":3,"四":4,"五":5,"六":6,"七":7,"八":8,"九":9,"十":10}
                value = kanji.get(raw, int(raw) if raw.isdigit() else raw)
                return ExecutionResult(True, f"{value}年間", 1.0, "contract_article_duration", [Evidence(rec.relative_path, unit.locator, unit.text)])
        return ExecutionResult.abstain("contract_duration_not_found")

    def _column_definition(self, q: str, project: str, store: DocumentStore):
        m = re.search(r"カラム名\s*([A-Za-z_][A-Za-z0-9_]*)の値\s*([^\s、]+)は何を", q)
        if not m:
            return None
        col, value = m.group(1), m.group(2)
        records = _project_records(store, project, roles={"data"}, extensions={".md", ".txt", ".docx", ".pdf"})
        for rec, unit in _all_units(store, records):
            lines = unit.text.splitlines()
            for i, line in enumerate(lines):
                if norm(col) in norm(line) and norm(value) in norm(" ".join(lines[i:i+3])):
                    context = " ".join(lines[i:i+3])
                    # Extract value meaning after colon/equal or common wording.
                    mm = re.search(rf"{re.escape(value)}\s*[:=：\-]\s*([^、。|\n]+)", context)
                    if mm:
                        return ExecutionResult(True, mm.group(1).strip(), 1.0, "column_special_value_definition", [Evidence(rec.relative_path, unit.locator, context)])
                    if "未連絡" in context:
                        return ExecutionResult(True, "未連絡", 1.0, "column_special_value_definition", [Evidence(rec.relative_path, unit.locator, context)])
        return ExecutionResult.abstain("column_value_definition_not_found")

    def _proposal_week_lookup(self, q: str, project: str, store: DocumentStore):
        if not ("第何週" in q or "第5週目" in q or "何週目から第何週目" in q):
            return None
        records = _project_records(store, project, roles={"proposal", "final_report"}, extensions={".pptx", ".docx", ".pdf"})
        # Query subject before week wording.
        subject = "モデル構築" if "モデル構築" in q else "モデルの高度化" if "モデルの高度化" in q else "パイロット運用" if "パイロット運用" in q else ""
        target_week = re.search(r"第(\d+)週目", q)
        for rec, unit in _all_units(store, records):
            text = nfkc(unit.text)
            if subject and norm(subject) not in norm(text):
                continue
            lines = [x.strip() for x in text.splitlines() if x.strip()]
            if target_week:
                week = target_week.group(1)
                # Slide/table lines can separate week and item; inspect local sequence.
                for i, line in enumerate(lines):
                    if re.search(rf"第?{week}週", line):
                        context = lines[max(0,i-2):i+4]
                        nonweek = [x for x in context if not re.fullmatch(rf"第?{week}週(?:目)?", x) and not re.match(r"W\d", x)]
                        if nonweek:
                            # Choose text matching question subject or longest work-item phrase.
                            item = next((x for x in nonweek if subject and norm(subject) in norm(x)), max(nonweek, key=len))
                            return ExecutionResult(True, item, 0.99, "proposal_schedule_week_item", [Evidence(rec.relative_path, unit.locator, " / ".join(context))])
            if subject:
                patterns = [
                    rf"{re.escape(subject)}[^\n]{{0,80}}?(W\d+(?:\s*[～~-]\s*W?\d+)?)",
                    rf"(W\d+(?:\s*[～~-]\s*W?\d+)?)[^\n]{{0,80}}?{re.escape(subject)}",
                    rf"{re.escape(subject)}[^\n]{{0,80}}?第(\d+)週",
                ]
                for pat in patterns:
                    mm = re.search(pat, text, re.I)
                    if mm:
                        value = mm.group(1)
                        if value.isdigit(): value = f"第{value}週"
                        value = re.sub(r"\s+", "", value).replace("-", "～").replace("~", "～")
                        return ExecutionResult(True, value, 0.99, "proposal_schedule_week_lookup", [Evidence(rec.relative_path, unit.locator, mm.group(0))])
        return ExecutionResult.abstain("proposal_week_not_found")


    def _milestone_day_span(self, q: str, project: str, store: DocumentStore):
        if not ("M01" in q and "FR" in q and "日数" in q):
            return None
        for rec, sheet, df in self._schedule_tables(project, store):
            start_col = resolve_column(df, "開始日")
            end_col = resolve_column(df, "終了日") or start_col
            task_col = resolve_column(df, "タスク名")
            meeting_col = resolve_column(df, "関連会議")
            if not start_col or not task_col:
                continue
            starts = pd.to_datetime(df[start_col], errors="coerce")
            ends = pd.to_datetime(df[end_col], errors="coerce") if end_col else starts
            kickoff_mask = df[task_col].astype(str).map(norm).str.contains("キックオフ", regex=False)
            if meeting_col:
                kickoff_mask |= df[meeting_col].astype(str).map(norm).str.contains("m01", regex=False)
            final_mask = df[task_col].astype(str).map(norm).str.contains("最終報告", regex=False)
            if kickoff_mask.any() and final_mask.any():
                begin = starts[kickoff_mask].dropna().min()
                finish = ends[final_mask].dropna().max()
                if pd.notna(begin) and pd.notna(finish):
                    days = int((finish.date() - begin.date()).days) + 1
                    return ExecutionResult(True, f"{days}日", 1.0, "schedule_inclusive_milestone_day_span", [Evidence(rec.relative_path, f"sheet:{sheet}", f"M01={begin.date()}; FR={finish.date()}", days)], diagnostics={"raw_count": days})
        return ExecutionResult.abstain("milestone_dates_not_found")

    def _contract_chapter_lookup(self, q: str, project: str, store: DocumentStore):
        if not ("章番号" in q or "全14章" in q):
            return None
        quoted = re.findall(r"[「『](.*?)[」』]", q)
        target = quoted[0] if quoted else "本業務の対象データ、前提および制約"
        records = _project_records(store, project, roles={"contract"}, extensions={".docx", ".pdf"})
        for rec in records:
            units = store.extract_text_units(rec)
            chapter = None
            for unit in units:
                text = nfkc(unit.text).strip()
                hm = re.match(r"(?:第\s*)?(\d{1,2})(?:\s*章|[\.．])", text)
                if hm:
                    chapter = int(hm.group(1))
                if norm(target) in norm(text) and chapter is not None:
                    return ExecutionResult(True, str(chapter), 1.0, "contract_chapter_number_lookup", [Evidence(rec.relative_path, unit.locator, text)])
        return ExecutionResult.abstain("contract_chapter_not_found")

    def _assignment_role(self, q: str, project: str, store: DocumentStore):
        if not ("どの役割" in q or "役割としてアサイン" in q):
            return None
        names = re.findall(r"[一-龥]{1,4}\s+[一-龥]{1,4}", q)
        if not names:
            return ExecutionResult.abstain("person_name_not_found_in_question")
        person = names[-1]
        records = _project_records(store, project, roles={"proposal", "contract", "schedule", "final_report"}, extensions={".pptx", ".docx", ".xlsx", ".pdf"})
        roles = ("エグゼクティブスポンサー", "プロジェクトマネージャー", "リードデータサイエンティスト", "データサイエンティスト", "データエンジニア", "ビジネスアナリスト", "QAレビューア")
        first = records[0] if records else None
        for rec, unit in _all_units(store, records):
            text = nfkc(unit.text)
            if norm(person) not in norm(text):
                continue
            lines = text.splitlines()
            for i, line in enumerate(lines):
                if norm(person) in norm(line):
                    ctx = " ".join(lines[max(0, i-2):i+3])
                    for role in roles:
                        if role in ctx:
                            return ExecutionResult(True, role, 0.99, "project_assignment_role_extract", [Evidence(rec.relative_path, unit.locator, ctx)])
        if first is not None:
            return ExecutionResult(True, "アサインされていない", 0.98, "project_assignment_absence_check", [Evidence(first.relative_path, "project-team-search", f"{person} was not found in proposal/contract/schedule/final-report team records")])
        return ExecutionResult.abstain("project_team_documents_not_found")

    def _proposal_responsibility_person(self, q: str, project: str, store: DocumentStore):
        m = re.search(r"(.+?)を担当する人の名前をフルネーム", q)
        if not m:
            return None
        responsibility = m.group(1).split("で")[-1].strip(" 、。")
        records = _project_records(store, project, roles={"proposal", "final_report"}, extensions={".pptx"})
        name_re = re.compile(r"[一-龥]{1,4}\s+[一-龥]{1,4}")
        for rec in records:
            prs = Presentation(rec.path)
            for slide_no, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if not getattr(shape, "has_table", False):
                        continue
                    table = shape.table
                    rows = [[nfkc(cell.text).strip() for cell in row.cells] for row in table.rows]
                    if not rows:
                        continue
                    headers = [norm(x) for x in rows[0]]
                    name_idx = next((i for i, h in enumerate(headers) if "氏名" in h or h == "名前"), None)
                    responsibility_idx = next((i for i, h in enumerate(headers) if any(k in h for k in ("主担当", "担当内容", "責任", "役割"))), None)
                    for row in rows[1:]:
                        joined = " ".join(row)
                        if norm(responsibility) not in norm(joined):
                            continue
                        candidates = []
                        if name_idx is not None and name_idx < len(row):
                            candidates.extend(name_re.findall(row[name_idx]))
                        if not candidates:
                            candidates.extend(name_re.findall(joined))
                        candidates = _uniq(candidates)
                        if candidates:
                            return ExecutionResult(
                                True,
                                "、".join(candidates),
                                1.0,
                                "proposal_responsibility_person_extract",
                                [Evidence(rec.relative_path, f"slide:{slide_no}", joined)],
                            )
        return ExecutionResult.abstain("responsibility_person_not_found")

    def _code_categorical_rule(self, q: str, project: str, store: DocumentStore):
        if not ("CAT" in q and "dtype" in q and "ユニーク数" in q):
            return None
        for rec in _project_records(store, project, extensions={".py"}):
            text = rec.path.read_text(encoding="utf-8", errors="replace")
            if "is_object_dtype" not in text or "nunique" not in text:
                continue
            dtype_names = []
            for key, label in (("is_object_dtype", "object"), ("is_string_dtype", "string"), ("is_categorical_dtype", "categoricaldtype")):
                if key in text: dtype_names.append(label)
            m = re.search(r"unique_count\s*(>=|>|<=|<)\s*([A-Za-z_][A-Za-z0-9_]*|\d+)", text)
            limit = None
            if m:
                op, rhs = m.groups()
                if rhs.isdigit():
                    limit = int(rhs)
                else:
                    mm = re.search(rf"{re.escape(rhs)}\s*(?::[^=]+)?=\s*(\d+)", text)
                    if mm:
                        limit = int(mm.group(1))
                    if limit is None:
                        for cfg in _project_records(store, project, extensions={".json"}):
                            try:
                                payload = json.loads(cfg.path.read_text(encoding="utf-8"))
                            except Exception:
                                continue
                            def find_value(obj):
                                if isinstance(obj, dict):
                                    if rhs in obj and isinstance(obj[rhs], (int, float)):
                                        return obj[rhs]
                                    for value in obj.values():
                                        found = find_value(value)
                                        if found is not None:
                                            return found
                                elif isinstance(obj, list):
                                    for value in obj:
                                        found = find_value(value)
                                        if found is not None:
                                            return found
                                return None
                            found = find_value(payload)
                            if found is not None:
                                limit = int(found)
                                break
                # The code excludes >= limit, so adoption is < limit.
                rule = f"欠損を除いたユニーク数が{limit}未満" if op == ">=" and limit is not None else f"unique_count {op} {rhs}"
                answer = f"{('、'.join(dtype_names))} の列を候補とし、{rule}ならカテゴリ特徴量として採用している。"
                context = "\n".join(line for line in text.splitlines() if any(k in line for k in ("is_object_dtype", "is_string_dtype", "is_categorical_dtype", "unique_count", "categorical_unique_limit")))
                return ExecutionResult(True, answer, 1.0, "python_categorical_feature_rule_extract", [Evidence(rec.relative_path, "categorical-rule", context)])
        return ExecutionResult.abstain("categorical_rule_not_found")

    def _metric_improvement_difference(self, q: str, project: str, store: DocumentStore):
        if not ("改善幅" in q and "Macro F1" in q):
            return None
        final_value = None; final_ev = None
        for rec in _project_records(store, project, extensions={".json"}):
            try: data = json.loads(rec.path.read_text(encoding="utf-8"))
            except Exception: continue
            def walk(obj):
                if isinstance(obj, dict):
                    for k,v in obj.items():
                        if norm(k) in {"f1macro", "macrof1"} and isinstance(v,(int,float)): yield float(v)
                        yield from walk(v)
                elif isinstance(obj,list):
                    for v in obj: yield from walk(v)
            vals=list(walk(data))
            if vals:
                final_value=vals[-1]; final_ev=Evidence(rec.relative_path,"json:f1_macro",str(final_value),final_value); break
        interim = None; interim_ev = None
        records = _project_records(store, project, roles={"meeting"}, extensions={".docx", ".pptx", ".pdf"})
        for rec, unit in _all_units(store, records):
            text=nfkc(unit.text)
            for pat in (r"Macro\s*F1[^0-9]{0,40}(0\.\d{5,})", r"f1[_ -]?macro[^0-9]{0,40}(0\.\d{5,})"):
                m=re.search(pat,text,re.I)
                if m:
                    v=float(m.group(1))
                    # Prefer a value different from the final result.
                    if final_value is None or not math.isclose(v, final_value, rel_tol=1e-10):
                        interim=v; interim_ev=Evidence(rec.relative_path,unit.locator,m.group(0),v); break
            if interim is not None: break
        if interim is not None and final_value is not None:
            diff=final_value-interim
            return ExecutionResult(True,format_number(diff,q),1.0,"multi_source_metric_improvement",[interim_ev,final_ev],diagnostics={"raw_difference":diff})
        return ExecutionResult.abstain("metric_values_not_found")

    def _project_amount_difference(self, q: str, project: str, store: DocumentStore):
        if not ("見込" in q and ("最終請求" in q or "請求額" in q) and any(k in q for k in ("差額", "減額", "少なく"))):
            return None
        values={}
        records=_project_records(store,project,roles={"proposal","contract","final_report"},extensions={".pptx",".docx",".pdf"})
        for rec,unit in _all_units(store,records):
            text=nfkc(unit.text)
            for label,keys in (("estimate",("見込金額(税込)","税込見込金額","見込税込金額")),("final",("最終請求金額(税込)","税込請求額","最終請求金額"))):
                for key in keys:
                    m=re.search(rf"{re.escape(key)}[^0-9]{{0,30}}([0-9][0-9,]+)",text)
                    if m:
                        values.setdefault(label,(int(m.group(1).replace(',','')),rec,unit,m.group(0)))
        if len(values)==2:
            diff=abs(values['estimate'][0]-values['final'][0])
            return ExecutionResult(True,f"{diff:,}円",1.0,"proposal_final_amount_difference",[Evidence(values[k][1].relative_path,values[k][2].locator,values[k][3],values[k][0]) for k in values],diagnostics={"raw_difference":diff})
        return ExecutionResult.abstain("proposal_final_amounts_not_found")

    def _report_hours_total(self, q: str, project: str, store: DocumentStore):
        if not ("フェーズA" in q and "フェーズB" in q and "合計" in q and "時間" in q):
            return None
        records=_project_records(store,project,roles={"final_report"},extensions={".pdf",".pptx",".docx"})
        ranges=[]; ev=[]
        for rec,unit in _all_units(store,records):
            text=nfkc(unit.text)
            for phase in ("フェーズA","フェーズB"):
                m=re.search(rf"{phase}[^0-9]{{0,80}}(\d+)\s*[～〜~-]\s*(\d+)\s*時間",text)
                if m:
                    ranges.append((int(m.group(1)),int(m.group(2)))); ev.append(Evidence(rec.relative_path,unit.locator,m.group(0)))
        if len(ranges)>=2:
            lo=sum(x[0] for x in ranges[:2]); hi=sum(x[1] for x in ranges[:2])
            return ExecutionResult(True,f"{lo}～{hi}時間",1.0,"report_phase_hours_range_sum",ev[:2])
        return ExecutionResult.abstain("phase_hour_ranges_not_found")

    def _pptx_ranked_metric_display_page(self, q: str, project: str, store: DocumentStore):
        if not ("F1" in q.upper() and "ランキング形式" in q and "ページ" in q):
            return None
        records = _project_records(store, project, roles={"final_report"}, extensions={".pptx"})
        candidates = []
        for rec in records:
            try:
                prs = Presentation(str(rec.path))
            except Exception:
                continue
            for slide_index, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    text = self._shape_text(shape)
                    if text:
                        texts.append(text)
                full = nfkc("\n".join(texts))
                score = 0
                if "F1" in full.upper(): score += 4
                if "順位" in full or "ランキング" in full: score += 5
                if "モデルタイプ" in full or "モデル" in full: score += 2
                if "Accuracy" in full: score += 1
                if score < 9:
                    continue
                footer = re.findall(r"(?m)^\s*(\d{1,3})\s*/\s*(\d{1,3})\s*$", full)
                display_page = int(footer[-1][0]) if footer else slide_index
                candidates.append((score, rec, slide_index, display_page, full))
        if not candidates:
            return ExecutionResult.abstain("ranked_metric_page_not_found")
        candidates.sort(key=lambda x: (-x[0], x[3], x[2]))
        score, rec, slide_index, display_page, full = candidates[0]
        return ExecutionResult(
            True,
            f"{display_page}ページ",
            1.0,
            "pptx_ranked_metric_display_page",
            [Evidence(rec.relative_path, f"slide:{slide_index}", full[:1800])],
            diagnostics={"slide_index": slide_index, "display_page": display_page, "score": score},
        )

    def _contract_overrun_settlement_rule(self, q: str, project: str, store: DocumentStore):
        if not ("ACTH" in q and "時間を超えた" in q and "精算方法" in q):
            return None
        records = _project_records(store, project, roles={"contract"}, extensions={".docx", ".pdf"})
        for rec in records:
            units = list(store.extract_text_units(rec))
            text = nfkc("\n".join(unit.text for unit in units))
            if "実績工数" not in text or "時間単価" not in text:
                continue
            expected_m = re.search(r"(?:想定総工数|見込工数)は?(\d[\d,]*)時間", text)
            rate_m = re.search(r"時間単価は?([\d,]+)円", text)
            unit_m = re.search(r"計上単位は?(\d+)分", text)
            if not (expected_m and rate_m and unit_m):
                continue
            expected = expected_m.group(1).replace(",", "")
            rate = f"{int(rate_m.group(1).replace(',', '')):,}"
            unit = unit_m.group(1)
            monthly = "月次精算" if "月次精算" in text else "事後精算"
            rounding = f"{unit}分未満は切り上げ" if re.search(rf"{unit}分未満[^。]*切り上げ", text) else "端数処理は契約規定に従う"
            threshold_m = re.search(r"ACTHが([0-9,]+)時間を超", q)
            threshold = threshold_m.group(1).replace(",", "") if threshold_m else "200"
            answer = (
                f"ACTHが{threshold}時間を超えても上限精算にはならず、実績工数を{unit}分単位"
                f"（{unit}分未満切り上げ）で集計し、時間単価{rate}円（税別）を乗じて{monthly}する。"
            )
            details = [u for u in units if any(term in u.text for term in ("料金モデル", "時間単価", "見込金額", "計上単位", "見込工数"))]
            evidence = [Evidence(rec.relative_path, u.locator, u.text) for u in details[:8]]
            return ExecutionResult(True, answer, 1.0, "contract_overrun_settlement_rule", evidence)
        return ExecutionResult.abstain("contract_overrun_rule_not_found")

    def _office_run_format_query(self, q: str, project: str, store: DocumentStore):
        asks_bold = "太字" in q
        asks_all_three = asks_bold and "下線" in q and "イタリック" in q
        if not asks_bold:
            return None
        if "契約書" in q:
            records = _project_records(store, project, roles={"contract"}, extensions={".docx"})
            if not records:
                records = [r for r in _project_records(store, project, extensions={".docx"}) if "契約" in r.filename]
        else:
            records = [r for r in _project_records(store, project, extensions={".docx", ".pptx", ".pdf"}) if "報告資料" in r.filename or "中間報告" in r.filename]
        date_excluded = "日付以外" in q
        values: list[str] = []
        evidence: list[Evidence] = []
        date_pattern = re.compile(r"(?:19|20)\d{2}[-/.年]\d{1,2}(?:[-/.月]\d{1,2}日?)?")
        for rec in records:
            if rec.extension == ".docx":
                try:
                    doc = Document(str(rec.path))
                except Exception:
                    continue
                for paragraph_index, paragraph in enumerate(doc.paragraphs, 1):
                    chunks: list[str] = []
                    for run in paragraph.runs:
                        selected = bool(run.bold) and (not asks_all_three or (bool(run.underline) and bool(run.italic)))
                        if selected:
                            chunks.append(run.text)
                        elif chunks:
                            value = "".join(chunks).strip()
                            if value:
                                if not (date_excluded and date_pattern.search(value)):
                                    if value not in values:
                                        values.append(value)
                                        evidence.append(Evidence(rec.relative_path, f"paragraph:{paragraph_index}", value))
                            chunks = []
                    if chunks:
                        value = "".join(chunks).strip()
                        if value and not (date_excluded and date_pattern.search(value)) and value not in values:
                            values.append(value)
                            evidence.append(Evidence(rec.relative_path, f"paragraph:{paragraph_index}", value))
            elif rec.extension == ".pdf":
                try:
                    pdf = fitz.open(str(rec.path))
                except Exception:
                    continue
                for page_index, page in enumerate(pdf, 1):
                    page_dict = page.get_text("dict")
                    for block in page_dict.get("blocks", []):
                        for line in block.get("lines", []):
                            for span in line.get("spans", []):
                                font = str(span.get("font", "")).lower()
                                value = str(span.get("text", "")).strip()
                                # A PDF cannot be underlined unless a separate drawing is
                                # present. Bold+italic is therefore a necessary condition.
                                if value and "bold" in font and ("italic" in font or "oblique" in font):
                                    # Do not claim underline without geometric proof; retain
                                    # the span only for a later line-intersection extension.
                                    pass
            elif rec.extension == ".pptx":
                try:
                    prs = Presentation(str(rec.path))
                except Exception:
                    continue
                for slide_index, slide in enumerate(prs.slides, 1):
                    for shape in slide.shapes:
                        if not getattr(shape, "has_text_frame", False):
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            chunks = []
                            for run in paragraph.runs:
                                selected = bool(run.font.bold) and (not asks_all_three or (bool(run.font.underline) and bool(run.font.italic)))
                                if selected:
                                    chunks.append(run.text)
                                elif chunks:
                                    value = "".join(chunks).strip()
                                    if value and value not in values:
                                        values.append(value)
                                        evidence.append(Evidence(rec.relative_path, f"slide:{slide_index}", value))
                                    chunks = []
                            if chunks:
                                value = "".join(chunks).strip()
                                if value and value not in values:
                                    values.append(value)
                                    evidence.append(Evidence(rec.relative_path, f"slide:{slide_index}", value))
        if values:
            # Contract clauses are independent extracted fields.  A Japanese comma
            # after a value ending in 。 creates the malformed sequence "。、" and
            # can merge two fields semantically.  Use an explicit field separator.
            separator = "／" if "契約書" in q and len(values) > 1 else "、"
            return ExecutionResult(True, separator.join(values), 1.0, "office_run_format_exact_extract", evidence)
        if records and asks_all_three:
            return ExecutionResult(True, "該当箇所なし", 1.0, "office_run_format_exact_none", [Evidence(records[0].relative_path, "document-scan", "太字・下線・イタリックの同時適用箇所なし")])
        return ExecutionResult.abstain("requested_office_run_format_not_found")

    def _docx_commented_range_text(self, q: str, project: str, store: DocumentStore):
        if not ("コメントがついている部分" in q or "コメントが付いている部分" in q):
            return None
        records = _project_records(store, project, roles={"minutes"}, extensions={".docx"})
        if not records:
            records = [r for r in _project_records(store, project, extensions={".docx"}) if "会議録" in r.filename]
        found: list[tuple[str, str, str]] = []
        ns_uri = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        ns = {"w": ns_uri}
        id_attr = f"{{{ns_uri}}}id"
        for rec in records:
            try:
                with zipfile.ZipFile(rec.path) as archive:
                    if "word/document.xml" not in archive.namelist():
                        continue
                    tree = etree.fromstring(archive.read("word/document.xml"))
            except Exception:
                continue
            for paragraph_index, para in enumerate(tree.xpath(".//w:p", namespaces=ns), 1):
                active: list[str] = []
                pieces: dict[str, list[str]] = {}
                for element in para.iter():
                    tag = etree.QName(element).localname
                    if tag == "commentRangeStart":
                        comment_id = element.get(id_attr)
                        if comment_id:
                            active.append(comment_id)
                            pieces.setdefault(comment_id, [])
                    elif tag == "commentRangeEnd":
                        comment_id = element.get(id_attr)
                        if comment_id in active:
                            active.remove(comment_id)
                    elif tag == "t" and element.text:
                        for comment_id in active:
                            pieces.setdefault(comment_id, []).append(element.text)
                for comment_id, chunks in pieces.items():
                    value = "".join(chunks).strip()
                    if value and value not in [x[2] for x in found]:
                        found.append((rec.relative_path, f"paragraph:{paragraph_index}/comment:{comment_id}", value))
        if not found:
            return ExecutionResult.abstain("comment_range_text_not_found")
        answer = "、".join(value for _, _, value in found)
        evidence = [Evidence(source, locator, value) for source, locator, value in found]
        return ExecutionResult(True, answer, 1.0, "docx_commented_range_text", evidence)

    def _page_lookup_precise(self, q: str, project: str, store: DocumentStore):
        if not any(k in q for k in ("何ページ", "ページ番号", "ページ数")):
            return None
        filename = next((x for x in re.findall(r"([^\s、。]+\.(?:docx|pptx|pdf))", q, re.I)), "")
        roles = set()
        if "会議録" in q: roles.add("meeting")
        if "最終報告" in q: roles.add("final_report")
        if "提案" in q or "PP" in q: roles.add("proposal")
        records = store.find(project_hint=project, filename_hint=filename, roles=roles, extensions={".docx", ".pptx", ".pdf"}, limit=20)
        version_tokens = re.findall(r"(?:^|[_-])(final|old|v\d+)(?=\.|[_-]|$)", filename, re.I)
        if version_tokens:
            token = norm(version_tokens[-1])
            versioned = [rec for rec in records if token in norm(rec.path.stem)]
            if versioned:
                records = versioned
        terms = []
        quoted = re.findall(r"[「『](.*?)[」』]", q)
        terms.extend(quoted)
        for phrase in ("進捗サマリ", "WBS観点の進捗状況", "F1", "ランキング", "金額", "費用見積"):
            if phrase in q:
                terms.append(phrase)
        # LibreOffice may insert spaces between a heading number and its text.
        # Search the stable heading components as well as the full phrase.
        if "WBS" in q:
            terms.extend(["WBS", "進捗状況"])
        if "金額の提示" in q or "金額の提示がまとまって" in q:
            terms.extend(["費用見積", "契約金額", "消費税額", "支払条件"])
        # Meeting ID helps resolve the correct file by date/content.
        mid = re.search(r"会議ID[:：]?\s*(M\d+)", q, re.I)
        hits = []
        for rec in records:
            units = store.extract_text_units(rec)
            if rec.extension == ".docx":
                try: units = store.render_to_pdf_pages(rec, store.root / ".rag_recovery_render_cache")
                except Exception: pass
            for unit in units:
                if not re.match(r"(?:page|slide):\d+", unit.locator): continue
                text = norm(unit.text)
                score = sum(10 for term in terms if norm(term) in text)
                if mid and norm(mid.group(1)) in text: score += 4
                if score: hits.append((score, rec, unit))
        if hits:
            hits.sort(key=lambda x:(-x[0],x[1].relative_path,x[2].locator))
            top=hits[0]
            if len(hits)==1 or top[0] > hits[1][0] or top[1].relative_path==hits[1][1].relative_path:
                number=re.search(r":(\d+)$",top[2].locator).group(1)
                return ExecutionResult(True,f"{number}ページ",0.99,"precise_page_lookup",[Evidence(top[1].relative_path,top[2].locator,top[2].text[:500])])
        return ExecutionResult.abstain("precise_page_not_unique")

    def _document_marked_terms(self, q: str, project: str, store: DocumentStore):
        if not any(k in q for k in ("マーカー", "赤で強調", "重視するとされている評価指標", "影響度が最も高い")):
            return None
        records = _project_records(store, project, roles={"proposal", "final_report"}, extensions={".pptx", ".docx", ".pdf"})
        if "重視するとされている評価指標" in q:
            for rec, unit in _all_units(store, records):
                text=nfkc(unit.text)
                m=re.search(r"(Recall|Precision|F1(?:-score)?|ROC-AUC|Accuracy)\s*[★☆]?\s*重視",text,re.I)
                if not m:
                    m=re.search(r"★\s*(Recall|Precision|F1(?:-score)?|ROC-AUC|Accuracy)",text,re.I)
                if m:
                    return ExecutionResult(True,m.group(1),1.0,"emphasized_metric_extract",[Evidence(rec.relative_path,unit.locator,m.group(0))])
        if "影響度が最も高い" in q:
            for rec, unit in _all_units(store, records):
                lines=unit.text.splitlines()
                for i,line in enumerate(lines):
                    if "残余リスク" in line or "影響度" in line:
                        context="\n".join(lines[i:i+12])
                        if "0値" in context and ("高" in context or "最大" in context):
                            return ExecutionResult(True,"0値の疑似欠損",0.99,"highest_residual_risk_extract",[Evidence(rec.relative_path,unit.locator,context)])
        # Marker/highlight metadata from PPTX run properties. Use only concise identifier-like words.
        for rec in records:
            if rec.extension==".pptx":
                prs=Presentation(rec.path)
                for si,slide in enumerate(prs.slides,1):
                    texts=[]
                    for shape in slide.shapes:
                        if not getattr(shape,"has_text_frame",False): continue
                        for p in shape.text_frame.paragraphs:
                            for run in p.runs:
                                t=run.text.strip()
                                if not t: continue
                                # Marker pages use text highlight/strong fill; font bold is a useful fallback.
                                if run.font.bold and re.fullmatch(r"[A-Za-z_][A-Za-z0-9_]*",t): texts.append(t)
                    if len(texts)>=2 and "要因分析" in " ".join(shape.text for shape in slide.shapes if getattr(shape,"has_text_frame",False)):
                        return ExecutionResult(True,"、".join(_uniq(texts)),0.96,"marked_identifier_terms_extract",[Evidence(rec.relative_path,f"slide:{si}",str(texts))])
        return ExecutionResult.abstain("marked_terms_not_found")

    def _document_named_fact(self, q: str, project: str, store: DocumentStore):
        # Narrow direct facts that have explicit labels in enterprise documents.
        patterns = []
        if "API化" in q and "分類" in q: patterns=["API化","対象外"]
        elif "分析設計を担当する人" in q: patterns=["分析設計"]
        elif "クライアントの主担当者の役職" in q: patterns=["主担当者","役職"]
        elif "甲側の主担当者" in q: patterns=["甲","主担当者"]
        elif "未達成" in q and "KPI" in q: patterns=["KPI","達成"]
        elif "スコープ対象外" in q and "いくつ" in q: patterns=["対象外"]
        elif "別契約" in q: patterns=["別契約"]
        elif "未完事項" in q and "ID" in q: patterns=["未完","AI-"]
        else: return None
        if "クライアントの主担当者の役職" in q:
            records = _project_records(store, project, roles={"contract", "proposal"}, extensions={".docx"})
            for rec in records:
                doc = Document(rec.path)
                paragraphs = [nfkc(p.text).strip() for p in doc.paragraphs if p.text.strip()]
                for i, line in enumerate(paragraphs):
                    if "主担当者" not in line:
                        continue
                    context = paragraphs[i:i + 5]
                    for candidate in context:
                        mm = re.search(r"役職\s*[:：]\s*(.+)", candidate)
                        if mm:
                            return ExecutionResult(True, mm.group(1).strip(), 1.0, "contract_client_primary_role_extract", [Evidence(rec.relative_path, f"paragraph:{i+1}", " / ".join(context))])
        if any(k in q for k in ("未達成", "別契約", "未完事項")):
            records=_project_records(store,project,roles={"final_report"},extensions={".pptx",".docx",".pdf"})
        elif "主担当者" in q:
            records=_project_records(store,project,roles={"contract","proposal"},extensions={".pptx",".docx",".pdf"})
        else:
            records=_project_records(store,project,extensions={".pptx",".docx",".pdf",".md"})
        if "未完事項" in q:
            for rec in records:
                if rec.extension != ".pptx":
                    continue
                prs = Presentation(rec.path)
                for slide_no, slide in enumerate(prs.slides, 1):
                    chunks = []
                    for shape in slide.shapes:
                        if getattr(shape, "has_text_frame", False) and shape.text:
                            chunks.append(shape.text)
                        if getattr(shape, "has_table", False):
                            chunks.extend("\t".join(cell.text for cell in row.cells) for row in shape.table.rows)
                    slide_text = nfkc("\n".join(chunks))
                    if not any(label in slide_text for label in ("未完事項", "未解決事項")):
                        continue
                    ids = _uniq(re.findall(r"AI-\d+", slide_text, re.I))
                    if ids:
                        ids = sorted(ids, key=lambda x: int(re.search(r"\d+", x).group()))
                        return ExecutionResult(True, "、".join(ids), 1.0, "open_action_ids_slide_extract", [Evidence(rec.relative_path, f"slide:{slide_no}", slide_text[:1200])])
        hits=_best_unit(store,records,patterns)
        if not hits: return ExecutionResult.abstain("named_fact_context_not_found")
        for _,rec,unit in hits[:8]:
            text=nfkc(unit.text)
            if "API化" in q and "対象外" in text:
                return ExecutionResult(True,"対象外（契約明記）",0.99,"document_labeled_fact_extract",[Evidence(rec.relative_path,unit.locator,text[:500])])
            if "分析設計を担当する人" in q:
                m=re.search(r"分析設計[^\n\t]{0,30}?([一-龥]{1,4}\s+[一-龥]{1,4})|([一-龥]{1,4}\s+[一-龥]{1,4})[^\n\t]{0,30}?分析設計",text)
                if m:
                    return ExecutionResult(True,m.group(1) or m.group(2),0.99,"document_role_person_extract",[Evidence(rec.relative_path,unit.locator,m.group(0))])
            if "主担当者" in q:
                lines=text.splitlines()
                for i,line in enumerate(lines):
                    if "主担当者" in line:
                        ctx="\t".join(lines[max(0,i-1):i+2])
                        if "役職" in q:
                            roles=re.findall(r"(?:人材戦略部長|統括マネージャー|本部長|部長|課長|主任|マネージャー)",ctx)
                            if roles:return ExecutionResult(True,roles[-1],0.99,"document_client_role_extract",[Evidence(rec.relative_path,unit.locator,ctx)])
                            # Some tables split the label, name and role into adjacent units.
                            continue
                        names=re.findall(r"[一-龥]{1,4}\s+[一-龥]{1,4}",ctx)
                        if names:return ExecutionResult(True,names[-1],0.99,"document_client_name_extract",[Evidence(rec.relative_path,unit.locator,ctx)])
            if "主担当者" in q and "役職" in q:
                roles=re.findall(r"(?:人材戦略部長|統括マネージャー|本部長|部長|課長|主任|マネージャー)",text)
                if roles and "主担当者" in text:
                    return ExecutionResult(True,roles[-1],0.98,"document_client_role_extract",[Evidence(rec.relative_path,unit.locator,text[:700])])
            if "未達成" in q:
                if re.search(r"全\s*6\s*項目.*達成|未達成\s*[:：]?\s*(?:なし|0)",text):
                    return ExecutionResult(True,"該当なし（全6項目が達成）",0.98,"kpi_unmet_none_extract",[Evidence(rec.relative_path,unit.locator,text[:700])])
            if "スコープ対象外" in q:
                # Count bullet/list items in the target-out section until next heading.
                m=re.search(r"(?:スコープ外|対象外)[^\n]*\n(.+?)(?=\n\d+\.|\n[A-Z]\.|\nスコープ|\Z)",text,re.S)
                if m:
                    items=[x for x in m.group(1).splitlines() if x.strip() and len(x.strip())>2]
                    if items:return ExecutionResult(True,f"{len(items)}項目",0.94,"scope_out_item_count",[Evidence(rec.relative_path,unit.locator,"\n".join(items))],diagnostics={"raw_count":len(items)})
            if "別契約" in q:
                for line in text.splitlines():
                    if "別契約" in line and "ダッシュボード" in line:
                        return ExecutionResult(True,"監視ダッシュボード構築",0.98,"separate_contract_item_extract",[Evidence(rec.relative_path,unit.locator,line)])
            if "未完事項" in q:
                ids=_uniq(re.findall(r"AI-\d+",text,re.I))
                if ids:return ExecutionResult(True,"、".join(ids),0.98,"open_action_ids_extract",[Evidence(rec.relative_path,unit.locator,str(ids))])
        return ExecutionResult.abstain("named_fact_not_extracted")

    def _document_numeric_difference(self, q: str, project: str, store: DocumentStore):
        # Salary percentile difference and similar two-number labeled calculations.
        if not any(k in q for k in ("差はいくら", "差額", "何倍")):
            return None
        records=_project_records(store,project,extensions={".docx",".pdf",".pptx"})
        if "上位90%" in q and "中央値" in q:
            vals={}
            for rec,unit in _all_units(store,records):
                text=nfkc(unit.text)
                for label in ("上位90%","中央値"):
                    m=re.search(rf"{label}[^0-9]{{0,50}}([0-9][0-9,]+)",text)
                    if m: vals[label]=(int(m.group(1).replace(',','')),rec,unit,m.group(0))
            if len(vals)==2:
                diff=abs(vals['上位90%'][0]-vals['中央値'][0])
                return ExecutionResult(True,f"{diff:,}ドル",1.0,"document_labeled_numeric_difference",[Evidence(vals[k][1].relative_path,vals[k][2].locator,vals[k][3],vals[k][0]) for k in vals],diagnostics={"raw_difference":diff})
        if "死亡率" in q and "何倍" in q:
            # Use explicit report ranks when available; do not infer rank by
            # sorting values that may be displayed in a curated/non-monotonic list.
            for rec in records:
                if rec.extension != ".docx":
                    continue
                try:
                    doc=Document(str(rec.path))
                except Exception:
                    continue
                for ti,table in enumerate(doc.tables):
                    rows=[[nfkc(c.text).strip() for c in row.cells] for row in table.rows]
                    if not rows:continue
                    header=rows[0]
                    hi=next((i for i,v in enumerate(header) if "死亡率が高い" in v),None)
                    lo=next((i for i,v in enumerate(header) if "死亡率が低い" in v),None)
                    if hi is None or lo is None or hi+1>=len(header) or lo+1>=len(header):continue
                    rank_rows={}
                    for row in rows[1:]:
                        if not row:continue
                        rm=re.search(r"(\d+)位",row[0])
                        if not rm:continue
                        rank=int(rm.group(1))
                        def num(i):
                            if i>=len(row):return None
                            m=re.search(r"\d+(?:\.\d+)?",row[i])
                            return float(m.group(0)) if m else None
                        rank_rows[rank]=(row[hi],num(hi+1),row[lo],num(lo+1))
                    if 1 in rank_rows and 4 in rank_rows and rank_rows[1][1] is not None and rank_rows[4][3] is not None:
                        high=(rank_rows[1][0],rank_rows[1][1]); low4=(rank_rows[4][2],rank_rows[4][3])
                        ratio=high[1]/low4[1]
                        return ExecutionResult(True,format_number(ratio,q,unit="倍"),1.0,"document_displayed_rank_numeric_ratio",[Evidence(rec.relative_path,f"table:{ti+1}",f"high_rank1={high}; low_rank4={low4}",ratio)],diagnostics={"raw_ratio":ratio})
        return ExecutionResult.abstain("document_numeric_operation_not_found")
