from __future__ import annotations

import re
import shutil
import subprocess
import tempfile
import unicodedata
from pathlib import Path
from datetime import datetime

OLE_MAGIC = bytes.fromhex("D0CF11E0A1B11AE1")
OFFICE_EXTS = {".docx", ".xlsx", ".pptx"}


def is_encrypted_office(path: Path) -> bool:
    try:
        return path.suffix.lower() in OFFICE_EXTS and path.read_bytes()[:8] == OLE_MAGIC
    except OSError:
        return False


def prepare_office_tree(source_root: Path, workspace: Path) -> tuple[Path, list[dict]]:
    """Copy source data, skip temp files, and decrypt protected Office documents.

    Passwords are derived from internal rules. No password brute force is performed.
    """
    prepared = workspace / "prepared_share"
    if prepared.exists():
        shutil.rmtree(prepared)
    prepared.mkdir(parents=True)
    events: list[dict] = []
    for src in source_root.rglob("*"):
        rel = src.relative_to(source_root)
        if src.is_dir():
            (prepared / rel).mkdir(parents=True, exist_ok=True)
            continue
        if _is_temp(src.name):
            events.append({"source": rel.as_posix(), "status": "ignored_office_temp"})
            continue
        dst = prepared / rel
        dst.parent.mkdir(parents=True, exist_ok=True)
        if not is_encrypted_office(src):
            shutil.copy2(src, dst)
            continue
        passwords = derive_password_candidates(src, source_root)
        ok = False
        errors: list[str] = []
        for password in passwords:
            try:
                _decrypt(src, dst, password)
                events.append({"source": rel.as_posix(), "status": "decrypted", "derivation": _mask(password)})
                ok = True
                break
            except Exception as exc:
                errors.append(f"{type(exc).__name__}:{exc}")
        if not ok:
            events.append({"source": rel.as_posix(), "status": "unresolved_encryption", "candidates": len(passwords), "errors": errors[-2:]})
            shutil.copy2(src, dst)
    return prepared, events


def derive_password_candidates(path: Path, root: Path) -> list[str]:
    """Derive candidates from explicit hints and the internal data-driven PW rule.

    No brute force is performed. The generic rule is discovered from the internal
    password-policy document and project glossary, then the contract start date is
    read from unprotected project documents.
    """
    name = unicodedata.normalize("NFKC", path.name)
    out: list[str] = []
    m = re.search(r"_pw-([A-Za-z0-9_-]+)", name, re.I)
    if m:
        out.append(m.group(1))

    policy_text = _read_internal_policy(root)
    if "DA-[案件略号]-[開始年月日8桁]-[拡張子コード]" in policy_text or ("案件略号" in policy_text and "開始年月日8桁" in policy_text):
        project_dir = _project_dir(path, root)
        if project_dir is not None:
            alias = _project_alias(project_dir.name, root)
            start = _contract_start_date(project_dir)
            if alias and start:
                out.append(f"DA-{alias}-{start}-{path.suffix.lower().lstrip('.')}")
    return list(dict.fromkeys(out))


def _project_dir(path: Path, root: Path) -> Path | None:
    try:
        rel = path.relative_to(root)
    except ValueError:
        return None
    parts = rel.parts
    if "プロジェクト" not in parts:
        return None
    idx = parts.index("プロジェクト")
    return root.joinpath(*parts[:idx + 2]) if idx + 1 < len(parts) else None


def _read_internal_policy(root: Path) -> str:
    candidates = []
    for candidate in root.rglob("*.docx"):
        normalized = unicodedata.normalize("NFKC", candidate.name).lower()
        if ("パスワード" in normalized and "規則" in normalized) or ("password" in normalized and "rule" in normalized):
            candidates.append(candidate)
    for candidate in candidates:
        try:
            from docx import Document
            doc = Document(candidate)
            parts = [p.text for p in doc.paragraphs]
            parts.extend("\n".join("\t".join(c.text for c in row.cells) for row in table.rows) for table in doc.tables)
            return unicodedata.normalize("NFKC", "\n".join(parts))
        except Exception:
            continue
    return ""


def _project_alias(project_name: str, root: Path) -> str:
    normalized_name = unicodedata.normalize("NFKC", project_name)
    for glossary in root.rglob("*用語集*.docx"):
        try:
            from docx import Document
            doc = Document(glossary)
            for table in doc.tables:
                for row in table.rows:
                    cells = [unicodedata.normalize("NFKC", c.text).strip() for c in row.cells]
                    if len(cells) >= 2 and cells[0] and (cells[0] in normalized_name or normalized_name in cells[0]):
                        alias = cells[1].strip()
                        if re.fullmatch(r"[A-Z][A-Z0-9_-]{1,15}", alias):
                            return alias
        except Exception:
            continue
    return ""


def _contract_start_date(project_dir: Path) -> str:
    date_patterns = [
        re.compile(r"(?:契約開始日|開始日|契約締結日兼効力発生日)\s*[:：]?\s*(20\d{2})[-/年](\d{1,2})[-/月](\d{1,2})日?"),
        re.compile(r"(?:案件開始|プロジェクト開始)\s*[:：]?\s*(20\d{2})[-/年](\d{1,2})[-/月](\d{1,2})日?"),
    ]
    role_priority = ("01.契約", "00.提案", "02.計画")
    files = sorted((p for p in project_dir.rglob("*") if p.is_file() and not is_encrypted_office(p)), key=lambda p: next((i for i,k in enumerate(role_priority) if k in p.as_posix()), 9))
    for candidate in files:
        text = _quick_office_text(candidate)
        for pattern in date_patterns:
            m = pattern.search(text)
            if m:
                return f"{int(m.group(1)):04d}{int(m.group(2)):02d}{int(m.group(3)):02d}"
    return ""


def _quick_office_text(path: Path) -> str:
    try:
        suffix = path.suffix.lower()
        if suffix == ".docx":
            from docx import Document
            doc = Document(path)
            parts = [p.text for p in doc.paragraphs]
            parts.extend("\n".join("\t".join(c.text for c in row.cells) for row in table.rows) for table in doc.tables)
            return unicodedata.normalize("NFKC", "\n".join(parts))
        if suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(path)
            return unicodedata.normalize("NFKC", "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False)))
        if suffix in {".md", ".txt", ".csv", ".json", ".py"}:
            return unicodedata.normalize("NFKC", path.read_text(encoding="utf-8", errors="ignore"))
    except Exception:
        return ""
    return ""


def _decrypt(src: Path, dst: Path, password: str) -> None:
    try:
        import msoffcrypto  # type: ignore
        with src.open("rb") as fh:
            office = msoffcrypto.OfficeFile(fh)
            office.load_key(password=password)
            with dst.open("wb") as out:
                office.decrypt(out)
        return
    except ImportError:
        pass
    helper = Path(__file__).with_name("uno_decrypt.py")
    subprocess.run(["/usr/bin/python3", str(helper), str(src), str(dst), password], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=120)


def _is_temp(name: str) -> bool:
    n = unicodedata.normalize("NFKC", name)
    return n.startswith("~$") or n.endswith(".pptx#") or n.endswith(".docx#") or n.endswith(".xlsx#")


def _mask(password: str) -> str:
    if len(password) <= 5:
        return "*" * len(password)
    return password[:2] + "*" * (len(password) - 4) + password[-2:]
