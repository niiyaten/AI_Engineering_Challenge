from __future__ import annotations

import csv
import json
from datetime import datetime
from pathlib import Path

import pandas as pd
import pytest
from docx import Document
from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.styles import PatternFill
from pptx import Presentation


@pytest.fixture()
def sample_share(tmp_path: Path) -> Path:
    root = tmp_path / "share"
    p1 = root / "共有ドライブ" / "プロジェクト" / "株式会社テスト分析"
    p2 = root / "共有ドライブ" / "プロジェクト" / "株式会社第二分析"
    for p in (p1, p2):
        (p / "00.提案").mkdir(parents=True)
        (p / "01.契約").mkdir(parents=True)
        (p / "02.計画").mkdir(parents=True)
        (p / "03.データ").mkdir(parents=True)
        (p / "04.分析" / "analysis_outputs").mkdir(parents=True)
        (p / "05.会議" / "報告資料").mkdir(parents=True)
        (p / "06.報告書").mkdir(parents=True)

    # Dataset for generic table operations and cross-project facts.
    pd.DataFrame({
        "id": [1, 2, 3, 4],
        "purpose": ["credit_card", "credit_card", "car", "credit_card"],
        "loan_amnt": [1000, 2000, 3000, 4000],
        "score": [1.0, None, 2.0, 3.0],
    }).to_csv(p1 / "03.データ" / "train.csv", index=False)
    pd.DataFrame({"id": [1, 2, 3], "score": [None, None, 1.0]}).to_csv(p2 / "03.データ" / "train.csv", index=False)

    # Schedule with dates, owner, colors and joinable IDs.
    wb = Workbook()
    ws = wb.active
    ws.title = "PLAN"
    ws.append(["タスクID", "タスク名", "担当者", "開始日", "終了日", "工数", "種別", "関連ID"])
    rows = [
        ["T01", "要件整理", "加藤 太郎", datetime(2025, 8, 10), datetime(2025, 8, 12), 10, "通常", "MS1"],
        ["T02", "中間報告会実施", "加藤 太郎 / 佐藤 花子", datetime(2025, 8, 20), datetime(2025, 8, 20), 5, "通常", "CP2"],
        ["T03", "モデル評価", "佐藤 花子", datetime(2025, 9, 8), datetime(2025, 9, 10), 7, "通常", "CP2"],
        ["T04", "最終報告", "田中 一郎", datetime(2025, 10, 1), datetime(2025, 10, 3), 4, "通常", "MS3"],
    ]
    for r in rows:
        ws.append(r)
    orange = PatternFill("solid", fgColor="F4B183")
    for cell in ws[3]:
        cell.fill = orange
    yellow = PatternFill("solid", fgColor="FFF2CC")
    ws[4][5].fill = yellow  # numeric work-hour cell
    wb.save(p1 / "02.計画" / "スケジュール_r2.xlsx")

    wb_old = Workbook()
    wo = wb_old.active
    wo.title = "PLAN"
    wo.append(["タスクID", "タスク名", "ステータス", "担当者"])
    wo.append(["T15", "モデル評価・重要特徴量整理", "未着手", "渡辺 遥"])
    wb_old.save(p1 / "02.計画" / "スケジュール_r1.xlsx")
    wb_new = Workbook()
    wn = wb_new.active
    wn.title = "PLAN"
    wn.append(["タスクID", "タスク名", "ステータス", "担当者"])
    wn.append(["T15", "モデル評価・重要特徴量整理", "完了", "渡辺 遥 / 小林 直樹"])
    wb_new.save(p1 / "02.計画" / "スケジュール_r3.xlsx")

    # Regression workbook.
    wb_reg = Workbook()
    data = wb_reg.active
    data.title = "train"
    data.append(["index", "x1", "x2"])
    data.append([1770, 2.0, 3.0])
    coef = wb_reg.create_sheet("coefficients")
    coef.append(["feature", "coefficient"])
    coef.append(["intercept", 1.0])
    coef.append(["x1", 2.0])
    coef.append(["x2", 4.0])
    wb_reg.save(p1 / "03.データ" / "regression.xlsx")

    # Embedded chart workbook.
    wb_chart = Workbook()
    cs = wb_chart.active
    cs.title = "Sheet1"
    cs.append(["day", "hum"])
    for day, hum in [(1, 10), (2, 20), (3, 15)]:
        cs.append([day, hum])
    chart = BarChart()
    chart.title = "Humidity"
    chart.add_data(Reference(cs, min_col=2, min_row=1, max_row=4), titles_from_data=True)
    chart.set_categories(Reference(cs, min_col=1, min_row=2, max_row=4))
    cs.add_chart(chart, "D2")
    wb_chart.save(p1 / "03.データ" / "chart.xlsx")

    # Document numeric lookup.
    doc = Document()
    doc.add_paragraph("機械学習（ML）エンジニア 米国平均給与：140,000ドル")
    doc.add_paragraph("データエンジニア 米国平均給与：125,256ドル")
    doc.save(p1 / "00.提案" / "データサイエンティスト調査.docx")

    # Run formatting intersection.
    report = Document()
    p = report.add_paragraph()
    r = p.add_run("重要な実施条件")
    r.bold = True
    r.italic = True
    r.underline = True
    report.add_paragraph("通常文")
    report.save(p1 / "05.会議" / "報告資料" / "報告資料.docx")

    # Contract facts.
    c1 = Document()
    c1.add_paragraph("固定金額契約。契約金額（税込）：5,500,000円")
    c1.add_paragraph("契約期間：2025年1月1日から2025年12月31日")
    c1.save(p1 / "01.契約" / "契約書.docx")
    c2 = Document()
    c2.add_paragraph("固定金額契約。契約金額（税込）：3,300,000円")
    c2.add_paragraph("契約期間：2025年2月1日から2025年8月31日")
    c2.save(p2 / "01.契約" / "契約書.docx")

    # Code + JSON.
    (p1 / "04.分析" / "features.py").write_text("def make(a, b):\n    return f'{a}__x__{b}'\n", encoding="utf-8")
    (p1 / "04.分析" / "analysis_outputs" / "metrics.json").write_text(
        json.dumps({"selected_columns": ["age", "income__x__region", "temp__x__hum"]}), encoding="utf-8"
    )

    # Final report to mark completion.
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "最終報告・案件完了"
    prs.save(p1 / "06.報告書" / "株式会社テスト分析_最終報告.pptx")
    prs2 = Presentation()
    slide = prs2.slides.add_slide(prs2.slide_layouts[5])
    slide.shapes.title.text = "最終報告・案件完了"
    prs2.save(p2 / "06.報告書" / "株式会社第二分析_最終報告.pptx")
    return root


@pytest.fixture()
def answers130(tmp_path: Path, sample_share: Path) -> Path:
    path = tmp_path / "answers.csv"
    fields = ["index", "split", "question_id", "question", "answer", "confidence_score", "confidence_level", "source_mode", "selected_sources", "source_count", "evidence_locator", "evidence", "method", "unknown_reason", "human_review_used"]
    rel = "共有ドライブ/プロジェクト/株式会社テスト分析/"
    unknowns = [
        (0, "テスト分析のデータサイエンティスト調査.docxにおいて、米国平均給与における機械学習（ML）エンジニアとデータエンジニアの差はいくらですか。", rel + "00.提案/データサイエンティスト調査.docx"),
        (1, "テスト分析のスケジュール_r2.xlsxにおいて、2025-08-11から2025-09-09の間に開始日または終了日が設定されているタスクIDをすべて挙げてください。", rel + "02.計画/スケジュール_r2.xlsx"),
        (2, "テスト分析のスケジュール_r2.xlsxにおいて、オレンジにハイライトされている行のタスク名をすべて答えてください。", rel + "02.計画/スケジュール_r2.xlsx"),
        (3, "テスト分析のregression.xlsxにおいて、回帰分析の結果として記載されている係数をindex=1770のデータに当てはめたときの予測値はいくつですか。小数第5位まで答えてください。", rel + "03.データ/regression.xlsx"),
        (4, "テスト分析の報告資料.docxの中で、太字、下線、イタリックのすべてに該当する箇所を抽出してください。", rel + "05.会議/報告資料/報告資料.docx"),
        (5, "テスト分析のスケジュール_r1.xlsxとスケジュール_r3.xlsxを比較したとき、未着手から完了への変更を除いて、案件遂行に関連する変更点を挙げてください。", rel + "02.計画/スケジュール_r1.xlsx\n" + rel + "02.計画/スケジュール_r3.xlsx"),
        (6, "テスト分析のmetrics.jsonとfeatures.pyから、生成された特徴量はいくつありますか。", rel + "04.分析/analysis_outputs/metrics.json\n" + rel + "04.分析/features.py"),
        (7, "全案件のtrain.csvにおいて、欠損行数が最も多い案件を案件名で答えてください。", ""),
    ]
    rows = []
    for i in range(130):
        if i < len(unknowns):
            qid, q, sources = unknowns[i]
            answer = "わからない"
        else:
            qid, q, sources = i, f"既存回答済み質問{i}", ""
            answer = "既存回答"
        rows.append({"index": i, "split": "test", "question_id": qid, "question": q, "answer": answer, "confidence_score": "0.2", "confidence_level": "low", "source_mode": "single_document", "selected_sources": sources, "source_count": "", "evidence_locator": "", "evidence": "", "method": "", "unknown_reason": "", "human_review_used": "False"})
    with path.open("w", encoding="utf-8-sig", newline="") as fh:
        w = csv.DictWriter(fh, fieldnames=fields)
        w.writeheader(); w.writerows(rows)
    return path
