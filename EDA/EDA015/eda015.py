from __future__ import annotations

import argparse
import csv
import hashlib
import json
import logging
import re
import shutil
import unicodedata
from datetime import datetime
from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# =============================================================================
# パス設定
# =============================================================================

# eda015.py は「プロジェクト直下 / EDA / EDA015 / eda015.py」に置く前提。
BASE_DIR = Path(__file__).resolve().parents[2]
RAW_SHARE_DIR = BASE_DIR / "data" / "raw" / "share"
RAW_DRIVE_DIR = RAW_SHARE_DIR / "share" / "共有ドライブ"
PROCESSED_SHARE_DIR = BASE_DIR / "data" / "processed" / "share"

OUTPUT_DIR = Path(__file__).resolve().parent
TABLE_DIR = OUTPUT_DIR / "tables"
REPORT_PATH = OUTPUT_DIR / "eda015_report.md"
MANIFEST_PATH = OUTPUT_DIR / "manifest.json"
LOG_PATH = OUTPUT_DIR / "eda015.log"


def setup() -> None:
    """出力フォルダとログ設定を準備する。"""
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    TABLE_DIR.mkdir(parents=True, exist_ok=True)
    PROCESSED_SHARE_DIR.mkdir(parents=True, exist_ok=True)
    logging.basicConfig(
        filename=LOG_PATH,
        level=logging.INFO,
        format="%(asctime)s [%(levelname)s] %(message)s",
        encoding="utf-8",
        force=True,
    )


def normalize_path_text(text: Any) -> str:
    """Windows上の濁点表記揺れなどをNFCへ寄せる。"""
    return unicodedata.normalize("NFC", str(text))


def relative(path: Path) -> str:
    """レポート用にプロジェクト相対パスを返す。"""
    try:
        return path.resolve().relative_to(BASE_DIR.resolve()).as_posix()
    except ValueError:
        return path.as_posix()


def file_sha1(path: Path) -> str:
    """入力ファイルの追跡用にSHA1を計算する。"""
    h = hashlib.sha1()
    with path.open("rb") as f:
        for block in iter(lambda: f.read(1024 * 1024), b""):
            h.update(block)
    return h.hexdigest()


def markdown_escape(text: Any) -> str:
    """Markdown表を壊しやすい文字だけを逃がす。"""
    value = "" if text is None else str(text)
    return value.replace("\\", "\\\\").replace("|", "\\|").replace("\n", " ")


def safe_filename(text: str) -> str:
    """スライド名や画像名に使える安全なファイル名を作る。"""
    value = normalize_path_text(text).strip()
    value = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "_", value)
    return value[:80] or "asset"


def markdown_table(rows: list[dict[str, Any]], max_rows: int = 20) -> str:
    """追加依存なしでMarkdown表を作る。"""
    if not rows:
        return "該当データはありません。"
    columns = list(rows[0].keys())
    lines = [
        "| " + " | ".join(columns) + " |",
        "| " + " | ".join("---" for _ in columns) + " |",
    ]
    for row in rows[:max_rows]:
        lines.append("| " + " | ".join(markdown_escape(row.get(col, "")) for col in columns) + " |")
    return "\n".join(lines)


def save_csv(rows: list[dict[str, Any]], path: Path) -> None:
    """Excelでも読みやすいUTF-8 BOM付きCSVを保存する。"""
    if not rows:
        path.write_text("", encoding="utf-8-sig")
        return
    columns = list(dict.fromkeys(col for row in rows for col in row.keys()))
    with path.open("w", encoding="utf-8-sig", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=columns)
        writer.writeheader()
        writer.writerows(rows)


def output_paths(raw_path: Path) -> tuple[Path, Path, Path]:
    """raw/shareからの相対構造を保ったMarkdown/JSON/assets出力先を返す。"""
    rel = raw_path.relative_to(RAW_SHARE_DIR)
    output_base = PROCESSED_SHARE_DIR / normalize_path_text(rel.as_posix())
    output_base.parent.mkdir(parents=True, exist_ok=True)
    md_path = output_base.with_suffix(output_base.suffix + ".md")
    json_path = output_base.with_suffix(output_base.suffix + ".structure.json")
    assets_dir = output_base.with_suffix(output_base.suffix + ".assets")
    return md_path, json_path, assets_dir


def is_temporary_office_file(path: Path) -> bool:
    """Officeの一時ファイルを除外する。"""
    return path.name.startswith("~$")


def emu_to_pt(value: Any) -> float | None:
    """PowerPointのEMU座標をptへ変換する。"""
    if value is None:
        return None
    return round(float(value) / 12700.0, 2)


def color_to_text(color: Any) -> str:
    """python-pptxの色オブジェクトを文字列化する。"""
    try:
        if color is not None and color.rgb is not None:
            return f"#{color.rgb}"
    except Exception:  # noqa: BLE001
        return ""
    return ""


def run_record(run: Any) -> dict[str, Any]:
    """テキストrunの本文と基本書式を保存する。"""
    font = run.font
    return {
        "text": run.text or "",
        "bold": bool(font.bold),
        "italic": bool(font.italic),
        "underline": bool(font.underline),
        "font_name": font.name or "",
        "font_size_pt": round(font.size.pt, 2) if font.size is not None else None,
        "font_color": color_to_text(font.color),
    }


def text_frame_record(shape: Any) -> dict[str, Any]:
    """図形内テキストを段落とrun単位で抽出する。"""
    paragraphs = []
    plain_lines = []
    for paragraph in shape.text_frame.paragraphs:
        runs = [run_record(run) for run in paragraph.runs]
        text = "".join(run["text"] for run in runs).strip()
        if text:
            plain_lines.append(text)
        paragraphs.append(
            {
                "text": text,
                "level": paragraph.level,
                "runs": runs,
            }
        )
    return {
        "text": "\n".join(plain_lines).strip(),
        "paragraphs": paragraphs,
    }


def table_record(shape: Any) -> dict[str, Any]:
    """PowerPoint表を行列の文字列として抽出する。"""
    rows: list[list[str]] = []
    table = shape.table
    for row in table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    return {
        "row_count": len(rows),
        "column_count": max((len(row) for row in rows), default=0),
        "cells": rows,
    }


def table_to_markdown(table: dict[str, Any]) -> str:
    """抽出済み表をMarkdown表へ変換する。"""
    rows = table.get("cells", [])
    if not rows:
        return ""
    header = rows[0]
    lines = [
        "| " + " | ".join(markdown_escape(cell) for cell in header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for row in rows[1:]:
        lines.append("| " + " | ".join(markdown_escape(cell) for cell in row) + " |")
    return "\n".join(lines)


def chart_record(shape: Any) -> dict[str, Any]:
    """PowerPointグラフの基本情報を抽出する。"""
    chart = shape.chart
    title = ""
    try:
        if chart.has_title:
            title = chart.chart_title.text_frame.text
    except Exception:  # noqa: BLE001
        title = ""
    def axis_title_text(axis: Any) -> str:
        """軸タイトルをJSON保存できる文字列として取り出す。"""
        try:
            if axis is not None and axis.has_title:
                return axis.axis_title.text_frame.text
        except Exception:  # noqa: BLE001
            return ""
        return ""

    return {
        "chart_type": str(chart.chart_type),
        "title": title,
        "series_count": len(chart.series),
        "category_axis_title": axis_title_text(getattr(chart, "category_axis", None)),
        "value_axis_title": axis_title_text(getattr(chart, "value_axis", None)),
    }


def save_picture(shape: Any, assets_dir: Path, slide_number: int, shape_index: int) -> dict[str, Any]:
    """スライド内画像をassetsへ保存し、パスとメタデータを返す。"""
    image = shape.image
    ext = image.ext or "bin"
    image_name = f"slide{slide_number:03d}_shape{shape_index:03d}.{safe_filename(ext)}"
    image_path = assets_dir / image_name
    image_path.write_bytes(image.blob)
    return {
        "image_path": relative(image_path),
        "image_ext": ext,
        "content_type": image.content_type,
        "bytes": image_path.stat().st_size,
        "sha1": hashlib.sha1(image.blob).hexdigest(),
    }


def shape_record(shape: Any, assets_dir: Path, slide_number: int, shape_index: int) -> dict[str, Any]:
    """1つのshapeから、テキスト、表、画像、グラフ、位置情報を抽出する。"""
    record: dict[str, Any] = {
        "shape_index": shape_index,
        "shape_id": getattr(shape, "shape_id", None),
        "name": getattr(shape, "name", ""),
        "shape_type": str(getattr(shape, "shape_type", "")),
        "left_pt": emu_to_pt(getattr(shape, "left", None)),
        "top_pt": emu_to_pt(getattr(shape, "top", None)),
        "width_pt": emu_to_pt(getattr(shape, "width", None)),
        "height_pt": emu_to_pt(getattr(shape, "height", None)),
        "text": "",
        "text_frame": None,
        "table": None,
        "image": None,
        "chart": None,
    }
    if getattr(shape, "has_text_frame", False):
        text_frame = text_frame_record(shape)
        record["text"] = text_frame["text"]
        record["text_frame"] = text_frame
    if getattr(shape, "has_table", False):
        record["table"] = table_record(shape)
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
        record["image"] = save_picture(shape, assets_dir, slide_number, shape_index)
    if getattr(shape, "has_chart", False):
        record["chart"] = chart_record(shape)
    return record


def notes_text(slide: Any) -> str:
    """スライドノートがあれば本文を取り出す。"""
    try:
        frame = slide.notes_slide.notes_text_frame
        return frame.text.strip()
    except Exception:  # noqa: BLE001
        return ""


def process_pptx(path: Path) -> tuple[dict[str, Any], str]:
    """PowerPointをスライド単位Markdownと構造JSONへ変換する。"""
    if is_temporary_office_file(path):
        raise RuntimeError("temporary_office_file")

    prs = Presentation(str(path))
    md_path, json_path, assets_dir = output_paths(path)
    if assets_dir.exists():
        shutil.rmtree(assets_dir)
    assets_dir.mkdir(parents=True, exist_ok=True)

    slide_records = []
    image_count = 0
    chart_count = 0
    table_count = 0
    text_shape_count = 0
    for slide_number, slide in enumerate(prs.slides, start=1):
        shapes = []
        for shape_index, shape in enumerate(slide.shapes, start=1):
            rec = shape_record(shape, assets_dir, slide_number, shape_index)
            if rec["text"]:
                text_shape_count += 1
            if rec["table"]:
                table_count += 1
            if rec["image"]:
                image_count += 1
            if rec["chart"]:
                chart_count += 1
            shapes.append(rec)
        slide_records.append(
            {
                "slide_number": slide_number,
                "notes": notes_text(slide),
                "shape_count": len(shapes),
                "shapes": shapes,
            }
        )

    record = {
        "raw_relative_path": normalize_path_text(path.relative_to(RAW_SHARE_DIR).as_posix()),
        "processed_markdown_path": relative(md_path),
        "processed_structure_path": relative(json_path),
        "assets_dir": relative(assets_dir),
        "file_name": path.name,
        "file_type": "pptx",
        "source_sha1": file_sha1(path),
        "slide_count": len(slide_records),
        "text_shape_count": text_shape_count,
        "table_count": table_count,
        "image_count": image_count,
        "chart_count": chart_count,
        "slides": slide_records,
    }
    json_path.write_text(json.dumps(record, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    md_path.write_text(markdown_document(record), encoding="utf-8")
    return record, "ok"


def markdown_document(record: dict[str, Any]) -> str:
    """抽出済みPowerPoint構造からLLM可読Markdownを作る。"""
    lines = [
        f"# PowerPoint Markdown: {record['file_name']}",
        "",
        "## Source",
        f"- raw_path: `{record['raw_relative_path']}`",
        f"- source_sha1: `{record['source_sha1']}`",
        f"- slide_count: {record['slide_count']}",
        f"- table_count: {record['table_count']}",
        f"- image_count: {record['image_count']}",
        f"- chart_count: {record['chart_count']}",
        "",
        "## Slides",
    ]
    for slide in record["slides"]:
        lines.extend(["", f"### Slide {slide['slide_number']}", ""])
        text_items = [shape["text"] for shape in slide["shapes"] if shape.get("text")]
        for text in text_items:
            lines.append(text)
            lines.append("")
        for shape in slide["shapes"]:
            if shape.get("table"):
                lines.extend(["#### Table", "", table_to_markdown(shape["table"]), ""])
            if shape.get("chart"):
                chart = shape["chart"]
                lines.extend(
                    [
                        "#### Chart",
                        f"- chart_type: `{chart.get('chart_type', '')}`",
                        f"- title: {chart.get('title', '')}",
                        f"- series_count: {chart.get('series_count', 0)}",
                        "",
                    ]
                )
            if shape.get("image"):
                image = shape["image"]
                lines.extend(
                    [
                        "#### Image",
                        f"- image_path: `{image.get('image_path', '')}`",
                        f"- content_type: `{image.get('content_type', '')}`",
                        "",
                    ]
                )
        if slide.get("notes"):
            lines.extend(["#### Notes", "", slide["notes"], ""])
    return "\n".join(lines).strip() + "\n"


def process_file(path: Path) -> dict[str, Any]:
    """1つのPPTXを処理し、ログ行を返す。"""
    row = {
        "raw_relative_path": normalize_path_text(path.relative_to(RAW_SHARE_DIR).as_posix()),
        "status": "",
        "error_type": "",
        "error_message": "",
        "processed_markdown_path": "",
        "processed_structure_path": "",
        "slide_count": "",
        "text_shape_count": "",
        "table_count": "",
        "image_count": "",
        "chart_count": "",
    }
    try:
        record, status = process_pptx(path)
        row.update(
            {
                "status": status,
                "processed_markdown_path": record["processed_markdown_path"],
                "processed_structure_path": record["processed_structure_path"],
                "slide_count": record["slide_count"],
                "text_shape_count": record["text_shape_count"],
                "table_count": record["table_count"],
                "image_count": record["image_count"],
                "chart_count": record["chart_count"],
            }
        )
    except Exception as exc:  # noqa: BLE001
        row["status"] = "error"
        row["error_type"] = exc.__class__.__name__
        row["error_message"] = str(exc)[:1000]
        logging.exception("failed: %s", path)
    return row


def collect_targets() -> list[Path]:
    """共有ドライブ配下からPowerPointを集める。"""
    return sorted(
        [path for path in RAW_DRIVE_DIR.rglob("*.pptx") if path.is_file()],
        key=lambda p: normalize_path_text(p.as_posix()).lower(),
    )


def write_report(log_rows: list[dict[str, Any]]) -> None:
    """EDA015の処理結果をMarkdownレポートに保存する。"""
    df = pd.DataFrame(log_rows)
    status_rows = df.groupby("status", dropna=False).size().reset_index(name="count").to_dict(orient="records")
    ok_df = df[df["status"] == "ok"].copy()
    totals = {
        "slide_count": int(pd.to_numeric(ok_df["slide_count"], errors="coerce").fillna(0).sum()),
        "text_shape_count": int(pd.to_numeric(ok_df["text_shape_count"], errors="coerce").fillna(0).sum()),
        "table_count": int(pd.to_numeric(ok_df["table_count"], errors="coerce").fillna(0).sum()),
        "image_count": int(pd.to_numeric(ok_df["image_count"], errors="coerce").fillna(0).sum()),
        "chart_count": int(pd.to_numeric(ok_df["chart_count"], errors="coerce").fillna(0).sum()),
    }
    error_rows = [row for row in log_rows if row["status"] != "ok"]
    lines = [
        "# EDA015: PowerPoint前処理",
        "",
        "## 目的",
        "",
        "PowerPointをスライド単位Markdownと構造JSONへ変換し、提案書、報告書、座席表、版違い比較で利用できる中間データを作る。",
        "",
        "## 出力",
        "",
        "- Markdown/JSON: `data/processed/share/**/*.pptx.md`, `*.pptx.structure.json`",
        "- 抽出画像: `data/processed/share/**/*.pptx.assets/*`",
        f"- 変換ログ: `{relative(TABLE_DIR / 'pptx_conversion_log.csv')}`",
        "",
        "## 処理結果",
        "",
        "凡例: `status` は処理状態、`count` は該当PowerPointファイル数を表します。",
        "",
        markdown_table(status_rows),
        "",
        "## 抽出総数",
        "",
        "凡例: 各項目は成功したPowerPointから抽出したスライド、テキスト図形、表、画像、グラフの合計件数を表します。",
        "",
        markdown_table([totals]),
        "",
        "## エラー・スキップ",
        "",
        "凡例: `raw_relative_path` は元ファイル、`status` は処理状態、`error_type` と `error_message` は失敗理由を表します。",
        "",
        markdown_table(error_rows, max_rows=30),
        "",
        "## 注意点",
        "",
        "- このEDAではスライド画像をOpenRouterへ送っていない。画像説明が必要な場合は、抽出assetsを対象に別工程で実行する。",
        "- グラフは基本メタデータを保存するが、グラフ内の数値系列の厳密抽出は次工程で扱う。",
        "- PowerPointの座標はpt単位に変換してJSONへ保存した。",
    ]
    REPORT_PATH.write_text("\n".join(lines) + "\n", encoding="utf-8")


def write_manifest(log_rows: list[dict[str, Any]]) -> None:
    """再現用の実行条件を保存する。"""
    manifest = {
        "eda": "EDA015",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "purpose": "Convert PowerPoint files into slide-level Markdown and structure JSON.",
        "target_count": len(log_rows),
        "outputs": {
            "report": relative(REPORT_PATH),
            "conversion_log": relative(TABLE_DIR / "pptx_conversion_log.csv"),
            "processed_root": relative(PROCESSED_SHARE_DIR),
        },
        "repro_steps": ["uv run python EDA/EDA015/eda015.py"],
    }
    MANIFEST_PATH.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def parse_args() -> argparse.Namespace:
    """コマンドライン引数を定義する。"""
    return argparse.ArgumentParser(description="Convert PowerPoint files into preprocessing artifacts.").parse_args()


def main() -> None:
    """EDA015を実行する。"""
    parse_args()
    setup()
    targets = collect_targets()
    log_rows = [process_file(path) for path in targets]
    save_csv(log_rows, TABLE_DIR / "pptx_conversion_log.csv")
    write_report(log_rows)
    write_manifest(log_rows)
    ok_count = sum(1 for row in log_rows if row["status"] == "ok")
    print(f"targets={len(log_rows)} ok={ok_count} errors={len(log_rows) - ok_count}")
    print(f"report={relative(REPORT_PATH)}")


if __name__ == "__main__":
    main()
