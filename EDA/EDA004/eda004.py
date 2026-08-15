from __future__ import annotations

import argparse
import hashlib
import json
import logging
import re
import unicodedata
import warnings
from collections import Counter
from collections.abc import Iterable
from pathlib import Path
from typing import Any

import matplotlib

# GUIがない環境でも画像保存できるようにする。
matplotlib.use("Agg")

import matplotlib.pyplot as plt
import pandas as pd
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

# =============================================================================
# パス設定
# =============================================================================

# eda004.py は「プロジェクト直下 / EDA / EDA004 / eda004.py」に置く前提。
BASE_DIR = Path(__file__).resolve().parents[2]
DATA_DIR = BASE_DIR / "data"
RAW_DIR = DATA_DIR / "raw"
INTERIM_DIR = DATA_DIR / "interim"
PROCESSED_DIR = DATA_DIR / "processed"

OUTPUT_DIR = Path(__file__).resolve().parent
FIG_DIR = OUTPUT_DIR / "figures"
TABLE_DIR = OUTPUT_DIR / "tables"
TEXT_DIR = OUTPUT_DIR / "texts"
REPORT_PATH = OUTPUT_DIR / "eda004_report.md"
LOG_PATH = OUTPUT_DIR / "eda004.log"

PROCESSED_TEXT_DIR = PROCESSED_DIR / "office_pdf_baseline"

TARGET_EXTENSIONS = {".docx", ".pptx", ".xlsx", ".pdf"}

DATASET_SEARCH_BASES = [
    RAW_DIR,
    INTERIM_DIR,
    DATA_DIR,
    BASE_DIR,
]

EXCLUDE_DIR_NAMES = {
    ".git",
    ".venv",
    "__pycache__",
    ".mypy_cache",
    ".pytest_cache",
    "node_modules",
}

DEFAULT_CHUNK_SIZE = 1_200
DEFAULT_CHUNK_OVERLAP = 200
DEFAULT_XLSX_MAX_ROWS_PER_SHEET = 120
DEFAULT_XLSX_MAX_COLS_PER_SHEET = 60
TEXT_PREVIEW_LENGTH = 500


# =============================================================================
# 基本ユーティリティ
# =============================================================================


def setup() -> None:
    """出力フォルダとログ設定を準備する。"""
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    TABLE_DIR.mkdir(parents=True, exist_ok=True)
    TEXT_DIR.mkdir(parents=True, exist_ok=True)
    PROCESSED_TEXT_DIR.mkdir(parents=True, exist_ok=True)

    logging.captureWarnings(True)
    warnings.simplefilter("always")
    warnings.filterwarnings("ignore", message="Glyph .* missing from font.*")
    logging.basicConfig(
        filename=LOG_PATH,
        level=logging.INFO,
        format="%(asctime)s [%(levelname)s] %(message)s",
        encoding="utf-8",
        force=True,
    )

    plt.rcParams["font.family"] = "sans-serif"
    plt.rcParams["font.sans-serif"] = ["Yu Gothic", "Meiryo", "MS Gothic", "DejaVu Sans"]
    plt.rcParams["axes.unicode_minus"] = False


_HASH_U_PATTERN = re.compile(r"#U([0-9a-fA-F]{4})")


def decode_hash_u_text(text: str) -> str:
    """#U5171 のように展開された日本語パスを通常の日本語へ戻す。"""

    def repl(match: re.Match[str]) -> str:
        return chr(int(match.group(1), 16))

    return unicodedata.normalize("NFC", _HASH_U_PATTERN.sub(repl, str(text)))


def normalize_text(text: Any) -> str:
    """日本語ファイル名や本文の濁点揺れを抑えるため、NFCに正規化する。"""
    return decode_hash_u_text(str(text))


def clean_text(text: Any) -> str:
    """抽出本文の余分な空白を軽く整える。"""
    text = normalize_text(text)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def clean_preview(text: Any, max_len: int = TEXT_PREVIEW_LENGTH) -> str:
    """CSVやMarkdownで確認しやすい短いプレビューを作る。"""
    text = re.sub(r"\s+", " ", normalize_text(text)).strip()
    if len(text) > max_len:
        return text[:max_len] + "..."
    return text


def safe_relative_to(path: Path, base: Path) -> str:
    """baseからの相対パスを返す。外側にある場合は絶対パスを返す。"""
    try:
        return normalize_text(path.relative_to(base).as_posix())
    except ValueError:
        return normalize_text(path.as_posix())


def is_excluded_path(path: Path) -> bool:
    """探索対象から除外すべきフォルダ配下かどうかを判定する。"""
    return any(part in EXCLUDE_DIR_NAMES for part in path.parts)


def iter_files(root: Path) -> Iterable[Path]:
    """root配下のファイルを再帰的に列挙する。"""
    for path in root.rglob("*"):
        if is_excluded_path(path):
            continue
        if path.is_file() and path.name != ".extracted":
            yield path


def iter_dirs(root: Path) -> Iterable[Path]:
    """root配下のフォルダを再帰的に列挙する。"""
    if not root.exists():
        return
    for path in root.rglob("*"):
        if is_excluded_path(path):
            continue
        if path.is_dir():
            yield path


def find_named_dirs(search_bases: list[Path], decoded_name: str) -> list[Path]:
    """指定した日本語フォルダ名に一致するフォルダを探す。"""
    hits: list[Path] = []
    seen: set[Path] = set()
    for base in search_bases:
        if not base.exists():
            continue
        for path in iter_dirs(base):
            if path in seen:
                continue
            if normalize_text(path.name) == decoded_name:
                hits.append(path)
                seen.add(path)
    return hits


def choose_best_path(paths: list[Path]) -> Path:
    """候補が複数ある場合、プロジェクト内のraw/interimを優先して選ぶ。"""
    if not paths:
        raise FileNotFoundError("候補パスが空です。")

    def score(path: Path) -> tuple[int, int]:
        text = path.as_posix()
        if "/data/raw/" in text:
            priority = 0
        elif "/data/interim/" in text:
            priority = 1
        else:
            priority = 2
        return (priority, len(path.parts))

    return sorted(paths, key=score)[0]


def find_drive_root() -> Path:
    """展開済みデータから `共有ドライブ` フォルダを探す。"""
    hits = find_named_dirs(DATASET_SEARCH_BASES, "共有ドライブ")
    if not hits:
        message = (
            "展開済みの `共有ドライブ` フォルダが見つかりません。\n"
            "share.zipを展開し、例えば `data/raw/share/share/共有ドライブ` になるように配置してください。"
        )
        raise FileNotFoundError(message)
    drive_root = choose_best_path(hits)
    logging.info("Use drive_root: %s", drive_root)
    return drive_root


def classify_shared_drive_path(rel_path: str) -> dict[str, str]:
    """共有ドライブ内のパスから、案件名・大分類フォルダなどを推定する。"""
    parts = Path(rel_path).parts
    parts_nfc = [normalize_text(p) for p in parts]

    project_name = ""
    major_folder = ""
    area = ""

    if "共有ドライブ" in parts_nfc:
        idx = parts_nfc.index("共有ドライブ")
        if len(parts_nfc) > idx + 1:
            area = parts_nfc[idx + 1]
        if area == "プロジェクト" and len(parts_nfc) > idx + 2:
            project_name = parts_nfc[idx + 2]
        if area == "プロジェクト" and len(parts_nfc) > idx + 3:
            major_folder = parts_nfc[idx + 3]
        elif len(parts_nfc) > idx + 2:
            major_folder = parts_nfc[idx + 2]

    return {
        "area": area,
        "project_name": project_name,
        "major_folder": major_folder,
    }


def make_document_id(rel_path: str) -> str:
    """パスから安定したdocument_idを作る。"""
    digest = hashlib.sha1(rel_path.encode("utf-8")).hexdigest()[:16]
    return f"doc_{digest}"


def save_csv(df: pd.DataFrame, path: Path) -> None:
    """Excelでも開きやすいようにUTF-8 BOM付きでCSV保存する。"""
    df.to_csv(path, index=False, encoding="utf-8-sig")


def df_to_markdown(df: pd.DataFrame, max_rows: int | None = None) -> str:
    """DataFrameをGitHub風Markdown表に変換する。"""
    if df.empty:
        return "該当データはありません。"
    view = df if max_rows is None else df.head(max_rows)
    columns = [str(col) for col in view.columns]

    def fmt(value: Any) -> str:
        if pd.isna(value):
            text = ""
        else:
            text = str(value)
        text = text.replace("\n", " ").replace("\r", " ")
        return text.replace("|", "\\|")

    lines = [
        "| " + " | ".join(columns) + " |",
        "| " + " | ".join("---" for _ in columns) + " |",
    ]
    for _, row in view.iterrows():
        lines.append("| " + " | ".join(fmt(row[col]) for col in view.columns) + " |")
    return "\n".join(lines)


def write_jsonl(records: list[dict[str, Any]], path: Path) -> None:
    """辞書リストをJSONLとして保存する。"""
    with path.open("w", encoding="utf-8") as f:
        for record in records:
            f.write(json.dumps(record, ensure_ascii=False) + "\n")


# =============================================================================
# ファイル棚卸し
# =============================================================================


def build_target_file_inventory(drive_root: Path) -> pd.DataFrame:
    """共有ドライブからEDA004対象形式のファイル一覧を作る。"""
    rows: list[dict[str, Any]] = []
    for path in iter_files(drive_root):
        extension = path.suffix.lower()
        if extension not in TARGET_EXTENSIONS:
            continue
        # Officeの一時ファイルは実体が通常のzip形式ではないため、抽出対象から外す。
        if path.name.startswith("~$"):
            continue
        rel_path = safe_relative_to(path, drive_root)
        attrs = classify_shared_drive_path(f"共有ドライブ/{rel_path}")
        rows.append(
            {
                "document_id": make_document_id(rel_path),
                "absolute_path": str(path),
                "relative_path": rel_path,
                "file_name": normalize_text(path.name),
                "extension": extension,
                "size_bytes": path.stat().st_size,
                "size_kb": round(path.stat().st_size / 1024, 2),
                **attrs,
            }
        )
    return pd.DataFrame(rows).sort_values(["extension", "relative_path"]).reset_index(drop=True)


# =============================================================================
# 形式別抽出
# =============================================================================


def bool_mark(value: Any) -> str:
    """True/False/Noneを本文に載せやすい短い表現にする。"""
    if value is True:
        return "yes"
    if value is False:
        return "no"
    return ""


def extract_docx(path: Path) -> tuple[str, dict[str, Any]]:
    """Word文書から段落、表、基本的な書式情報を抽出する。"""
    doc = Document(path)
    lines = [f"# DOCXファイル: {normalize_text(path.name)}"]
    paragraph_count = 0
    table_count = len(doc.tables)
    styled_runs = Counter()
    highlight_count = 0

    for idx, para in enumerate(doc.paragraphs, start=1):
        text = clean_text(para.text)
        if not text:
            continue
        paragraph_count += 1
        style_name = getattr(para.style, "name", "")
        lines.append(f"\n## paragraph_{idx:03d}")
        if style_name:
            lines.append(f"- style: {style_name}")
        lines.append(text)
        run_notes = []
        for run in para.runs:
            run_text = clean_text(run.text)
            if not run_text:
                continue
            flags = []
            if run.bold:
                flags.append("bold")
                styled_runs["bold"] += 1
            if run.italic:
                flags.append("italic")
                styled_runs["italic"] += 1
            if run.underline:
                flags.append("underline")
                styled_runs["underline"] += 1
            if run.font.highlight_color is not None:
                flags.append(f"highlight={run.font.highlight_color}")
                styled_runs["highlight"] += 1
                highlight_count += 1
            if run.font.color is not None and run.font.color.rgb is not None:
                flags.append(f"font_color={run.font.color.rgb}")
                styled_runs["font_color"] += 1
            if flags:
                run_notes.append(f"- {'/'.join(flags)}: {run_text}")
        if run_notes:
            lines.append("### run_styles")
            lines.extend(run_notes)

    for t_idx, table in enumerate(doc.tables, start=1):
        lines.append(f"\n## table_{t_idx:03d}")
        for r_idx, row in enumerate(table.rows, start=1):
            cells = [clean_text(cell.text).replace("\n", " / ") for cell in row.cells]
            lines.append(f"row_{r_idx:03d}: " + " | ".join(cells))

    metadata = {
        "extraction_method": "docx_paragraphs_tables_basic_styles",
        "paragraph_count": paragraph_count,
        "table_count": table_count,
        "slide_count": 0,
        "sheet_count": 0,
        "page_count": 0,
        "formula_count": 0,
        "highlight_count": highlight_count,
        "styled_run_count": sum(styled_runs.values()),
        "note": dict(styled_runs),
    }
    return clean_text("\n".join(lines)), metadata


def iter_shape_text(shape: Any) -> list[str]:
    """PowerPointの図形から、本文と表セルのテキストを取り出す。"""
    texts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = clean_text(shape.text)
        if text:
            texts.append(text)
    if getattr(shape, "has_table", False):
        for r_idx, row in enumerate(shape.table.rows, start=1):
            cells = [clean_text(cell.text).replace("\n", " / ") for cell in row.cells]
            if any(cells):
                texts.append(f"table_row_{r_idx:03d}: " + " | ".join(cells))
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            texts.extend(iter_shape_text(child))
    return texts


def extract_pptx(path: Path) -> tuple[str, dict[str, Any]]:
    """PowerPointからスライド単位のテキストと表セルを抽出する。"""
    prs = Presentation(path)
    lines = [f"# PPTXファイル: {normalize_text(path.name)}"]
    slide_count = len(prs.slides)
    shape_count = 0
    table_shape_count = 0

    for s_idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"\n## slide_{s_idx:03d}")
        for shape in slide.shapes:
            shape_count += 1
            if getattr(shape, "has_table", False):
                table_shape_count += 1
            texts = iter_shape_text(shape)
            for text in texts:
                lines.append(text)

    metadata = {
        "extraction_method": "pptx_slide_text_tables",
        "paragraph_count": 0,
        "table_count": table_shape_count,
        "slide_count": slide_count,
        "sheet_count": 0,
        "page_count": 0,
        "formula_count": 0,
        "highlight_count": 0,
        "styled_run_count": 0,
        "note": {"shape_count": shape_count},
    }
    return clean_text("\n".join(lines)), metadata


def cell_style_flags(cell: Any) -> list[str]:
    """Excelセルの主な書式フラグを抽出する。"""
    flags: list[str] = []
    if cell.font is not None and cell.font.bold:
        flags.append("bold")
    if cell.font is not None and cell.font.italic:
        flags.append("italic")
    if cell.font is not None and cell.font.underline:
        flags.append("underline")
    if cell.fill is not None and cell.fill.fill_type:
        flags.append(f"fill={cell.fill.fill_type}")
    if cell.font is not None and cell.font.color is not None and cell.font.color.type == "rgb":
        flags.append(f"font_color={cell.font.color.rgb}")
    return flags


def extract_xlsx(
    path: Path,
    max_rows_per_sheet: int,
    max_cols_per_sheet: int,
) -> tuple[str, dict[str, Any], list[dict[str, Any]]]:
    """Excelブックからシート概要、セル値、数式、書式情報を抽出する。"""
    wb = load_workbook(path, data_only=False, read_only=False)
    lines = [f"# XLSXファイル: {normalize_text(path.name)}"]
    formula_count = 0
    styled_cell_count = 0
    hidden_sheet_count = 0
    sheet_rows: list[dict[str, Any]] = []

    for ws in wb.worksheets:
        max_row = ws.max_row or 0
        max_col = ws.max_column or 0
        used_rows = min(max_row, max_rows_per_sheet)
        used_cols = min(max_col, max_cols_per_sheet)
        if ws.sheet_state != "visible":
            hidden_sheet_count += 1

        filter_ref = getattr(ws.auto_filter, "ref", None) or ""
        freeze_panes = str(ws.freeze_panes or "")
        sheet_formula_count = 0
        sheet_styled_count = 0
        non_empty_count = 0
        truncated = max_row > max_rows_per_sheet or max_col > max_cols_per_sheet

        lines.append(f"\n## sheet: {ws.title}")
        lines.append(f"- size: rows={max_row}, cols={max_col}")
        lines.append(f"- visible_state: {ws.sheet_state}")
        if filter_ref:
            lines.append(f"- auto_filter: {filter_ref}")
        if freeze_panes:
            lines.append(f"- freeze_panes: {freeze_panes}")
        if truncated:
            lines.append(f"- note: preview limited to rows={used_rows}, cols={used_cols}")

        for row in ws.iter_rows(min_row=1, max_row=used_rows, min_col=1, max_col=used_cols):
            values: list[str] = []
            style_notes: list[str] = []
            for cell in row:
                value = cell.value
                if value is None:
                    values.append("")
                else:
                    non_empty_count += 1
                    text_value = normalize_text(value)
                    values.append(text_value.replace("\n", " / "))
                    if isinstance(value, str) and value.startswith("="):
                        formula_count += 1
                        sheet_formula_count += 1

                flags = cell_style_flags(cell)
                if flags:
                    styled_cell_count += 1
                    sheet_styled_count += 1
                    style_notes.append(f"{cell.coordinate}:{'/'.join(flags)}")

            if any(values):
                lines.append(f"row_{row[0].row:05d}: " + " | ".join(values))
            if style_notes:
                lines.append("styles: " + "; ".join(style_notes[:20]))

        sheet_rows.append(
            {
                "file_name": normalize_text(path.name),
                "sheet_name": ws.title,
                "max_row": max_row,
                "max_column": max_col,
                "preview_rows": used_rows,
                "preview_columns": used_cols,
                "non_empty_preview_cells": non_empty_count,
                "formula_count_preview": sheet_formula_count,
                "styled_cell_count_preview": sheet_styled_count,
                "auto_filter_ref": filter_ref,
                "freeze_panes": freeze_panes,
                "sheet_state": ws.sheet_state,
                "truncated": truncated,
            }
        )

    metadata = {
        "extraction_method": "xlsx_sheet_values_formulas_basic_styles",
        "paragraph_count": 0,
        "table_count": 0,
        "slide_count": 0,
        "sheet_count": len(wb.worksheets),
        "page_count": 0,
        "formula_count": formula_count,
        "highlight_count": 0,
        "styled_run_count": styled_cell_count,
        "note": {"hidden_sheet_count": hidden_sheet_count},
    }
    return clean_text("\n".join(lines)), metadata, sheet_rows


def extract_pdf(path: Path) -> tuple[str, dict[str, Any]]:
    """PDFからページ単位のテキストを抽出する。"""
    reader = PdfReader(str(path))
    lines = [f"# PDFファイル: {normalize_text(path.name)}"]
    page_count = len(reader.pages)
    empty_pages = 0

    for p_idx, page in enumerate(reader.pages, start=1):
        text = clean_text(page.extract_text() or "")
        lines.append(f"\n## page_{p_idx:03d}")
        if text:
            lines.append(text)
        else:
            empty_pages += 1
            lines.append("[テキスト抽出なし]")

    metadata = {
        "extraction_method": "pdf_page_text_pypdf",
        "paragraph_count": 0,
        "table_count": 0,
        "slide_count": 0,
        "sheet_count": 0,
        "page_count": page_count,
        "formula_count": 0,
        "highlight_count": 0,
        "styled_run_count": 0,
        "note": {"empty_pages": empty_pages},
    }
    return clean_text("\n".join(lines)), metadata


def extract_file(
    row: pd.Series,
    max_rows_per_sheet: int,
    max_cols_per_sheet: int,
) -> tuple[dict[str, Any], list[dict[str, Any]], list[dict[str, Any]]]:
    """1ファイルを抽出し、文書レコードとシート補助レコードを返す。"""
    path = Path(row["absolute_path"])
    extension = str(row["extension"])
    sheet_rows: list[dict[str, Any]] = []

    if extension == ".docx":
        text, metadata = extract_docx(path)
    elif extension == ".pptx":
        text, metadata = extract_pptx(path)
    elif extension == ".xlsx":
        text, metadata, sheet_rows = extract_xlsx(path, max_rows_per_sheet, max_cols_per_sheet)
    elif extension == ".pdf":
        text, metadata = extract_pdf(path)
    else:
        raise ValueError(f"未対応の拡張子です: {extension}")

    record = row.to_dict()
    record.update(metadata)
    record["text"] = text
    record["text_length"] = len(text)
    record["line_count"] = text.count("\n") + 1 if text else 0
    record["text_preview"] = clean_preview(text)

    for sheet_row in sheet_rows:
        sheet_row.update(
            {
                "document_id": row["document_id"],
                "relative_path": row["relative_path"],
                "project_name": row["project_name"],
                "major_folder": row["major_folder"],
            }
        )

    return record, sheet_rows, []


# =============================================================================
# チャンク化と集計
# =============================================================================


def chunk_text(text: str, chunk_size: int, chunk_overlap: int) -> list[str]:
    """長い抽出本文を検索に使いやすい固定長チャンクへ分割する。"""
    text = clean_text(text)
    if not text:
        return []
    if chunk_overlap >= chunk_size:
        raise ValueError("chunk_overlap は chunk_size より小さくしてください。")

    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end])
        if end >= len(text):
            break
        start = max(0, end - chunk_overlap)
    return chunks


def run_extraction(
    inventory_df: pd.DataFrame,
    chunk_size: int,
    chunk_overlap: int,
    max_rows_per_sheet: int,
    max_cols_per_sheet: int,
) -> tuple[pd.DataFrame, pd.DataFrame, pd.DataFrame, pd.DataFrame]:
    """対象ファイル全体の抽出とチャンク化を実行する。"""
    documents: list[dict[str, Any]] = []
    chunks: list[dict[str, Any]] = []
    sheet_rows: list[dict[str, Any]] = []
    errors: list[dict[str, Any]] = []

    for _, row in inventory_df.iterrows():
        try:
            doc, sheets, _ = extract_file(row, max_rows_per_sheet, max_cols_per_sheet)
            documents.append(doc)
            sheet_rows.extend(sheets)

            for chunk_idx, chunk in enumerate(chunk_text(doc["text"], chunk_size, chunk_overlap)):
                chunk_id = f"{doc['document_id']}_chunk_{chunk_idx:04d}"
                chunks.append(
                    {
                        "chunk_id": chunk_id,
                        "document_id": doc["document_id"],
                        "chunk_index": chunk_idx,
                        "extension": doc["extension"],
                        "relative_path": doc["relative_path"],
                        "file_name": doc["file_name"],
                        "project_name": doc["project_name"],
                        "major_folder": doc["major_folder"],
                        "chunk_text": chunk,
                        "chunk_length": len(chunk),
                    }
                )
        except Exception as exc:  # noqa: BLE001 - EDAでは失敗ファイルを一覧化して処理を継続する。
            logging.exception("Failed to extract: %s", row.get("relative_path"))
            errors.append(
                {
                    "document_id": row.get("document_id", ""),
                    "relative_path": row.get("relative_path", ""),
                    "extension": row.get("extension", ""),
                    "error_type": type(exc).__name__,
                    "error_message": str(exc),
                }
            )

    doc_df = pd.DataFrame(documents)
    chunk_df = pd.DataFrame(chunks)
    sheet_df = pd.DataFrame(sheet_rows)
    error_df = pd.DataFrame(errors)

    write_jsonl(documents, TEXT_DIR / "extracted_documents.jsonl")
    write_jsonl(chunks, TEXT_DIR / "text_chunks.jsonl")
    write_jsonl(documents, PROCESSED_TEXT_DIR / "extracted_documents.jsonl")
    write_jsonl(chunks, PROCESSED_TEXT_DIR / "text_chunks.jsonl")

    return doc_df, chunk_df, sheet_df, error_df


def make_summary_tables(
    inventory_df: pd.DataFrame,
    doc_df: pd.DataFrame,
    chunk_df: pd.DataFrame,
    sheet_df: pd.DataFrame,
    error_df: pd.DataFrame,
) -> dict[str, pd.DataFrame]:
    """レポートと確認用CSVに使う集計表を作る。"""
    target_extension_counts = (
        inventory_df.groupby("extension", as_index=False)
        .agg(file_count=("document_id", "count"), total_size_kb=("size_kb", "sum"))
        .sort_values(["file_count", "extension"], ascending=[False, True])
    )
    target_extension_counts["total_size_kb"] = target_extension_counts["total_size_kb"].round(2)

    success_by_ext = doc_df.groupby("extension").size().rename("success_count") if not doc_df.empty else pd.Series(dtype=int)
    error_by_ext = error_df.groupby("extension").size().rename("error_count") if not error_df.empty else pd.Series(dtype=int)
    extraction_status = target_extension_counts.set_index("extension").join(success_by_ext).join(error_by_ext)
    extraction_status[["success_count", "error_count"]] = extraction_status[["success_count", "error_count"]].fillna(0).astype(int)
    extraction_status["success_rate"] = (
        extraction_status["success_count"] / extraction_status["file_count"].replace(0, pd.NA)
    ).fillna(0).round(4)
    extraction_status = extraction_status.reset_index()

    extraction_summary = (
        doc_df.groupby(["extension", "extraction_method"], as_index=False)
        .agg(
            extracted_files=("document_id", "count"),
            total_text_length=("text_length", "sum"),
            mean_text_length=("text_length", "mean"),
            total_line_count=("line_count", "sum"),
            total_formula_count=("formula_count", "sum"),
            total_styled_count=("styled_run_count", "sum"),
        )
        if not doc_df.empty
        else pd.DataFrame()
    )
    if not extraction_summary.empty:
        extraction_summary["mean_text_length"] = extraction_summary["mean_text_length"].round(1)

    chunk_summary = (
        chunk_df.groupby("extension", as_index=False)
        .agg(
            chunk_count=("chunk_id", "count"),
            mean_chunk_length=("chunk_length", "mean"),
            max_chunk_length=("chunk_length", "max"),
        )
        if not chunk_df.empty
        else pd.DataFrame()
    )
    if not chunk_summary.empty:
        chunk_summary["mean_chunk_length"] = chunk_summary["mean_chunk_length"].round(1)

    file_text_length_ranking = (
        doc_df[
            [
                "document_id",
                "extension",
                "relative_path",
                "extraction_method",
                "size_bytes",
                "text_length",
                "line_count",
                "text_preview",
            ]
        ]
        .sort_values("text_length", ascending=False)
        .reset_index(drop=True)
        if not doc_df.empty
        else pd.DataFrame()
    )

    sample_documents = (
        doc_df[
            [
                "document_id",
                "extension",
                "relative_path",
                "extraction_method",
                "text_length",
                "paragraph_count",
                "table_count",
                "slide_count",
                "sheet_count",
                "page_count",
                "text_preview",
            ]
        ]
        .sort_values(["extension", "relative_path"])
        .head(20)
        .reset_index(drop=True)
        if not doc_df.empty
        else pd.DataFrame()
    )

    project_extension_counts = (
        inventory_df.groupby(["project_name", "extension"], as_index=False)
        .agg(file_count=("document_id", "count"))
        .sort_values(["project_name", "extension"])
    )

    return {
        "target_file_inventory": inventory_df,
        "target_extension_counts": target_extension_counts,
        "extraction_status": extraction_status,
        "extraction_summary": extraction_summary,
        "chunk_summary": chunk_summary,
        "file_text_length_ranking": file_text_length_ranking,
        "sample_documents": sample_documents,
        "sheet_summary": sheet_df,
        "project_extension_counts": project_extension_counts,
        "extraction_errors": error_df,
    }


def save_tables(tables: dict[str, pd.DataFrame]) -> None:
    """集計表をCSVとして保存する。"""
    for name, df in tables.items():
        save_csv(df, TABLE_DIR / f"{name}.csv")


def plot_extension_counts(df: pd.DataFrame) -> None:
    """対象拡張子別ファイル数を可視化する。"""
    if df.empty:
        return
    fig, ax = plt.subplots(figsize=(8, 4))
    ax.bar(df["extension"], df["file_count"], color="#3B82F6")
    ax.set_title("EDA004 target file counts by extension")
    ax.set_xlabel("extension")
    ax.set_ylabel("file_count")
    fig.tight_layout()
    fig.savefig(FIG_DIR / "01_target_extension_counts.png", dpi=160)
    plt.close(fig)


def plot_text_lengths(doc_df: pd.DataFrame) -> None:
    """抽出テキスト長の分布を可視化する。"""
    if doc_df.empty:
        return
    fig, ax = plt.subplots(figsize=(8, 4))
    for extension, group in doc_df.groupby("extension"):
        ax.hist(group["text_length"], bins=20, alpha=0.55, label=extension)
    ax.set_title("Extracted text length distribution")
    ax.set_xlabel("text_length")
    ax.set_ylabel("file_count")
    ax.legend()
    fig.tight_layout()
    fig.savefig(FIG_DIR / "02_text_length_distribution.png", dpi=160)
    plt.close(fig)


def plot_chunk_counts(chunk_summary: pd.DataFrame) -> None:
    """拡張子別チャンク数を可視化する。"""
    if chunk_summary.empty:
        return
    fig, ax = plt.subplots(figsize=(8, 4))
    ax.bar(chunk_summary["extension"], chunk_summary["chunk_count"], color="#10B981")
    ax.set_title("Chunk counts by extension")
    ax.set_xlabel("extension")
    ax.set_ylabel("chunk_count")
    fig.tight_layout()
    fig.savefig(FIG_DIR / "03_chunk_counts_by_extension.png", dpi=160)
    plt.close(fig)


def plot_top_text_files(file_text_length_ranking: pd.DataFrame) -> None:
    """抽出テキスト量が大きいファイルを可視化する。"""
    if file_text_length_ranking.empty:
        return
    view = file_text_length_ranking.head(15).copy()
    labels = view["relative_path"].map(lambda x: Path(str(x)).name)
    fig, ax = plt.subplots(figsize=(9, 5))
    ax.barh(labels, view["text_length"], color="#F59E0B")
    ax.invert_yaxis()
    ax.set_title("Top extracted text length files")
    ax.set_xlabel("text_length")
    fig.tight_layout()
    fig.savefig(FIG_DIR / "04_top_text_length_files.png", dpi=160)
    plt.close(fig)


def make_report(
    drive_root: Path,
    tables: dict[str, pd.DataFrame],
    chunk_size: int,
    chunk_overlap: int,
    max_rows_per_sheet: int,
    max_cols_per_sheet: int,
) -> None:
    """EDA004のMarkdownレポートを生成する。"""
    inventory_df = tables["target_file_inventory"]
    doc_df = pd.read_csv(TABLE_DIR / "file_text_length_ranking.csv") if (TABLE_DIR / "file_text_length_ranking.csv").exists() else pd.DataFrame()
    status_df = tables["extraction_status"]
    summary_df = tables["extraction_summary"]
    chunk_summary = tables["chunk_summary"]
    file_ranking = tables["file_text_length_ranking"]
    sheet_summary = tables["sheet_summary"]
    error_df = tables["extraction_errors"]
    sample_documents = tables["sample_documents"]

    total_files = len(inventory_df)
    success_count = int(status_df["success_count"].sum()) if not status_df.empty else 0
    error_count = int(status_df["error_count"].sum()) if not status_df.empty else 0
    total_chunks = int(chunk_summary["chunk_count"].sum()) if not chunk_summary.empty else 0
    total_text_length = int(doc_df["text_length"].sum()) if not doc_df.empty and "text_length" in doc_df else 0

    lines: list[str] = []
    lines.append("# EDA004: Office文書・PDFの抽出ベースライン")
    lines.append("")
    lines.append("## 目的・背景")
    lines.append("")
    lines.append(
        "EDA003では、EDA002で抽出した `.md`, `.csv`, `.json`, `.py`, `.ipynb` の検索ベースラインを確認した。"
        "その結果、検索に失敗したvalid質問の多くは、Word、PowerPoint、Excel、PDF、書式情報、セル情報、"
        "グラフや画像など、EDA002の対象外形式を必要としていることが分かった。"
    )
    lines.append("")
    lines.append(
        "EDA004では、次の段階として `.docx`, `.pptx`, `.xlsx`, `.pdf` を対象に、本文と一部メタ情報を抽出する。"
        "ここでは最終回答生成までは行わず、後続の検索・集計・照合処理に渡せるテキスト表現を作ることを目的とする。"
    )
    lines.append("")
    lines.append("## 入力データ")
    lines.append("")
    lines.append(f"- 共有ドライブ: `{safe_relative_to(drive_root, BASE_DIR)}`")
    lines.append(f"- 対象拡張子: `{', '.join(sorted(TARGET_EXTENSIONS))}`")
    lines.append(f"- チャンクサイズ: {chunk_size} 文字")
    lines.append(f"- チャンクオーバーラップ: {chunk_overlap} 文字")
    lines.append(f"- Excelプレビュー上限: 1シートあたり {max_rows_per_sheet} 行 x {max_cols_per_sheet} 列")
    lines.append("")
    lines.append("## 出力ファイル")
    lines.append("")
    lines.append("- `EDA/EDA004/texts/extracted_documents.jsonl`: 1ファイル1レコードの抽出本文")
    lines.append("- `EDA/EDA004/texts/text_chunks.jsonl`: 検索インデックス投入用のチャンク")
    lines.append("- `data/processed/office_pdf_baseline/extracted_documents.jsonl`: 後続処理向けコピー")
    lines.append("- `data/processed/office_pdf_baseline/text_chunks.jsonl`: 後続処理向けコピー")
    lines.append("- `EDA/EDA004/tables/*.csv`: 抽出結果、シート概要、エラー一覧、プレビュー")
    lines.append("")
    lines.append("## 全体サマリ")
    lines.append("")
    lines.append(f"- 対象ファイル数: {total_files}")
    lines.append(f"- 抽出成功: {success_count}")
    lines.append(f"- 抽出失敗: {error_count}")
    lines.append(f"- 作成チャンク数: {total_chunks}")
    lines.append(f"- 抽出テキスト総文字数: {total_text_length}")
    lines.append("")
    lines.append("## 拡張子別の抽出状況")
    lines.append("")
    lines.append(df_to_markdown(status_df))
    lines.append("")
    lines.append("凡例: `extension` は拡張子、`file_count` は対象ファイル数、`total_size_kb` は合計サイズ、`success_count` は抽出成功数、`error_count` は抽出失敗数、`success_rate` は成功率を表します。")
    lines.append("")
    lines.append("## 抽出方法別サマリ")
    lines.append("")
    lines.append(df_to_markdown(summary_df))
    lines.append("")
    lines.append("凡例: `extraction_method` は抽出方法、`extracted_files` は抽出済みファイル数、`total_text_length` は抽出文字数、`total_formula_count` はExcelプレビュー内の数式数、`total_styled_count` は検出した書式付き要素数を表します。")
    lines.append("")
    lines.append("## チャンク数サマリ")
    lines.append("")
    lines.append(df_to_markdown(chunk_summary))
    lines.append("")
    lines.append("凡例: `chunk_count` は作成した検索用チャンク数、`mean_chunk_length` は平均チャンク長、`max_chunk_length` は最大チャンク長を表します。")
    lines.append("")
    lines.append("## Excelシート概要")
    lines.append("")
    if sheet_summary.empty:
        lines.append("Excelシート概要はありません。")
    else:
        lines.append(df_to_markdown(sheet_summary.head(30)))
        lines.append("")
        lines.append("凡例: `max_row` と `max_column` はシートの使用範囲、`preview_rows` と `preview_columns` は抽出対象にした範囲、`formula_count_preview` はプレビュー範囲内の数式数、`auto_filter_ref` はフィルター範囲を表します。")
    lines.append("")
    lines.append("## ファイル別テキスト量ランキング")
    lines.append("")
    lines.append(df_to_markdown(file_ranking.head(20)))
    lines.append("")
    lines.append("凡例: `text_length` は抽出本文の文字数、`line_count` は抽出本文の行数、`text_preview` は本文冒頭の確認用プレビューを表します。")
    lines.append("")
    lines.append("## 抽出本文サンプル")
    lines.append("")
    lines.append(df_to_markdown(sample_documents, max_rows=20))
    lines.append("")
    lines.append("## 考察")
    lines.append("")
    lines.append(
        "EDA004により、EDA003で不足していたOffice文書とPDFを機械的にテキスト化できる状態になった。"
        "特にWordとPowerPointは、提案書、契約書、最終報告書などの本文検索に使える可能性が高い。"
    )
    lines.append("")
    lines.append(
        "Excelについては、セル値、シート名、フィルター範囲、数式、基本的なセル書式を抽出している。"
        "ただし、PivotTable、条件付き書式、グラフ、非表示フィルターの詳細などは、この初期抽出だけでは完全には扱えない。"
        "表データに対する条件抽出や集計は、RAGチャンク検索とは別に、openpyxlやpandasで直接処理する方針が必要である。"
    )
    lines.append("")
    lines.append(
        "PDFはpypdfでページ単位のテキスト抽出を行っている。画像化されたPDFや複雑な表は本文が欠落する可能性があるため、"
        "抽出なしページや短すぎるページは後続でOCR候補として扱う。"
    )
    lines.append("")
    lines.append("## 抽出失敗")
    lines.append("")
    if error_df.empty:
        lines.append("抽出失敗はありません。")
    else:
        lines.append(df_to_markdown(error_df, max_rows=30))
    lines.append("")
    lines.append("## 次にやること")
    lines.append("")
    lines.append("1. EDA002とEDA004のチャンクを統合し、valid質問で検索ヒット率を再評価する。")
    lines.append("2. Excel/CSVの集計問題は、チャンク検索ではなく表データを直接読んで計算する処理を作る。")
    lines.append("3. 書式が問われる問題では、Word/PPTのrunや図形単位の書式抽出をさらに細かくする。")
    lines.append("4. 画像化PDF、pngグラフ、Office内画像はOCRまたは画像理解の対象として別EDAで扱う。")

    REPORT_PATH.write_text("\n".join(lines), encoding="utf-8")


def validate_report_artifact() -> None:
    """生成したMarkdown成果物に必要な章が含まれるか確認する。"""
    report_text = REPORT_PATH.read_text(encoding="utf-8")
    required_sections = [
        "## 目的・背景",
        "## 全体サマリ",
        "## 拡張子別の抽出状況",
        "## Excelシート概要",
        "## 考察",
        "## 次にやること",
    ]
    missing = [section for section in required_sections if section not in report_text]
    if missing:
        raise RuntimeError(f"eda004_report.mdに必要な章が不足しています: {missing}")


# =============================================================================
# main
# =============================================================================


def parse_args() -> argparse.Namespace:
    """コマンドライン引数を読む。"""
    parser = argparse.ArgumentParser(description="EDA004: Office and PDF extraction baseline.")
    parser.add_argument("--chunk-size", type=int, default=DEFAULT_CHUNK_SIZE)
    parser.add_argument("--chunk-overlap", type=int, default=DEFAULT_CHUNK_OVERLAP)
    parser.add_argument("--xlsx-max-rows-per-sheet", type=int, default=DEFAULT_XLSX_MAX_ROWS_PER_SHEET)
    parser.add_argument("--xlsx-max-cols-per-sheet", type=int, default=DEFAULT_XLSX_MAX_COLS_PER_SHEET)
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    setup()

    drive_root = find_drive_root()
    inventory_df = build_target_file_inventory(drive_root)
    doc_df, chunk_df, sheet_df, error_df = run_extraction(
        inventory_df,
        chunk_size=args.chunk_size,
        chunk_overlap=args.chunk_overlap,
        max_rows_per_sheet=args.xlsx_max_rows_per_sheet,
        max_cols_per_sheet=args.xlsx_max_cols_per_sheet,
    )
    tables = make_summary_tables(inventory_df, doc_df, chunk_df, sheet_df, error_df)
    save_tables(tables)
    plot_extension_counts(tables["target_extension_counts"])
    plot_text_lengths(doc_df)
    plot_chunk_counts(tables["chunk_summary"])
    plot_top_text_files(tables["file_text_length_ranking"])
    make_report(
        drive_root,
        tables,
        chunk_size=args.chunk_size,
        chunk_overlap=args.chunk_overlap,
        max_rows_per_sheet=args.xlsx_max_rows_per_sheet,
        max_cols_per_sheet=args.xlsx_max_cols_per_sheet,
    )
    validate_report_artifact()

    print(f"EDA004 finished: {REPORT_PATH}")
    print(f"tables: {TABLE_DIR}")
    print(f"texts: {TEXT_DIR}")
    print(f"processed copy: {PROCESSED_TEXT_DIR}")


if __name__ == "__main__":
    main()
