from __future__ import annotations

import csv
import importlib.metadata
import re
import time
import unicodedata
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from .io_utils import assert_formal_input_allowed, relative_to_root, sha1_text, write_csv, write_jsonl
from .schemas import DecryptionAttempt, FileRecord, PasswordCandidate, ProtectedFileRecord, SCHEMA_VERSION, to_dict


OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}
RESOLUTION_AMBIGUOUS = {"ambiguous_project", "ambiguous_alias", "ambiguous_contract_date"}


@dataclass
class CandidateSecret:
    record: PasswordCandidate
    password: str


@dataclass
class ContractDateCandidate:
    date: str
    source_path: str
    location: str
    matched_label: str
    confidence: float


@dataclass
class ProjectAliasResolution:
    project_name: str
    alias: str
    source: str
    confidence: float
    candidates: list[dict[str, Any]] = field(default_factory=list)
    error_status: str = ""


@dataclass
class PasswordRule:
    valid: bool
    source_path: str
    reason: str
    prefix: str = "DA"


@dataclass
class ProtectedResolutionResult:
    records: list[ProtectedFileRecord]
    candidates: list[PasswordCandidate]
    attempts: list[DecryptionAttempt]
    input_path_overrides: dict[str, Path]


def normalize_text(value: str) -> str:
    return unicodedata.normalize("NFKC", value).replace("\xa0", " ")


def masked_path(path: str) -> str:
    text = normalize_text(path)
    text = re.sub(r"(?i)(pw|password)\s*=\s*([^/\\\s()（）\[\]{}]+)", r"\1=***", text)
    text = re.sub(r"(?i)(pw)-([A-Za-z0-9][A-Za-z0-9_.-]*)", r"\1-***", text)
    return text


def mask_password(password: str) -> str:
    return f"***len{len(password)}***"


def password_hash(password: str) -> str:
    return sha1_text(password)


def safe_name(path: Path) -> str:
    return re.sub(r"[^0-9A-Za-z_.-]+", "_", path.stem).strip("_") or "decrypted"


def msoffcrypto_version() -> str:
    try:
        return importlib.metadata.version("msoffcrypto-tool")
    except importlib.metadata.PackageNotFoundError:
        return "unknown"


def find_management_doc(raw_root: Path, project_root: Path, name_part: str) -> Path | None:
    normalized_part = normalize_text(name_part)
    for path in raw_root.rglob("*.docx"):
        assert_formal_input_allowed(path, project_root)
        normalized_name = normalize_text(path.name)
        parent_text = normalize_text(path.parent.as_posix())
        if normalized_part in normalized_name and "社内管理" in parent_text:
            return path
    for path in raw_root.rglob("*.docx"):
        assert_formal_input_allowed(path, project_root)
        if normalized_part in normalize_text(path.name):
            return path
    return None


def read_docx_text(path: Path) -> str:
    from docx import Document

    document = Document(path)
    chunks: list[str] = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
    for table in document.tables:
        for row in table.rows:
            chunks.append("\t".join(cell.text.strip() for cell in row.cells))
    return "\n".join(chunks)


def read_password_rule(raw_root: Path, project_root: Path) -> PasswordRule:
    path = find_management_doc(raw_root, project_root, "パスワード導出規則")
    if path is None:
        return PasswordRule(False, "", "rule_document_not_found")
    text = normalize_text(read_docx_text(path))
    required = ["DA-[案件略号]-[開始年月日8桁]-[拡張子コード]", "主略称", "契約開始日", "YYYYMMDD"]
    missing = [item for item in required if item not in text]
    return PasswordRule(
        valid=not missing,
        source_path=masked_path(relative_to_root(path, project_root)),
        reason="ok" if not missing else "missing:" + ",".join(missing),
    )


def read_project_aliases(raw_root: Path, project_root: Path) -> list[dict[str, Any]]:
    path = find_management_doc(raw_root, project_root, "社内用語集")
    if path is None:
        return []
    from docx import Document

    aliases: list[dict[str, Any]] = []
    document = Document(path)
    source = masked_path(relative_to_root(path, project_root))
    for table in document.tables:
        if not table.rows:
            continue
        headers = [normalize_text(cell.text.strip()) for cell in table.rows[0].cells]
        if "案件名" not in headers or "主略称" not in headers:
            continue
        name_index = headers.index("案件名")
        alias_index = headers.index("主略称")
        alt_index = headers.index("別名候補") if "別名候補" in headers else -1
        for row in table.rows[1:]:
            cells = [normalize_text(cell.text.strip()) for cell in row.cells]
            if len(cells) <= max(name_index, alias_index):
                continue
            project_name = cells[name_index]
            alias = cells[alias_index]
            if not project_name or not alias:
                continue
            alternatives = []
            if alt_index >= 0 and len(cells) > alt_index:
                alternatives = [item.strip() for item in re.split(r"[,、/]", cells[alt_index]) if item.strip()]
            aliases.append({"project_name": project_name, "alias": alias, "alternatives": alternatives, "source": source})
    return aliases


def resolve_project_alias(project_name: str, file: FileRecord, aliases: list[dict[str, Any]]) -> ProjectAliasResolution:
    normalized_project = normalize_text(project_name)
    candidates: list[dict[str, Any]] = []
    for row in aliases:
        official = normalize_text(row["project_name"])
        alternatives = [normalize_text(item) for item in row.get("alternatives", [])]
        confidence = 0.0
        reason = ""
        if normalized_project == official:
            confidence = 0.99
            reason = "project_name_exact"
        elif normalized_project in official or official in normalized_project:
            confidence = 0.85
            reason = "project_name_partial"
        elif any(alt and alt in normalize_text(file.relative_path + " " + file.file_name) for alt in alternatives):
            confidence = 0.7
            reason = "alias_in_path"
        if confidence:
            candidates.append({"alias": row["alias"], "project_name": row["project_name"], "source": row["source"], "confidence": confidence, "reason": reason})

    candidates.sort(key=lambda item: item["confidence"], reverse=True)
    if not candidates:
        return ProjectAliasResolution(project_name, "", "", 0.0, [], "ambiguous_alias")
    top = candidates[0]
    tied = [item for item in candidates if item["confidence"] == top["confidence"]]
    if len({item["alias"] for item in tied}) > 1:
        return ProjectAliasResolution(project_name, "", "", top["confidence"], candidates, "ambiguous_alias")
    return ProjectAliasResolution(project_name, top["alias"], f"{top['source']}:{top['reason']}", top["confidence"], candidates)


def is_office_temporary(file: FileRecord) -> bool:
    name = normalize_text(file.file_name)
    return file.is_temp_office_file or name.startswith("~$") or name.startswith(".~lock.")


def detect_protection(file: FileRecord, source_path: Path, project_root: Path) -> tuple[str, str, bool, str]:
    if is_office_temporary(file):
        return "temporary_file", "office_temp_name", False, ""
    if file.extension not in OFFICE_EXTENSIONS:
        return "not_protected", "not_office_extension", False, ""
    assert_formal_input_allowed(source_path, project_root)
    try:
        import msoffcrypto

        with source_path.open("rb") as handle:
            office = msoffcrypto.OfficeFile(handle)
            if office.is_encrypted():
                return "password_protected", "msoffcrypto", True, ""
            return "not_protected", "msoffcrypto", False, ""
    except Exception as exc:
        if zipfile.is_zipfile(source_path):
            return "not_protected", "zipfile_fallback", False, ""
        return "corrupted_file", "msoffcrypto_error", False, f"{type(exc).__name__}: {exc}"


def extract_filename_passwords(file: FileRecord, project_root: Path) -> list[CandidateSecret]:
    stem = normalize_text(Path(file.file_name).stem)
    found: list[tuple[str, str]] = []
    for match in re.finditer(r"(?i)(?:^|[_\s(（])(?P<key>password|pw)\s*=\s*(?P<password>[^\s()（）\[\]{}「」]+)", stem):
        key = match.group("key").lower()
        password = match.group("password").strip(" _-.,，。")
        if password:
            found.append((f"filename_{key}", password))
    for match in re.finditer(r"(?i)(?:^|[_\s(（])pw-(?P<password>[A-Za-z0-9][A-Za-z0-9_.-]*)", stem):
        password = match.group("password").strip(" _-.,，。")
        if password:
            found.append(("filename_pw_legacy", password))

    candidates: list[CandidateSecret] = []
    for method, password in list(dict.fromkeys(found)):
        candidate_id = "pwd_" + sha1_text(f"{file.file_id}:{method}:{password_hash(password)}")[:16]
        record = PasswordCandidate(
            schema_version=SCHEMA_VERSION,
            candidate_id=candidate_id,
            file_id=file.file_id,
            derivation_method=method,
            project_name=file.project_name,
            project_alias="",
            contract_start_date="",
            extension_code=file.extension.lstrip("."),
            source_paths=[masked_path(relative_to_root(project_root / file.raw_path, project_root))],
            confidence=0.95,
            masked_password=mask_password(password),
            password_hash=password_hash(password),
        )
        candidates.append(CandidateSecret(record, password))
    return candidates


def validate_decrypted_file(file: FileRecord, path: Path) -> tuple[bool, str]:
    if not path.exists() or path.stat().st_size <= 0:
        return False, "decrypted_file_empty"
    if file.extension in OFFICE_EXTENSIONS:
        if not zipfile.is_zipfile(path):
            return False, "decrypted_file_is_not_zip_ooxml"
        with zipfile.ZipFile(path) as archive:
            names = set(archive.namelist())
        required = {
            ".docx": "word/document.xml",
            ".xlsx": "xl/workbook.xml",
            ".pptx": "ppt/presentation.xml",
        }[file.extension]
        if required not in names:
            return False, f"missing_ooxml_part:{required}"
    try:
        if file.extension == ".docx":
            from docx import Document

            Document(path)
        elif file.extension == ".xlsx":
            import openpyxl

            workbook = openpyxl.load_workbook(path, read_only=True, data_only=False)
            workbook.close()
        elif file.extension == ".pptx":
            from pptx import Presentation

            Presentation(path)
    except Exception as exc:
        return False, f"parser_validation_failed:{type(exc).__name__}: {exc}"
    return True, ""


def decrypt_with_candidate(
    file: FileRecord,
    source_path: Path,
    candidate: CandidateSecret,
    decrypted_root: Path,
    attempt_order: int,
) -> tuple[DecryptionAttempt, Path | None]:
    import msoffcrypto

    started = time.perf_counter()
    out_dir = decrypted_root / file.file_id
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"decrypted{file.extension}"
    error_type = ""
    error_message = ""
    success = False
    validation_success = False
    try:
        with source_path.open("rb") as src:
            office = msoffcrypto.OfficeFile(src)
            try:
                office.load_key(password=candidate.password, verify_password=True)
            except TypeError:
                office.load_key(password=candidate.password)
            with out_path.open("wb") as dst:
                office.decrypt(dst)
        success = True
        validation_success, validation_error = validate_decrypted_file(file, out_path)
        if not validation_success:
            error_type = "validation_failed"
            error_message = validation_error
    except Exception as exc:
        error_type = type(exc).__name__
        error_message = str(exc)[:300]
        if out_path.exists():
            out_path.unlink()
    elapsed = round(time.perf_counter() - started, 4)
    attempt = DecryptionAttempt(
        schema_version=SCHEMA_VERSION,
        file_id=file.file_id,
        candidate_id=candidate.record.candidate_id,
        attempt_order=attempt_order,
        success=success,
        validation_success=validation_success,
        library=f"msoffcrypto-tool {msoffcrypto_version()}",
        elapsed_seconds=elapsed,
        error_type=error_type,
        error_message=error_message,
    )
    return attempt, out_path if success and validation_success else None


DATE_RE = re.compile(r"(20\d{2})\s*[年/\-.]\s*(\d{1,2})\s*[月/\-.]\s*(\d{1,2})\s*日?")
DATE_LABELS = [
    ("契約開始日", 0.98),
    ("契約期間", 0.9),
    ("契約開始", 0.86),
    ("開始日", 0.72),
]


def normalize_date(match: re.Match[str]) -> str:
    year, month, day = match.groups()
    return f"{int(year):04d}{int(month):02d}{int(day):02d}"


def read_text_for_date(file: FileRecord, path: Path) -> str:
    if file.extension == ".docx":
        return read_docx_text(path)
    if file.extension == ".pptx":
        from pptx import Presentation

        prs = Presentation(path)
        chunks: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(shape.text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        chunks.append("\t".join(cell.text for cell in row.cells))
        return "\n".join(chunks)
    if file.extension == ".xlsx":
        import openpyxl

        workbook = openpyxl.load_workbook(path, read_only=True, data_only=False)
        chunks: list[str] = []
        try:
            for ws in workbook.worksheets:
                for row in ws.iter_rows(values_only=True):
                    values = [str(value) for value in row if value is not None]
                    if values:
                        chunks.append("\t".join(values))
        finally:
            workbook.close()
        return "\n".join(chunks)
    if file.extension in {".md", ".txt", ".csv", ".tsv"}:
        return path.read_text(encoding="utf-8", errors="replace")
    if file.extension == ".pdf":
        import fitz

        doc = fitz.open(path)
        return "\n".join(page.get_text("text") for page in doc)
    return ""


def find_contract_dates(
    project_name: str,
    files: list[FileRecord],
    project_root: Path,
    input_overrides: dict[str, Path],
) -> list[ContractDateCandidate]:
    candidates: list[ContractDateCandidate] = []
    normalized_project = normalize_text(project_name)
    for file in files:
        if normalize_text(file.project_name) != normalized_project:
            continue
        if file.extension not in {".docx", ".pptx", ".xlsx", ".pdf", ".md", ".txt", ".csv", ".tsv"}:
            continue
        source_path = input_overrides.get(file.file_id, project_root / file.raw_path)
        if file.extension in OFFICE_EXTENSIONS and not zipfile.is_zipfile(source_path):
            continue
        try:
            text = normalize_text(read_text_for_date(file, source_path))
        except Exception:
            continue
        if not text:
            continue
        file_bonus = 0.0
        path_text = normalize_text(file.relative_path + " " + file.file_name)
        if "契約" in path_text or file.document_kind == "contract":
            file_bonus += 0.08
        if "提案" in path_text or "計画" in path_text or "報告" in path_text:
            file_bonus += 0.02
        for match in DATE_RE.finditer(text):
            context = text[max(0, match.start() - 80) : min(len(text), match.end() + 40)]
            label = ""
            confidence = 0.2 + file_bonus
            for label_text, label_score in DATE_LABELS:
                if label_text in context:
                    label = label_text
                    confidence = max(confidence, label_score + file_bonus)
            if not label:
                continue
            candidates.append(
                ContractDateCandidate(
                    date=normalize_date(match),
                    source_path=masked_path(file.raw_path),
                    location=f"text_offset:{match.start()}",
                    matched_label=label,
                    confidence=round(min(confidence, 0.99), 4),
                )
            )
    best_by_date: dict[str, ContractDateCandidate] = {}
    for item in candidates:
        if item.date not in best_by_date or item.confidence > best_by_date[item.date].confidence:
            best_by_date[item.date] = item
    return sorted(best_by_date.values(), key=lambda item: item.confidence, reverse=True)[:5]


def make_rule_candidates(
    file: FileRecord,
    alias: ProjectAliasResolution,
    dates: list[ContractDateCandidate],
    rule: PasswordRule,
) -> list[CandidateSecret]:
    candidates: list[CandidateSecret] = []
    extension_code = file.extension.lstrip(".").lower()
    for date in dates:
        password = f"{rule.prefix}-{alias.alias}-{date.date}-{extension_code}"
        confidence = round(min(alias.confidence * 0.45 + date.confidence * 0.45 + 0.1, 0.99), 4)
        candidate_id = "pwd_" + sha1_text(f"{file.file_id}:rule:{password_hash(password)}")[:16]
        record = PasswordCandidate(
            schema_version=SCHEMA_VERSION,
            candidate_id=candidate_id,
            file_id=file.file_id,
            derivation_method="rule_derived",
            project_name=file.project_name,
            project_alias=alias.alias,
            contract_start_date=date.date,
            extension_code=extension_code,
            source_paths=[alias.source, date.source_path, rule.source_path],
            confidence=confidence,
            masked_password=mask_password(password),
            password_hash=password_hash(password),
        )
        candidates.append(CandidateSecret(record, password))
    return candidates


def protected_record_for_file(file: FileRecord, status: str, method: str, requires_password: bool, error: str) -> ProtectedFileRecord:
    if status == "temporary_file":
        resolution_status = "temporary_file"
    elif status == "not_protected":
        resolution_status = "not_protected"
    elif status == "corrupted_file":
        resolution_status = "corrupted_file"
    else:
        resolution_status = "error" if error else "not_protected"
    return ProtectedFileRecord(
        schema_version=SCHEMA_VERSION,
        file_id=file.file_id,
        source_path=masked_path(file.raw_path),
        protection_status=status,
        detection_method=method,
        is_temporary_file=is_office_temporary(file),
        requires_password=requires_password,
        resolution_status=resolution_status,
        warnings=[],
        error=error,
        project_name=file.project_name,
        file_type=file.extension.lstrip("."),
        failure_reason=error,
    )


def update_record_success(record: ProtectedFileRecord, candidate: PasswordCandidate, out_path: Path, status: str) -> None:
    record.resolution_status = status
    record.selected_derivation_method = candidate.derivation_method
    record.decrypted_work_path = out_path.as_posix()
    record.derivation_method = candidate.derivation_method
    record.project_alias = candidate.project_alias
    record.contract_start_date = candidate.contract_start_date
    if candidate.derivation_method == "rule_derived" and len(candidate.source_paths) >= 2:
        record.contract_date_source = candidate.source_paths[1]
    else:
        record.contract_date_source = next((path for path in candidate.source_paths if "契約" in path), "")
    record.decryption_success = True
    record.validation_success = True
    record.failure_reason = ""


def write_protected_outputs(output_dir: Path, records: list[ProtectedFileRecord], candidates: list[PasswordCandidate], attempts: list[DecryptionAttempt]) -> None:
    write_jsonl(output_dir / "protected_file_inventory.jsonl", [to_dict(item) for item in records])
    write_jsonl(output_dir / "password_candidates.jsonl", [to_dict(item) for item in candidates])
    write_jsonl(output_dir / "decryption_attempts.jsonl", [to_dict(item) for item in attempts])
    fields = [
        "source_path",
        "project_name",
        "file_type",
        "protection_status",
        "filename_password_hint_found",
        "project_alias",
        "project_alias_source",
        "contract_start_date",
        "contract_date_source",
        "derivation_method",
        "decryption_success",
        "validation_success",
        "extraction_success",
        "failure_reason",
    ]
    rows = []
    for item in records:
        if not item.requires_password and not item.is_temporary_file:
            continue
        row = to_dict(item)
        rows.append({field: row.get(field, "") for field in fields})
    write_csv(output_dir / "decryption_summary.csv", rows, fields)


def resolve_protected_files(
    files: list[FileRecord],
    project_root: Path,
    raw_root: Path,
    run_dir: Path,
    strict: bool = False,
) -> ProtectedResolutionResult:
    output_dir = run_dir / "protected_files"
    decrypted_root = run_dir / "decrypted"
    records: list[ProtectedFileRecord] = []
    all_candidates: list[PasswordCandidate] = []
    attempts: list[DecryptionAttempt] = []
    input_overrides: dict[str, Path] = {}
    protected_by_file: dict[str, ProtectedFileRecord] = {}

    office_files = [file for file in files if file.extension in OFFICE_EXTENSIONS]
    for file in office_files:
        source_path = project_root / file.raw_path
        status, method, requires_password, error = detect_protection(file, source_path, project_root)
        record = protected_record_for_file(file, status, method, requires_password, error)
        records.append(record)
        protected_by_file[file.file_id] = record
        if not requires_password:
            continue
        filename_candidates = extract_filename_passwords(file, project_root)
        if filename_candidates:
            record.filename_password_hint_found = True
        for secret in filename_candidates:
            all_candidates.append(secret.record)
            attempt, out_path = decrypt_with_candidate(file, source_path, secret, decrypted_root, len(attempts) + 1)
            attempts.append(attempt)
            record.attempt_count += 1
            if out_path is not None:
                input_overrides[file.file_id] = out_path
                update_record_success(record, secret.record, out_path, "resolved_from_filename")
                break
        if record.requires_password and not record.decryption_success and filename_candidates:
            record.failure_reason = "filename_password_candidates_failed"

    rule = read_password_rule(raw_root, project_root)
    aliases = read_project_aliases(raw_root, project_root)

    for file in office_files:
        record = protected_by_file[file.file_id]
        if not record.requires_password or record.decryption_success:
            continue
        if not rule.valid:
            record.resolution_status = "error"
            record.failure_reason = f"password_rule_invalid:{rule.reason}"
            continue
        alias = resolve_project_alias(file.project_name, file, aliases)
        record.project_alias_source = alias.source
        if alias.error_status:
            record.resolution_status = alias.error_status
            record.failure_reason = alias.error_status
            continue
        record.project_alias = alias.alias
        dates = find_contract_dates(file.project_name, files, project_root, input_overrides)
        if not dates:
            record.resolution_status = "contract_date_not_found"
            record.failure_reason = "contract_date_not_found"
            continue
        record.contract_start_date = dates[0].date
        record.contract_date_source = dates[0].source_path
        source_path = project_root / file.raw_path
        for secret in make_rule_candidates(file, alias, dates, rule):
            all_candidates.append(secret.record)
            attempt, out_path = decrypt_with_candidate(file, source_path, secret, decrypted_root, len(attempts) + 1)
            attempts.append(attempt)
            record.attempt_count += 1
            if out_path is not None:
                input_overrides[file.file_id] = out_path
                update_record_success(record, secret.record, out_path, "resolved_from_rule")
                break
        if not record.decryption_success:
            record.resolution_status = "password_candidates_failed"
            record.failure_reason = "password_candidates_failed"

    if strict:
        unresolved = [record for record in records if record.requires_password and not record.decryption_success and record.resolution_status not in RESOLUTION_AMBIGUOUS]
        if unresolved:
            names = ", ".join(record.source_path for record in unresolved[:5])
            raise RuntimeError(f"protected Office files could not be resolved: {names}")

    write_protected_outputs(output_dir, records, all_candidates, attempts)
    return ProtectedResolutionResult(records, all_candidates, attempts, input_overrides)


def update_extraction_status(
    output_dir: Path,
    records: list[ProtectedFileRecord],
    candidates: list[PasswordCandidate],
    attempts: list[DecryptionAttempt],
    extraction_status_by_file: dict[str, tuple[bool, str]],
) -> None:
    for record in records:
        success, error = extraction_status_by_file.get(record.file_id, (False, ""))
        record.extraction_success = success
        if error and not record.failure_reason:
            record.failure_reason = error
    write_protected_outputs(output_dir, records, candidates, attempts)
