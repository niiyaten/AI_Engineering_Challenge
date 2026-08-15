from __future__ import annotations

import ast
import csv
import json
import re
from collections import Counter
from pathlib import Path
from typing import Any

from .io_utils import assert_formal_input_allowed, sha1_text, write_json, write_jsonl
from .pptx_colors import slide_color_map
from .schemas import CompactFileProfile, EXTRACTOR_VERSION, ExtractionResult, FileRecord, SearchRecord, to_dict


def compact_text(value: object, limit: int = 1200) -> str:
    text = re.sub(r"\s+", " ", "" if value is None else str(value)).strip()
    return text[:limit]


def record_id(file_id: str, record_type: str, key: str) -> str:
    return f"{record_type}_{sha1_text(file_id + ':' + key)[:16]}"


def make_search_record(file: FileRecord, record_type: str, key: str, text: str, metadata: dict[str, Any]) -> SearchRecord:
    return SearchRecord(
        record_id=record_id(file.file_id, record_type, key),
        file_id=file.file_id,
        record_type=record_type,
        raw_path=file.raw_path,
        text=compact_text(text, 6000),
        metadata=metadata,
    )


def safe_read_text(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp932"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def csv_preview(path: Path, delimiter: str = ",", max_rows: int = 40) -> tuple[list[str], list[list[str]]]:
    rows: list[list[str]] = []
    with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as handle:
        reader = csv.reader(handle, delimiter=delimiter)
        for row in reader:
            rows.append(row)
            if len(rows) >= max_rows:
                break
    header = rows[0] if rows else []
    return header, rows


def write_rows_csv(path: Path, rows: list[list[Any]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8-sig", newline="") as handle:
        writer = csv.writer(handle)
        writer.writerows(rows)


def extract_docx(file: FileRecord, raw_path: Path) -> tuple[dict[str, Any], list[SearchRecord], list[str], list[str]]:
    from docx import Document

    document = Document(raw_path)
    blocks: list[dict[str, Any]] = []
    warnings: list[str] = []
    for i, paragraph in enumerate(document.paragraphs):
        runs = [
            {
                "text": run.text,
                "bold": run.bold if run.bold is not None else (paragraph.style.font.bold if paragraph.style and paragraph.style.font.bold is not None else False),
                "italic": run.italic if run.italic is not None else (paragraph.style.font.italic if paragraph.style and paragraph.style.font.italic is not None else False),
                "underline": run.underline if run.underline is not None else (paragraph.style.font.underline if paragraph.style and paragraph.style.font.underline is not None else False),
                "bold_source": "explicit" if run.bold is not None else "inherited" if paragraph.style and paragraph.style.font.bold is not None else "default",
                "italic_source": "explicit" if run.italic is not None else "inherited" if paragraph.style and paragraph.style.font.italic is not None else "default",
                "underline_source": "explicit" if run.underline is not None else "inherited" if paragraph.style and paragraph.style.font.underline is not None else "default",
                "font_color": str(run.font.color.rgb) if run.font.color and run.font.color.rgb else "",
                "highlight_color": str(run.font.highlight_color) if run.font.highlight_color else "",
                "run_index": run_index,
            }
            for run_index, run in enumerate(paragraph.runs)
            if run.text
        ]
        text = paragraph.text.strip()
        if text or runs:
            style_name = paragraph.style.name if paragraph.style else ""
            heading_level = int(style_name.split()[-1]) if style_name.lower().startswith("heading") and style_name.split()[-1].isdigit() else None
            blocks.append({"type": "paragraph", "index": i, "style": style_name, "heading_level": heading_level, "text": text, "runs": runs})

    tables: list[dict[str, Any]] = []
    for table_index, table in enumerate(document.tables):
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        tables.append({"table_index": table_index, "rows": rows, "section_type": "table"})

    structure = {
        "file_type": "docx",
        "raw_path": file.raw_path,
        "paragraph_count": len(document.paragraphs),
        "table_count": len(document.tables),
        "blocks": blocks,
        "tables": tables,
        "document_title": file.file_name,
        "intermediate_representation": "document_ir_v1",
    }
    records = [
        make_search_record(file, "docx_paragraphs", "paragraphs", "\n".join(block["text"] for block in blocks), {"paragraph_count": len(blocks)})
    ]
    for table in tables:
        text = "\n".join(" | ".join(row) for row in table["rows"][:40])
        records.append(make_search_record(file, "docx_table", str(table["table_index"]), text, {"table_index": table["table_index"]}))
    warnings.append("DOCXコメントと埋め込み画像の詳細抽出は未実装")
    return structure, records, [], warnings


def extract_pptx(file: FileRecord, raw_path: Path) -> tuple[dict[str, Any], list[SearchRecord], list[str], list[str]]:
    from pptx import Presentation

    presentation = Presentation(raw_path)
    raw_colors = slide_color_map(raw_path)
    slides: list[dict[str, Any]] = []
    warnings: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        shapes: list[dict[str, Any]] = []
        tables: list[list[list[str]]] = []
        image_count = 0
        raw_shape_cursor = 0
        for shape_index, shape in enumerate(slide.shapes):
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text.strip())
                shape_runs: list[dict[str, Any]] = []
                if getattr(shape, "has_text_frame", False):
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                        for run_index, run in enumerate(paragraph.runs):
                            color_info = {}
                            slide_colors = raw_colors.get(slide_index, [])
                            if raw_shape_cursor < len(slide_colors):
                                raw_paragraphs = slide_colors[raw_shape_cursor]
                                if paragraph_index < len(raw_paragraphs) and run_index < len(raw_paragraphs[paragraph_index]):
                                    color_info = raw_paragraphs[paragraph_index][run_index]
                            shape_runs.append({
                                "text": run.text,
                                "bold": run.font.bold,
                                "italic": run.font.italic,
                                "underline": run.font.underline,
                                "run_index": run_index,
                                "paragraph_index": paragraph_index,
                                "font_color": color_info.get("normalized_color_name", "unknown"),
                                "font_color_raw_type": color_info.get("raw_color_type", ""),
                                "font_color_raw_value": color_info.get("raw_color_value", ""),
                                "font_color_scheme": color_info.get("scheme_color_name", ""),
                                "font_color_transforms": color_info.get("color_transforms", []),
                                "font_color_resolved_rgb": color_info.get("resolved_rgb", ""),
                                "font_color_resolved_argb": color_info.get("resolved_argb", ""),
                                "font_color_normalized_name": color_info.get("normalized_color_name", "unknown"),
                                "font_color_source": color_info.get("color_source", "unknown"),
                                "font_color_resolution_status": color_info.get("resolution_status", "not_specified"),
                                "fill_color": color_info.get("shape_fill_normalized_name", "unknown"),
                                "fill_color_resolved_rgb": color_info.get("shape_fill_rgb", ""),
                                "fill_color_source": color_info.get("shape_fill_source", "unknown"),
                            })
                    raw_shape_cursor += 1
                shapes.append({"shape_index": shape_index, "text": shape.text.strip(), "runs": shape_runs, "left": shape.left, "top": shape.top, "width": shape.width, "height": shape.height})
            if getattr(shape, "has_table", False):
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                tables.append(rows)
            if "PICTURE" in str(getattr(shape, "shape_type", "")):
                image_count += 1
        slides.append({"slide_number": slide_index, "texts": texts, "tables": tables, "shapes": shapes, "image_count": image_count})

    structure = {
        "file_type": "pptx",
        "raw_path": file.raw_path,
        "slide_count": len(slides),
        "slides": slides,
        "document_title": file.file_name,
        "intermediate_representation": "document_ir_v1",
    }
    records = []
    for slide in slides:
        text = "\n".join(slide["texts"])
        for table in slide["tables"]:
            text += "\n" + "\n".join(" | ".join(row) for row in table)
        records.append(make_search_record(file, "pptx_slide", str(slide["slide_number"]), text, {"slide_number": slide["slide_number"]}))
    warnings.append("PPTXグラフの系列値と画像実体の抽出は限定的")
    return structure, records, [], warnings


def color_to_text(color: Any) -> str:
    if color is None:
        return ""
    parts: list[str] = []
    for name in ("type", "rgb", "indexed", "theme", "tint"):
        value = getattr(color, name, None)
        if value in (None, ""):
            continue
        parts.append(f"{name}={value}")
    return ";".join(parts)


def cell_style(cell: Any) -> dict[str, Any]:
    fill = getattr(cell.fill, "fgColor", None)
    font_color = getattr(cell.font, "color", None)
    font = cell.font
    return {
        "coordinate": cell.coordinate,
        "fill_color": color_to_text(fill),
        "bold": bool(font.bold),
        "italic": bool(font.italic),
        "underline": bool(font.underline),
        "font_color": color_to_text(font_color),
        "number_format": cell.number_format,
        "comment": cell.comment.text if cell.comment else "",
        "formula": str(cell.value) if isinstance(cell.value, str) and cell.value.startswith("=") else "",
    }


def useful_style(style: dict[str, Any]) -> bool:
    return any(style.get(key) for key in ("fill_color", "bold", "italic", "underline", "comment", "formula"))


def extract_xlsx(
    file: FileRecord,
    raw_path: Path,
    table_dir: Path,
) -> tuple[dict[str, Any], list[SearchRecord], list[str], list[str]]:
    import openpyxl

    workbook = openpyxl.load_workbook(raw_path, data_only=False)
    sheets: list[dict[str, Any]] = []
    table_paths: list[str] = []
    for ws in workbook.worksheets:
        rows = [[cell.value for cell in row] for row in ws.iter_rows()]
        csv_path = table_dir / f"{file.file_id}_{safe_filename(ws.title)}.csv"
        write_rows_csv(csv_path, rows)
        table_paths.append(csv_path.as_posix())
        styles = []
        for row in ws.iter_rows():
            for cell in row:
                style = cell_style(cell)
                if useful_style(style):
                    styles.append(style)
        sheets.append(
            {
                "sheet_name": ws.title,
                "max_row": ws.max_row,
                "max_column": ws.max_column,
                "csv_path": csv_path.as_posix(),
                "used_range": f"A1:{ws.cell(ws.max_row, ws.max_column).coordinate}" if ws.max_row and ws.max_column else "",
                "merged_cells": [str(rng) for rng in ws.merged_cells.ranges],
                "styled_cells": styles,
                "auto_filter": str(ws.auto_filter.ref or "") if ws.auto_filter else "",
                "auto_filter_columns": [
                    {
                        "col_id": item.colId,
                        "filters": list(item.filters.filter) if item.filters else [],
                        "custom_filters": [
                            {"operator": custom.operator, "val": custom.val}
                            for custom in (item.customFilters.customFilter if item.customFilters else [])
                        ],
                    }
                    for item in (ws.auto_filter.filterColumn if ws.auto_filter else [])
                ],
                "chart_count": len(getattr(ws, "_charts", [])),
            }
        )

    structure = {"file_type": "xlsx", "raw_path": file.raw_path, "sheet_count": len(sheets), "sheets": sheets}
    records = []
    for sheet in sheets:
        header, rows = csv_preview(Path(sheet["csv_path"]))
        preview = "\n".join(" | ".join(map(lambda v: "" if v is None else str(v), row)) for row in rows[:40])
        style_text = "\n".join(json.dumps(item, ensure_ascii=False) for item in sheet["styled_cells"][:80])
        records.append(
            make_search_record(
                file,
                "xlsx_sheet",
                sheet["sheet_name"],
                f"sheet={sheet['sheet_name']}\ncolumns={header}\n{preview}\nstyles={style_text}",
                {"sheet_name": sheet["sheet_name"], "csv_path": sheet["csv_path"], "styled_cell_count": len(sheet["styled_cells"])},
            )
        )
    return structure, records, table_paths, []


def safe_filename(text: str) -> str:
    ascii_part = re.sub(r"[^0-9A-Za-z_.-]+", "_", text).strip("_") or "sheet"
    # 日本語だけのシート名がすべて sheet.csv へ衝突しないよう、元の名前のハッシュを残す。
    return f"{ascii_part}_{sha1_text(text)[:10]}"


def extract_pdf(
    file: FileRecord,
    raw_path: Path,
    image_dir: Path,
    render_pages: bool,
    max_render_pages: int,
) -> tuple[dict[str, Any], list[SearchRecord], list[str], list[str]]:
    import fitz

    doc = fitz.open(raw_path)
    pages: list[dict[str, Any]] = []
    image_paths: list[str] = []
    for page_index, page in enumerate(doc, start=1):
        text = page.get_text("text")
        raw_blocks = page.get_text("dict").get("blocks", [])
        blocks = []
        for raw_block in raw_blocks:
            safe_block = {"bbox": raw_block.get("bbox"), "lines": []}
            for raw_line in raw_block.get("lines", []) or []:
                safe_line = {"bbox": raw_line.get("bbox"), "spans": []}
                for raw_span in raw_line.get("spans", []) or []:
                    safe_line["spans"].append({key: raw_span.get(key) for key in ("text", "bbox", "font", "size", "flags", "color")})
                safe_block["lines"].append(safe_line)
            blocks.append(safe_block)
        image_path = ""
        if render_pages and (max_render_pages <= 0 or page_index <= max_render_pages):
            pix = page.get_pixmap(matrix=fitz.Matrix(1.0, 1.0), alpha=False)
            out = image_dir / f"{file.file_id}_page{page_index:03d}.png"
            out.parent.mkdir(parents=True, exist_ok=True)
            pix.save(out)
            image_path = out.as_posix()
            image_paths.append(image_path)
        pages.append({"page_number": page_index, "text": text, "blocks": blocks, "image_path": image_path})
    structure = {"file_type": "pdf", "raw_path": file.raw_path, "page_count": len(pages), "pages": pages, "document_title": file.file_name, "intermediate_representation": "document_ir_v1"}
    records = [
        make_search_record(file, "pdf_page", str(page["page_number"]), page["text"], {"page_number": page["page_number"], "image_path": page["image_path"]})
        for page in pages
    ]
    return structure, records, image_paths, []


def extract_image(file: FileRecord, raw_path: Path) -> tuple[dict[str, Any], list[SearchRecord], list[str], list[str]]:
    from PIL import Image

    with Image.open(raw_path) as image:
        width, height = image.size
    structure = {"file_type": "image", "raw_path": file.raw_path, "width": width, "height": height}
    record = make_search_record(file, "image", "metadata", f"画像ファイル: {file.file_name} width={width} height={height}", {"width": width, "height": height})
    return structure, [record], [raw_path.as_posix()], []


def extract_delimited(file: FileRecord, raw_path: Path, table_dir: Path, delimiter: str) -> tuple[dict[str, Any], list[SearchRecord], list[str], list[str]]:
    header, preview_rows = csv_preview(raw_path, delimiter=delimiter, max_rows=80)
    with raw_path.open("r", encoding="utf-8-sig", newline="") as handle:
        rows = list(csv.reader(handle, delimiter=delimiter))
    out = table_dir / f"{file.file_id}.csv"
    write_rows_csv(out, rows)
    text = f"columns={header}\n" + "\n".join(" | ".join(row) for row in preview_rows[:60])
    structure = {
        "file_type": file.extension.lstrip("."),
        "raw_path": file.raw_path,
        "columns": header,
        "row_count": max(0, len(rows) - 1),
        "table_data_path": out.as_posix(),
    }
    return structure, [make_search_record(file, "table_file", "preview", text, {"csv_path": out.as_posix(), "columns": header})], [out.as_posix()], []


def extract_json(file: FileRecord, raw_path: Path) -> tuple[dict[str, Any], list[SearchRecord], list[str], list[str]]:
    obj = json.loads(safe_read_text(raw_path))
    text = json.dumps(obj, ensure_ascii=False, indent=2)[:6000]
    structure = {"file_type": "json", "raw_path": file.raw_path, "top_level_type": type(obj).__name__, "preview": text[:1000]}
    return structure, [make_search_record(file, "json_chunk", "root", text, {})], [], []


def extract_markdown_or_text(file: FileRecord, raw_path: Path) -> tuple[dict[str, Any], list[SearchRecord], list[str], list[str]]:
    text = safe_read_text(raw_path)
    headings = [line.strip() for line in text.splitlines() if line.strip().startswith("#")]
    structure = {"file_type": file.extension.lstrip(".") or "txt", "raw_path": file.raw_path, "char_count": len(text), "headings": headings}
    chunks = [text[i : i + 5000] for i in range(0, len(text), 4500)] or [""]
    records = [make_search_record(file, "markdown_chunk", str(i), chunk, {"chunk_index": i}) for i, chunk in enumerate(chunks)]
    return structure, records, [], []


def extract_python(file: FileRecord, raw_path: Path) -> tuple[dict[str, Any], list[SearchRecord], list[str], list[str]]:
    source = safe_read_text(raw_path)
    tree = ast.parse(source)
    functions = []
    imports = []
    assignments: list[str] = []
    calls: list[str] = []
    for node in ast.walk(tree):
        if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef)):
            functions.append({
                "name": node.name,
                "lineno": node.lineno,
                "end_lineno": getattr(node, "end_lineno", node.lineno),
                "args": [arg.arg for arg in node.args.args],
                "source": ast.get_source_segment(source, node) or "",
            })
        elif isinstance(node, (ast.Import, ast.ImportFrom)):
            imports.append(ast.get_source_segment(source, node) or "")
        elif isinstance(node, (ast.Assign, ast.AnnAssign)):
            targets = node.targets if isinstance(node, ast.Assign) else [node.target]
            assignments.extend(target.id for target in targets if isinstance(target, ast.Name))
        elif isinstance(node, ast.Call):
            call_name = ast.get_source_segment(source, node.func) or ""
            if call_name:
                calls.append(call_name)
    assignments = list(dict.fromkeys(assignments))
    calls = list(dict.fromkeys(calls))
    structure = {
        "file_type": "py",
        "raw_path": file.raw_path,
        "imports": imports,
        "functions": functions,
        "assignments": assignments,
        "calls": calls,
    }
    # 変数名と呼出名を先頭に置き、Compact File Profileでも検索条件を失わないようにする。
    summary = (
        "assignments: " + " ".join(assignments[:80])
        + "\ncalls: " + " ".join(calls[:80])
        + "\nfunctions: " + " ".join(fn["name"] for fn in functions)
        + "\nimports:\n" + "\n".join(imports)
    )
    records = [make_search_record(file, "python_summary", "summary", summary, {"function_count": len(functions)})]
    for fn in functions:
        function_text = fn["source"][:6000] or json.dumps(fn, ensure_ascii=False)
        metadata = {key: value for key, value in fn.items() if key != "source"}
        records.append(make_search_record(file, "python_function", fn["name"], function_text, metadata))
    return structure, records, [], []


def extract_notebook(file: FileRecord, raw_path: Path) -> tuple[dict[str, Any], list[SearchRecord], list[str], list[str]]:
    obj = json.loads(safe_read_text(raw_path))
    cells = []
    records = []
    for i, cell in enumerate(obj.get("cells", [])):
        source = "".join(cell.get("source", []))
        outputs = []
        for output in cell.get("outputs", []):
            outputs.append(str(output.get("text", ""))[:1000])
        cells.append({"cell_index": i, "cell_type": cell.get("cell_type", ""), "source": source, "outputs_preview": outputs})
        records.append(make_search_record(file, "notebook_cell", str(i), source + "\n" + "\n".join(outputs), {"cell_index": i, "cell_type": cell.get("cell_type", "")}))
    return {"file_type": "ipynb", "raw_path": file.raw_path, "cell_count": len(cells), "cells": cells}, records, [], []


def extract_file(
    file: FileRecord,
    project_root: Path,
    extracted_dir: Path,
    table_dir: Path,
    image_dir: Path,
    render_pdf_pages: bool,
    max_pdf_render_pages: int,
    input_path_override: Path | None = None,
) -> tuple[ExtractionResult, list[SearchRecord], CompactFileProfile]:
    """1つのrawファイルから構造情報、検索レコード、短いファイル要約を生成する。"""
    raw_path = input_path_override or project_root / file.raw_path
    if input_path_override is None:
        assert_formal_input_allowed(raw_path, project_root)
    extractor = file.extension.lstrip(".") or "unknown"
    warnings: list[str] = []
    records: list[SearchRecord] = []
    table_paths: list[str] = []
    image_paths: list[str] = []
    structure: dict[str, Any] = {}
    status = "success"
    error = ""
    try:
        if file.is_temp_office_file:
            status = "skipped"
            warnings.append("一時Officeファイルは抽出対象外")
        elif file.extension == ".docx":
            structure, records, image_paths, warnings = extract_docx(file, raw_path)
        elif file.extension == ".pptx":
            structure, records, image_paths, warnings = extract_pptx(file, raw_path)
        elif file.extension == ".xlsx":
            structure, records, table_paths, warnings = extract_xlsx(file, raw_path, table_dir)
        elif file.extension == ".pdf":
            structure, records, image_paths, warnings = extract_pdf(file, raw_path, image_dir, render_pdf_pages, max_pdf_render_pages)
        elif file.extension in {".png", ".jpg", ".jpeg"}:
            structure, records, image_paths, warnings = extract_image(file, raw_path)
        elif file.extension == ".csv":
            structure, records, table_paths, warnings = extract_delimited(file, raw_path, table_dir, ",")
        elif file.extension == ".tsv":
            structure, records, table_paths, warnings = extract_delimited(file, raw_path, table_dir, "\t")
        elif file.extension == ".json":
            structure, records, table_paths, warnings = extract_json(file, raw_path)
        elif file.extension == ".py":
            structure, records, table_paths, warnings = extract_python(file, raw_path)
        elif file.extension == ".ipynb":
            structure, records, table_paths, warnings = extract_notebook(file, raw_path)
        elif file.extension in {".md", ".txt"}:
            structure, records, table_paths, warnings = extract_markdown_or_text(file, raw_path)
        else:
            status = "skipped"
            warnings.append(f"未対応拡張子: {file.extension}")
    except Exception as exc:
        status = "error"
        error = f"{type(exc).__name__}: {exc}"
        structure = {"file_type": file.extension.lstrip("."), "raw_path": file.raw_path, "error": error}
        records = [make_search_record(file, "metadata", "fallback", f"{file.file_name} {file.relative_path}", {"error": error})]

    extracted_path = extracted_dir / f"{file.file_id}.json"
    write_json(extracted_path, structure)
    type_counts = Counter(record.record_type for record in records)
    joined = " ".join(record.text for record in records[:5])
    profile = CompactFileProfile(
        file_id=file.file_id,
        raw_path=file.raw_path,
        file_name=file.file_name,
        extension=file.extension,
        project_name=file.project_name,
        major_folder=file.major_folder,
        document_kind=file.document_kind,
        version_label=file.version_label,
        summary=compact_text(joined, 1200),
        record_type_counts=dict(type_counts),
        keywords=[file.project_name, file.major_folder, file.document_kind, file.version_label, file.file_name],
    )
    result = ExtractionResult(
        file_id=file.file_id,
        raw_path=file.raw_path,
        status=status,
        extractor=extractor,
        extracted_path=extracted_path.as_posix(),
        search_record_count=len(records),
        table_data_paths=table_paths,
        image_paths=image_paths,
        warnings=warnings,
        error=error,
    )
    return result, records, profile


def extract_all(
    files: list[FileRecord],
    project_root: Path,
    output_dir: Path,
    render_pdf_pages: bool = True,
    max_pdf_render_pages: int = 0,
    input_path_overrides: dict[str, Path] | None = None,
) -> tuple[list[ExtractionResult], list[SearchRecord], list[CompactFileProfile]]:
    extracted_dir = output_dir / "extracted"
    table_dir = output_dir / "table_data"
    image_dir = output_dir / "images"
    results: list[ExtractionResult] = []
    search_records: list[SearchRecord] = []
    profiles: list[CompactFileProfile] = []
    overrides = input_path_overrides or {}
    for file in files:
        result, records, profile = extract_file(
            file,
            project_root,
            extracted_dir,
            table_dir,
            image_dir,
            render_pdf_pages,
            max_pdf_render_pages,
            input_path_override=overrides.get(file.file_id),
        )
        results.append(result)
        search_records.extend(records)
        profiles.append(profile)

    write_jsonl(output_dir / "extraction_results.jsonl", [to_dict(item) for item in results])
    write_jsonl(output_dir / "search_records.jsonl", [to_dict(item) for item in search_records])
    write_jsonl(output_dir / "compact_file_profiles.jsonl", [to_dict(item) for item in profiles])
    return results, search_records, profiles
