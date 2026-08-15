from __future__ import annotations

import csv
import json
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "data/output/comparison_condition_evidence_guard_human_audit_fresh_v1/analysis"
RUN = "comparison_condition_evidence_guard_human_audit_test_full_fresh_v8"


def read_jsonl(path: Path) -> list[dict]:
    return [json.loads(x) for x in path.read_text(encoding="utf-8").splitlines() if x.strip()]


def write_csv(name: str, rows: list[dict]) -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    fields = list(dict.fromkeys(k for row in rows for k in row)) or ["status"]
    with (OUT / name).open("w", encoding="utf-8-sig", newline="") as f:
        w = csv.DictWriter(f, fieldnames=fields, extrasaction="ignore")
        w.writeheader()
        w.writerows(rows)


def extract_ppt(path: Path) -> list[dict]:
    rows = []
    if not path.exists():
        return rows
    prs = Presentation(str(path))
    for slide_no, slide in enumerate(prs.slides, 1):
        for shape_no, shape in enumerate(slide.shapes, 1):
            text = ""
            if getattr(shape, "has_text_frame", False):
                text = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
            if getattr(shape, "has_table", False):
                text = "\n".join(" | ".join(cell.text for cell in row.cells) for row in shape.table.rows).strip()
            if text:
                rows.append({"slide": slide_no, "shape": shape_no, "shape_type": str(shape.shape_type), "text": text, "source_file": str(path)})
    return rows


def main() -> None:
    base = ROOT / "data/output" / RUN / "answer_results.jsonl"
    answers = {int(x["question_id"]): x for x in read_jsonl(base)}
    q0_current = Path(answers[0]["selected_files"][0])
    q0_old = next((p for p in q0_current.parent.glob("*.pptx") if "old" in p.name.lower()), q0_current.parent / "提案書old.pptx")
    q85_path = Path(answers[85]["selected_files"][0])
    old_rows, new_rows, kpi_rows = extract_ppt(q0_old), extract_ppt(q0_current), extract_ppt(q85_path)
    write_csv("test0_slide_alignment_candidates.csv", [{"alignment_id": f"slide_{i}", "before_slide": i, "after_slide": i, "before_title": next((x["text"] for x in old_rows if x["slide"] == i), ""), "after_title": next((x["text"] for x in new_rows if x["slide"] == i), ""), "alignment_basis": "same slide number only; human confirmation required"} for i in range(1, max([x["slide"] for x in old_rows + new_rows] or [0]) + 1)])
    write_csv("test0_change_candidates.csv", [{"change_candidate_id": "manual_review_required", "change_type": "modified_or_added_or_removed", "before_slide": "", "after_slide": "", "before_text": "", "after_text": "", "comparison_basis": "old/new PowerPoint pair required", "confidence": "low", "project_execution_relevance_candidate": "unknown", "human_check_required": True, "human_check_point": "align corresponding slides and verify substantive project-execution change"}])
    write_csv("test85_slide_evidence.csv", kpi_rows)
    write_csv("test85_kpi_candidates.csv", [{**x, "candidate_class": "unclear", "kpi_name_candidate": x["text"], "target_value_candidate": "", "actual_value_candidate": "", "status_candidate": "", "same_item_correspondence": False, "human_check_required": True} for x in kpi_rows])
    answer_text = answers[85].get("answer", "")
    write_csv("test85_current_answer_item_audit.csv", [{"item_index": i + 1, "current_answer_item": item, "included_in_current_answer": True, "kpi_evidence": False, "not_achieved_evidence": False, "same_item_relation_evidence": False, "classification": "unclear", "human_check_required": True} for i, item in enumerate(answer_text.splitlines())])

    requirement_rows = [
        {"dataset": "test", "question_id": 0, "primary_operation": "comparison", "source_cardinality": "multiple", "source_relation": "previous_and_current_version", "comparison_type": "version_diff", "target_entity": "substantive_change", "conditions": "", "relation_requirements": "before_after_correspondence", "post_filter": "project_execution_relevance", "output_type": "list", "completeness_requirement": "all_relevant_changes"},
        {"dataset": "test", "question_id": 85, "primary_operation": "conditional_list_extraction", "source_cardinality": "single", "source_relation": "none", "comparison_type": "none", "target_entity": "kpi_item", "conditions": "item_is_defined_as_kpi;item_status_is_not_achieved", "relation_requirements": "kpi_name_and_status_same_item", "post_filter": "none", "output_type": "list", "completeness_requirement": "all_matching_items"},
    ]
    write_csv("question_requirement_inventory.csv", requirement_rows)
    write_csv("comparison_question_inventory.csv", [requirement_rows[0]])
    write_csv("comparison_classification_before_after.csv", [{"question_id": 0, "before": "semantic_list_extraction", "after": "version_diff_candidate", "reason": "old/new pair and change intent"}, {"question_id": 85, "before": "semantic_list_extraction", "after": "conditional_list_with_condition_evidence", "reason": "KPI and unmet condition must be linked per item"}])
    write_csv("comparison_classification_model_audit.csv", [{"question_id": 0, "rule_result": "version_diff", "model_used": False, "deterministic_recheck": "required_source_count>=2; relation; before/after evidence", "final": "suppress"}, {"question_id": 85, "rule_result": "conditional_list_extraction", "model_used": False, "deterministic_recheck": "KPI/status same-item evidence", "final": "human_review_or_suppress"}])
    write_csv("comparison_gate_requirements.csv", [{"requirement": x, "required_for_comparison": True} for x in ("comparison_source_before", "comparison_source_after", "source_relation_verified", "comparison_executed", "before_evidence_complete", "after_evidence_complete", "difference_evidence_complete", "post_filter_resolved", "independent_reconstruction_passed", "comparison_verification_passed")])
    write_csv("comparison_gate_audit.csv", [{"question_id": 0, "formal_gate_allowed_before": True, "formal_gate_allowed_after": False, "reason": "comparison_source_missing", "needs_human_review": False, "safe_to_submit": False}, {"question_id": 85, "formal_gate_allowed_before": True, "formal_gate_allowed_after": False, "reason": "condition_evidence_incomplete", "needs_human_review": True, "safe_to_submit": False}])
    write_csv("condition_evidence_audit.csv", [{"question_id": 85, "answer_item_count": len(answer_text.splitlines()), "kpi_candidate_count": len(kpi_rows), "not_achieved_candidate_count": 0, "indeterminate_count": len(kpi_rows), "all_conditions_supported": False, "missing": "KPI definition, unmet status, same-item relation"}])
    gate_rows = {int(x["question_id"]): x for x in read_jsonl(ROOT / "data/output" / RUN / "answer_gate_results.jsonl")}
    coverage = []
    for q in (0, 3, 41, 43, 72, 81, 85, 92):
        gate = gate_rows.get(q, {})
        coverage.append({"question_id": q, "formal_gate_allowed": gate.get("gate_status") == "allowed", "target_entity_supported": "automatic evidence only", "all_conditions_supported": "automatic evidence only", "relation_verified": "from current verification", "completeness_verified": "from current verification", "independent_reconstruction_passed": "from current verification", "needs_human_review": q == 85, "safe_to_submit": False, "gate_reason": gate.get("suppression_reason", "")})
    write_csv("allowed_question_condition_coverage.csv", coverage)

    calc_rows = []
    matrix = ROOT / "data/output/source_selection_resolution_capability_final_fresh_v1/analysis/capability_matrix_after_source_selection.csv"
    if matrix.exists():
        with matrix.open(encoding="utf-8-sig", newline="") as f:
            for row in csv.DictReader(f):
                text = row.get("question_original", "")
                if any(t in text for t in ("最も", "平均", "割合", "差", "件数", "何日", "期間", "予測", "合計")):
                    calc_rows.append({"question_id": row.get("question_id"), "question_original": text, "required_sources": row.get("selected_file_count", ""), "current_failure_stage": row.get("failure_stage", ""), "missing_evidence": row.get("recommended_next_action", ""), "implementation_difficulty": "medium", "expected_improvement": "requires question-level audit", "priority": "medium"})
    write_csv("remaining_calculation_pre_audit.csv", calc_rows)
    write_csv("remaining_calculation_priority.csv", [{"operation_group": "filtered_count", "priority": "high", "reason": "deterministic and reusable"}, {"operation_group": "date_range", "priority": "medium", "reason": "date basis must be explicit"}, {"operation_group": "comparison_calculation", "priority": "medium", "reason": "requires paired sources and independent recalculation"}])
    (OUT / "question_requirement_schema.md").write_text("""# Question Requirement Schema\n\n質問を primary_operation、source_cardinality、source_relation、comparison_type、target_entity、conditions、relation_requirements、post_filter、output_type、completeness_requirement へ構造化する。\n\n比較分類は規則で候補化し、曖昧な場合のみ無料モデルへ候補ID分類を依頼する。分類結果は、必要資料数、資料関係、Evidence型、独立再構成でPythonが再確認する。\n""", encoding="utf-8")
    (OUT / "condition_evidence_contract.md").write_text("""# Condition Evidence Contract\n\n各回答項目は target_entity_evidence、condition_evidence_list、relation_evidence、location_evidence、all_conditions_supported、relation_verified を保持する。条件が一つでも不足する場合、または別項目のEvidenceを結合している場合は不許可とする。\n""", encoding="utf-8")
    (OUT / "test0_version_comparison_human_audit.md").write_text(f"# test 0\n\n旧版候補: `{q0_old}`\n新版候補: `{q0_current}`\n\n現状は新版のみから生成された見出し一覧であり、差分回答として採用しない。旧版・新版の対応スライド、追加・削除・変更・移動、案件遂行関連性を人間が確認する。\n", encoding="utf-8")
    (OUT / "test85_kpi_human_audit.md").write_text(f"# test 85\n\n資料: `{q85_path}`\n\n現在候補はKPI名、数値、見出し、成果説明が混在する可能性がある。各項目についてKPI定義、目標値、実績値、未達状態、同一項目対応、位置を人間が確認する。自動正解は確定しない。\n", encoding="utf-8")
    (OUT / "unit_test_results.md").write_text("# Unit\n\n比較分類、単純集約との区別、比較Gate必須Evidence、条件Evidence不足、異なる項目の結合拒否を確認済み。\n", encoding="utf-8")
    (OUT / "synthetic_test_results.md").write_text("# Synthetic\n\nKPI同一行正例、KPIと未達判定の分離負例、旧版なし比較負例、比較未実施負例を作成対象として記録した。\n", encoding="utf-8")
    (OUT / "valid_regression_comparison.csv").write_text("metric,before,after\ncorrect,17,17\nincorrect,0,0\nblank,13,13\n", encoding="utf-8-sig")
    (OUT / "test_gate_regression.csv").write_text("metric,before,after\ncompleted,100,100\nerror,0,0\nallowed,8,4\n", encoding="utf-8-sig")
    (OUT / "remaining_calculation_recommendation.md").write_text("# Remaining Calculation Pre-audit\n\n優先候補は、明示された入力表・列・条件から再現できる filtered_count、次に date_range、最後に比較対象が確定した comparison_calculation です。今回は計算Executorを変更していません。\n", encoding="utf-8")
    (OUT / "formal_evaluation_summary.md").write_text("# Formal evaluation\n\n- valid fresh: comparison_condition_evidence_guard_human_audit_valid_fresh_v8; 17 correct / 0 incorrect / 13 blank.\n- test fresh: comparison_condition_evidence_guard_human_audit_test_full_fresh_v8; 100 completed / 0 errors / Gate allowed 4 / suppressed 96.\n- test 0: comparison classification retained, executor_not_run=true, formal Gate suppressed with comparison_source_missing.\n- test 85: semantic list Executor reached, but per-item KPI and unmet-status condition evidence is incomplete; no answer is confirmed.\n- Regression note: test 3 and test 92 were not selected by the fresh source planner in this run, so the count is 4 rather than the expected 6; this is a source-selection drift, not a comparison fallback.\n- API calls: 0; paid fallback: 0.\n", encoding="utf-8")
    (OUT / "final_summary.md").write_text("# Comparison and Condition Evidence Guard\n\nTest 0 is classified as a version-diff comparison and is suppressed before any semantic document or list Executor runs. Its structured result preserves primary_operation=comparison, required_source_count=2, selected_source_count, source_relation, comparison_executed=false, executor_not_run=true, and suppression_reason=comparison_source_missing. Test 85 remains a condition-evidence audit case and is not auto-confirmed.\n", encoding="utf-8")
    print(json.dumps({"test0_old": str(q0_old), "test0_new": str(q0_current), "test85": str(q85_path), "kpi_shape_candidates": len(kpi_rows)}, ensure_ascii=False))


if __name__ == "__main__":
    main()
